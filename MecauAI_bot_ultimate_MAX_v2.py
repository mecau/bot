import asyncio
import logging
import base64
import io
import json
import os
import re
import urllib.parse
import aiohttp
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import CommandStart, Command
from aiogram.exceptions import TelegramRetryAfter, TelegramForbiddenError, TelegramBadRequest
from aiogram.types import (
    InlineKeyboardMarkup, InlineKeyboardButton,
    ReplyKeyboardMarkup, KeyboardButton, Message, BotCommand, BufferedInputFile
)
from openai import AsyncOpenAI
from docx import Document
from docx.shared import Pt, Inches, Mm
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from openpyxl.utils import get_column_letter

# Для вставки картинок в Excel дополнительно нужен пакет Pillow:
#   pip install pillow
# Для чтения PDF в разделе «Анализ документа» нужен пакет pypdf:
#   pip install pypdf
# Оба импортируются лениво (внутри try/except), чтобы их отсутствие
# не ломало весь бот, а просто отключало соответствующую функцию.

from config import (
    BOT_TOKEN, GROQ_API_KEY,
    CHANNEL_1_USERNAME, CHANNEL_1_URL,
    CHANNEL_2_USERNAME, CHANNEL_2_URL,
    TEXT_MODEL, VISION_MODEL, PROMPTS, AD_FOOTER
)

MY_ADMIN_ID = 1184589026

logging.basicConfig(level=logging.INFO)
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher()

groq_client = AsyncOpenAI(
    api_key=GROQ_API_KEY,
    base_url="https://api.groq.com/openai/v1"
)

USERS_FILE = "users.json"
FAV_FILE = "favorites.json"
MODES_FILE = "user_modes.json"

all_users_cache = set()

def load_user_ids() -> set:
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list):
                    return set(data)
        except Exception:
            pass
    return set()

def save_user_id(user_id: int):
    global all_users_cache
    all_users_cache.add(user_id)
    try:
        with open(USERS_FILE, "w", encoding="utf-8") as f:
            json.dump(list(all_users_cache), f, ensure_ascii=False)
    except Exception as e:
        logging.error(f"Не удалось сохранить пользователя: {e}")

all_users_cache = load_user_ids()

def load_favorites() -> dict:
    if os.path.exists(FAV_FILE):
        try:
            with open(FAV_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def save_favorites(favs: dict):
    with open(FAV_FILE, "w", encoding="utf-8") as f:
        json.dump(favs, f, ensure_ascii=False, indent=2)

def add_to_favorites(user_id: int, text: str):
    favs = load_favorites()
    uid_str = str(user_id)
    if uid_str not in favs:
        favs[uid_str] = []
    if text not in favs[uid_str]:
        favs[uid_str].append(text)
        save_favorites(favs)

def load_user_modes() -> dict:
    if os.path.exists(MODES_FILE):
        try:
            with open(MODES_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, dict):
                    return data
        except Exception:
            pass
    return {}

def save_user_mode(user_id: int, mode: str):
    """Сохраняем выбранный режим на диск, чтобы он не сбрасывался при перезапуске бота."""
    modes = load_user_modes()
    modes[str(user_id)] = mode
    try:
        with open(MODES_FILE, "w", encoding="utf-8") as f:
            json.dump(modes, f, ensure_ascii=False)
    except Exception as e:
        logging.error(f"Не удалось сохранить режим пользователя: {e}")

users_db = {}
saved_modes_cache = load_user_modes()
broadcast_states = set()
ppt_states = set()
excel_states = set()
doc_analysis_states = set()
busy_users = set()
user_ppt_images = {}
user_ppt_design = {}
user_excel_mode = {}
# Lightweight per-user workflow state. Cleared after a completed/cancelled flow.
user_task_state = {}
title_states = {}
user_last_request = {}
user_custom_images = {}  # temporary presentation assets: user_id -> list Telegram file_id
ppt_creation_users = set()

def set_task_state(user_id: int, **kwargs):
    state = user_task_state.setdefault(user_id, {})
    state.update(kwargs)
    return state

def get_task_state(user_id: int):
    return user_task_state.get(user_id, {})

def clear_task_state(user_id: int):
    user_task_state.pop(user_id, None)
    user_custom_images.pop(user_id, None)
    ppt_creation_users.discard(user_id)

def ppt_assets(user_id: int):
    return user_custom_images.get(user_id, [])

def save_presentation_image(user_id: int, file_id: str):
    if user_id not in ppt_creation_users:
        return False
    images = user_custom_images.setdefault(user_id, [])
    if file_id not in images:
        images.append(file_id)
    user_custom_images[user_id] = images[-20:]
    return True

def add_custom_image(user_id: int, file_id: str):
    images = user_custom_images.setdefault(user_id, [])
    if file_id not in images:
        images.append(file_id)
    user_custom_images[user_id] = images[-20:]

user_doc_context = {}

# Лёгкая аналитика поведения — без хранения содержимого сообщений.
user_stats = {}

def get_user_stats(user_id: int):
    if user_id not in user_stats:
        user_stats[user_id] = {
            "messages": 0,
            "voice": 0,
            "files": 0,
            "images": 0,
            "exports": 0,
            "last_action": None,
        }
    return user_stats[user_id]

def remember_action(user_id: int, action: str):
    stats = get_user_stats(user_id)
    stats["last_action"] = action
    stats["messages"] += 1

def get_help_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="🎙 Голос", callback_data="help_voice"),
            InlineKeyboardButton(text="📎 Файл", callback_data="help_file"),
        ],
        [
            InlineKeyboardButton(text="📈 Презентация", callback_data="help_ppt"),
            InlineKeyboardButton(text="📊 Excel", callback_data="help_excel"),
        ],
        [
            InlineKeyboardButton(text="🧠 Режим ИИ", callback_data="help_mode"),
        ]
    ])

def get_settings_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🧠 Режим ИИ", callback_data="settings_mode")],
        [InlineKeyboardButton(text="🧹 Очистить контекст", callback_data="settings_clear")],
    ])

def get_user_data(user_id: int):
    if user_id not in users_db:
        users_db[user_id] = {
            "mode": saved_modes_cache.get(str(user_id), "ai"),
            "history": [],
            "last_output": "Здесь пока нет ответов."
        }
    return users_db[user_id]

def clear_pending_states(user_id: int):
    """
    Сбрасывает все 'ожидающие ввода' состояния пользователя (создание презентации,
    таблицы, анализ документа). Нужно, чтобы если пользователь передумал и нажал
    другую кнопку меню, бот не принял её текст за тему презентации/таблицы —
    раньше это приводило к путанице и «зависшим» состояниям.
    """
    ppt_states.discard(user_id)
    excel_states.discard(user_id)
    doc_analysis_states.discard(user_id)
    user_ppt_images.pop(user_id, None)
    user_ppt_design.pop(user_id, None)
    user_excel_mode.pop(user_id, None)
    title_states.pop(user_id, None)

async def call_groq_with_retry(messages, model: str = None, temperature: float = 0.7, max_retries: int = 2, timeout: int = 45):
    """
    Обёртка над вызовом Groq API с повторными попытками и таймаутом.
    Раньше любой сетевой сбой или подвисание запроса приводило к ошибке
    «Попробуй ещё раз» — теперь бот сам делает до 2 повторных попыток
    с небольшой паузой, прежде чем сдаться.
    """
    used_model = model or TEXT_MODEL
    last_err = None
    for attempt in range(max_retries + 1):
        try:
            return await asyncio.wait_for(
                groq_client.chat.completions.create(model=used_model, messages=messages, temperature=temperature),
                timeout=timeout
            )
        except Exception as e:
            last_err = e
            logging.warning(f"Groq API — попытка {attempt + 1} не удалась: {e}")
            if attempt < max_retries:
                await asyncio.sleep(1.5 * (attempt + 1))
    raise last_err


async def transcribe_voice(file_bytes: bytes, filename: str = "voice.ogg") -> str:
    """Надёжное распознавание Telegram voice через Groq Whisper."""
    last_err = None
    for attempt in range(3):
        try:
            result = await asyncio.wait_for(
                groq_client.audio.transcriptions.create(
                    model="whisper-large-v3-turbo",
                    file=(filename, file_bytes),
                    response_format="text",
                    language="ru",
                    temperature=0
                ),
                timeout=60
            )
            transcript = result if isinstance(result, str) else getattr(result, "text", "")
            transcript = re.sub(r"\s+", " ", (transcript or "")).strip()
            if transcript:
                return transcript
            raise ValueError("Пустая расшифровка")
        except Exception as e:
            last_err = e
            logging.warning(f"Whisper attempt {attempt + 1}/3 failed: {e}")
            if attempt < 2:
                await asyncio.sleep(1.5 * (attempt + 1))
    raise last_err

async def check_subscription(user_id: int) -> bool:
    if user_id == MY_ADMIN_ID:
        return True
    channels = [CHANNEL_1_USERNAME, CHANNEL_2_USERNAME]
    for ch in channels:
        try:
            member = await bot.get_chat_member(chat_id=f"@{ch}", user_id=user_id)
            if member.status not in ["creator", "administrator", "member"]:
                return False
        except Exception as e:
            logging.error(f"Ошибка проверки подписки на канал @{ch}: {e}")
            return False
    return True

def get_sub_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📢 Подписаться на канал 1", url=CHANNEL_1_URL)],
        [InlineKeyboardButton(text="📢 Подписаться на канал 2", url=CHANNEL_2_URL)],
        [InlineKeyboardButton(text="✅ Проверить подписку", callback_data="check_sub")]
    ])

def get_main_keyboard():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="🤖 Помощник"), KeyboardButton(text="🛠 Создать")],
            [KeyboardButton(text="📁 Документы"), KeyboardButton(text="⭐ Избранное")],
            [KeyboardButton(text="⚙️ Настройки")],
        ],
        resize_keyboard=True,
        is_persistent=True
    )

def get_admin_keyboard():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="🤖 Помощник"), KeyboardButton(text="📁 Документы")],
            [KeyboardButton(text="🛠 Создать"), KeyboardButton(text="⭐ Избранное")],
            [KeyboardButton(text="⚙️ Настройки")],
            [KeyboardButton(text="📊 Статистика"), KeyboardButton(text="📢 Рассылка")],
        ],
        resize_keyboard=True,
        is_persistent=True
    )

def keyboard_for(user_id: int):
    return get_admin_keyboard() if user_id == MY_ADMIN_ID else get_main_keyboard()

def get_capabilities_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎙 Голос", callback_data="guide_voice"),
         InlineKeyboardButton(text="🖼 Фото", callback_data="guide_image")],
        [InlineKeyboardButton(text="📁 Документы", callback_data="guide_doc"),
         InlineKeyboardButton(text="📈 Презентации", callback_data="guide_ppt")],
        [InlineKeyboardButton(text="📊 Excel", callback_data="guide_excel"),
         InlineKeyboardButton(text="📄 Word", callback_data="guide_word")],
    ])

def get_cancel_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="↩️ Начать заново", callback_data="flow_restart")]
    ])

def get_quick_actions_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎙 Голос", callback_data="guide_voice"),
         InlineKeyboardButton(text="🖼 Фото", callback_data="guide_image")],
        [InlineKeyboardButton(text="📁 Документ", callback_data="guide_doc"),
         InlineKeyboardButton(text="📈 Презентация", callback_data="guide_ppt")],
        [InlineKeyboardButton(text="📊 Excel", callback_data="guide_excel"),
         InlineKeyboardButton(text="📄 Word", callback_data="guide_word")],
    ])

def get_create_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📈 Презентация", callback_data="create_ppt"),
         InlineKeyboardButton(text="📊 Excel", callback_data="create_excel")],
        [InlineKeyboardButton(text="📄 Word", callback_data="create_word"),
         InlineKeyboardButton(text="📑 Титульный лист", callback_data="create_title")],
    ])

def get_helper_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎙 Голос", callback_data="guide_voice"),
         InlineKeyboardButton(text="🖼 Картинка", callback_data="guide_image")],
        [InlineKeyboardButton(text="📁 Документ", callback_data="guide_doc"),
         InlineKeyboardButton(text="📈 Презентация", callback_data="guide_ppt")],
        [InlineKeyboardButton(text="📊 Excel", callback_data="guide_excel"),
         InlineKeyboardButton(text="📄 Word", callback_data="guide_word")],
    ])

def get_answer_inline_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="💡 Проще", callback_data="btn_simplify"),
            InlineKeyboardButton(text="⭐ Сохранить", callback_data="btn_save_fav"),
        ],
        [
            InlineKeyboardButton(text="🔁 Доработать", callback_data="btn_continue"),
            InlineKeyboardButton(text="📄 В Word", callback_data="btn_word"),
        ],
    ])

def clean_text_for_html(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
    return text.strip()

async def safe_answer(message: Message, text: str, parse_mode: str = None, reply_markup=None, disable_web_page_preview: bool = True):
    """
    Безопасная отправка сообщения: если Telegram не может распарсить
    HTML/Markdown в ответе ИИ (например, из-за случайных символов < > _ *),
    бот отправит тот же текст обычным сообщением вместо падения с ошибкой.
    """
    try:
        return await message.answer(
            text, parse_mode=parse_mode, reply_markup=reply_markup,
            disable_web_page_preview=disable_web_page_preview
        )
    except Exception as e:
        logging.warning(f"Не удалось отправить с parse_mode={parse_mode}, отправляю без форматирования: {e}")
        plain_text = re.sub(r'<[^>]+>', '', text) if parse_mode == "HTML" else text
        try:
            return await message.answer(
                plain_text, reply_markup=reply_markup,
                disable_web_page_preview=disable_web_page_preview
            )
        except Exception as e2:
            logging.error(f"Повторная ошибка при отправке сообщения: {e2}")
            return None

def extract_json(raw: str):
    """
    Надёжное извлечение JSON из ответа ИИ. Модель иногда добавляет
    markdown-обёртку (```json ... ```) или лишние пояснения до/после JSON —
    эта функция вытаскивает валидный JSON в любом из этих случаев.
    """
    if not raw:
        raise ValueError("Пустой ответ от модели")
    raw = raw.strip()

    if "```" in raw:
        parts = raw.split("```")
        for part in parts:
            candidate = part.strip()
            if candidate.lower().startswith("json"):
                candidate = candidate[4:].strip()
            if candidate.startswith("{") or candidate.startswith("["):
                raw = candidate
                break

    start_obj = raw.find("{")
    start_arr = raw.find("[")

    if start_obj == -1 and start_arr == -1:
        raise ValueError("В ответе модели не найден JSON")

    if start_arr != -1 and (start_obj == -1 or start_arr < start_obj):
        start = start_arr
        end = raw.rfind("]")
    else:
        start = start_obj
        end = raw.rfind("}")

    if end == -1 or end < start:
        raise ValueError("Не удалось найти конец JSON в ответе модели")

    return json.loads(raw[start:end + 1])

async def generate_ai_image(session: aiohttp.ClientSession, prompt: str, width: int = 800, height: int = 600):
    """Генерация иллюстрации через Pollinations.ai. Возвращает байты картинки или None."""
    try:
        encoded_prompt = urllib.parse.quote(prompt)
        img_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width={width}&height={height}&nologo=true&model=flux"
        async with session.get(img_url, timeout=aiohttp.ClientTimeout(total=20)) as resp:
            if resp.status == 200:
                img_bytes = await resp.read()
                if len(img_bytes) > 1000:
                    return img_bytes
    except Exception as e:
        logging.error(f"Не удалось сгенерировать ИИ-картинку: {e}")
    return None

@dp.message(CommandStart())
async def cmd_start(message: Message):
    user_id = message.from_user.id
    save_user_id(user_id)

    if not await check_subscription(user_id):
        await message.answer(
            f"🔒 Доступ заблокирован!\n\n"
            f"Чтобы пользоваться MecauAI, подпишись на оба наших канала:\n"
            f" 👉 {CHANNEL_1_URL}\n"
            f" 👉 {CHANNEL_2_URL}\n\n"
            f"После подписки нажми кнопку ниже 👇",
            reply_markup=get_sub_keyboard()
        )
        return

    start_text = (
        f"Привет, {message.from_user.first_name}! Ты активировал MecauAI 🚀\n\n"
        "Я твой карманный помощник для учебы и разработки. Можешь отправлять вопросы, задачи, код, изображения, документы и голосовые сообщения — я распознаю запрос, разберусь в нём и дам готовый ответ. Также могу подготовить Word, презентацию, Excel или помочь разобрать документ. Если задача большая, просто опиши её своими словами — я сам выберу подходящий способ обработки."
    )
    await message.answer(start_text, reply_markup=keyboard_for(user_id))

@dp.callback_query(F.data == "check_sub")
async def cb_check_sub(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    if await check_subscription(user_id):
        save_user_id(user_id)
        try:
            await callback.message.delete()
        except Exception:
            pass
        await callback.message.answer(
            "🎉 Подписка на оба канала подтверждена! Добро пожаловать в MecauAI 🚀",
            reply_markup=keyboard_for(user_id)
        )
    else:
        await callback.answer("❌ Ты подписался еще не на все каналы!", show_alert=True)


@dp.message(F.text == "✨ Возможности")
async def cmd_capabilities(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    clear_pending_states(message.from_user.id)
    await message.answer(
        "✨ <b>MecauAI умеет больше, чем кажется</b>\n\n"
        "💬 <b>Общение</b> — вопросы, объяснения, тексты, код, идеи.\n"
        "🎙 <b>Голос</b> — отправь голосовое, я сам превращу его в задачу.\n"
        "🖼 <b>Изображения</b> — распознаю текст, задачи, схемы и содержимое фото.\n"
        "📚 <b>Документы</b> — PDF, DOCX, TXT и фото страниц с вопросами по содержимому.\n"
        "📄 <b>Экспорт</b> — ответ можно сохранить в Word.\n"
        "📈 <b>Презентации</b> — попроси сделать презентацию на нужную тему.\n"
        "📊 <b>Excel</b> — попроси создать таблицу с данными и диаграммой.\n\n"
        "💡 Главное: <b>не нужно искать правильную кнопку</b>. Напиши или скажи задачу своими словами.",
        parse_mode="HTML",
        reply_markup=get_help_keyboard()
    )

@dp.message(F.text == "📂 Файлы")
async def cmd_files_help(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    clear_pending_states(message.from_user.id)
    await message.answer(
        "📂 <b>Работа с файлами</b>\n\n"
        "Просто отправь файл в чат.\n\n"
        "📚 Для PDF/DOCX/TXT я могу прочитать содержание и отвечать на вопросы по нему.\n"
        "🖼 Фото страницы тоже можно отправить — я попробую распознать текст.\n\n"
        "Если хочешь создать файл, просто напиши, например:\n"
        "«Сделай Word из последнего ответа»\n"
        "«Создай презентацию про искусственный интеллект»\n"
        "«Сделай Excel с оценками студентов»",
        parse_mode="HTML"
    )

@dp.message(Command("help"))
async def cmd_help(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🤖 <b>MecauAI умеет работать не только с текстом.</b>\n\n"
        "🎙 Голос — говори задачу обычными словами.\n"
        "🖼 Фото — отправляй задачи, скриншоты и изображения.\n"
        "📁 Документы — загружай PDF/DOCX/TXT и задавай вопросы.\n"
        "📈 Презентации — создание с выбором содержания и оформления.\n"
        "📊 Excel — таблицы, формулы, расчёты и визуализация.\n"
        "📄 Word — оформление готового материала в документ.\n\n"
        "💡 Самый простой способ: просто напиши, что тебе нужно.",
        parse_mode="HTML",
        reply_markup=get_quick_actions_keyboard()
    )

@dp.message(F.text == "🤖 Помощник")
async def cmd_helper_menu(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🤖 <b>Что умеет MecauAI?</b>\n\n"
        "🎙 <b>Голос</b> — отправь голосовое, я распознаю речь и отвечу.\n"
        "🖼 <b>Картинка</b> — отправь фото, я разберу его.\n"
        "📁 <b>Документ</b> — отправь PDF/DOCX/TXT, я изучу его.\n"
        "📈 <b>Презентация</b> — я проведу тебя по оформлению шаг за шагом.\n"
        "📊 <b>Excel</b> — создам таблицу и диаграммы.\n"
        "📄 <b>Word</b> — превращу готовый ответ в документ.\n\n"
        "💡 Не нужно знать команды: выбирай кнопку или просто отправляй задачу.\n\nЕсли не знаешь, что выбрать — просто напиши или скажи, что хочешь получить. Я сам подберу нужный сценарий.",
        parse_mode="HTML",
        reply_markup=get_quick_actions_keyboard()
    )

@dp.message(F.text == "📁 Документы")
async def cmd_documents_menu(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "📁 <b>Работа с документами</b>\n\n"
        "Просто отправь PDF, DOCX или TXT — дальше можно:\n"
        "• сделать конспект;\n"
        "• найти нужную информацию;\n"
        "• задать вопросы по документу;\n"
        "• подготовить материал к экзамену;\n"
        "• сравнить документы;\n"
        "• сделать основу для презентации.\n\n"
        "💡 Можно сразу написать запрос вместе с файлом.",
        parse_mode="HTML",
        reply_markup=get_cancel_keyboard()
    )

@dp.message(F.text == "🛠 Создать")
async def cmd_create_menu(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🛠 <b>Что создать?</b>\n\n"
        "Выбери готовый вариант. Если нужного нет — просто опиши задачу обычными словами.\n\n"
        "📈 В презентации я отдельно помогу с темой, структурой, стилем, титульником и твоими картинками.",
        parse_mode="HTML",
        reply_markup=get_create_keyboard()
    )

@dp.message(F.text == "⚙️ Настройки")
async def cmd_settings(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    clear_pending_states(message.from_user.id)
    user_data = get_user_data(message.from_user.id)
    names = {
        "coder": "💻 ИИ-Программист",
        "ai": "🧠 Академический ассистент",
        "friend": "🫂 Лучший друг"
    }
    await message.answer(
        f"⚙️ <b>Настройки</b>\n\n"
        f"Текущий режим: <b>{names.get(user_data['mode'], 'Стандартный')}</b>\n\n"
        "Можно сменить стиль общения или очистить контекст текущего диалога.",
        parse_mode="HTML",
        reply_markup=get_settings_keyboard()
    )

@dp.callback_query(F.data.in_({"help_voice", "help_file", "help_ppt", "help_excel", "help_mode",
                               "settings_mode", "settings_clear"}))
async def cb_help_and_settings(callback: types.CallbackQuery):
    user_id = callback.from_user.id

    if callback.data == "help_voice":
        await callback.message.answer("🎙 Просто отправь мне голосовое сообщение. Ничего включать не нужно.")
    elif callback.data == "help_file":
        await callback.message.answer("📎 Просто отправь PDF, DOCX, TXT или фото страницы — я предложу, что можно сделать.")
    elif callback.data == "help_ppt":
        await callback.message.answer("📈 Напиши: «Сделай презентацию на тему ...» — бот сам запустит создание презентации.")
    elif callback.data == "help_excel":
        await callback.message.answer("📊 Напиши: «Сделай Excel-таблицу ...» — бот предложит нужный формат и создаст файл.")
    elif callback.data == "help_mode":
        await cmd_mode(callback.message)
    elif callback.data == "settings_mode":
        await cmd_mode(callback.message)
    elif callback.data == "settings_clear":
        user_data = get_user_data(user_id)
        user_data["history"] = []
        user_data["last_output"] = "Здесь пока нет ответов."
        user_doc_context.pop(user_id, None)
        await callback.message.answer("🧹 Контекст очищен. Начинаем с чистого листа.")

    await callback.answer()

@dp.message(F.text == "ℹ️ О MecauAI")
@dp.message(Command("about"))
async def cmd_about(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    about_text = (
            "🧠 <b>Возможности MecauAI</b>\n\n"
            "MecauAI задуман как помощник, которому не нужно объяснять, какую кнопку нажать. "
            "Просто отправь задачу — текстом, голосом, фото или файлом.\n\n"
            "<b>Что уже работает:</b>\n"
            "• 🎙 голосовые сообщения → расшифровка → ответ ИИ;\n"
            "• 🖼 анализ изображений и распознавание текста;\n"
            "• 📚 анализ PDF/DOCX/TXT и вопросы по документу;\n"
            "• 📄 экспорт ответа в Word;\n"
            "• 📈 презентации с дизайном и иллюстрациями;\n"
            "• 📊 Excel-таблицы с диаграммами;\n"
            "• ⭐ избранное;\n"
            "• 🧠 три режима общения.\n\n"
            "<b>Главное улучшение интерфейса:</b> функций много, но кнопок мало. "
            "Сложные действия запускаются обычной фразой — так бот ощущается как помощник, а не как меню."
        )
    await message.answer(about_text, parse_mode="HTML", reply_markup=keyboard_for(message.from_user.id))

@dp.message(F.text == "🔄 Режим ИИ")
@dp.message(Command("mode"))
async def cmd_mode(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(message.from_user.id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💻 ИИ-Программист",       callback_data="set_mode_coder")],
        [InlineKeyboardButton(text="🧠 Академический ассистент", callback_data="set_mode_ai")],
        [InlineKeyboardButton(text="🫂 Лучший друг",             callback_data="set_mode_friend")]
    ])
    await message.answer("Выбери режим работы бота:", reply_markup=kb)

@dp.callback_query(F.data.startswith("set_mode_"))
async def cb_set_mode(callback: types.CallbackQuery):
    mode = callback.data.split("_")[-1]
    user_id = callback.from_user.id
    user_data = get_user_data(user_id)
    user_data["mode"] = mode
    save_user_mode(user_id, mode)
    clear_pending_states(user_id)

    names = {
        "coder": "ИИ-Программист 💻",
        "ai": "Академический ассистент 🧠",
        "friend": "Лучший друг 🫂"
    }
    await callback.message.edit_text(f"Режим переключен на: {names.get(mode, 'Стандартный')}")
    await callback.answer()

@dp.message(F.text == "📊 Статистика")
async def cmd_stats(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    users = load_user_ids()
    total_messages = sum(v.get("messages", 0) for v in user_stats.values())
    total_voice = sum(v.get("voice", 0) for v in user_stats.values())
    total_files = sum(v.get("files", 0) for v in user_stats.values())
    total_images = sum(v.get("images", 0) for v in user_stats.values())
    total_exports = sum(v.get("exports", 0) for v in user_stats.values())
    await message.answer(
        f"📊 <b>Статистика бота</b>\n\n"
        f"👥 Всего пользователей: {len(users)}\n"
        f"💬 Сообщений: {total_messages}\n"
        f"🎙 Голосовых: {total_voice}\n"
        f"📎 Файлов: {total_files}\n"
        f"🖼 Изображений: {total_images}\n"
        f"📄 Экспортов: {total_exports}\n"
        f"⏳ Сейчас заняты генерацией: {len(busy_users)}",
        parse_mode="HTML"
    )

@dp.message(F.text == "📢 Рассылка")
async def cmd_broadcast_prompt(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    clear_pending_states(message.from_user.id)
    broadcast_states.add(message.from_user.id)
    await message.answer("📢 Отправь текст рассылки следующим сообщением (все пользователи бота получат его).")

# ======================= ПРЕЗЕНТАЦИИ =======================

PPT_THEMES = {
    "ppt_blue": {
        "label": "🔵 Классический синий",
        "primary": RGBColor(0x1F, 0x4E, 0x79),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xD6, 0xE4, 0xF0),
        "body_text": RGBColor(0x21, 0x21, 0x21),
    },
    "ppt_dark": {
        "label": "⚫ Тёмный модерн",
        "primary": RGBColor(0x18, 0x18, 0x18),
        "bg": RGBColor(0x2B, 0x2B, 0x2B),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xCC, 0xCC, 0xCC),
        "body_text": RGBColor(0xF2, 0xF2, 0xF2),
    },
    "ppt_green": {
        "label": "🟢 Изумрудный",
        "primary": RGBColor(0x0E, 0x6B, 0x4F),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xD2, 0xF0, 0xE4),
        "body_text": RGBColor(0x1A, 0x1A, 0x1A),
    },
    "ppt_orange": {
        "label": "🟠 Тёплый оранжевый",
        "primary": RGBColor(0xC1, 0x5A, 0x11),
        "bg": RGBColor(0xFF, 0xF8, 0xF2),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xFC, 0xE2, 0xC8),
        "body_text": RGBColor(0x2B, 0x1A, 0x0D),
    },
}

def build_title_slide(prs, theme, topic_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme["primary"]

    tb = slide.shapes.add_textbox(PptxInches(1), PptxInches(2.7), PptxInches(11.3), PptxInches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Презентация проекта"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = PptxPt(40)
    p.font.bold = True
    p.font.color.rgb = theme["title_text"]

    sub_tb = slide.shapes.add_textbox(PptxInches(1), PptxInches(4.3), PptxInches(11.3), PptxInches(1.5))
    sub_tb.text_frame.word_wrap = True
    sp = sub_tb.text_frame.paragraphs[0]
    sp.text = topic_text[:300]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = PptxPt(22)
    sp.font.color.rgb = theme["subtitle_text"]
    return slide

def build_content_slide(prs, theme, idx, item, img_stream):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = theme["bg"]

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), prs.slide_width, PptxInches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["primary"]
    bar.line.fill.background()
    bar.shadow.inherit = False

    tb_title = s.shapes.add_textbox(PptxInches(0.6), PptxInches(0.25), PptxInches(12.1), PptxInches(0.9))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = item.get("title", f"Слайд {idx + 1}")
    p_title.font.size = PptxPt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = theme["title_text"]

    has_img = img_stream is not None
    content_width = PptxInches(6.8) if has_img else PptxInches(11.9)
    tb_content = s.shapes.add_textbox(PptxInches(0.6), PptxInches(1.7), content_width, PptxInches(5.3))
    tf = tb_content.text_frame
    tf.word_wrap = True
    points = item.get("points", []) or ["Нет данных"]
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"●  {pt}"
        p.font.name = 'Calibri'
        p.font.size = PptxPt(18)
        p.font.color.rgb = theme["body_text"]
        p.space_after = PptxPt(10)

    if has_img:
        try:
            s.shapes.add_picture(img_stream, left=PptxInches(7.8), top=PptxInches(1.8), width=PptxInches(4.9))
        except Exception as img_err:
            logging.error(f"Не удалось вставить картинку на слайд {idx}: {img_err}")
    return s

@dp.message(F.text == "📈 Презентация")
async def cmd_ppt_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    clear_pending_states(user_id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=v["label"], callback_data=k)] for k, v in PPT_THEMES.items()
    ] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
    await message.answer("🎨 Выбери дизайн презентации:", reply_markup=kb)

@dp.callback_query(F.data.in_(set(PPT_THEMES.keys())))
async def cb_ppt_design(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    user_ppt_design[user_id] = callback.data
    ppt_states.add(user_id)
    theme_label = PPT_THEMES[callback.data]["label"]
    try:
        await callback.message.edit_text(f"🎨 Дизайн выбран: {theme_label}")
    except Exception:
        pass
    await callback.message.answer(
        "📈 Отправь тему презентации следующим сообщением.\n\n"
        "💡 Хочешь добавить свои картинки? Отправь фото сейчас (можно несколько), а затем — тему презентации.",
    )
    await callback.answer()

@dp.callback_query(F.data == "cancel_flow")
async def cb_cancel_flow(callback: types.CallbackQuery):
    clear_pending_states(callback.from_user.id)
    try:
        await callback.message.edit_text("❌ Действие отменено.")
    except Exception:
        pass
    await callback.answer("Отменено")

# ======================= EXCEL-ТАБЛИЦЫ =======================

EXCEL_MODES = {
    "excel_plain": "📋 Только таблица",
    "excel_chart": "📊 Таблица + диаграмма",
    "excel_full": "🖼 Таблица + диаграмма + ИИ-картинка",
}

@dp.message(F.text == "📊 Excel-таблица")
async def cmd_excel_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    clear_pending_states(user_id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=v, callback_data=k)] for k, v in EXCEL_MODES.items()
    ] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
    await message.answer("📊 Выбери, что должно быть в файле Excel:", reply_markup=kb)

@dp.callback_query(F.data.in_(set(EXCEL_MODES.keys())))
async def cb_excel_mode(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    user_excel_mode[user_id] = callback.data
    excel_states.add(user_id)
    mode_label = EXCEL_MODES[callback.data]
    try:
        await callback.message.edit_text(f"📊 Выбрано: {mode_label}")
    except Exception:
        pass
    await callback.message.answer(
        "Опиши задачу или тему для таблицы следующим сообщением "
        "(например: «Таблица успеваемости студентов с оценками»), "
        "и я сформирую готовый Excel-файл (.xlsx)!"
    )
    await callback.answer()

# ======================= АНАЛИЗ ДОКУМЕНТОВ =======================

MAX_DOC_CONTEXT_CHARS = 12000
MAX_DOC_SIZE_BYTES = 20 * 1024 * 1024  # ограничение Telegram Bot API на скачивание файла

@dp.message(F.text == "📚 Анализ документа")
async def cmd_doc_analysis_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    clear_pending_states(user_id)
    doc_analysis_states.add(user_id)
    await message.answer(
        "📚 Пришли файл в формате <b>.pdf</b>, <b>.docx</b> или <b>.txt</b> (можно также прислать фото страницы) — "
        "я прочитаю его, сделаю краткое содержание, и дальше сможешь задавать вопросы по документу прямо в чате.",
        parse_mode="HTML"
    )

async def extract_text_from_document(ext: str, file_bytes: bytes):
    """Извлекает текст из .pdf / .docx / .txt. Возвращает (текст, ошибка_или_None)."""
    if ext == "txt":
        return file_bytes.decode("utf-8", errors="ignore"), None
    if ext == "docx":
        d = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in d.paragraphs), None
    if ext == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return None, "⚠️ Для чтения PDF на сервере не установлена библиотека pypdf. Установи её командой: pip install pypdf"
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
            return text, None
        except Exception as e:
            return None, f"⚠️ Не удалось прочитать PDF: {e}"
    return None, "⚠️ Неподдерживаемый формат файла."

async def process_document_context(message: Message, filename: str, text: str, status_msg: Message):
    user_id = message.from_user.id
    text = (text or "").strip()
    if not text:
        await status_msg.edit_text("⚠️ Не удалось извлечь текст из документа (возможно, это скан без текстового слоя).")
        doc_analysis_states.discard(user_id)
        return

    truncated = text[:MAX_DOC_CONTEXT_CHARS]
    user_doc_context[user_id] = truncated
    doc_analysis_states.discard(user_id)

    try:
        summary_resp = await call_groq_with_retry(messages=[
            {"role": "system", "content": "Сделай краткое структурированное содержание документа на русском языке (5-8 предложений)."},
            {"role": "user", "content": truncated}
        ])
        summary = clean_text_for_html(summary_resp.choices[0].message.content)
    except Exception as e:
        logging.error(f"Не удалось получить содержание документа: {e}")
        summary = "(не удалось сформировать краткое содержание, но документ загружен — можешь задавать вопросы)"

    await status_msg.edit_text(f"✅ Документ «{filename}» прочитан ({len(text)} символов).")
    await safe_answer(
        message,
        f"📋 <b>Краткое содержание:</b>\n\n{summary}\n\n"
        f"💬 Теперь можешь задавать вопросы по документу, попросить краткий конспект, выделить главное, сравнить части документа или подготовить ответ/структуру на его основе.",
        parse_mode="HTML"
    )

@dp.message(F.document)
async def handle_document(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id not in doc_analysis_states:
        await message.answer("📚 Чтобы проанализировать документ, сначала нажми кнопку «📚 Анализ документа», а затем пришли файл.")
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    get_user_stats(user_id)["files"] += 1
    doc = message.document
    filename = doc.file_name or "документ"
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in ("pdf", "docx", "txt"):
        await message.answer("⚠️ Поддерживаются только форматы .pdf, .docx и .txt")
        return
    if doc.file_size and doc.file_size > MAX_DOC_SIZE_BYTES:
        await message.answer("⚠️ Файл слишком большой (лимит Telegram на скачивание — 20 МБ).")
        return

    busy_users.add(user_id)
    status_msg = await message.answer("📚 Читаю документ...")
    try:
        file_info = await bot.get_file(doc.file_id)
        downloaded = await bot.download_file(file_info.file_path)
        file_bytes = downloaded.read()

        text, err = await extract_text_from_document(ext, file_bytes)
        if err:
            await status_msg.edit_text(err)
            doc_analysis_states.discard(user_id)
            return

        await process_document_context(message, filename, text, status_msg)
    except Exception as e:
        logging.error(f"Ошибка при обработке документа: {e}", exc_info=True)
        try:
            await status_msg.edit_text("⚠️ Ошибка при чтении документа. Попробуй ещё раз.")
        except Exception:
            pass
        doc_analysis_states.discard(user_id)
    finally:
        busy_users.discard(user_id)

@dp.message(F.text == "⭐ Избранное")
async def cmd_favorites(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(user_id)

    favs = load_favorites()
    uid_str = str(user_id)
    user_favs = favs.get(uid_str, [])

    if not user_favs:
        await message.answer("⭐ У тебя пока нет сохраненных ответов в избранном.")
        return

    await message.answer(f"⭐ Твои сохраненные ответы ({len(user_favs)}):")
    for idx, fav_text in enumerate(user_favs, 1):
        display_text = f"Сохранение #{idx}:\n\n{fav_text}"
        if len(display_text) > 4000:
            display_text = display_text[:3997] + "..."
        await message.answer(display_text, disable_web_page_preview=True)

@dp.callback_query(F.data.in_({
    "guide_voice", "guide_image", "guide_doc", "guide_ppt", "guide_excel", "guide_word", "create_ppt", "create_excel", "create_word", "create_title"
}))
async def cb_product_navigation(callback: types.CallbackQuery):
    data = callback.data
    messages = {
        "guide_voice": "🎙 Отправь голосовое — ничего включать не нужно.",
        "guide_image": "🖼 Отправь фото — я могу разобрать задачу, текст или изображение. Для презентации свои картинки можно добавить прямо в процессе её создания.",
        "guide_doc": "📁 Отправь PDF, DOCX или TXT — после загрузки задавай вопросы.",
        "guide_ppt": "📈 Нажми «Создать → Презентация». Я буду задавать вопросы по одному.",
        "guide_excel": "📊 Нажми «Создать → Excel» и опиши, какую таблицу хочешь.",
        "guide_word": "📄 Под ответом нажми «В Word», чтобы получить файл.",
        "create_ppt": "📈 Хорошо. Я помогу выбрать назначение, стиль и титульный лист — без сложных терминов.",
        "create_excel": "📊 Опиши нужную таблицу или пришли пример.",
        "create_word": "📄 Получи ответ и нажми «В Word».",
        "create_title": "📑 Я сначала уточню, для чего титульник нужен, и предложу понятные варианты."
    }
    await callback.message.answer(messages.get(data, "Готово."))
    await callback.answer()

@dp.callback_query(F.data == "flow_restart")
async def cb_flow_restart(callback: types.CallbackQuery):
    clear_task_state(callback.from_user.id)
    await callback.message.answer(
        "🔄 Начинаем заново.\n\n"
        "Просто напиши, что хочешь сделать, или выбери действие ниже.",
        reply_markup=get_quick_actions_keyboard()
    )
    await callback.answer()

@dp.callback_query(F.data.in_({"set_mode_ai", "set_mode_coder", "set_mode_friend"}))
async def cb_modes(callback: types.CallbackQuery):
    modes = {
        "set_mode_ai": ("🤖 AI-помощник", "универсальный режим"),
        "set_mode_coder": ("💻 Кодер", "режим для программирования и разбора кода"),
        "set_mode_friend": ("🤝 Дружеский режим", "более простой и неформальный стиль"),
    }
    title, mode = modes.get(callback.data, ("Режим", "универсальный режим"))
    set_task_state(callback.from_user.id, preferred_mode=mode)
    await callback.message.answer(
        f"✅ <b>{title}</b>\nНастройка сохранена. Я буду учитывать её в следующих запросах.",
        parse_mode="HTML"
    )
    await callback.answer()

@dp.callback_query(F.data.in_({"btn_simplify", "btn_save_fav", "btn_word", "btn_continue"}))
async def cb_answer_actions(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    msg_text = callback.message.text or callback.message.caption or ""

    clean_msg = msg_text.split("—\n⚡")[0].strip() if "—\n⚡" in msg_text else msg_text.strip()
    if not clean_msg:
        await callback.answer("⚠️ Нечего обрабатывать!", show_alert=True)
        return

    if callback.data == "btn_word":
        user_data = get_user_data(user_id)
        user_data["last_output"] = clean_msg
        user_data_stats = get_user_stats(user_id)
        user_data_stats["exports"] += 1

        doc = Document()
        for section in doc.sections:
            section.top_margin = Mm(20)
            section.bottom_margin = Mm(20)
            section.left_margin = Mm(30)
            section.right_margin = Mm(15)
        p = doc.add_paragraph()
        run = p.add_run(clean_msg)
        run.font.name = "Times New Roman"
        run.font.size = Pt(14)
        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        await callback.message.answer_document(
            BufferedInputFile(bio.read(), filename="MecauAI_Answer.docx"),
            caption="📄 Готово — сохранил этот ответ в Word."
        )
        await callback.answer()
        return

    if callback.data == "btn_continue":
        await callback.message.answer(
            "🔁 Продолжаем. Напиши следующим сообщением, что изменить, добавить или уточнить — контекст сохранён."
        )
        await callback.answer()
        return

    if callback.data == "btn_save_fav":
        add_to_favorites(user_id, clean_msg)
        await callback.answer("⭐ Успешно сохранено в избранное!", show_alert=True)
    elif callback.data == "btn_simplify":
        await callback.answer("💡 Сжимаю до сути...")
        try:
            response = await call_groq_with_retry(messages=[
                {"role": "system", "content": "Объясни предельно коротко и просто в 3-5 предложениях."},
                {"role": "user", "content": clean_msg}
            ])
            simplified_reply = clean_text_for_html(response.choices[0].message.content)
            full_reply = f"💡 <b>Коротко на пальцах:</b>\n\n{simplified_reply}{AD_FOOTER}"
            await safe_answer(callback.message, full_reply, parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
        except Exception:
            await callback.message.answer("⚠️ Ошибка при обработке. Попробуй ещё раз. Если повторится, пиши - @mecau")

@dp.message(F.text == "📄 Ответ в Word")
async def cmd_download_word(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(user_id)

    user_data = get_user_data(user_id)
    text_to_save = user_data.get("last_output", "").strip()
    if not text_to_save or text_to_save == "Здесь пока нет ответов.":
        await message.answer("⚠️ Нет данных для сохранения. Попробуй ещё раз. Если повторится, пиши - @mecau")
        return

    doc = Document()
    for section in doc.sections:
        section.top_margin = Mm(20)
        section.bottom_margin = Mm(20)
        section.left_margin = Mm(30)
        section.right_margin = Mm(15)

    for block in re.split(r"\n\s*\n", text_to_save):
        if not block.strip(): continue
        p=doc.add_paragraph(); p.paragraph_format.first_line_indent=Mm(12); p.paragraph_format.line_spacing=1.15
        for i,line in enumerate(block.splitlines()):
            if i: p.add_run().add_break()
            r=p.add_run(re.sub(r"<[^>]+>","",line).strip()); r.font.name='Times New Roman'; r.font.size=Pt(14)

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)

    file_doc = BufferedInputFile(bio.read(), filename="MecauAI_Answer.docx")
    await message.answer_document(file_doc, caption="📄 Вот твой ответ в формате Word!")

@dp.message(F.text == "📑 Титульник ГОСТ")
async def cmd_gost_title(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(message.from_user.id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📘 Индивидуальный проект", callback_data="gost_project")],
        [InlineKeyboardButton(text="📗 Курсовая работа",        callback_data="gost_coursework")],
        [InlineKeyboardButton(text="📕 Дипломная работа (ВКР)", callback_data="gost_diploma")],
        [InlineKeyboardButton(text="📙 Отчёт по практике",      callback_data="gost_practice")],
    ])
    await message.answer("📑 Выбери тип работы для титульника:", reply_markup=kb)

@dp.callback_query(F.data.startswith("gost_"))
async def cb_gost_start(callback: types.CallbackQuery):
    work_type_map = {
        "gost_project": ("ИНДИВИДУАЛЬНЫЙ ПРОЕКТ", "Индивидуальный_проект"),
        "gost_coursework": ("КУРСОВАЯ РАБОТА", "Курсовая_работа"),
        "gost_diploma": ("ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА (ВКР)", "ВКР"),
        "gost_practice": ("ОТЧЁТ ПО ПРАКТИКЕ", "Отчет_по_практике"),
    }
    item=work_type_map.get(callback.data)
    if not item: return await callback.answer("Неизвестный тип", show_alert=True)
    label, filename_base=item; uid=callback.from_user.id
    title_states[uid]={"step":"institution","work_label":label,"filename_base":filename_base}
    await callback.message.answer("📑 <b>Шаг 1/5 — полное название колледжа / вуза</b>\n\nПример: <code>Московский государственный технический университет</code>",parse_mode="HTML")
    await callback.answer()

def _title_preview(d):
    e=lambda x: html.escape(str(x))
    return (f"📑 <b>Проверь данные титульника</b>\n\n🏫 <b>Учебное заведение:</b> {e(d['institution'])}\n👨‍🎓 <b>Студент:</b> {e(d['student'])}\n👨‍🏫 <b>Преподаватель:</b> {e(d['teacher'])}\n📚 <b>Тема:</b> {e(d['topic'])}\n📍 <b>Город, год:</b> {e(d['city'])}, {e(d['year'])}")

@dp.callback_query(F.data.in_({"title_confirm","title_edit"}))
async def cb_title_confirm(callback: types.CallbackQuery):
    uid=callback.from_user.id; d=title_states.get(uid)
    if not d: return await callback.answer("Запусти титульник заново",show_alert=True)
    if callback.data=="title_edit":
        d["step"]="institution"
        await callback.message.answer("✏️ Начинаем заново.\n\n🏫 <b>Шаг 1/5 — полное название колледжа / вуза</b>\nПример: <code>Московский государственный технический университет</code>",parse_mode="HTML")
        return await callback.answer()
    doc=Document()
    for sec in doc.sections: sec.top_margin=Mm(20); sec.bottom_margin=Mm(20); sec.left_margin=Mm(30); sec.right_margin=Mm(15)
    def addc(t,b=False):
        p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run(t); r.font.name="Times New Roman"; r.font.size=Pt(14); r.bold=b
    def addr(t):
        p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.RIGHT; r=p.add_run(t); r.font.name="Times New Roman"; r.font.size=Pt(14)
    addc(d["institution"]); [doc.add_paragraph() for _ in range(4)]; addc(d["work_label"],True); doc.add_paragraph(); addc(f"на тему:\n«{d['topic']}»",True); [doc.add_paragraph() for _ in range(6)]; addr(f"Выполнил(а):\n{d['student']}\n\nРуководитель:\n{d['teacher']}"); [doc.add_paragraph() for _ in range(5)]; addc(f"{d['city']}, {d['year']}")
    bio=io.BytesIO(); doc.save(bio); bio.seek(0)
    await callback.message.answer_document(BufferedInputFile(bio.read(),filename=f"Titulnik_{d['filename_base']}.docx"),caption=f"📑 Титульник «{d['work_label']}» готов!")
    title_states.pop(uid,None); await callback.answer("Готово!")

# Голосовой режим намеренно не добавлен в меню: достаточно отправить voice-сообщение.
# Это сохраняет интерфейс компактным и не перегружает пользователя кнопками.
MENU_BUTTONS = {
    "✨ Возможности", "📂 Файлы", "⭐ Избранное", "⚙️ Настройки",
    "📢 Рассылка", "📊 Статистика",
    # Старые кнопки тоже игнорируем, если они остались у уже открытого клиента.
    "📄 Ответ в Word", "📑 Титульник ГОСТ", "📈 Презентация",
    "📊 Excel-таблица", "📚 Анализ документа", "🔄 Режим ИИ", "ℹ️ О MecauAI"
}


@dp.message(F.voice)
async def handle_voice(message: Message):
    """
    Голосовой режим без отдельной кнопки:
    пользователь отправляет voice -> бот распознаёт речь -> использует
    полученный текст как обычный запрос ИИ.
    """
    user_id = message.from_user.id

    if not await check_subscription(user_id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return

    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    save_user_id(user_id)
    get_user_stats(user_id)["voice"] += 1
    status_msg = await message.answer("🎙 Распознаю голосовое сообщение...")

    busy_users.add(user_id)
    try:
        file_info = await bot.get_file(message.voice.file_id)
        downloaded = await bot.download_file(file_info.file_path)
        audio_bytes = downloaded.read()

        if len(audio_bytes) > MAX_DOC_SIZE_BYTES:
            await status_msg.edit_text("⚠️ Голосовое сообщение слишком большое. Пришли более короткую запись.")
            return

        transcript = await transcribe_voice(audio_bytes)
        if not transcript:
            await status_msg.edit_text("⚠️ Не удалось разобрать речь. Попробуй записать сообщение ещё раз.")
            return

        user_data = get_user_data(user_id)
        current_mode = user_data["mode"]
        system_prompt = PROMPTS.get(current_mode, PROMPTS["ai"])

        doc_context = user_doc_context.get(user_id)
        if doc_context:
            system_prompt = (
                f"{system_prompt}\n\n"
                f"Пользователь ранее загрузил документ. Используй его содержимое, "
                f"если вопрос с ним связан:\n---\n{doc_context}\n---"
            )

        # Голос обрабатывается как обычный текстовый запрос, поэтому история,
        # режим ИИ и экспорт последнего ответа продолжают работать одинаково.
        user_data["history"].append({"role": "user", "content": transcript})
        if len(user_data["history"]) > 6:
            user_data["history"] = user_data["history"][-6:]

        messages_payload = [{"role": "system", "content": system_prompt}] + user_data["history"]

        await status_msg.edit_text(f"🎙 Распознал:\n\n«{transcript[:3500]}»\n\n🤖 Формирую ответ...")
        await bot.send_chat_action(chat_id=message.chat.id, action="typing")

        response = await call_groq_with_retry(messages=messages_payload)
        ai_reply = clean_text_for_html(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply

        full_message = (
            f"🎙 <b>Твой запрос:</b> {transcript}\n\n"
            f"🤖 <b>Ответ:</b>\n\n{ai_reply}{AD_FOOTER}"
        )

        try:
            await status_msg.delete()
        except Exception:
            pass

        parse_mode = "Markdown" if current_mode == "coder" else "HTML"
        await safe_answer(
            message,
            full_message,
            parse_mode=parse_mode,
            reply_markup=get_answer_inline_keyboard()
        )

    except Exception as e:
        logging.error(f"Ошибка обработки голосового сообщения: {e}", exc_info=True)
        try:
            await status_msg.edit_text(
                "⚠️ Не удалось обработать голосовое сообщение. "
                "Попробуй записать его ещё раз чуть короче и без сильного фонового шума."
            )
        except Exception:
            pass
    finally:
        busy_users.discard(user_id)

@dp.message(F.photo)
async def handle_photo(message: Message):
    user_id = message.from_user.id
    if save_presentation_image(user_id, message.photo[-1].file_id):
        await message.answer("🖼 Картинка добавлена в текущую презентацию. Можешь отправить ещё или продолжить.")
        return
    try:
        add_custom_image(user_id, message.photo[-1].file_id)
    except Exception:
        pass
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    save_user_id(user_id)
    get_user_stats(user_id)["images"] += 1
    photo = message.photo[-1]

    # Скачивание фото вынесено в отдельный try/except: раньше ошибка сети или
    # слишком большого файла на этом шаге не перехватывалась и «ломала» обработку.
    try:
        if photo.file_size and photo.file_size > MAX_DOC_SIZE_BYTES:
            await message.answer("⚠️ Фото слишком большое (лимит Telegram на скачивание — 20 МБ). Пришли файл меньшего размера.")
            return
        file_info = await bot.get_file(photo.file_id)
        downloaded_file = await bot.download_file(file_info.file_path)
        img_bytes = downloaded_file.read()
    except Exception as e:
        logging.error(f"Не удалось скачать фото от пользователя {user_id}: {e}")
        await message.answer("⚠️ Не удалось загрузить фото. Проверь соединение и попробуй отправить ещё раз.")
        return

    # Если пользователь сейчас создает презентацию, сохраняем картинку для слайдов
    if user_id in ppt_states:
        if user_id not in user_ppt_images:
            user_ppt_images[user_id] = []
        user_ppt_images[user_id].append(img_bytes)
        await message.answer(f"🖼 Картинка сохранена для презентации ({len(user_ppt_images[user_id])} шт.)! Теперь отправь тему презентации.")
        return

    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    # Если пользователь в режиме анализа документа — считаем фото сканом страницы
    # и распознаём текст через vision-модель вместо обычного анализа изображения.
    if user_id in doc_analysis_states:
        busy_users.add(user_id)
        status_msg = await message.answer("📚 Распознаю текст на фото...")
        try:
            base64_image = base64.b64encode(img_bytes).decode('utf-8')
            response = await call_groq_with_retry(
                model=VISION_MODEL,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Распознай и выведи весь текст, присутствующий на этом изображении, без комментариев."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                    ]
                }]
            )
            recognized_text = clean_text_for_html(response.choices[0].message.content)
            await process_document_context(message, "фото документа", recognized_text, status_msg)
        except Exception as e:
            logging.error(f"Ошибка распознавания текста с фото: {e}", exc_info=True)
            try:
                await status_msg.edit_text("⚠️ Не удалось распознать текст на фото. Попробуй ещё раз.")
            except Exception:
                pass
            doc_analysis_states.discard(user_id)
        finally:
            busy_users.discard(user_id)
        return

    # Обычный режим анализа картинки через Vision ИИ
    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        base64_image = base64.b64encode(img_bytes).decode('utf-8')
        response = await call_groq_with_retry(
            model=VISION_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": message.caption or "Подробно проанализируй это изображение, распознай текст если он есть и ответь по сути."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }]
        )
        reply = clean_text_for_html(response.choices[0].message.content)
        user_data = get_user_data(user_id)
        user_data["last_output"] = reply
        await safe_answer(message, f"{reply}{AD_FOOTER}", parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
    except Exception as e:
        logging.error(f"Ошибка обработки фото: {e}")
        await message.answer("⚠️ Ошибка при обработке изображения. Попробуй ещё раз. Если повторится, пиши - @mecau")

@dp.message(F.text)
async def handle_text(message: Message):
    user_id = message.from_user.id

    ts=title_states.get(user_id)
    if ts:
        raw=(message.text or "").strip(); step=ts.get("step")
        if not raw: return await message.answer("⚠️ Поле не может быть пустым.")
        if step=="institution": ts["institution"]=raw[:300]; ts["step"]="student"; return await message.answer("👨‍🎓 <b>Шаг 2/5 — ФИО студента</b>\nПример: <code>Иванов Иван Иванович</code>",parse_mode="HTML")
        if step=="student": ts["student"]=raw[:150]; ts["step"]="teacher"; return await message.answer("👨‍🏫 <b>Шаг 3/5 — ФИО преподавателя</b>\nПример: <code>Петров Пётр Петрович</code>",parse_mode="HTML")
        if step=="teacher": ts["teacher"]=raw[:150]; ts["step"]="topic"; return await message.answer("📚 <b>Шаг 4/5 — тема работы</b>",parse_mode="HTML")
        if step=="topic": ts["topic"]=raw[:500]; ts["step"]="city_year"; return await message.answer("📍 <b>Шаг 5/5 — город и год</b>\nПример: <code>Москва, 2026</code>",parse_mode="HTML")
        if step=="city_year":
            m=re.match(r"^\s*(.+?)\s*[,;]\s*(\d{4})\s*$",raw) or re.match(r"^\s*(.+?)\s+(\d{4})\s*$",raw)
            if not m: return await message.answer("⚠️ Формат: <code>Москва, 2026</code>",parse_mode="HTML")
            city,year=m.group(1).strip(" ,;.-"),m.group(2)
            if not city or not 1900<=int(year)<=2200: return await message.answer("⚠️ Проверь город и год.")
            ts.update(city=city,year=year,step="confirm")
            kb=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="✅ Всё верно",callback_data="title_confirm"),InlineKeyboardButton(text="✏️ Изменить",callback_data="title_edit")]])
            return await message.answer(_title_preview(ts),parse_mode="HTML",reply_markup=kb)
        return await message.answer("Нажми кнопку под предпросмотром.")

    # Естественные команды: пользователь может запускать функции словами,
    # не изучая меню.
    normalized = (message.text or "").lower().strip()

    if normalized in {"/clear", "очисти контекст", "забудь предыдущий диалог", "начать заново"}:
        user_data = get_user_data(user_id)
        user_data["history"] = []
        user_data["last_output"] = "Здесь пока нет ответов."
        user_doc_context.pop(user_id, None)
        title_states.pop(user_id, None)
        await message.answer("🧹 Готово. Контекст очищен — можем начать заново.")
        return

    if normalized in {"/help", "помощь", "что ты умеешь", "что умеешь"}:
        await cmd_capabilities(message)
        return

    if normalized in {"сделай презентацию", "создай презентацию"} or normalized.startswith(("сделай презентацию ", "создай презентацию ")):
        # Сохраняем привычный сценарий, но позволяем запускать его текстом.
        await cmd_ppt_prompt(message)
        if normalized not in {"сделай презентацию", "создай презентацию"}:
            topic = re.sub(r"^(сделай|создай) презентацию\s*", "", message.text, flags=re.I).strip()
            if topic and user_id not in busy_users:
                # Запоминаем тему как следующий шаг, чтобы пользователь не вводил её второй раз.
                ppt_states.discard(user_id)
                user_ppt_design[user_id] = "ppt_blue"
                ppt_states.add(user_id)
                await message.answer(f"📈 Тема сохранена: «{topic}»\nТеперь выбери дизайн — или нажми отмену.")
        return

    if normalized.startswith(("сделай excel ", "создай excel ", "сделай таблицу ", "создай таблицу ")):
        await cmd_excel_prompt(message)
        return

    if user_id == MY_ADMIN_ID and user_id in broadcast_states:
        broadcast_states.remove(user_id)
        users = load_user_ids()
        success, failed = 0, 0
        status_msg = await message.answer("📢 Начинаю рассылку...")
        for uid in users:
            try:
                await bot.send_message(uid, message.text, disable_web_page_preview=True)
                success += 1
            except TelegramRetryAfter as e:
                # Telegram просит подождать из-за лимита скорости — ждём и пробуем ещё раз
                await asyncio.sleep(e.retry_after)
                try:
                    await bot.send_message(uid, message.text, disable_web_page_preview=True)
                    success += 1
                except Exception:
                    failed += 1
            except TelegramForbiddenError:
                # Пользователь заблокировал бота — это не сбой, просто пропускаем
                failed += 1
            except Exception:
                failed += 1
            await asyncio.sleep(0.05)
        await status_msg.edit_text(f"✅ Рассылка завершена!\n\n👥 Доставлено: {success}\n❌ Ошибок: {failed}")
        return

    # ----------------- ГЕНЕРАЦИЯ ПРЕЗЕНТАЦИИ -----------------
    if user_id in ppt_states:
        ppt_states.remove(user_id)
        if user_id in busy_users:
            await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
            return
        busy_users.add(user_id)
        theme_key = user_ppt_design.pop(user_id, "ppt_blue")
        theme = PPT_THEMES.get(theme_key, PPT_THEMES["ppt_blue"])
        status_msg = await message.answer("📈 Генерирую презентацию и иллюстрации, подожди немного...")
        try:
            await bot.send_chat_action(chat_id=message.chat.id, action="typing")
            user_images = user_ppt_images.pop(user_id, [])
            num_slides = max(len(user_images), 5) if user_images else 5

            response = await call_groq_with_retry(
                messages=[
                    {
                        "role": "system",
                        "content": f"Составь презентацию ровно из {num_slides} слайдов. Ответ выдай СТРОГО в формате валидного JSON без markdown-оформления (без ```json), в виде списка объектов: [{{\"title\": \"Заголовок слайда\", \"points\": [\"Тезис 1\", \"Тезис 2\"], \"image_prompt\": \"Detailed professional visual illustration of the slide topic in English\"}}]. В поле image_prompt ВСЕГДА пиши качественный промпт на английском языке для генерации картинки."
                    },
                    {"role": "user", "content": f"Тема: {message.text}"}
                ],
                temperature=0.7
            )
            raw_content = clean_text_for_html(response.choices[0].message.content)
            slides_data = extract_json(raw_content)
            if not isinstance(slides_data, list) or not slides_data:
                raise ValueError("Модель вернула пустой или некорректный список слайдов")

            prs = Presentation()
            prs.slide_width = PptxInches(13.333)
            prs.slide_height = PptxInches(7.5)

            build_title_slide(prs, theme, message.text)

            async with aiohttp.ClientSession() as session:
                for idx, item in enumerate(slides_data):
                    img_stream = None
                    if user_images and idx < len(user_images):
                        img_stream = io.BytesIO(user_images[idx])
                    elif item.get("image_prompt"):
                        img_bytes = await generate_ai_image(session, item["image_prompt"], width=800, height=600)
                        if img_bytes:
                            img_stream = io.BytesIO(img_bytes)

                    build_content_slide(prs, theme, idx, item, img_stream)

            bio = io.BytesIO()
            prs.save(bio)
            bio.seek(0)
            file_doc = BufferedInputFile(bio.read(), filename="Presentation.pptx")

            await status_msg.delete()
            await message.answer_document(file_doc, caption=f"📈 Презентация готова! Дизайн: {theme['label']}")

        except Exception as e:
            logging.error(f"Критическая ошибка при генерации презентации: {e}", exc_info=True)
            try:
                await status_msg.delete()
            except Exception:
                pass
            await message.answer("⚠️ Произошла ошибка при создании презентации. Попробуй еще раз.")
        finally:
            busy_users.discard(user_id)
        return

    # ----------------- ГЕНЕРАЦИЯ EXCEL -----------------
    if user_id in excel_states:
        excel_states.remove(user_id)
        if user_id in busy_users:
            await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
            return
        busy_users.add(user_id)
        excel_mode = user_excel_mode.pop(user_id, "excel_plain")
        status_msg = await message.answer("📊 Генерирую таблицу Excel, подожди секунду...")
        try:
            await bot.send_chat_action(chat_id=message.chat.id, action="typing")
            response = await call_groq_with_retry(
                messages=[
                    {
                        "role": "system",
                        "content": "Создай структуру таблицы на основе запроса пользователя. Ответ выдай СТРОГО в формате валидного JSON без markdown-оформления (без ```json), в виде объекта с ключом 'headers' (список строк-заголовков) и ключом 'rows' (список списков с данными для строк таблицы, где числовые значения — это числа, а не строки). Пример: {\"headers\": [\"№\", \"Название\", \"Значение\"], \"rows\": [[1, \"Пример 1\", 100], [2, \"Пример 2\", 200]]}"
                    },
                    {"role": "user", "content": f"Задача для таблицы: {message.text}"}
                ],
                temperature=0.7
            )
            raw_content = clean_text_for_html(response.choices[0].message.content)
            table_data = extract_json(raw_content)

            headers = table_data.get("headers") or ["Колонка 1", "Колонка 2"]
            raw_rows = table_data.get("rows") or [["Данные 1", "Данные 2"]]

            # Нормализуем строки: выравниваем количество столбцов под заголовки,
            # чтобы "кривые" ответы от ИИ не ломали генерацию файла.
            num_cols = len(headers)
            norm_rows = []
            for row in raw_rows:
                row = list(row) if isinstance(row, (list, tuple)) else [row]
                if len(row) < num_cols:
                    row = row + [""] * (num_cols - len(row))
                elif len(row) > num_cols:
                    row = row[:num_cols]
                norm_rows.append(row)

            wb = Workbook()
            ws = wb.active
            ws.title = "Таблица"

            header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="4F81BD", end_color="4F81BD", fill_type="solid")
            align_center = Alignment(horizontal="center", vertical="center")
            align_left = Alignment(horizontal="left", vertical="center")
            align_right = Alignment(horizontal="right", vertical="center")
            border_thin = Border(
                left=Side(style='thin', color='D9D9D9'),
                right=Side(style='thin', color='D9D9D9'),
                top=Side(style='thin', color='D9D9D9'),
                bottom=Side(style='thin', color='D9D9D9')
            )

            ws.append(headers)
            for col_num in range(1, num_cols + 1):
                cell = ws.cell(row=1, column=col_num)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = align_center
                cell.border = border_thin

            def is_number(val):
                if isinstance(val, (int, float)):
                    return True
                try:
                    float(str(val).replace(",", "."))
                    return True
                except (TypeError, ValueError):
                    return False

            for row_idx, row_data in enumerate(norm_rows, start=2):
                ws.append(row_data)
                for col_num in range(1, num_cols + 1):
                    cell = ws.cell(row=row_idx, column=col_num)
                    cell.font = Font(name="Calibri", size=11)
                    cell.alignment = align_right if is_number(cell.value) else align_left
                    cell.border = border_thin

            for col in ws.columns:
                max_length = 0
                column_letter = col[0].column_letter
                for cell in col:
                    try:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    except Exception:
                        pass
                ws.column_dimensions[column_letter].width = max(max_length + 4, 12)

            # ---- Диаграмма по числовым данным (если выбран нужный режим) ----
            if excel_mode in ("excel_chart", "excel_full") and num_cols >= 2:
                try:
                    numeric_cols = []
                    for col_idx in range(1, num_cols):
                        total, numeric = 0, 0
                        for row in norm_rows:
                            val = row[col_idx]
                            if val not in ("", None):
                                total += 1
                                if is_number(val):
                                    numeric += 1
                        if total > 0 and numeric / total >= 0.6:
                            numeric_cols.append(col_idx)

                    if numeric_cols:
                        chart = BarChart()
                        chart.title = "Диаграмма по данным таблицы"
                        chart.style = 10
                        chart.x_axis.title = headers[0]
                        min_col = min(numeric_cols) + 1
                        max_col = max(numeric_cols) + 1
                        data = Reference(ws, min_col=min_col, max_col=max_col, min_row=1, max_row=len(norm_rows) + 1)
                        cats = Reference(ws, min_col=1, min_row=2, max_row=len(norm_rows) + 1)
                        chart.add_data(data, titles_from_data=True)
                        chart.set_categories(cats)
                        chart.width = 18
                        chart.height = 10
                        anchor_col_letter = get_column_letter(num_cols + 2)
                        ws.add_chart(chart, f"{anchor_col_letter}2")
                except Exception as chart_err:
                    logging.error(f"Не удалось построить диаграмму в Excel: {chart_err}")

            # ---- ИИ-картинка по теме (если выбран полный режим) ----
            if excel_mode == "excel_full":
                try:
                    async with aiohttp.ClientSession() as session:
                        img_bytes = await generate_ai_image(
                            session,
                            f"Professional business infographic illustration about: {message.text}",
                            width=768, height=512
                        )
                    if img_bytes:
                        from openpyxl.drawing.image import Image as XLImage
                        xl_img = XLImage(io.BytesIO(img_bytes))
                        xl_img.width = 480
                        xl_img.height = 320
                        img_row = len(norm_rows) + 4
                        ws.add_image(xl_img, f"A{img_row}")
                except Exception as img_err:
                    logging.error(f"Не удалось вставить ИИ-картинку в Excel (проверь, установлен ли Pillow): {img_err}")

            bio = io.BytesIO()
            wb.save(bio)
            bio.seek(0)
            file_doc = BufferedInputFile(bio.read(), filename="MecauAI_Table.xlsx")

            await status_msg.delete()
            await message.answer_document(file_doc, caption=f"📊 Твоя таблица Excel готова! ({EXCEL_MODES.get(excel_mode)})")
        except Exception as e:
            logging.error(f"Ошибка при создании Excel: {e}", exc_info=True)
            try:
                await status_msg.delete()
            except Exception:
                pass
            await message.answer("⚠️ Произошла ошибка при создании таблицы Excel. Попробуй еще раз.")
        finally:
            busy_users.discard(user_id)
        return

    if message.text in MENU_BUTTONS:
        return

    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    save_user_id(user_id)
    user_data = get_user_data(user_id)
    current_mode = user_data["mode"]

    system_prompt = PROMPTS.get(current_mode, PROMPTS["ai"]) + (
        "\n\nПРАВИЛА ОТВЕТА: "
        "отвечай естественно и полезно. Сначала дай прямой ответ, затем при необходимости "
        "короткое объяснение, пример или конкретные следующие шаги. "
        "Не повторяй вопрос пользователя. Не добавляй пустые вступления. "
        "Если запрос можно решить практическим способом — предложи готовый результат, а не только совет. "
        "Если информации недостаточно — задай один точный уточняющий вопрос."
    ) + "\n\nОтвечай содержательно: если вопрос требует объяснения, дай краткий контекст, основной ответ и 1-2 практических примера или следующих шага. Не растягивай ответ без необходимости."

    # Если ранее был загружен документ через «📚 Анализ документа» — добавляем
    # его содержимое в системный промпт, чтобы отвечать с учётом контекста.
    doc_context = user_doc_context.get(user_id)
    if doc_context:
        system_prompt = (
            f"{system_prompt}\n\n"
            f"Пользователь ранее загрузил документ. Используй его содержимое, если вопрос с ним связан:\n"
            f"---\n{doc_context}\n---"
        )

    user_data["history"].append({"role": "user", "content": message.text})
    if len(user_data["history"]) > 6:
        user_data["history"] = user_data["history"][-6:]

    messages_payload = [{"role": "system", "content": system_prompt}] + user_data["history"]

    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        response = await call_groq_with_retry(messages=messages_payload)
        ai_reply = clean_text_for_html(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply

        if current_mode == "coder":
            full_message = f"{ai_reply}\n\n{AD_FOOTER}"
            await safe_answer(message, full_message, parse_mode="Markdown", reply_markup=get_answer_inline_keyboard())
        else:
            full_message = f"{ai_reply}{AD_FOOTER}"
            await safe_answer(message, full_message, parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
    except Exception:
        await message.answer("⚠️ Ошибка при обработке запроса. Попробуй ещё раз. Если повторится, пиши - @mecau")

@dp.errors()
async def global_error_handler(event: types.ErrorEvent):
    """
    Глобальный перехватчик непредвиденных ошибок. Раньше необработанное
    исключение в любом хэндлере могло привести к тому, что пользователь
    просто не получал ответа и не понимал, что произошло. Теперь бот
    логирует ошибку и сообщает пользователю, что что-то пошло не так.
    """
    logging.error(f"Необработанная ошибка: {event.exception}", exc_info=True)
    try:
        update = event.update
        if update.message:
            await update.message.answer("⚠️ Произошла непредвиденная ошибка. Попробуй ещё раз позже. Если повторится — пиши @mecau")
        elif update.callback_query:
            await update.callback_query.answer("⚠️ Произошла ошибка, попробуй ещё раз.", show_alert=True)
    except Exception:
        pass
    return True

async def main():
    await bot.set_my_commands([
        BotCommand(command="start", description="Перезапустить бота"),
        BotCommand(command="mode",  description="Сменить режим ассистента"),
        BotCommand(command="about", description="О возможностях"),
    ])
    await bot.delete_webhook()
    print("🚀 Бот MecauAI запущен и слушает обновления...")
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
