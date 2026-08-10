import asyncio
import logging
import base64
import io
import json
import os
import re
import urllib.parse
import aiohttp
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import CommandStart, Command
from aiogram.exceptions import TelegramRetryAfter, TelegramForbiddenError, TelegramBadRequest
from aiogram.types import (
    InlineKeyboardMarkup, InlineKeyboardButton,
    ReplyKeyboardMarkup, KeyboardButton, Message, BotCommand, BufferedInputFile
)
from openai import AsyncOpenAI
from docx import Document
from docx.shared import Pt, Inches, Mm
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from openpyxl.utils import get_column_letter

# Для вставки картинок в Excel дополнительно нужен пакет Pillow:
#   pip install pillow
# Для чтения PDF в разделе «Анализ документа» нужен пакет pypdf:
#   pip install pypdf
# Оба импортируются лениво (внутри try/except), чтобы их отсутствие
# не ломало весь бот, а просто отключало соответствующую функцию.

from config import (
    BOT_TOKEN, GROQ_API_KEY,
    CHANNEL_1_USERNAME, CHANNEL_1_URL,
    CHANNEL_2_USERNAME, CHANNEL_2_URL,
    TEXT_MODEL, VISION_MODEL, PROMPTS, AD_FOOTER
)

MY_ADMIN_ID = 1184589026

logging.basicConfig(level=logging.INFO)
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher()

groq_client = AsyncOpenAI(
    api_key=GROQ_API_KEY,
    base_url="https://api.groq.com/openai/v1"
)

USERS_FILE = "users.json"
FAV_FILE = "favorites.json"
MODES_FILE = "user_modes.json"

all_users_cache = set()

def load_user_ids() -> set:
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list):
                    return set(data)
        except Exception:
            pass
    return set()

def save_user_id(user_id: int):
    global all_users_cache
    all_users_cache.add(user_id)
    try:
        with open(USERS_FILE, "w", encoding="utf-8") as f:
            json.dump(list(all_users_cache), f, ensure_ascii=False)
    except Exception as e:
        logging.error(f"Не удалось сохранить пользователя: {e}")

all_users_cache = load_user_ids()

def load_favorites() -> dict:
    if os.path.exists(FAV_FILE):
        try:
            with open(FAV_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def save_favorites(favs: dict):
    with open(FAV_FILE, "w", encoding="utf-8") as f:
        json.dump(favs, f, ensure_ascii=False, indent=2)

def add_to_favorites(user_id: int, text: str):
    favs = load_favorites()
    uid_str = str(user_id)
    if uid_str not in favs:
        favs[uid_str] = []
    if text not in favs[uid_str]:
        favs[uid_str].append(text)
        save_favorites(favs)

def load_user_modes() -> dict:
    if os.path.exists(MODES_FILE):
        try:
            with open(MODES_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, dict):
                    return data
        except Exception:
            pass
    return {}

def save_user_mode(user_id: int, mode: str):
    """Сохраняем выбранный режим на диск, чтобы он не сбрасывался при перезапуске бота."""
    modes = load_user_modes()
    modes[str(user_id)] = mode
    try:
        with open(MODES_FILE, "w", encoding="utf-8") as f:
            json.dump(modes, f, ensure_ascii=False)
    except Exception as e:
        logging.error(f"Не удалось сохранить режим пользователя: {e}")

users_db = {}
saved_modes_cache = load_user_modes()
broadcast_states = set()
ppt_states = set()
excel_states = set()
doc_analysis_states = set()
busy_users = set()
user_ppt_images = {}
user_ppt_design = {}
user_excel_mode = {}
user_doc_context = {}

def get_user_data(user_id: int):
    if user_id not in users_db:
        users_db[user_id] = {
            "mode": saved_modes_cache.get(str(user_id), "ai"),
            "history": [],
            "last_output": "Здесь пока нет ответов."
        }
    return users_db[user_id]

def clear_pending_states(user_id: int):
    """
    Сбрасывает все 'ожидающие ввода' состояния пользователя (создание презентации,
    таблицы, анализ документа). Нужно, чтобы если пользователь передумал и нажал
    другую кнопку меню, бот не принял её текст за тему презентации/таблицы —
    раньше это приводило к путанице и «зависшим» состояниям.
    """
    ppt_states.discard(user_id)
    excel_states.discard(user_id)
    doc_analysis_states.discard(user_id)
    user_ppt_images.pop(user_id, None)
    user_ppt_design.pop(user_id, None)
    user_excel_mode.pop(user_id, None)

async def call_groq_with_retry(messages, model: str = None, temperature: float = 0.7, max_retries: int = 2, timeout: int = 45):
    """
    Обёртка над вызовом Groq API с повторными попытками и таймаутом.
    Раньше любой сетевой сбой или подвисание запроса приводило к ошибке
    «Попробуй ещё раз» — теперь бот сам делает до 2 повторных попыток
    с небольшой паузой, прежде чем сдаться.
    """
    used_model = model or TEXT_MODEL
    last_err = None
    for attempt in range(max_retries + 1):
        try:
            return await asyncio.wait_for(
                groq_client.chat.completions.create(model=used_model, messages=messages, temperature=temperature),
                timeout=timeout
            )
        except Exception as e:
            last_err = e
            logging.warning(f"Groq API — попытка {attempt + 1} не удалась: {e}")
            if attempt < max_retries:
                await asyncio.sleep(1.5 * (attempt + 1))
    raise last_err

async def check_subscription(user_id: int) -> bool:
    if user_id == MY_ADMIN_ID:
        return True
    channels = [CHANNEL_1_USERNAME, CHANNEL_2_USERNAME]
    for ch in channels:
        try:
            member = await bot.get_chat_member(chat_id=f"@{ch}", user_id=user_id)
            if member.status not in ["creator", "administrator", "member"]:
                return False
        except Exception as e:
            logging.error(f"Ошибка проверки подписки на канал @{ch}: {e}")
            return False
    return True

def get_sub_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📢 Подписаться на канал 1", url=CHANNEL_1_URL)],
        [InlineKeyboardButton(text="📢 Подписаться на канал 2", url=CHANNEL_2_URL)],
        [InlineKeyboardButton(text="✅ Проверить подписку", callback_data="check_sub")]
    ])

def get_main_keyboard():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="📄 Ответ в Word"), KeyboardButton(text="📑 Титульник ГОСТ")],
            [KeyboardButton(text="📈 Презентация"), KeyboardButton(text="📊 Excel-таблица")],
            [KeyboardButton(text="📚 Анализ документа"), KeyboardButton(text="⭐ Избранное")],
            [KeyboardButton(text="🔄 Режим ИИ"), KeyboardButton(text="ℹ️ О MecauAI")]
        ],
        resize_keyboard=True
    )

def get_admin_keyboard():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="📄 Ответ в Word"), KeyboardButton(text="📑 Титульник ГОСТ")],
            [KeyboardButton(text="📈 Презентация"), KeyboardButton(text="📊 Excel-таблица")],
            [KeyboardButton(text="📚 Анализ документа"), KeyboardButton(text="⭐ Избранное")],
            [KeyboardButton(text="🔄 Режим ИИ"), KeyboardButton(text="ℹ️ О MecauAI")],
            [KeyboardButton(text="📊 Статистика"), KeyboardButton(text="📢 Рассылка")]
        ],
        resize_keyboard=True
    )

def keyboard_for(user_id: int):
    return get_admin_keyboard() if user_id == MY_ADMIN_ID else get_main_keyboard()

def get_answer_inline_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="💡 Объяснить проще", callback_data="btn_simplify"),
            InlineKeyboardButton(text="⭐ В избранное", callback_data="btn_save_fav")
        ]
    ])

def clean_text_for_html(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
    return text.strip()

async def safe_answer(message: Message, text: str, parse_mode: str = None, reply_markup=None, disable_web_page_preview: bool = True):
    """
    Безопасная отправка сообщения: если Telegram не может распарсить
    HTML/Markdown в ответе ИИ (например, из-за случайных символов < > _ *),
    бот отправит тот же текст обычным сообщением вместо падения с ошибкой.
    """
    try:
        return await message.answer(
            text, parse_mode=parse_mode, reply_markup=reply_markup,
            disable_web_page_preview=disable_web_page_preview
        )
    except Exception as e:
        logging.warning(f"Не удалось отправить с parse_mode={parse_mode}, отправляю без форматирования: {e}")
        plain_text = re.sub(r'<[^>]+>', '', text) if parse_mode == "HTML" else text
        try:
            return await message.answer(
                plain_text, reply_markup=reply_markup,
                disable_web_page_preview=disable_web_page_preview
            )
        except Exception as e2:
            logging.error(f"Повторная ошибка при отправке сообщения: {e2}")
            return None

def extract_json(raw: str):
    """
    Надёжное извлечение JSON из ответа ИИ. Модель иногда добавляет
    markdown-обёртку (```json ... ```) или лишние пояснения до/после JSON —
    эта функция вытаскивает валидный JSON в любом из этих случаев.
    """
    if not raw:
        raise ValueError("Пустой ответ от модели")
    raw = raw.strip()

    if "```" in raw:
        parts = raw.split("```")
        for part in parts:
            candidate = part.strip()
            if candidate.lower().startswith("json"):
                candidate = candidate[4:].strip()
            if candidate.startswith("{") or candidate.startswith("["):
                raw = candidate
                break

    start_obj = raw.find("{")
    start_arr = raw.find("[")

    if start_obj == -1 and start_arr == -1:
        raise ValueError("В ответе модели не найден JSON")

    if start_arr != -1 and (start_obj == -1 or start_arr < start_obj):
        start = start_arr
        end = raw.rfind("]")
    else:
        start = start_obj
        end = raw.rfind("}")

    if end == -1 or end < start:
        raise ValueError("Не удалось найти конец JSON в ответе модели")

    return json.loads(raw[start:end + 1])

async def generate_ai_image(session: aiohttp.ClientSession, prompt: str, width: int = 800, height: int = 600):
    """Генерация иллюстрации через Pollinations.ai. Возвращает байты картинки или None."""
    try:
        encoded_prompt = urllib.parse.quote(prompt)
        img_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width={width}&height={height}&nologo=true&model=flux"
        async with session.get(img_url, timeout=aiohttp.ClientTimeout(total=20)) as resp:
            if resp.status == 200:
                img_bytes = await resp.read()
                if len(img_bytes) > 1000:
                    return img_bytes
    except Exception as e:
        logging.error(f"Не удалось сгенерировать ИИ-картинку: {e}")
    return None

@dp.message(CommandStart())
async def cmd_start(message: Message):
    user_id = message.from_user.id
    save_user_id(user_id)

    if not await check_subscription(user_id):
        await message.answer(
            f"🔒 Доступ заблокирован!\n\n"
            f"Чтобы пользоваться MecauAI, подпишись на оба наших канала:\n"
            f" 👉 {CHANNEL_1_URL}\n"
            f" 👉 {CHANNEL_2_URL}\n\n"
            f"После подписки нажми кнопку ниже 👇",
            reply_markup=get_sub_keyboard()
        )
        return

    start_text = (
        f"Привет, {message.from_user.first_name}! Ты активировал MecauAI 🚀\n\n"
        "Я твой карманный помощник для учебы и разработки. Отправляй любые вопросы, задачи, картинки, код или создавай презентации с таблицами!"
    )
    await message.answer(start_text, reply_markup=keyboard_for(user_id))

@dp.callback_query(F.data == "check_sub")
async def cb_check_sub(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    if await check_subscription(user_id):
        save_user_id(user_id)
        try:
            await callback.message.delete()
        except Exception:
            pass
        await callback.message.answer(
            "🎉 Подписка на оба канала подтверждена! Добро пожаловать в MecauAI 🚀",
            reply_markup=keyboard_for(user_id)
        )
    else:
        await callback.answer("❌ Ты подписался еще не на все каналы!", show_alert=True)

@dp.message(F.text == "ℹ️ О MecauAI")
@dp.message(Command("about"))
async def cmd_about(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    about_text = (
        "🧠 <b>Возможности MecauAI</b>\n\n"
        "• 💻 Режим «ИИ-Программист»: написание чистого кода с копированием в 1 клик.\n"
        "• 🧠 Академический ассистент и 🫂 Лучший друг — разные стили общения.\n"
        "• 📄 Экспорт ответов в .docx и титульники по ГОСТу (проект, курсовая, ВКР, практика).\n"
        "• 📈 Генерация презентаций (.pptx) с выбором дизайна и твоими или ИИ-картинками.\n"
        "• 📊 Автоматические таблицы Excel (.xlsx) с диаграммами и ИИ-иллюстрациями.\n"
        "• 📚 Анализ документов (.pdf/.docx/.txt/фото) с ответами на вопросы по содержимому.\n"
        "• ⭐ Избранное — сохраняй лучшие ответы бота.\n\n"

        "🆕 <b>Что нового в этом обновлении:</b>\n\n"

        "<u>🔧 Исправлены ошибки:</u>\n"
        "• Устранена критическая ошибка в рассылке сообщений, из-за которой бот вообще не запускался.\n"
        "• Исправлена генерация картинок для презентаций — раньше ссылка на сервис была битой, и картинки не подставлялись в слайды.\n"
        "• Переработана обработка ответов ИИ в формате JSON (для таблиц и презентаций) — теперь бот корректно распознаёт данные, даже если модель добавила лишний текст вокруг JSON.\n"
        "• Бот больше не падает с ошибкой при отправке сообщений с разметкой (HTML/Markdown) — если форматирование некорректно, сообщение отправится обычным текстом.\n"
        "• Таблицы Excel стали устойчивы к «кривым» данным от ИИ (когда в разных строках разное число столбцов).\n"
        "• Исправлена ошибка обработки фото: скачивание изображения теперь отдельно защищено от сбоев сети и слишком больших файлов (лимит 20 МБ), с понятным сообщением вместо тихого сбоя.\n"
        "• Исправлена путаница с «зависшими» состояниями: если во время создания презентации/таблицы нажать другую кнопку меню, прежнее ожидание темы корректно сбрасывается.\n"
        "• Добавлен общий обработчик непредвиденных ошибок — бот не «зависает» и всегда отвечает пользователю.\n\n"

        "<u>✨ Новые функции:</u>\n"
        "• 📚 <b>Анализ документов</b> — загрузи .pdf, .docx, .txt или фото страницы, бот прочитает файл, пришлёт краткое содержание и будет отвечать на вопросы по документу прямо в чате.\n"
        "• 🎨 <b>Выбор дизайна презентации</b> — 4 темы оформления (🔵 синяя, ⚫ тёмная, 🟢 изумрудная, 🟠 оранжевая) с фирменными цветами, заголовками и акцентами на каждом слайде.\n"
        "• 📊 <b>Расширенное создание Excel</b> — выбор формата: 📋 только таблица / 📊 таблица + диаграмма (строится автоматически по числовым данным) / 🖼 таблица + диаграмма + ИИ-картинка по теме.\n"
        "• ❌ Кнопка «Отмена» в диалогах выбора дизайна/формата.\n\n"

        "<u>🛡 Повышена надёжность:</u>\n"
        "• Запросы к ИИ теперь автоматически повторяются при сетевых сбоях (до 2 повторов) и имеют таймаут — бот не «висит» бесконечно.\n"
        "• Добавлена защита от одновременного запуска нескольких тяжёлых задач одним пользователем («Подожди, предыдущая задача ещё выполняется»).\n"
        "• Выбранный режим общения (ИИ-Программист / Академический / Друг) теперь сохраняется на диске и не сбрасывается при перезапуске бота.\n"
        "• Рассылка сообщений админом теперь корректно обрабатывает лимиты Telegram (flood control) и заблокировавших бота пользователей, не прерываясь на ошибках.\n"
    )
    await message.answer(about_text, parse_mode="HTML", reply_markup=keyboard_for(message.from_user.id))

@dp.message(F.text == "🔄 Режим ИИ")
@dp.message(Command("mode"))
async def cmd_mode(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(message.from_user.id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💻 ИИ-Программист",       callback_data="set_mode_coder")],
        [InlineKeyboardButton(text="🧠 Академический ассистент", callback_data="set_mode_ai")],
        [InlineKeyboardButton(text="🫂 Лучший друг",             callback_data="set_mode_friend")]
    ])
    await message.answer("Выбери режим работы бота:", reply_markup=kb)

@dp.callback_query(F.data.startswith("set_mode_"))
async def cb_set_mode(callback: types.CallbackQuery):
    mode = callback.data.split("_")[-1]
    user_id = callback.from_user.id
    user_data = get_user_data(user_id)
    user_data["mode"] = mode
    save_user_mode(user_id, mode)
    clear_pending_states(user_id)

    names = {
        "coder": "ИИ-Программист 💻",
        "ai": "Академический ассистент 🧠",
        "friend": "Лучший друг 🫂"
    }
    await callback.message.edit_text(f"Режим переключен на: {names.get(mode, 'Стандартный')}")
    await callback.answer()

@dp.message(F.text == "📊 Статистика")
async def cmd_stats(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    users = load_user_ids()
    await message.answer(
        f"📊 Статистика бота:\n\n"
        f"👥 Всего пользователей: {len(users)}\n"
        f"⏳ Сейчас заняты генерацией: {len(busy_users)}"
    )

@dp.message(F.text == "📢 Рассылка")
async def cmd_broadcast_prompt(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    clear_pending_states(message.from_user.id)
    broadcast_states.add(message.from_user.id)
    await message.answer("📢 Отправь текст рассылки следующим сообщением (все пользователи бота получат его).")

# ======================= ПРЕЗЕНТАЦИИ =======================

PPT_THEMES = {
    "ppt_blue": {
        "label": "🔵 Классический синий",
        "primary": RGBColor(0x1F, 0x4E, 0x79),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xD6, 0xE4, 0xF0),
        "body_text": RGBColor(0x21, 0x21, 0x21),
    },
    "ppt_dark": {
        "label": "⚫ Тёмный модерн",
        "primary": RGBColor(0x18, 0x18, 0x18),
        "bg": RGBColor(0x2B, 0x2B, 0x2B),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xCC, 0xCC, 0xCC),
        "body_text": RGBColor(0xF2, 0xF2, 0xF2),
    },
    "ppt_green": {
        "label": "🟢 Изумрудный",
        "primary": RGBColor(0x0E, 0x6B, 0x4F),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xD2, 0xF0, 0xE4),
        "body_text": RGBColor(0x1A, 0x1A, 0x1A),
    },
    "ppt_orange": {
        "label": "🟠 Тёплый оранжевый",
        "primary": RGBColor(0xC1, 0x5A, 0x11),
        "bg": RGBColor(0xFF, 0xF8, 0xF2),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xFC, 0xE2, 0xC8),
        "body_text": RGBColor(0x2B, 0x1A, 0x0D),
    },
}

def build_title_slide(prs, theme, topic_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme["primary"]

    tb = slide.shapes.add_textbox(PptxInches(1), PptxInches(2.7), PptxInches(11.3), PptxInches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Презентация проекта"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = PptxPt(40)
    p.font.bold = True
    p.font.color.rgb = theme["title_text"]

    sub_tb = slide.shapes.add_textbox(PptxInches(1), PptxInches(4.3), PptxInches(11.3), PptxInches(1.5))
    sub_tb.text_frame.word_wrap = True
    sp = sub_tb.text_frame.paragraphs[0]
    sp.text = topic_text[:300]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = PptxPt(22)
    sp.font.color.rgb = theme["subtitle_text"]
    return slide

def build_content_slide(prs, theme, idx, item, img_stream):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = theme["bg"]

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), prs.slide_width, PptxInches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["primary"]
    bar.line.fill.background()
    bar.shadow.inherit = False

    tb_title = s.shapes.add_textbox(PptxInches(0.6), PptxInches(0.25), PptxInches(12.1), PptxInches(0.9))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = item.get("title", f"Слайд {idx + 1}")
    p_title.font.size = PptxPt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = theme["title_text"]

    has_img = img_stream is not None
    content_width = PptxInches(6.8) if has_img else PptxInches(11.9)
    tb_content = s.shapes.add_textbox(PptxInches(0.6), PptxInches(1.7), content_width, PptxInches(5.3))
    tf = tb_content.text_frame
    tf.word_wrap = True
    points = item.get("points", []) or ["Нет данных"]
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"●  {pt}"
        p.font.name = 'Calibri'
        p.font.size = PptxPt(18)
        p.font.color.rgb = theme["body_text"]
        p.space_after = PptxPt(10)

    if has_img:
        try:
            s.shapes.add_picture(img_stream, left=PptxInches(7.8), top=PptxInches(1.8), width=PptxInches(4.9))
        except Exception as img_err:
            logging.error(f"Не удалось вставить картинку на слайд {idx}: {img_err}")
    return s

@dp.message(F.text == "📈 Презентация")
async def cmd_ppt_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    clear_pending_states(user_id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=v["label"], callback_data=k)] for k, v in PPT_THEMES.items()
    ] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
    await message.answer("🎨 Выбери дизайн презентации:", reply_markup=kb)

@dp.callback_query(F.data.in_(set(PPT_THEMES.keys())))
async def cb_ppt_design(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    user_ppt_design[user_id] = callback.data
    ppt_states.add(user_id)
    theme_label = PPT_THEMES[callback.data]["label"]
    try:
        await callback.message.edit_text(f"🎨 Дизайн выбран: {theme_label}")
    except Exception:
        pass
    await callback.message.answer(
        "📈 Отправь тему презентации следующим сообщением.\n\n"
        "💡 Хочешь добавить свои картинки? Отправь фото сейчас (можно несколько), а затем — тему презентации.",
    )
    await callback.answer()

@dp.callback_query(F.data == "cancel_flow")
async def cb_cancel_flow(callback: types.CallbackQuery):
    clear_pending_states(callback.from_user.id)
    try:
        await callback.message.edit_text("❌ Действие отменено.")
    except Exception:
        pass
    await callback.answer("Отменено")

# ======================= EXCEL-ТАБЛИЦЫ =======================

EXCEL_MODES = {
    "excel_plain": "📋 Только таблица",
    "excel_chart": "📊 Таблица + диаграмма",
    "excel_full": "🖼 Таблица + диаграмма + ИИ-картинка",
}

@dp.message(F.text == "📊 Excel-таблица")
async def cmd_excel_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    clear_pending_states(user_id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=v, callback_data=k)] for k, v in EXCEL_MODES.items()
    ] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
    await message.answer("📊 Выбери, что должно быть в файле Excel:", reply_markup=kb)

@dp.callback_query(F.data.in_(set(EXCEL_MODES.keys())))
async def cb_excel_mode(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    user_excel_mode[user_id] = callback.data
    excel_states.add(user_id)
    mode_label = EXCEL_MODES[callback.data]
    try:
        await callback.message.edit_text(f"📊 Выбрано: {mode_label}")
    except Exception:
        pass
    await callback.message.answer(
        "Опиши задачу или тему для таблицы следующим сообщением "
        "(например: «Таблица успеваемости студентов с оценками»), "
        "и я сформирую готовый Excel-файл (.xlsx)!"
    )
    await callback.answer()

# ======================= АНАЛИЗ ДОКУМЕНТОВ =======================

MAX_DOC_CONTEXT_CHARS = 12000
MAX_DOC_SIZE_BYTES = 20 * 1024 * 1024  # ограничение Telegram Bot API на скачивание файла

@dp.message(F.text == "📚 Анализ документа")
async def cmd_doc_analysis_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    clear_pending_states(user_id)
    doc_analysis_states.add(user_id)
    await message.answer(
        "📚 Пришли файл в формате <b>.pdf</b>, <b>.docx</b> или <b>.txt</b> (можно также прислать фото страницы) — "
        "я прочитаю его, сделаю краткое содержание, и дальше сможешь задавать вопросы по документу прямо в чате.",
        parse_mode="HTML"
    )

async def extract_text_from_document(ext: str, file_bytes: bytes):
    """Извлекает текст из .pdf / .docx / .txt. Возвращает (текст, ошибка_или_None)."""
    if ext == "txt":
        return file_bytes.decode("utf-8", errors="ignore"), None
    if ext == "docx":
        d = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in d.paragraphs), None
    if ext == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return None, "⚠️ Для чтения PDF на сервере не установлена библиотека pypdf. Установи её командой: pip install pypdf"
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
            return text, None
        except Exception as e:
            return None, f"⚠️ Не удалось прочитать PDF: {e}"
    return None, "⚠️ Неподдерживаемый формат файла."

async def process_document_context(message: Message, filename: str, text: str, status_msg: Message):
    user_id = message.from_user.id
    text = (text or "").strip()
    if not text:
        await status_msg.edit_text("⚠️ Не удалось извлечь текст из документа (возможно, это скан без текстового слоя).")
        doc_analysis_states.discard(user_id)
        return

    truncated = text[:MAX_DOC_CONTEXT_CHARS]
    user_doc_context[user_id] = truncated
    doc_analysis_states.discard(user_id)

    try:
        summary_resp = await call_groq_with_retry(messages=[
            {"role": "system", "content": "Сделай краткое структурированное содержание документа на русском языке (5-8 предложений)."},
            {"role": "user", "content": truncated}
        ])
        summary = clean_text_for_html(summary_resp.choices[0].message.content)
    except Exception as e:
        logging.error(f"Не удалось получить содержание документа: {e}")
        summary = "(не удалось сформировать краткое содержание, но документ загружен — можешь задавать вопросы)"

    await status_msg.edit_text(f"✅ Документ «{filename}» прочитан ({len(text)} символов).")
    await safe_answer(
        message,
        f"📋 <b>Краткое содержание:</b>\n\n{summary}\n\n"
        f"💬 Теперь просто пиши вопросы по документу в чат — отвечу с учётом его содержимого.",
        parse_mode="HTML"
    )

@dp.message(F.document)
async def handle_document(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id not in doc_analysis_states:
        await message.answer("📚 Чтобы проанализировать документ, сначала нажми кнопку «📚 Анализ документа», а затем пришли файл.")
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    doc = message.document
    filename = doc.file_name or "документ"
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in ("pdf", "docx", "txt"):
        await message.answer("⚠️ Поддерживаются только форматы .pdf, .docx и .txt")
        return
    if doc.file_size and doc.file_size > MAX_DOC_SIZE_BYTES:
        await message.answer("⚠️ Файл слишком большой (лимит Telegram на скачивание — 20 МБ).")
        return

    busy_users.add(user_id)
    status_msg = await message.answer("📚 Читаю документ...")
    try:
        file_info = await bot.get_file(doc.file_id)
        downloaded = await bot.download_file(file_info.file_path)
        file_bytes = downloaded.read()

        text, err = await extract_text_from_document(ext, file_bytes)
        if err:
            await status_msg.edit_text(err)
            doc_analysis_states.discard(user_id)
            return

        await process_document_context(message, filename, text, status_msg)
    except Exception as e:
        logging.error(f"Ошибка при обработке документа: {e}", exc_info=True)
        try:
            await status_msg.edit_text("⚠️ Ошибка при чтении документа. Попробуй ещё раз.")
        except Exception:
            pass
        doc_analysis_states.discard(user_id)
    finally:
        busy_users.discard(user_id)

@dp.message(F.text == "⭐ Избранное")
async def cmd_favorites(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(user_id)

    favs = load_favorites()
    uid_str = str(user_id)
    user_favs = favs.get(uid_str, [])

    if not user_favs:
        await message.answer("⭐ У тебя пока нет сохраненных ответов в избранном.")
        return

    await message.answer(f"⭐ Твои сохраненные ответы ({len(user_favs)}):")
    for idx, fav_text in enumerate(user_favs, 1):
        display_text = f"Сохранение #{idx}:\n\n{fav_text}"
        if len(display_text) > 4000:
            display_text = display_text[:3997] + "..."
        await message.answer(display_text, disable_web_page_preview=True)

@dp.callback_query(F.data.in_({"btn_simplify", "btn_save_fav"}))
async def cb_answer_actions(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    msg_text = callback.message.text or callback.message.caption or ""

    clean_msg = msg_text.split("—\n⚡")[0].strip() if "—\n⚡" in msg_text else msg_text.strip()
    if not clean_msg:
        await callback.answer("⚠️ Нечего обрабатывать!", show_alert=True)
        return

    if callback.data == "btn_save_fav":
        add_to_favorites(user_id, clean_msg)
        await callback.answer("⭐ Успешно сохранено в избранное!", show_alert=True)
    elif callback.data == "btn_simplify":
        await callback.answer("💡 Сжимаю до сути...")
        try:
            response = await call_groq_with_retry(messages=[
                {"role": "system", "content": "Объясни предельно коротко и просто в 3-5 предложениях."},
                {"role": "user", "content": clean_msg}
            ])
            simplified_reply = clean_text_for_html(response.choices[0].message.content)
            full_reply = f"💡 <b>Коротко на пальцах:</b>\n\n{simplified_reply}{AD_FOOTER}"
            await safe_answer(callback.message, full_reply, parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
        except Exception:
            await callback.message.answer("⚠️ Ошибка при обработке. Попробуй ещё раз. Если повторится, пиши - @mecau")

@dp.message(F.text == "📄 Ответ в Word")
async def cmd_download_word(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(user_id)

    user_data = get_user_data(user_id)
    text_to_save = user_data.get("last_output", "").strip()
    if not text_to_save or text_to_save == "Здесь пока нет ответов.":
        await message.answer("⚠️ Нет данных для сохранения. Попробуй ещё раз. Если повторится, пиши - @mecau")
        return

    doc = Document()
    for section in doc.sections:
        section.top_margin = Mm(20)
        section.bottom_margin = Mm(20)
        section.left_margin = Mm(30)
        section.right_margin = Mm(15)

    p = doc.add_paragraph()
    run = p.add_run(text_to_save)
    run.font.name = 'Times New Roman'
    run.font.size = Pt(14)

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)

    file_doc = BufferedInputFile(bio.read(), filename="MecauAI_Answer.docx")
    await message.answer_document(file_doc, caption="📄 Вот твой ответ в формате Word!")

@dp.message(F.text == "📑 Титульник ГОСТ")
async def cmd_gost_title(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(message.from_user.id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📘 Индивидуальный проект", callback_data="gost_project")],
        [InlineKeyboardButton(text="📗 Курсовая работа",        callback_data="gost_coursework")],
        [InlineKeyboardButton(text="📕 Дипломная работа (ВКР)", callback_data="gost_diploma")],
        [InlineKeyboardButton(text="📙 Отчёт по практике",      callback_data="gost_practice")],
    ])
    await message.answer("📑 Выбери тип работы для титульника:", reply_markup=kb)

@dp.callback_query(F.data.startswith("gost_"))
async def cb_gost_generate(callback: types.CallbackQuery):
    work_type_map = {
        "gost_project":    ("ИНДИВИДУАЛЬНЫЙ ПРОЕКТ",                   "Индивидуальный_проект"),
        "gost_coursework": ("КУРСОВАЯ РАБОТА",                         "Курсовая_работа"),
        "gost_diploma":    ("ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА (ВКР)", "ВКР"),
        "gost_practice":   ("ОТЧЁТ ПО ПРАКТИКЕ",                      "Отчет_по_практике"),
    }
    work_label, filename_base = work_type_map[callback.data]

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Mm(20)
        section.bottom_margin = Mm(20)
        section.left_margin   = Mm(30)
        section.right_margin  = Mm(15)

    def add_centered(text, size=14, bold=False):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.name = 'Times New Roman'
        run.font.size = Pt(size)
        run.bold = bold
        return p

    def add_right(text, size=14):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = p.add_run(text)
        run.font.name = 'Times New Roman'
        run.font.size = Pt(size)
        return p

    def add_empty(count=1):
        for _ in range(count):
            p = doc.add_paragraph()
            p.add_run("").font.name = 'Times New Roman'

    add_centered("ФЕДЕРАЛЬНОЕ ГОСУДАРСТВЕННОЕ БЮДЖЕТНОЕ ПРОФЕССИОНАЛЬНОЕ\nОБРАЗОВАТЕЛЬНОЕ УЧРЕЖДЕНИЕ\n«НАЗВАНИЕ КОЛЛЕДЖА»", size=14)
    add_empty(4)
    add_centered(work_label, size=14, bold=True)
    add_empty(1)
    add_centered("на тему:\n«Введи тему работы здесь»", size=14, bold=True)
    add_empty(6)
    add_right("Выполнил(а): студент(ка) группы ГРУППА\nФамилия Имя Отчество\n\nРуководитель:\nДолжность, Фамилия И.О.")
    add_empty(5)
    add_centered("Город — 2026", size=14)

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)

    file_doc = BufferedInputFile(bio.read(), filename=f"Titulnik_{filename_base}.docx")
    await callback.message.answer_document(file_doc, caption=f"📑 Титульник «{work_label}» готов!")
    await callback.answer()

MENU_BUTTONS = {
    "📄 Ответ в Word", "📑 Титульник ГОСТ",
    "📈 Презентация", "📊 Excel-таблица", "📚 Анализ документа",
    "⭐ Избранное", "🔄 Режим ИИ", "ℹ️ О MecauAI",
    "📢 Рассылка", "📊 Статистика"
}

@dp.message(F.photo)
async def handle_photo(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    save_user_id(user_id)
    photo = message.photo[-1]

    # Скачивание фото вынесено в отдельный try/except: раньше ошибка сети или
    # слишком большого файла на этом шаге не перехватывалась и «ломала» обработку.
    try:
        if photo.file_size and photo.file_size > MAX_DOC_SIZE_BYTES:
            await message.answer("⚠️ Фото слишком большое (лимит Telegram на скачивание — 20 МБ). Пришли файл меньшего размера.")
            return
        file_info = await bot.get_file(photo.file_id)
        downloaded_file = await bot.download_file(file_info.file_path)
        img_bytes = downloaded_file.read()
    except Exception as e:
        logging.error(f"Не удалось скачать фото от пользователя {user_id}: {e}")
        await message.answer("⚠️ Не удалось загрузить фото. Проверь соединение и попробуй отправить ещё раз.")
        return

    # Если пользователь сейчас создает презентацию, сохраняем картинку для слайдов
    if user_id in ppt_states:
        if user_id not in user_ppt_images:
            user_ppt_images[user_id] = []
        user_ppt_images[user_id].append(img_bytes)
        await message.answer(f"🖼 Картинка сохранена для презентации ({len(user_ppt_images[user_id])} шт.)! Теперь отправь тему презентации.")
        return

    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    # Если пользователь в режиме анализа документа — считаем фото сканом страницы
    # и распознаём текст через vision-модель вместо обычного анализа изображения.
    if user_id in doc_analysis_states:
        busy_users.add(user_id)
        status_msg = await message.answer("📚 Распознаю текст на фото...")
        try:
            base64_image = base64.b64encode(img_bytes).decode('utf-8')
            response = await call_groq_with_retry(
                model=VISION_MODEL,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Распознай и выведи весь текст, присутствующий на этом изображении, без комментариев."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                    ]
                }]
            )
            recognized_text = clean_text_for_html(response.choices[0].message.content)
            await process_document_context(message, "фото документа", recognized_text, status_msg)
        except Exception as e:
            logging.error(f"Ошибка распознавания текста с фото: {e}", exc_info=True)
            try:
                await status_msg.edit_text("⚠️ Не удалось распознать текст на фото. Попробуй ещё раз.")
            except Exception:
                pass
            doc_analysis_states.discard(user_id)
        finally:
            busy_users.discard(user_id)
        return

    # Обычный режим анализа картинки через Vision ИИ
    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        base64_image = base64.b64encode(img_bytes).decode('utf-8')
        response = await call_groq_with_retry(
            model=VISION_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": message.caption or "Подробно проанализируй это изображение, распознай текст если он есть и ответь по сути."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }]
        )
        reply = clean_text_for_html(response.choices[0].message.content)
        user_data = get_user_data(user_id)
        user_data["last_output"] = reply
        await safe_answer(message, f"{reply}{AD_FOOTER}", parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
    except Exception as e:
        logging.error(f"Ошибка обработки фото: {e}")
        await message.answer("⚠️ Ошибка при обработке изображения. Попробуй ещё раз. Если повторится, пиши - @mecau")

@dp.message(F.text)
async def handle_text(message: Message):
    user_id = message.from_user.id

    if user_id == MY_ADMIN_ID and user_id in broadcast_states:
        broadcast_states.remove(user_id)
        users = load_user_ids()
        success, failed = 0, 0
        status_msg = await message.answer("📢 Начинаю рассылку...")
        for uid in users:
            try:
                await bot.send_message(uid, message.text, disable_web_page_preview=True)
                success += 1
            except TelegramRetryAfter as e:
                # Telegram просит подождать из-за лимита скорости — ждём и пробуем ещё раз
                await asyncio.sleep(e.retry_after)
                try:
                    await bot.send_message(uid, message.text, disable_web_page_preview=True)
                    success += 1
                except Exception:
                    failed += 1
            except TelegramForbiddenError:
                # Пользователь заблокировал бота — это не сбой, просто пропускаем
                failed += 1
            except Exception:
                failed += 1
            await asyncio.sleep(0.05)
        await status_msg.edit_text(f"✅ Рассылка завершена!\n\n👥 Доставлено: {success}\n❌ Ошибок: {failed}")
        return

    # ----------------- ГЕНЕРАЦИЯ ПРЕЗЕНТАЦИИ -----------------
    if user_id in ppt_states:
        ppt_states.remove(user_id)
        if user_id in busy_users:
            await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
            return
        busy_users.add(user_id)
        theme_key = user_ppt_design.pop(user_id, "ppt_blue")
        theme = PPT_THEMES.get(theme_key, PPT_THEMES["ppt_blue"])
        status_msg = await message.answer("📈 Генерирую презентацию и иллюстрации, подожди немного...")
        try:
            await bot.send_chat_action(chat_id=message.chat.id, action="typing")
            user_images = user_ppt_images.pop(user_id, [])
            num_slides = max(len(user_images), 5) if user_images else 5

            response = await call_groq_with_retry(
                messages=[
                    {
                        "role": "system",
                        "content": f"Составь презентацию ровно из {num_slides} слайдов. Ответ выдай СТРОГО в формате валидного JSON без markdown-оформления (без ```json), в виде списка объектов: [{{\"title\": \"Заголовок слайда\", \"points\": [\"Тезис 1\", \"Тезис 2\"], \"image_prompt\": \"Detailed professional visual illustration of the slide topic in English\"}}]. В поле image_prompt ВСЕГДА пиши качественный промпт на английском языке для генерации картинки."
                    },
                    {"role": "user", "content": f"Тема: {message.text}"}
                ],
                temperature=0.7
            )
            raw_content = clean_text_for_html(response.choices[0].message.content)
            slides_data = extract_json(raw_content)
            if not isinstance(slides_data, list) or not slides_data:
                raise ValueError("Модель вернула пустой или некорректный список слайдов")

            prs = Presentation()
            prs.slide_width = PptxInches(13.333)
            prs.slide_height = PptxInches(7.5)

            build_title_slide(prs, theme, message.text)

            async with aiohttp.ClientSession() as session:
                for idx, item in enumerate(slides_data):
                    img_stream = None
                    if user_images and idx < len(user_images):
                        img_stream = io.BytesIO(user_images[idx])
                    elif item.get("image_prompt"):
                        img_bytes = await generate_ai_image(session, item["image_prompt"], width=800, height=600)
                        if img_bytes:
                            img_stream = io.BytesIO(img_bytes)

                    build_content_slide(prs, theme, idx, item, img_stream)

            bio = io.BytesIO()
            prs.save(bio)
            bio.seek(0)
            file_doc = BufferedInputFile(bio.read(), filename="Presentation.pptx")

            await status_msg.delete()
            await message.answer_document(file_doc, caption=f"📈 Презентация готова! Дизайн: {theme['label']}")

        except Exception as e:
            logging.error(f"Критическая ошибка при генерации презентации: {e}", exc_info=True)
            try:
                await status_msg.delete()
            except Exception:
                pass
            await message.answer("⚠️ Произошла ошибка при создании презентации. Попробуй еще раз.")
        finally:
            busy_users.discard(user_id)
        return

    # ----------------- ГЕНЕРАЦИЯ EXCEL -----------------
    if user_id in excel_states:
        excel_states.remove(user_id)
        if user_id in busy_users:
            await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
            return
        busy_users.add(user_id)
        excel_mode = user_excel_mode.pop(user_id, "excel_plain")
        status_msg = await message.answer("📊 Генерирую таблицу Excel, подожди секунду...")
        try:
            await bot.send_chat_action(chat_id=message.chat.id, action="typing")
            response = await call_groq_with_retry(
                messages=[
                    {
                        "role": "system",
                        "content": "Создай структуру таблицы на основе запроса пользователя. Ответ выдай СТРОГО в формате валидного JSON без markdown-оформления (без ```json), в виде объекта с ключом 'headers' (список строк-заголовков) и ключом 'rows' (список списков с данными для строк таблицы, где числовые значения — это числа, а не строки). Пример: {\"headers\": [\"№\", \"Название\", \"Значение\"], \"rows\": [[1, \"Пример 1\", 100], [2, \"Пример 2\", 200]]}"
                    },
                    {"role": "user", "content": f"Задача для таблицы: {message.text}"}
                ],
                temperature=0.7
            )
            raw_content = clean_text_for_html(response.choices[0].message.content)
            table_data = extract_json(raw_content)

            headers = table_data.get("headers") or ["Колонка 1", "Колонка 2"]
            raw_rows = table_data.get("rows") or [["Данные 1", "Данные 2"]]

            # Нормализуем строки: выравниваем количество столбцов под заголовки,
            # чтобы "кривые" ответы от ИИ не ломали генерацию файла.
            num_cols = len(headers)
            norm_rows = []
            for row in raw_rows:
                row = list(row) if isinstance(row, (list, tuple)) else [row]
                if len(row) < num_cols:
                    row = row + [""] * (num_cols - len(row))
                elif len(row) > num_cols:
                    row = row[:num_cols]
                norm_rows.append(row)

            wb = Workbook()
            ws = wb.active
            ws.title = "Таблица"

            header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="4F81BD", end_color="4F81BD", fill_type="solid")
            align_center = Alignment(horizontal="center", vertical="center")
            align_left = Alignment(horizontal="left", vertical="center")
            align_right = Alignment(horizontal="right", vertical="center")
            border_thin = Border(
                left=Side(style='thin', color='D9D9D9'),
                right=Side(style='thin', color='D9D9D9'),
                top=Side(style='thin', color='D9D9D9'),
                bottom=Side(style='thin', color='D9D9D9')
            )

            ws.append(headers)
            for col_num in range(1, num_cols + 1):
                cell = ws.cell(row=1, column=col_num)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = align_center
                cell.border = border_thin

            def is_number(val):
                if isinstance(val, (int, float)):
                    return True
                try:
                    float(str(val).replace(",", "."))
                    return True
                except (TypeError, ValueError):
                    return False

            for row_idx, row_data in enumerate(norm_rows, start=2):
                ws.append(row_data)
                for col_num in range(1, num_cols + 1):
                    cell = ws.cell(row=row_idx, column=col_num)
                    cell.font = Font(name="Calibri", size=11)
                    cell.alignment = align_right if is_number(cell.value) else align_left
                    cell.border = border_thin

            for col in ws.columns:
                max_length = 0
                column_letter = col[0].column_letter
                for cell in col:
                    try:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    except Exception:
                        pass
                ws.column_dimensions[column_letter].width = max(max_length + 4, 12)

            # ---- Диаграмма по числовым данным (если выбран нужный режим) ----
            if excel_mode in ("excel_chart", "excel_full") and num_cols >= 2:
                try:
                    numeric_cols = []
                    for col_idx in range(1, num_cols):
                        total, numeric = 0, 0
                        for row in norm_rows:
                            val = row[col_idx]
                            if val not in ("", None):
                                total += 1
                                if is_number(val):
                                    numeric += 1
                        if total > 0 and numeric / total >= 0.6:
                            numeric_cols.append(col_idx)

                    if numeric_cols:
                        chart = BarChart()
                        chart.title = "Диаграмма по данным таблицы"
                        chart.style = 10
                        chart.x_axis.title = headers[0]
                        min_col = min(numeric_cols) + 1
                        max_col = max(numeric_cols) + 1
                        data = Reference(ws, min_col=min_col, max_col=max_col, min_row=1, max_row=len(norm_rows) + 1)
                        cats = Reference(ws, min_col=1, min_row=2, max_row=len(norm_rows) + 1)
                        chart.add_data(data, titles_from_data=True)
                        chart.set_categories(cats)
                        chart.width = 18
                        chart.height = 10
                        anchor_col_letter = get_column_letter(num_cols + 2)
                        ws.add_chart(chart, f"{anchor_col_letter}2")
                except Exception as chart_err:
                    logging.error(f"Не удалось построить диаграмму в Excel: {chart_err}")

            # ---- ИИ-картинка по теме (если выбран полный режим) ----
            if excel_mode == "excel_full":
                try:
                    async with aiohttp.ClientSession() as session:
                        img_bytes = await generate_ai_image(
                            session,
                            f"Professional business infographic illustration about: {message.text}",
                            width=768, height=512
                        )
                    if img_bytes:
                        from openpyxl.drawing.image import Image as XLImage
                        xl_img = XLImage(io.BytesIO(img_bytes))
                        xl_img.width = 480
                        xl_img.height = 320
                        img_row = len(norm_rows) + 4
                        ws.add_image(xl_img, f"A{img_row}")
                except Exception as img_err:
                    logging.error(f"Не удалось вставить ИИ-картинку в Excel (проверь, установлен ли Pillow): {img_err}")

            bio = io.BytesIO()
            wb.save(bio)
            bio.seek(0)
            file_doc = BufferedInputFile(bio.read(), filename="MecauAI_Table.xlsx")

            await status_msg.delete()
            await message.answer_document(file_doc, caption=f"📊 Твоя таблица Excel готова! ({EXCEL_MODES.get(excel_mode)})")
        except Exception as e:
            logging.error(f"Ошибка при создании Excel: {e}", exc_info=True)
            try:
                await status_msg.delete()
            except Exception:
                pass
            await message.answer("⚠️ Произошла ошибка при создании таблицы Excel. Попробуй еще раз.")
        finally:
            busy_users.discard(user_id)
        return

    if message.text in MENU_BUTTONS:
        return

    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    save_user_id(user_id)
    user_data = get_user_data(user_id)
    current_mode = user_data["mode"]

    system_prompt = PROMPTS.get(current_mode, PROMPTS["ai"])

    # Если ранее был загружен документ через «📚 Анализ документа» — добавляем
    # его содержимое в системный промпт, чтобы отвечать с учётом контекста.
    doc_context = user_doc_context.get(user_id)
    if doc_context:
        system_prompt = (
            f"{system_prompt}\n\n"
            f"Пользователь ранее загрузил документ. Используй его содержимое, если вопрос с ним связан:\n"
            f"---\n{doc_context}\n---"
        )

    user_data["history"].append({"role": "user", "content": message.text})
    if len(user_data["history"]) > 6:
        user_data["history"] = user_data["history"][-6:]

    messages_payload = [{"role": "system", "content": system_prompt}] + user_data["history"]

    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        response = await call_groq_with_retry(messages=messages_payload)
        ai_reply = clean_text_for_html(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply

        if current_mode == "coder":
            full_message = f"{ai_reply}\n\n{AD_FOOTER}"
            await safe_answer(message, full_message, parse_mode="Markdown", reply_markup=get_answer_inline_keyboard())
        else:
            full_message = f"{ai_reply}{AD_FOOTER}"
            await safe_answer(message, full_message, parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
    except Exception:
        await message.answer("⚠️ Ошибка при обработке запроса. Попробуй ещё раз. Если повторится, пиши - @mecau")

@dp.errors()
async def global_error_handler(event: types.ErrorEvent):
    """
    Глобальный перехватчик непредвиденных ошибок. Раньше необработанное
    исключение в любом хэндлере могло привести к тому, что пользователь
    просто не получал ответа и не понимал, что произошло. Теперь бот
    логирует ошибку и сообщает пользователю, что что-то пошло не так.
    """
    logging.error(f"Необработанная ошибка: {event.exception}", exc_info=True)
    try:
        update = event.update
        if update.message:
            await update.message.answer("⚠️ Произошла непредвиденная ошибка. Попробуй ещё раз позже. Если повторится — пиши @mecau")
        elif update.callback_query:
            await update.callback_query.answer("⚠️ Произошла ошибка, попробуй ещё раз.", show_alert=True)
    except Exception:
        pass
    return True

async def main():
    await bot.set_my_commands([
        BotCommand(command="start", description="Перезапустить бота"),
        BotCommand(command="mode",  description="Сменить режим ассистента"),
        BotCommand(command="about", description="О возможностях"),
    ])
    await bot.delete_webhook()
    print("🚀 Бот MecauAI запущен и слушает обновления...")
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
