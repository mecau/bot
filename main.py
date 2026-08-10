import asyncio
import logging
import base64
import io
import json
import os
import re
import time
import gc
from datetime import datetime, timezone
import urllib.parse
import aiohttp
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import CommandStart, Command
from aiogram.exceptions import TelegramRetryAfter, TelegramForbiddenError, TelegramBadRequest
from aiogram.types import (
    InlineKeyboardMarkup, InlineKeyboardButton,
    ReplyKeyboardMarkup, KeyboardButton, Message, BotCommand, BufferedInputFile
)
from openai import AsyncOpenAI
from docx import Document
from docx.shared import Pt, Inches, Mm
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from openpyxl.worksheet.table import Table, TableStyleInfo
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.utils import get_column_letter

# Для вставки картинок в Excel дополнительно нужен пакет Pillow:
#   pip install pillow
# Для чтения PDF в разделе «Анализ документа» нужен пакет pypdf:
#   pip install pypdf
# Оба импортируются лениво (внутри try/except), чтобы их отсутствие
# не ломало весь бот, а просто отключало соответствующую функцию.

from config import (
    BOT_TOKEN, GROQ_API_KEY,
    CHANNEL_1_USERNAME, CHANNEL_1_URL,
    CHANNEL_2_USERNAME, CHANNEL_2_URL,
    TEXT_MODEL, VISION_MODEL, PROMPTS, AD_FOOTER
)

MY_ADMIN_ID = int(os.getenv("MECAUAI_ADMIN_ID", "1184589026"))

logging.basicConfig(level=logging.INFO)
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher()

groq_client = AsyncOpenAI(
    api_key=GROQ_API_KEY,
    base_url="https://api.groq.com/openai/v1"
)

USERS_FILE = "users.json"
FAV_FILE = "favorites.json"
MODES_FILE = "user_modes.json"

all_users_cache = set()

def load_user_ids() -> set:
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list):
                    return set(data)
        except Exception:
            pass
    return set()

def save_user_id(user_id: int):
    global all_users_cache
    if user_id in all_users_cache:
        return
    all_users_cache.add(user_id)
    try:
        tmp = USERS_FILE + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(sorted(all_users_cache), f, ensure_ascii=False)
        os.replace(tmp, USERS_FILE)
    except Exception as e:
        logging.error(f"Не удалось сохранить пользователя: {e}")

all_users_cache = load_user_ids()

def load_favorites() -> dict:
    if os.path.exists(FAV_FILE):
        try:
            with open(FAV_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def save_favorites(favs: dict):
    with open(FAV_FILE, "w", encoding="utf-8") as f:
        json.dump(favs, f, ensure_ascii=False, indent=2)

def add_to_favorites(user_id: int, text: str):
    favs = load_favorites()
    uid_str = str(user_id)
    if uid_str not in favs:
        favs[uid_str] = []
    if text not in favs[uid_str]:
        favs[uid_str].append(text)
        save_favorites(favs)

def load_user_modes() -> dict:
    if os.path.exists(MODES_FILE):
        try:
            with open(MODES_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, dict):
                    return data
        except Exception:
            pass
    return {}

def save_user_mode(user_id: int, mode: str):
    """Сохраняем выбранный режим на диск, чтобы он не сбрасывался при перезапуске бота."""
    modes = load_user_modes()
    modes[str(user_id)] = mode
    try:
        with open(MODES_FILE, "w", encoding="utf-8") as f:
            json.dump(modes, f, ensure_ascii=False)
    except Exception as e:
        logging.error(f"Не удалось сохранить режим пользователя: {e}")

users_db = {}
saved_modes_cache = load_user_modes()
broadcast_states = set()
ppt_states = set()
excel_states = set()
doc_analysis_states = set()
busy_users = set()
user_ppt_images = {}
user_ppt_design = {}
user_excel_mode = {}
user_ppt_topic = {}
user_excel_topic = {}
# Lightweight per-user workflow state. Cleared after a completed/cancelled flow.
user_task_state = {}
user_last_request = {}
user_custom_images = {}  # temporary presentation assets: user_id -> list Telegram file_id
ppt_creation_users = set()

# Ограничители нагрузки: не даём слабому VPS одновременно запускать несколько
# тяжёлых генераций и десятки обращений к AI.
AI_SEMAPHORE = None
HEAVY_JOB_SEMAPHORE = None
MAX_AI_CONCURRENCY = 2
MAX_HEAVY_CONCURRENCY = 1
MAX_PPT_IMAGES = 8
MAX_PPT_AI_IMAGES = 4

# ===== MecauAI 2.0: единый registry артефактов + лёгкий persistence =====
ARTIFACT_DIR = os.getenv("MECAUAI_ARTIFACT_DIR", "generated_files")
ARTIFACTS_FILE = os.getenv("MECAUAI_ARTIFACTS_FILE", "artifacts.json")
ARTIFACT_TTL_SECONDS = int(os.getenv("MECAUAI_ARTIFACT_TTL_SECONDS", str(24 * 3600)))
MAX_ARTIFACTS_PER_USER = 6
MAX_ARTIFACT_BYTES = 25 * 1024 * 1024
RESOURCE_CHECK_INTERVAL = 60
RAM_SOFT_LIMIT_MB = int(os.getenv("MECAUAI_RAM_SOFT_LIMIT_MB", "700"))
RAM_HARD_LIMIT_MB = int(os.getenv("MECAUAI_RAM_HARD_LIMIT_MB", "900"))
artifact_registry = {}
artifact_lock = None

def _artifact_lock():
    global artifact_lock
    if artifact_lock is None:
        artifact_lock = asyncio.Lock()
    return artifact_lock

def _load_artifacts():
    global artifact_registry
    try:
        if os.path.exists(ARTIFACTS_FILE):
            with open(ARTIFACTS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            artifact_registry = data if isinstance(data, dict) else {}
    except Exception as e:
        logging.warning("Не удалось загрузить registry артефактов: %s", e)
        artifact_registry = {}

def _save_artifacts():
    try:
        tmp = ARTIFACTS_FILE + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(artifact_registry, f, ensure_ascii=False, indent=2)
        os.replace(tmp, ARTIFACTS_FILE)
    except Exception as e:
        logging.warning("Не удалось сохранить registry артефактов: %s", e)

_load_artifacts()
os.makedirs(ARTIFACT_DIR, exist_ok=True)

def _artifact_path(user_id, ext, stem):
    safe = re.sub(r"[^A-Za-zА-Яа-я0-9_-]+", "_", str(stem))[:45].strip("_") or "file"
    uid_dir = os.path.join(ARTIFACT_DIR, str(user_id))
    os.makedirs(uid_dir, exist_ok=True)
    return os.path.join(uid_dir, f"{int(time.time())}_{safe}.{ext}")

def register_artifact(user_id, kind, path, meta=None):
    if not path or not os.path.exists(path):
        return None
    entry = {"kind": kind, "path": os.path.abspath(path), "created": time.time(), "meta": meta or {}}
    items = artifact_registry.setdefault(str(user_id), [])
    items.append(entry)
    artifact_registry[str(user_id)] = items[-MAX_ARTIFACTS_PER_USER:]
    _save_artifacts()
    return entry

def last_artifact(user_id, kind=None):
    items = artifact_registry.get(str(user_id), [])
    for item in reversed(items):
        if kind is None or item.get("kind") == kind:
            if os.path.exists(item.get("path", "")):
                return item
    return None

def _memory_mb():
    try:
        import resource
        # Linux reports KiB for ru_maxrss.
        return resource.getrusage(resource.RUSAGE_SELF).ru_maxrss / 1024
    except Exception:
        return 0

def cleanup_artifacts():
    now = time.time()
    changed = False
    for uid, items in list(artifact_registry.items()):
        kept = []
        for item in items:
            path = item.get("path", "")
            age = now - float(item.get("created", now))
            try:
                size = os.path.getsize(path) if os.path.exists(path) else 0
                if path and os.path.exists(path) and (age > ARTIFACT_TTL_SECONDS or size > MAX_ARTIFACT_BYTES):
                    os.remove(path); changed = True
                    continue
            except Exception:
                pass
            if os.path.exists(path):
                kept.append(item)
        if kept:
            artifact_registry[uid] = kept[-MAX_ARTIFACTS_PER_USER:]
        else:
            artifact_registry.pop(uid, None)
    # Empty user directories are cheap to remove.
    try:
        for name in os.listdir(ARTIFACT_DIR):
            full = os.path.join(ARTIFACT_DIR, name)
            if os.path.isdir(full) and not os.listdir(full):
                os.rmdir(full)
    except Exception:
        pass
    if changed:
        _save_artifacts()

async def maintenance_loop():
    while True:
        try:
            cleanup_artifacts()
            ram = _memory_mb()
            if ram >= RAM_HARD_LIMIT_MB:
                logging.error("RAM watchdog: %s MB — принудительная очистка временных данных", round(ram))
                gc.collect()
                await asyncio.sleep(2)
            elif ram >= RAM_SOFT_LIMIT_MB:
                logging.warning("RAM watchdog: %s MB — мягкая очистка", round(ram))
                gc.collect()
        except asyncio.CancelledError:
            raise
        except Exception as e:
            logging.warning("Maintenance loop: %s", e)
        await asyncio.sleep(RESOURCE_CHECK_INTERVAL)

def _file_safe_size(path):
    try:
        return os.path.getsize(path) <= MAX_ARTIFACT_BYTES
    except Exception:
        return False


def set_task_state(user_id: int, **kwargs):
    state = user_task_state.setdefault(user_id, {})
    state.update(kwargs)
    return state

def get_task_state(user_id: int):
    return user_task_state.get(user_id, {})

def clear_task_state(user_id: int):
    user_task_state.pop(user_id, None)
    user_custom_images.pop(user_id, None)
    user_ppt_images.pop(user_id, None)
    ppt_creation_users.discard(user_id)

def ppt_assets(user_id: int):
    return user_custom_images.get(user_id, [])

def save_presentation_image(user_id: int, file_id: str):
    if user_id not in ppt_creation_users:
        return False
    images = user_custom_images.setdefault(user_id, [])
    if file_id not in images:
        images.append(file_id)
    user_custom_images[user_id] = images[-20:]
    return True

def add_custom_image(user_id: int, file_id: str):
    images = user_custom_images.setdefault(user_id, [])
    if file_id not in images:
        images.append(file_id)
    user_custom_images[user_id] = images[-20:]

user_doc_context = {}

# Лёгкая аналитика поведения — без хранения содержимого сообщений.
user_stats = {}
user_rate_buckets = {}
SUBSCRIPTION_CACHE_TTL = 300
subscription_cache = {}

def allow_request(user_id: int, limit: int = 15, window: int = 30) -> bool:
    """Мягкий антиспам: защищает VPS и API, не мешая обычному использованию."""
    now = time.monotonic()
    bucket = [t for t in user_rate_buckets.get(user_id, []) if now - t < window]
    if len(bucket) >= limit:
        user_rate_buckets[user_id] = bucket
        return False
    bucket.append(now)
    user_rate_buckets[user_id] = bucket
    if len(user_rate_buckets) > 5000:
        # Не даём служебному словарю расти бесконечно. Удаляем самые старые записи.
        oldest = sorted(user_rate_buckets.items(), key=lambda kv: kv[1][-1] if kv[1] else 0)[:1000]
        for uid, _ in oldest:
            user_rate_buckets.pop(uid, None)
    return True

def get_user_stats(user_id: int):
    if user_id not in user_stats:
        user_stats[user_id] = {
            "messages": 0,
            "voice": 0,
            "files": 0,
            "images": 0,
            "exports": 0,
            "last_action": None,
        }
    return user_stats[user_id]

def remember_action(user_id: int, action: str):
    stats = get_user_stats(user_id)
    stats["last_action"] = action
    stats["messages"] += 1

def get_help_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="🎙 Голос", callback_data="help_voice"),
            InlineKeyboardButton(text="📎 Файл", callback_data="help_file"),
        ],
        [
            InlineKeyboardButton(text="📈 Презентация", callback_data="help_ppt"),
            InlineKeyboardButton(text="📊 Excel", callback_data="help_excel"),
        ],
        [
            InlineKeyboardButton(text="🧠 Режим ИИ", callback_data="help_mode"),
        ]
    ])

def get_settings_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🧠 Режим ИИ", callback_data="settings_mode")],
        [InlineKeyboardButton(text="🧹 Очистить контекст", callback_data="settings_clear")],
    ])

def get_user_data(user_id: int):
    if user_id not in users_db:
        users_db[user_id] = {
            "mode": saved_modes_cache.get(str(user_id), "ai"),
            "history": [],
            "last_output": "Здесь пока нет ответов."
        }
    return users_db[user_id]

def clear_pending_states(user_id: int):
    """
    Сбрасывает все 'ожидающие ввода' состояния пользователя (создание презентации,
    таблицы, анализ документа). Нужно, чтобы если пользователь передумал и нажал
    другую кнопку меню, бот не принял её текст за тему презентации/таблицы —
    раньше это приводило к путанице и «зависшим» состояниям.
    """
    ppt_states.discard(user_id)
    excel_states.discard(user_id)
    doc_analysis_states.discard(user_id)
    user_ppt_images.pop(user_id, None)
    user_ppt_design.pop(user_id, None)
    user_excel_mode.pop(user_id, None)
    user_ppt_topic.pop(user_id, None)
    user_excel_topic.pop(user_id, None)

def _get_semaphores():
    global AI_SEMAPHORE, HEAVY_JOB_SEMAPHORE
    if AI_SEMAPHORE is None:
        AI_SEMAPHORE = asyncio.Semaphore(MAX_AI_CONCURRENCY)
    if HEAVY_JOB_SEMAPHORE is None:
        HEAVY_JOB_SEMAPHORE = asyncio.Semaphore(MAX_HEAVY_CONCURRENCY)
    return AI_SEMAPHORE, HEAVY_JOB_SEMAPHORE


async def acquire_heavy_job(status_msg=None):
    _, heavy = _get_semaphores()
    if heavy.locked() and status_msg is not None:
        try:
            await status_msg.edit_text("⏳ Сервер занят другой тяжёлой задачей. Я поставил твою задачу в очередь — бот не перегружается.")
        except Exception:
            pass
    await heavy.acquire()
    return heavy


async def call_groq_with_retry(messages, model: str = None, temperature: float = 0.7, max_retries: int = 2, timeout: int = 45):
    """
    Обёртка над вызовом Groq API с повторными попытками и таймаутом.
    Раньше любой сетевой сбой или подвисание запроса приводило к ошибке
    «Попробуй ещё раз» — теперь бот сам делает до 2 повторных попыток
    с небольшой паузой, прежде чем сдаться.
    """
    used_model = model or TEXT_MODEL
    last_err = None
    for attempt in range(max_retries + 1):
        try:
            ai_sem, _ = _get_semaphores()
            async with ai_sem:
                return await asyncio.wait_for(
                    groq_client.chat.completions.create(model=used_model, messages=messages, temperature=temperature),
                    timeout=timeout
                )
        except Exception as e:
            last_err = e
            logging.warning(f"Groq API — попытка {attempt + 1} не удалась: {e}")
            if attempt < max_retries:
                await asyncio.sleep(1.5 * (attempt + 1))
    raise last_err


async def transcribe_voice(file_bytes: bytes, filename: str = "voice.ogg") -> str:
    """Надёжное распознавание Telegram voice через Groq Whisper."""
    last_err = None
    for attempt in range(3):
        try:
            ai_sem, _ = _get_semaphores()
            async with ai_sem:
                result = await asyncio.wait_for(
                    groq_client.audio.transcriptions.create(
                        model="whisper-large-v3-turbo",
                        file=(filename, file_bytes),
                        response_format="text",
                        language="ru",
                        temperature=0
                    ),
                    timeout=60
                )
            transcript = result if isinstance(result, str) else getattr(result, "text", "")
            transcript = re.sub(r"\s+", " ", (transcript or "")).strip()
            if transcript:
                return transcript
            raise ValueError("Пустая расшифровка")
        except Exception as e:
            last_err = e
            logging.warning(f"Whisper attempt {attempt + 1}/3 failed: {e}")
            if attempt < 2:
                await asyncio.sleep(1.5 * (attempt + 1))
    raise last_err

async def check_subscription(user_id: int) -> bool:
    if user_id == MY_ADMIN_ID:
        return True
    now = time.monotonic()
    cached = subscription_cache.get(user_id)
    if cached and now - cached[1] < SUBSCRIPTION_CACHE_TTL:
        return cached[0]
    channels = [CHANNEL_1_USERNAME, CHANNEL_2_USERNAME]
    for ch in channels:
        try:
            member = await bot.get_chat_member(chat_id=f"@{ch}", user_id=user_id)
            if member.status not in ["creator", "administrator", "member"]:
                subscription_cache[user_id] = (False, now)
                return False
        except Exception as e:
            logging.error(f"Ошибка проверки подписки на канал @{ch}: {e}")
            # Если Telegram временно недоступен, используем свежий старый результат;
            # иначе безопасно запрещаем доступ.
            if cached and now - cached[1] < 900:
                return cached[0]
            return False
    subscription_cache[user_id] = (True, now)
    return True

def get_sub_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📢 Подписаться на канал 1", url=CHANNEL_1_URL)],
        [InlineKeyboardButton(text="📢 Подписаться на канал 2", url=CHANNEL_2_URL)],
        [InlineKeyboardButton(text="✅ Проверить подписку", callback_data="check_sub")]
    ])

def get_main_keyboard():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="🤖 Помощник"), KeyboardButton(text="🛠 Создать")],
            [KeyboardButton(text="📁 Документы"), KeyboardButton(text="⭐ Избранное")],
            [KeyboardButton(text="⚙️ Настройки")],
        ],
        resize_keyboard=True,
        is_persistent=True
    )

def get_admin_keyboard():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="🤖 Помощник"), KeyboardButton(text="📁 Документы")],
            [KeyboardButton(text="🛠 Создать"), KeyboardButton(text="⭐ Избранное")],
            [KeyboardButton(text="⚙️ Настройки")],
            [KeyboardButton(text="📊 Статистика"), KeyboardButton(text="📢 Рассылка")],
        ],
        resize_keyboard=True,
        is_persistent=True
    )

def keyboard_for(user_id: int):
    return get_admin_keyboard() if user_id == MY_ADMIN_ID else get_main_keyboard()

def get_capabilities_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎙 Голос", callback_data="guide_voice"),
         InlineKeyboardButton(text="🖼 Фото", callback_data="guide_image")],
        [InlineKeyboardButton(text="📁 Документы", callback_data="guide_doc"),
         InlineKeyboardButton(text="📈 Презентации", callback_data="guide_ppt")],
        [InlineKeyboardButton(text="📊 Excel", callback_data="guide_excel"),
         InlineKeyboardButton(text="📄 Word", callback_data="guide_word")],
    ])

def get_cancel_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="↩️ Начать заново", callback_data="flow_restart")]
    ])

def get_quick_actions_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎙 Голос", callback_data="guide_voice"),
         InlineKeyboardButton(text="🖼 Фото", callback_data="guide_image")],
        [InlineKeyboardButton(text="📁 Документ", callback_data="guide_doc"),
         InlineKeyboardButton(text="📈 Презентация", callback_data="guide_ppt")],
        [InlineKeyboardButton(text="📊 Excel", callback_data="guide_excel"),
         InlineKeyboardButton(text="📄 Word", callback_data="guide_word")],
    ])

def get_create_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📈 Презентация", callback_data="create_ppt"),
         InlineKeyboardButton(text="📊 Excel", callback_data="create_excel")],
        [InlineKeyboardButton(text="📄 Word", callback_data="create_word"),
         InlineKeyboardButton(text="📑 Титульный лист", callback_data="create_title")],
    ])

def get_helper_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎙 Голос", callback_data="guide_voice"),
         InlineKeyboardButton(text="🖼 Картинка", callback_data="guide_image")],
        [InlineKeyboardButton(text="📁 Документ", callback_data="guide_doc"),
         InlineKeyboardButton(text="📈 Презентация", callback_data="guide_ppt")],
        [InlineKeyboardButton(text="📊 Excel", callback_data="guide_excel"),
         InlineKeyboardButton(text="📄 Word", callback_data="guide_word")],
    ])

def get_answer_inline_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="💡 Проще", callback_data="btn_simplify"),
            InlineKeyboardButton(text="⭐ Сохранить", callback_data="btn_save_fav"),
        ],
        [
            InlineKeyboardButton(text="🔁 Доработать", callback_data="btn_continue"),
            InlineKeyboardButton(text="📄 В Word", callback_data="btn_word"),
        ],
    ])

def clean_text_for_html(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
    return text.strip()

async def safe_answer(message: Message, text: str, parse_mode: str = None, reply_markup=None, disable_web_page_preview: bool = True):
    """
    Безопасная отправка сообщения: если Telegram не может распарсить
    HTML/Markdown в ответе ИИ (например, из-за случайных символов < > _ *),
    бот отправит тот же текст обычным сообщением вместо падения с ошибкой.
    """
    try:
        return await message.answer(
            text, parse_mode=parse_mode, reply_markup=reply_markup,
            disable_web_page_preview=disable_web_page_preview
        )
    except Exception as e:
        logging.warning(f"Не удалось отправить с parse_mode={parse_mode}, отправляю без форматирования: {e}")
        plain_text = re.sub(r'<[^>]+>', '', text) if parse_mode == "HTML" else text
        try:
            return await message.answer(
                plain_text, reply_markup=reply_markup,
                disable_web_page_preview=disable_web_page_preview
            )
        except Exception as e2:
            logging.error(f"Повторная ошибка при отправке сообщения: {e2}")
            return None

def extract_json(raw: str):
    """
    Надёжное извлечение JSON из ответа ИИ. Модель иногда добавляет
    markdown-обёртку (```json ... ```) или лишние пояснения до/после JSON —
    эта функция вытаскивает валидный JSON в любом из этих случаев.
    """
    if not raw:
        raise ValueError("Пустой ответ от модели")
    raw = raw.strip()

    if "```" in raw:
        parts = raw.split("```")
        for part in parts:
            candidate = part.strip()
            if candidate.lower().startswith("json"):
                candidate = candidate[4:].strip()
            if candidate.startswith("{") or candidate.startswith("["):
                raw = candidate
                break

    start_obj = raw.find("{")
    start_arr = raw.find("[")

    if start_obj == -1 and start_arr == -1:
        raise ValueError("В ответе модели не найден JSON")

    if start_arr != -1 and (start_obj == -1 or start_arr < start_obj):
        start = start_arr
        end = raw.rfind("]")
    else:
        start = start_obj
        end = raw.rfind("}")

    if end == -1 or end < start:
        raise ValueError("Не удалось найти конец JSON в ответе модели")

    return json.loads(raw[start:end + 1])

async def generate_ai_image(session: aiohttp.ClientSession, prompt: str, width: int = 800, height: int = 600):
    """Генерация иллюстрации через Pollinations.ai. Возвращает байты картинки или None."""
    try:
        encoded_prompt = urllib.parse.quote(prompt)
        img_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width={width}&height={height}&nologo=true&model=flux"
        async with session.get(img_url, timeout=aiohttp.ClientTimeout(total=20)) as resp:
            if resp.status == 200:
                img_bytes = await resp.read()
                if len(img_bytes) > 1000:
                    return img_bytes
    except Exception as e:
        logging.error(f"Не удалось сгенерировать ИИ-картинку: {e}")
    return None

@dp.message(CommandStart())
async def cmd_start(message: Message):
    user_id = message.from_user.id
    save_user_id(user_id)

    if not await check_subscription(user_id):
        await message.answer(
            f"🔒 Доступ заблокирован!\n\n"
            f"Чтобы пользоваться MecauAI, подпишись на оба наших канала:\n"
            f" 👉 {CHANNEL_1_URL}\n"
            f" 👉 {CHANNEL_2_URL}\n\n"
            f"После подписки нажми кнопку ниже 👇",
            reply_markup=get_sub_keyboard()
        )
        return

    start_text = (
        f"Привет, {message.from_user.first_name}! Ты активировал MecauAI 🚀\n\n"
        "Я твой карманный помощник для учебы и разработки. Можешь отправлять вопросы, задачи, код, изображения, документы и голосовые сообщения — я распознаю запрос, разберусь в нём и дам готовый ответ. Также могу подготовить Word, презентацию, Excel или помочь разобрать документ. Если задача большая, просто опиши её своими словами — я сам выберу подходящий способ обработки."
    )
    await message.answer(start_text, reply_markup=keyboard_for(user_id))

@dp.callback_query(F.data == "check_sub")
async def cb_check_sub(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    if await check_subscription(user_id):
        save_user_id(user_id)
        try:
            await callback.message.delete()
        except Exception:
            pass
        await callback.message.answer(
            "🎉 Подписка на оба канала подтверждена! Добро пожаловать в MecauAI 🚀",
            reply_markup=keyboard_for(user_id)
        )
    else:
        await callback.answer("❌ Ты подписался еще не на все каналы!", show_alert=True)


@dp.message(F.text == "✨ Возможности")
async def cmd_capabilities(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    clear_pending_states(message.from_user.id)
    await message.answer(
        "✨ <b>MecauAI умеет больше, чем кажется</b>\n\n"
        "💬 <b>Общение</b> — вопросы, объяснения, тексты, код, идеи.\n"
        "🎙 <b>Голос</b> — отправь голосовое, я сам превращу его в задачу.\n"
        "🖼 <b>Изображения</b> — распознаю текст, задачи, схемы и содержимое фото.\n"
        "📚 <b>Документы</b> — PDF, DOCX, TXT и фото страниц с вопросами по содержимому.\n"
        "📄 <b>Экспорт</b> — ответ можно сохранить в Word.\n"
        "📈 <b>Презентации</b> — попроси сделать презентацию на нужную тему.\n"
        "📊 <b>Excel</b> — попроси создать таблицу с данными и диаграммой.\n\n"
        "💡 Главное: <b>не нужно искать правильную кнопку</b>. Напиши или скажи задачу своими словами.",
        parse_mode="HTML",
        reply_markup=get_help_keyboard()
    )

@dp.message(F.text == "📂 Файлы")
async def cmd_files_help(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    clear_pending_states(message.from_user.id)
    await message.answer(
        "📂 <b>Работа с файлами</b>\n\n"
        "Просто отправь файл в чат.\n\n"
        "📚 Для PDF/DOCX/TXT я могу прочитать содержание и отвечать на вопросы по нему.\n"
        "🖼 Фото страницы тоже можно отправить — я попробую распознать текст.\n\n"
        "Если хочешь создать файл, просто напиши, например:\n"
        "«Сделай Word из последнего ответа»\n"
        "«Создай презентацию про искусственный интеллект»\n"
        "«Сделай Excel с оценками студентов»",
        parse_mode="HTML"
    )

@dp.message(Command("help"))
async def cmd_help(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🤖 <b>MecauAI умеет работать не только с текстом.</b>\n\n"
        "🎙 Голос — говори задачу обычными словами.\n"
        "🖼 Фото — отправляй задачи, скриншоты и изображения.\n"
        "📁 Документы — загружай PDF/DOCX/TXT и задавай вопросы.\n"
        "📈 Презентации — создание с выбором содержания и оформления.\n"
        "📊 Excel — таблицы, формулы, расчёты и визуализация.\n"
        "📄 Word — оформление готового материала в документ.\n\n"
        "💡 Самый простой способ: просто напиши, что тебе нужно.",
        parse_mode="HTML",
        reply_markup=get_quick_actions_keyboard()
    )

@dp.message(F.text == "🤖 Помощник")
async def cmd_helper_menu(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🤖 <b>Что умеет MecauAI?</b>\n\n"
        "🎙 <b>Голос</b> — отправь голосовое, я распознаю речь и отвечу.\n"
        "🖼 <b>Картинка</b> — отправь фото, я разберу его.\n"
        "📁 <b>Документ</b> — отправь PDF/DOCX/TXT, я изучу его.\n"
        "📈 <b>Презентация</b> — я проведу тебя по оформлению шаг за шагом.\n"
        "📊 <b>Excel</b> — создам таблицу и диаграммы.\n"
        "📄 <b>Word</b> — превращу готовый ответ в документ.\n\n"
        "💡 Не нужно знать команды: выбирай кнопку или просто отправляй задачу.\n\nЕсли не знаешь, что выбрать — просто напиши или скажи, что хочешь получить. Я сам подберу нужный сценарий.",
        parse_mode="HTML",
        reply_markup=get_quick_actions_keyboard()
    )

@dp.message(F.text == "📁 Документы")
async def cmd_documents_menu(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "📁 <b>Работа с документами</b>\n\n"
        "Просто отправь PDF, DOCX или TXT — дальше можно:\n"
        "• сделать конспект;\n"
        "• найти нужную информацию;\n"
        "• задать вопросы по документу;\n"
        "• подготовить материал к экзамену;\n"
        "• сравнить документы;\n"
        "• сделать основу для презентации.\n\n"
        "💡 Можно сразу написать запрос вместе с файлом.",
        parse_mode="HTML",
        reply_markup=get_cancel_keyboard()
    )

@dp.message(F.text == "🛠 Создать")
async def cmd_create_menu(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🛠 <b>Что создать?</b>\n\n"
        "Выбери готовый вариант. Если нужного нет — просто опиши задачу обычными словами.\n\n"
        "📈 В презентации я отдельно помогу с темой, структурой, стилем, титульником и твоими картинками.",
        parse_mode="HTML",
        reply_markup=get_create_keyboard()
    )

@dp.message(F.text == "⚙️ Настройки")
async def cmd_settings(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    clear_pending_states(message.from_user.id)
    user_data = get_user_data(message.from_user.id)
    names = {
        "coder": "💻 ИИ-Программист",
        "ai": "🧠 Академический ассистент",
        "friend": "🫂 Лучший друг"
    }
    await message.answer(
        f"⚙️ <b>Настройки</b>\n\n"
        f"Текущий режим: <b>{names.get(user_data['mode'], 'Стандартный')}</b>\n\n"
        "Можно сменить стиль общения или очистить контекст текущего диалога.",
        parse_mode="HTML",
        reply_markup=get_settings_keyboard()
    )

@dp.callback_query(F.data.in_({"help_voice", "help_file", "help_ppt", "help_excel", "help_mode",
                               "settings_mode", "settings_clear"}))
async def cb_help_and_settings(callback: types.CallbackQuery):
    user_id = callback.from_user.id

    if callback.data == "help_voice":
        await callback.message.answer("🎙 Просто отправь мне голосовое сообщение. Ничего включать не нужно.")
    elif callback.data == "help_file":
        await callback.message.answer("📎 Просто отправь PDF, DOCX, TXT или фото страницы — я предложу, что можно сделать.")
    elif callback.data == "help_ppt":
        await callback.message.answer("📈 Напиши: «Сделай презентацию на тему ...» — бот сам запустит создание презентации.")
    elif callback.data == "help_excel":
        await callback.message.answer("📊 Напиши: «Сделай Excel-таблицу ...» — бот предложит нужный формат и создаст файл.")
    elif callback.data == "help_mode":
        await cmd_mode(callback.message)
    elif callback.data == "settings_mode":
        await cmd_mode(callback.message)
    elif callback.data == "settings_clear":
        user_data = get_user_data(user_id)
        user_data["history"] = []
        user_data["last_output"] = "Здесь пока нет ответов."
        user_doc_context.pop(user_id, None)
        await callback.message.answer("🧹 Контекст очищен. Начинаем с чистого листа.")

    await callback.answer()

@dp.message(F.text == "ℹ️ О MecauAI")
@dp.message(Command("about"))
async def cmd_about(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    about_text = (
            "🧠 <b>Возможности MecauAI</b>\n\n"
            "MecauAI задуман как помощник, которому не нужно объяснять, какую кнопку нажать. "
            "Просто отправь задачу — текстом, голосом, фото или файлом.\n\n"
            "<b>Что уже работает:</b>\n"
            "• 🎙 голосовые сообщения → расшифровка → ответ ИИ;\n"
            "• 🖼 анализ изображений и распознавание текста;\n"
            "• 📚 анализ PDF/DOCX/TXT и вопросы по документу;\n"
            "• 📄 экспорт ответа в Word;\n"
            "• 📈 презентации с дизайном и иллюстрациями;\n"
            "• 📊 Excel-таблицы с диаграммами;\n"
            "• ⭐ избранное;\n"
            "• 🧠 три режима общения.\n\n"
            "<b>Главное улучшение интерфейса:</b> функций много, но кнопок мало. "
            "Сложные действия запускаются обычной фразой — так бот ощущается как помощник, а не как меню."
        )
    await message.answer(about_text, parse_mode="HTML", reply_markup=keyboard_for(message.from_user.id))

@dp.message(F.text == "🔄 Режим ИИ")
@dp.message(Command("mode"))
async def cmd_mode(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(message.from_user.id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💻 ИИ-Программист",       callback_data="set_mode_coder")],
        [InlineKeyboardButton(text="🧠 Академический ассистент", callback_data="set_mode_ai")],
        [InlineKeyboardButton(text="🫂 Лучший друг",             callback_data="set_mode_friend")]
    ])
    await message.answer("Выбери режим работы бота:", reply_markup=kb)

@dp.callback_query(F.data.startswith("set_mode_"))
async def cb_set_mode(callback: types.CallbackQuery):
    mode = callback.data.split("_")[-1]
    user_id = callback.from_user.id
    user_data = get_user_data(user_id)
    user_data["mode"] = mode
    save_user_mode(user_id, mode)
    clear_pending_states(user_id)

    names = {
        "coder": "ИИ-Программист 💻",
        "ai": "Академический ассистент 🧠",
        "friend": "Лучший друг 🫂"
    }
    await callback.message.edit_text(f"Режим переключен на: {names.get(mode, 'Стандартный')}")
    await callback.answer()

@dp.message(F.text == "📊 Статистика")
async def cmd_stats(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    users = load_user_ids()
    total_messages = sum(v.get("messages", 0) for v in user_stats.values())
    total_voice = sum(v.get("voice", 0) for v in user_stats.values())
    total_files = sum(v.get("files", 0) for v in user_stats.values())
    total_images = sum(v.get("images", 0) for v in user_stats.values())
    total_exports = sum(v.get("exports", 0) for v in user_stats.values())
    await message.answer(
        f"📊 <b>Статистика бота</b>\n\n"
        f"👥 Всего пользователей: {len(users)}\n"
        f"💬 Сообщений: {total_messages}\n"
        f"🎙 Голосовых: {total_voice}\n"
        f"📎 Файлов: {total_files}\n"
        f"🖼 Изображений: {total_images}\n"
        f"📄 Экспортов: {total_exports}\n"
        f"⏳ Сейчас заняты генерацией: {len(busy_users)}\n"
        f"🧠 AI одновременно: до {MAX_AI_CONCURRENCY}\n"
        f"🛡 Тяжёлых задач одновременно: до {MAX_HEAVY_CONCURRENCY}",
        parse_mode="HTML"
    )

@dp.message(F.text == "📢 Рассылка")
async def cmd_broadcast_prompt(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    clear_pending_states(message.from_user.id)
    broadcast_states.add(message.from_user.id)
    await message.answer("📢 Отправь текст рассылки следующим сообщением (все пользователи бота получат его).")

# ======================= ПРЕЗЕНТАЦИИ =======================

PPT_THEMES = {
    "ppt_blue": {
        "label": "🔵 Классический синий",
        "primary": RGBColor(0x1F, 0x4E, 0x79),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xD6, 0xE4, 0xF0),
        "body_text": RGBColor(0x21, 0x21, 0x21),
    },
    "ppt_dark": {
        "label": "⚫ Тёмный модерн",
        "primary": RGBColor(0x18, 0x18, 0x18),
        "bg": RGBColor(0x2B, 0x2B, 0x2B),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xCC, 0xCC, 0xCC),
        "body_text": RGBColor(0xF2, 0xF2, 0xF2),
    },
    "ppt_green": {
        "label": "🟢 Изумрудный",
        "primary": RGBColor(0x0E, 0x6B, 0x4F),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xD2, 0xF0, 0xE4),
        "body_text": RGBColor(0x1A, 0x1A, 0x1A),
    },
    "ppt_orange": {
        "label": "🟠 Тёплый оранжевый",
        "primary": RGBColor(0xC1, 0x5A, 0x11),
        "bg": RGBColor(0xFF, 0xF8, 0xF2),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xFC, 0xE2, 0xC8),
        "body_text": RGBColor(0x2B, 0x1A, 0x0D),
    },
}

def _set_text_style(paragraph, size, color, bold=False, font="Aptos"):
    paragraph.font.name = font
    paragraph.font.size = PptxPt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold


def _add_footer(slide, theme, slide_no=None):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0.55), PptxInches(7.08), PptxInches(12.2), PptxInches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = theme["primary"]; line.line.fill.background()
    if slide_no is not None:
        tb = slide.shapes.add_textbox(PptxInches(11.9), PptxInches(7.1), PptxInches(0.7), PptxInches(0.25))
        p = tb.text_frame.paragraphs[0]; p.text = str(slide_no); p.alignment = PP_ALIGN.RIGHT
        _set_text_style(p, 9, theme["body_text"])


def build_title_slide(prs, theme, topic_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = theme["primary"]

    # Декоративные блоки делают титульник похожим на современный шаблон, а не на обычный текстовый лист.
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(0.22), PptxInches(7.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = theme["subtitle_text"]; accent.line.fill.background()

    label = slide.shapes.add_textbox(PptxInches(0.9), PptxInches(1.0), PptxInches(3.0), PptxInches(0.4))
    p = label.text_frame.paragraphs[0]; p.text = "MECAUAI • ПРЕЗЕНТАЦИЯ"
    _set_text_style(p, 11, theme["subtitle_text"], True)

    title = slide.shapes.add_textbox(PptxInches(0.9), PptxInches(2.0), PptxInches(11.2), PptxInches(2.2))
    tf = title.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(topic_text).strip()[:180]; p.alignment = PP_ALIGN.LEFT
    _set_text_style(p, 34 if len(str(topic_text)) < 80 else 28, theme["title_text"], True)

    sub = slide.shapes.add_textbox(PptxInches(0.92), PptxInches(4.55), PptxInches(8.8), PptxInches(0.8))
    p = sub.text_frame.paragraphs[0]; p.text = "Структурировано • Визуально • Готово к защите"
    _set_text_style(p, 16, theme["subtitle_text"])

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(9.9), PptxInches(5.55), PptxInches(2.25), PptxInches(0.7))
    badge.fill.solid(); badge.fill.fore_color.rgb = theme["bg"]; badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]; p.text = "AI • 2026"; p.alignment = PP_ALIGN.CENTER
    _set_text_style(p, 13, theme["body_text"], True)
    return slide


def build_content_slide(prs, theme, idx, item, img_stream):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = theme["bg"]

    # Небольшая верхняя полоса вместо огромного синего заголовка.
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), prs.slide_width, PptxInches(0.18))
    top.fill.solid(); top.fill.fore_color.rgb = theme["primary"]; top.line.fill.background()

    num = s.shapes.add_shape(MSO_SHAPE.OVAL, PptxInches(0.62), PptxInches(0.55), PptxInches(0.55), PptxInches(0.55))
    num.fill.solid(); num.fill.fore_color.rgb = theme["primary"]; num.line.fill.background()
    p = num.text_frame.paragraphs[0]; p.text = str(idx + 1); p.alignment = PP_ALIGN.CENTER
    _set_text_style(p, 13, theme["title_text"], True)

    tb_title = s.shapes.add_textbox(PptxInches(1.35), PptxInches(0.47), PptxInches(10.9), PptxInches(0.75))
    p_title = tb_title.text_frame.paragraphs[0]; p_title.text = str(item.get("title", f"Слайд {idx + 1}"))[:120]
    _set_text_style(p_title, 27 if len(p_title.text) < 65 else 22, theme["body_text"], True)

    points = [str(x).strip() for x in (item.get("points") or []) if str(x).strip()]
    points = points[:6] or ["Ключевая информация по теме будет добавлена автоматически."]
    has_img = img_stream is not None
    left_w = 7.05 if has_img else 11.7

    # Карточки с тезисами вместо длинного списка маркеров.
    card_h = 0.72 if len(points) <= 4 else 0.62
    y = 1.55
    for i, pt in enumerate(points):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.62), PptxInches(y), PptxInches(left_w), PptxInches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA) if theme["bg"] == RGBColor(0xFF,0xFF,0xFF) else theme["bg"]
        card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xEC)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, PptxInches(0.82), PptxInches(y + 0.16), PptxInches(0.32), PptxInches(0.32))
        badge.fill.solid(); badge.fill.fore_color.rgb = theme["primary"]; badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]; bp.text = str(i + 1); bp.alignment = PP_ALIGN.CENTER
        _set_text_style(bp, 8, theme["title_text"], True)
        tb = s.shapes.add_textbox(PptxInches(1.3), PptxInches(y + 0.10), PptxInches(left_w - 0.85), PptxInches(card_h - 0.15))
        tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
        pp = tf.paragraphs[0]; pp.text = pt[:220]
        _set_text_style(pp, 14 if len(pt) < 130 else 12, theme["body_text"])
        y += card_h + 0.13

    if has_img:
        frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(8.05), PptxInches(1.55), PptxInches(4.65), PptxInches(4.95))
        frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(0xEE,0xF1,0xF5); frame.line.color.rgb = RGBColor(0xD9,0xDE,0xE5)
        try:
            s.shapes.add_picture(img_stream, left=PptxInches(8.18), top=PptxInches(1.68), width=PptxInches(4.39), height=PptxInches(4.69))
        except Exception as img_err:
            logging.error(f"Не удалось вставить картинку на слайд {idx}: {img_err}")

    _add_footer(s, theme, idx + 2)
    return s


async def generate_presentation_file(message: Message, topic: str, user_id: int = None):
    """Генерирует презентацию из нормализованной структуры, не теряя исходную тему."""
    user_id = user_id if user_id is not None else message.from_user.id
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    busy_users.add(user_id)
    theme_key = user_ppt_design.pop(user_id, "ppt_blue")
    theme = PPT_THEMES.get(theme_key, PPT_THEMES["ppt_blue"])
    user_images = user_ppt_images.pop(user_id, [])
    status_msg = await message.answer("📈 Готовлю структуру, дизайн и визуалы презентации…")
    heavy_lock = None
    try:
        heavy_lock = await acquire_heavy_job(status_msg)
        await bot.send_chat_action(chat_id=message.chat.id, action="typing")
        num_content_slides = max(len(user_images), 6)
        response = await call_groq_with_retry(
            messages=[
                {"role": "system", "content": (
                    f"Создай профессиональную учебную/деловую презентацию на русском языке. "
                    f"Нужно ровно {num_content_slides} содержательных слайдов ПОСЛЕ титульного. "
                    "Верни только валидный JSON-массив без markdown. Каждый объект: "
                    "{title: string, points: [3-5 коротких тезисов], image_prompt: string}. "
                    "Структура должна быть логичной: введение, ключевые факты/механизмы, данные или примеры, "
                    "проблемы/риски, перспективы, вывод. Не пиши заголовок вроде 'Презентация проекта': "
                    "title должен быть конкретным. Тезисы — не длиннее 180 символов, без повторов. "
                    "image_prompt — подробный англоязычный промпт для современной деловой/учебной иллюстрации, "
                    "без текста на изображении."
                )},
                {"role": "user", "content": f"Тема презентации: {topic}"}
            ],
            temperature=0.45
        )
        slides_data = extract_json(clean_text_for_html(response.choices[0].message.content))
        if not isinstance(slides_data, list) or not slides_data:
            raise ValueError("ИИ не вернул структуру слайдов")
        slides_data = slides_data[:num_content_slides]
        while len(slides_data) < num_content_slides:
            slides_data.append({"title": "Выводы", "points": ["Ключевые результаты темы.", "Практическое значение.", "Основной вывод."], "image_prompt": "Minimal professional abstract business illustration"})

        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        build_title_slide(prs, theme, topic)
        async with aiohttp.ClientSession() as session:
            for idx, item in enumerate(slides_data):
                img_stream = None
                if user_images and idx < len(user_images):
                    try:
                        file_info = await bot.get_file(user_images[idx])
                        downloaded = await bot.download_file(file_info.file_path)
                        img_bytes = downloaded.read()
                        if len(img_bytes) <= MAX_DOC_SIZE_BYTES:
                            img_stream = io.BytesIO(img_bytes)
                    except Exception as img_err:
                        logging.warning(f"Не удалось скачать пользовательскую картинку для слайда {idx + 1}: {img_err}")
                elif item.get("image_prompt") and idx < MAX_PPT_AI_IMAGES:
                    img_bytes = await generate_ai_image(session, item["image_prompt"], width=900, height=650)
                    if img_bytes:
                        img_stream = io.BytesIO(img_bytes)
                build_content_slide(prs, theme, idx, item, img_stream)

        bio = io.BytesIO(); prs.save(bio); bio.seek(0)
        filename = re.sub(r"[^A-Za-zА-Яа-я0-9_-]+", "_", topic[:55]).strip("_") or "Presentation"
        await status_msg.delete()
        await message.answer_document(
            BufferedInputFile(bio.read(), filename=f"MecauAI_{filename}.pptx"),
            caption=f"📈 Готово: {len(slides_data)+1} слайдов • {theme['label']}"
        )
        get_user_stats(user_id)["exports"] += 1
    except Exception as e:
        logging.error(f"Ошибка генерации презентации: {e}", exc_info=True)
        try: await status_msg.delete()
        except Exception: pass
        await message.answer("⚠️ Не удалось собрать презентацию. Я сохранил запрос; попробуй ещё раз — если ошибка повторится, пришли тему короче.")
    finally:
        busy_users.discard(user_id)
        user_ppt_topic.pop(user_id, None)
        ppt_states.discard(user_id)
        user_ppt_design.pop(user_id, None)
        user_ppt_images.pop(user_id, None)
        if heavy_lock is not None:
            heavy_lock.release()
        gc.collect()


async def generate_excel_file(message: Message, topic: str, user_id: int = None):
    """Генерирует устойчивый XLSX: данные + оформление + фильтр + диаграмма при наличии чисел."""
    user_id = user_id if user_id is not None else message.from_user.id
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    busy_users.add(user_id)
    excel_mode = user_excel_mode.pop(user_id, "excel_chart")
    status_msg = await message.answer("📊 Создаю Excel: данные, оформление и аналитику…")
    heavy_lock = None
    try:
        heavy_lock = await acquire_heavy_job(status_msg)
        response = await call_groq_with_retry(
            messages=[
                {"role": "system", "content": (
                    "Сформируй данные для Excel. Верни ТОЛЬКО валидный JSON-объект без markdown: "
                    "{title, sheet_name, headers, rows}. headers — 2-12 строк. rows — 5-30 строк. "
                    "Числа возвращай числами, даты — строками YYYY-MM-DD. Не придумывай сложные формулы. "
                    "Если пользователь просит статистику/динамику, включи числовые колонки, подходящие для диаграммы. "
                    "Все значения должны соответствовать теме и быть реалистично согласованными между собой."
                )},
                {"role": "user", "content": f"Задача: {topic}"}
            ],
            temperature=0.25
        )
        data = extract_json(clean_text_for_html(response.choices[0].message.content))
        headers = [str(x).strip() for x in (data.get("headers") or []) if str(x).strip()][:12]
        rows = data.get("rows") or []
        if not headers: raise ValueError("Пустые заголовки")
        num_cols = len(headers); norm_rows=[]
        for row in rows[:30]:
            row=list(row) if isinstance(row,(list,tuple)) else [row]
            row=(row+[""]*num_cols)[:num_cols]
            norm_rows.append(row)
        if not norm_rows: norm_rows=[["" for _ in headers]]

        wb=Workbook(); ws=wb.active
        sheet_name=re.sub(r"[\\/*?:\[\]]", "_", str(data.get("sheet_name") or "Данные"))[:31] or "Данные"
        ws.title=sheet_name
        title=str(data.get("title") or topic)[:120]
        ws.merge_cells(start_row=1,start_column=1,end_row=1,end_column=num_cols)
        c=ws.cell(1,1,title); c.font=Font(name="Aptos Display",size=18,bold=True,color="FFFFFF"); c.fill=PatternFill("solid",fgColor="1F4E78"); c.alignment=Alignment(horizontal="left",vertical="center")
        ws.row_dimensions[1].height=30
        ws.append(headers)
        header_row=2
        for cell in ws[header_row]:
            cell.font=Font(name="Aptos",size=11,bold=True,color="FFFFFF"); cell.fill=PatternFill("solid",fgColor="2F75B5"); cell.alignment=Alignment(horizontal="center",vertical="center",wrap_text=True)
        for row in norm_rows:
            ws.append(row)
        thin=Side(style="thin",color="D9E2F3")
        for row in ws.iter_rows(min_row=2,max_row=ws.max_row,min_col=1,max_col=num_cols):
            for cell in row:
                cell.border=Border(bottom=thin); cell.alignment=Alignment(vertical="center",wrap_text=True)
                if isinstance(cell.value,(int,float)): cell.number_format='#,##0.00' if isinstance(cell.value,float) else '#,##0'
        ws.freeze_panes="A3"; ws.auto_filter.ref=f"A2:{get_column_letter(num_cols)}{ws.max_row}"
        for i in range(1,num_cols+1):
            vals=[str(ws.cell(r,i).value or "") for r in range(2,ws.max_row+1)]
            ws.column_dimensions[get_column_letter(i)].width=min(max(max([len(str(headers[i-1]))]+[len(v) for v in vals])+2,12),34)
        if ws.max_row >= 3:
            tab=Table(displayName="DataTable",ref=f"A2:{get_column_letter(num_cols)}{ws.max_row}")
            tab.tableStyleInfo=TableStyleInfo(name="TableStyleMedium2",showFirstColumn=False,showLastColumn=False,showRowStripes=True,showColumnStripes=False)
            ws.add_table(tab)

        # Отдельный лист с диаграммой — не накладываем график на данные.
        if excel_mode in ("excel_chart","excel_full") and num_cols>=2:
            numeric=[]
            for ci in range(2,num_cols+1):
                vals=[ws.cell(r,ci).value for r in range(3,ws.max_row+1)]
                if vals and sum(isinstance(v,(int,float)) for v in vals)/len(vals)>=0.6: numeric.append(ci)
            if numeric:
                ch=BarChart(); ch.type="col"; ch.style=10; ch.title=f"Аналитика: {headers[numeric[0]-1]}"; ch.y_axis.title=headers[numeric[0]-1]; ch.x_axis.title=headers[0]
                ch.add_data(Reference(ws,min_col=numeric[0],max_col=numeric[0],min_row=2,max_row=ws.max_row),titles_from_data=True)
                ch.set_categories(Reference(ws,min_col=1,min_row=3,max_row=ws.max_row)); ch.width=18; ch.height=9
                a=wb.create_sheet("Аналитика"); a["A1"]=title; a["A1"].font=Font(size=18,bold=True); a.add_chart(ch,"A3")
                # Простые живые KPI: Excel пересчитает их при открытии файла.
                data_col = get_column_letter(numeric[0])
                first_data = 3; last_data = ws.max_row
                a["A15"] = "Показатель"; a["B15"] = "Значение"
                for cell in a[15]:
                    cell.font = Font(bold=True, color="FFFFFF"); cell.fill = PatternFill("solid", fgColor="2F75B5")
                a["A16"] = "Всего"; a["B16"] = f"=SUM(Данные!{data_col}{first_data}:{data_col}{last_data})"
                a["A17"] = "Среднее"; a["B17"] = f"=AVERAGE(Данные!{data_col}{first_data}:{data_col}{last_data})"
                a["A18"] = "Минимум"; a["B18"] = f"=MIN(Данные!{data_col}{first_data}:{data_col}{last_data})"
                a["A19"] = "Максимум"; a["B19"] = f"=MAX(Данные!{data_col}{first_data}:{data_col}{last_data})"
                a.column_dimensions["A"].width = 20; a.column_dimensions["B"].width = 20

        if excel_mode=="excel_full":
            try:
                async with aiohttp.ClientSession() as session:
                    img_bytes=await generate_ai_image(session,f"Clean modern business infographic about {topic}",width=900,height=500)
                if img_bytes:
                    from openpyxl.drawing.image import Image as XLImage
                    img=XLImage(io.BytesIO(img_bytes)); img.width=540; img.height=300
                    ws.add_image(img,f"A{ws.max_row+3}")
            except Exception as img_err: logging.warning(f"Excel image skipped: {img_err}")

        bio=io.BytesIO(); wb.save(bio); bio.seek(0)
        filename=re.sub(r"[^A-Za-zА-Яа-я0-9_-]+","_",topic[:55]).strip("_") or "Table"
        await status_msg.delete()
        await message.answer_document(BufferedInputFile(bio.read(),filename=f"MecauAI_{filename}.xlsx"),caption=f"📊 Excel готов • {len(norm_rows)} строк • {num_cols} столбцов")
        get_user_stats(user_id)["exports"] += 1
    except Exception as e:
        logging.error(f"Ошибка генерации Excel: {e}",exc_info=True)
        try: await status_msg.delete()
        except Exception: pass
        await message.answer("⚠️ Excel не удалось собрать. Попробуй описать таблицу чуть конкретнее: какие столбцы и за какой период нужны.")
    finally:
        busy_users.discard(user_id); user_excel_topic.pop(user_id,None); excel_states.discard(user_id)
        user_excel_mode.pop(user_id, None)
        if heavy_lock is not None:
            heavy_lock.release()
        gc.collect()

@dp.message(F.text == "📈 Презентация")
async def cmd_ppt_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    # Если тема уже передана естественной фразой, не заставляем вводить её второй раз.
    topic = user_ppt_topic.get(user_id)
    clear_pending_states(user_id)
    if topic:
        user_ppt_topic[user_id] = topic
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=v["label"], callback_data=k)] for k, v in PPT_THEMES.items()
    ] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
    prompt = "🎨 Выбери дизайн презентации:" if not topic else f"🎨 Тема: «{topic[:180]}»\nВыбери дизайн:"
    await message.answer(prompt, reply_markup=kb)

@dp.callback_query(F.data.in_(set(PPT_THEMES.keys())))
async def cb_ppt_design(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    user_ppt_design[user_id] = callback.data
    topic = user_ppt_topic.get(user_id)
    ppt_states.add(user_id)
    theme_label = PPT_THEMES[callback.data]["label"]
    try:
        await callback.message.edit_text(f"🎨 Дизайн выбран: {theme_label}")
    except Exception:
        pass
    if topic:
        await callback.message.answer(f"✅ Дизайн выбран: {theme_label}. Начинаю генерацию по сохранённой теме…")
        await generate_presentation_file(callback.message, topic, user_id=user_id)
    else:
        await callback.message.answer(
            "📈 Отправь тему презентации следующим сообщением.\n\n"
            "💡 Свои картинки можно отправить перед темой — я распределю их по слайдам.",
        )
    await callback.answer()

@dp.callback_query(F.data == "cancel_flow")
async def cb_cancel_flow(callback: types.CallbackQuery):
    clear_pending_states(callback.from_user.id)
    try:
        await callback.message.edit_text("❌ Действие отменено.")
    except Exception:
        pass
    await callback.answer("Отменено")

# ======================= EXCEL-ТАБЛИЦЫ =======================

EXCEL_MODES = {
    "excel_plain": "📋 Только таблица",
    "excel_chart": "📊 Таблица + диаграмма",
    "excel_full": "🖼 Таблица + диаграмма + ИИ-картинка",
}

@dp.message(F.text == "📊 Excel-таблица")
async def cmd_excel_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    topic = user_excel_topic.get(user_id)
    clear_pending_states(user_id)
    if topic:
        user_excel_topic[user_id] = topic
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=v, callback_data=k)] for k, v in EXCEL_MODES.items()
    ] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
    prompt = "📊 Выбери формат Excel:" if not topic else f"📊 Запрос: «{topic[:180]}»\nВыбери формат:"
    await message.answer(prompt, reply_markup=kb)

@dp.callback_query(F.data.in_(set(EXCEL_MODES.keys())))
async def cb_excel_mode(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    user_excel_mode[user_id] = callback.data
    topic = user_excel_topic.get(user_id)
    excel_states.add(user_id)
    mode_label = EXCEL_MODES[callback.data]
    try:
        await callback.message.edit_text(f"📊 Выбрано: {mode_label}")
    except Exception:
        pass
    if topic:
        await callback.message.answer(f"✅ Формат выбран: {mode_label}. Начинаю генерацию Excel…")
        await generate_excel_file(callback.message, topic, user_id=user_id)
    else:
        await callback.message.answer(
            "Опиши задачу или тему для таблицы следующим сообщением "
            "(например: «Таблица успеваемости студентов с оценками»), "
            "и я сформирую готовый Excel-файл (.xlsx)!"
        )
    await callback.answer()

# ======================= АНАЛИЗ ДОКУМЕНТОВ =======================

MAX_DOC_CONTEXT_CHARS = 12000
MAX_DOC_SIZE_BYTES = 20 * 1024 * 1024  # ограничение Telegram Bot API на скачивание файла

@dp.message(F.text == "📚 Анализ документа")
async def cmd_doc_analysis_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    clear_pending_states(user_id)
    doc_analysis_states.add(user_id)
    await message.answer(
        "📚 Пришли файл в формате <b>.pdf</b>, <b>.docx</b> или <b>.txt</b> (можно также прислать фото страницы) — "
        "я прочитаю его, сделаю краткое содержание, и дальше сможешь задавать вопросы по документу прямо в чате.",
        parse_mode="HTML"
    )

async def extract_text_from_document(ext: str, file_bytes: bytes):
    """Извлекает текст из .pdf / .docx / .txt. Возвращает (текст, ошибка_или_None)."""
    if ext == "txt":
        return file_bytes.decode("utf-8", errors="ignore"), None
    if ext == "docx":
        d = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in d.paragraphs), None
    if ext == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return None, "⚠️ Для чтения PDF на сервере не установлена библиотека pypdf. Установи её командой: pip install pypdf"
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
            return text, None
        except Exception as e:
            return None, f"⚠️ Не удалось прочитать PDF: {e}"
    return None, "⚠️ Неподдерживаемый формат файла."

async def process_document_context(message: Message, filename: str, text: str, status_msg: Message):
    user_id = message.from_user.id
    text = (text or "").strip()
    if not text:
        await status_msg.edit_text("⚠️ Не удалось извлечь текст из документа (возможно, это скан без текстового слоя).")
        doc_analysis_states.discard(user_id)
        return

    truncated = text[:MAX_DOC_CONTEXT_CHARS]
    user_doc_context[user_id] = truncated
    doc_analysis_states.discard(user_id)

    try:
        summary_resp = await call_groq_with_retry(messages=[
            {"role": "system", "content": "Сделай краткое структурированное содержание документа на русском языке (5-8 предложений)."},
            {"role": "user", "content": truncated}
        ])
        summary = clean_text_for_html(summary_resp.choices[0].message.content)
    except Exception as e:
        logging.error(f"Не удалось получить содержание документа: {e}")
        summary = "(не удалось сформировать краткое содержание, но документ загружен — можешь задавать вопросы)"

    await status_msg.edit_text(f"✅ Документ «{filename}» прочитан ({len(text)} символов).")
    await safe_answer(
        message,
        f"📋 <b>Краткое содержание:</b>\n\n{summary}\n\n"
        f"💬 Теперь можешь задавать вопросы по документу, попросить краткий конспект, выделить главное, сравнить части документа или подготовить ответ/структуру на его основе.",
        parse_mode="HTML"
    )

@dp.message(F.document)
async def handle_document(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if not allow_request(user_id, limit=5, window=60):
        await message.answer("🛡 Слишком много файлов подряд. Подожди немного и попробуй снова.")
        return
    if user_id not in doc_analysis_states:
        await message.answer("📚 Чтобы проанализировать документ, сначала нажми кнопку «📚 Анализ документа», а затем пришли файл.")
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    get_user_stats(user_id)["files"] += 1
    doc = message.document
    filename = doc.file_name or "документ"
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in ("pdf", "docx", "txt"):
        await message.answer("⚠️ Поддерживаются только форматы .pdf, .docx и .txt")
        return
    if doc.file_size and doc.file_size > MAX_DOC_SIZE_BYTES:
        await message.answer("⚠️ Файл слишком большой (лимит Telegram на скачивание — 20 МБ).")
        return

    busy_users.add(user_id)
    status_msg = await message.answer("📚 Читаю документ...")
    try:
        file_info = await bot.get_file(doc.file_id)
        downloaded = await bot.download_file(file_info.file_path)
        file_bytes = downloaded.read()

        text, err = await extract_text_from_document(ext, file_bytes)
        if err:
            await status_msg.edit_text(err)
            doc_analysis_states.discard(user_id)
            return

        await process_document_context(message, filename, text, status_msg)
    except Exception as e:
        logging.error(f"Ошибка при обработке документа: {e}", exc_info=True)
        try:
            await status_msg.edit_text("⚠️ Ошибка при чтении документа. Попробуй ещё раз.")
        except Exception:
            pass
        doc_analysis_states.discard(user_id)
    finally:
        busy_users.discard(user_id)

@dp.message(F.text == "⭐ Избранное")
async def cmd_favorites(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(user_id)

    favs = load_favorites()
    uid_str = str(user_id)
    user_favs = favs.get(uid_str, [])

    if not user_favs:
        await message.answer("⭐ У тебя пока нет сохраненных ответов в избранном.")
        return

    await message.answer(f"⭐ Твои сохраненные ответы ({len(user_favs)}):")
    for idx, fav_text in enumerate(user_favs, 1):
        display_text = f"Сохранение #{idx}:\n\n{fav_text}"
        if len(display_text) > 4000:
            display_text = display_text[:3997] + "..."
        await message.answer(display_text, disable_web_page_preview=True)

@dp.callback_query(F.data.in_({
    "guide_voice", "guide_image", "guide_doc", "guide_ppt", "guide_excel", "guide_word", "create_ppt", "create_excel", "create_word", "create_title"
}))
async def cb_product_navigation(callback: types.CallbackQuery):
    data = callback.data
    messages = {
        "guide_voice": "🎙 Отправь голосовое — ничего включать не нужно.",
        "guide_image": "🖼 Отправь фото — я могу разобрать задачу, текст или изображение. Для презентации свои картинки можно добавить прямо в процессе её создания.",
        "guide_doc": "📁 Отправь PDF, DOCX или TXT — после загрузки задавай вопросы.",
        "guide_ppt": "📈 Нажми «Создать → Презентация». Я буду задавать вопросы по одному.",
        "guide_excel": "📊 Нажми «Создать → Excel» и опиши, какую таблицу хочешь.",
        "guide_word": "📄 Под ответом нажми «В Word», чтобы получить файл.",
        "create_ppt": "📈 Хорошо. Я помогу выбрать назначение, стиль и титульный лист — без сложных терминов.",
        "create_excel": "📊 Опиши нужную таблицу или пришли пример.",
        "create_word": "📄 Получи ответ и нажми «В Word».",
        "create_title": "📑 Я сначала уточню, для чего титульник нужен, и предложу понятные варианты."
    }
    if data == "create_ppt":
        user_id = callback.from_user.id
        clear_pending_states(user_id)
        kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text=v["label"], callback_data=k)] for k, v in PPT_THEMES.items()] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
        await callback.message.answer(
            "📈 <b>Создание презентации</b>\n\n"
            "Сначала выбери стиль. После этого я попрошу тему. Можно также сразу написать тему обычным сообщением.",
            parse_mode="HTML", reply_markup=kb
        )
    elif data == "create_excel":
        user_id = callback.from_user.id
        clear_pending_states(user_id)
        kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text=v, callback_data=k)] for k, v in EXCEL_MODES.items()] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
        await callback.message.answer(
            "📊 <b>Создание Excel</b>\n\n"
            "Выбери формат. После выбора просто напиши, какую таблицу нужно сделать.",
            parse_mode="HTML", reply_markup=kb
        )
    elif data == "create_title":
        await callback.message.answer("📑 Выбери тип титульного листа:", reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📘 Индивидуальный проект", callback_data="gost_project")],
            [InlineKeyboardButton(text="📗 Курсовая работа", callback_data="gost_coursework")],
            [InlineKeyboardButton(text="📕 Дипломная работа (ВКР)", callback_data="gost_diploma")],
            [InlineKeyboardButton(text="📙 Отчёт по практике", callback_data="gost_practice")]
        ]))
    else:
        await callback.message.answer(messages.get(data, "Готово."))
    await callback.answer()

@dp.callback_query(F.data == "flow_restart")
async def cb_flow_restart(callback: types.CallbackQuery):
    clear_task_state(callback.from_user.id)
    await callback.message.answer(
        "🔄 Начинаем заново.\n\n"
        "Просто напиши, что хочешь сделать, или выбери действие ниже.",
        reply_markup=get_quick_actions_keyboard()
    )
    await callback.answer()

@dp.callback_query(F.data.in_({"set_mode_ai", "set_mode_coder", "set_mode_friend"}))
async def cb_modes(callback: types.CallbackQuery):
    modes = {
        "set_mode_ai": ("🤖 AI-помощник", "универсальный режим"),
        "set_mode_coder": ("💻 Кодер", "режим для программирования и разбора кода"),
        "set_mode_friend": ("🤝 Дружеский режим", "более простой и неформальный стиль"),
    }
    title, mode = modes.get(callback.data, ("Режим", "универсальный режим"))
    set_task_state(callback.from_user.id, preferred_mode=mode)
    await callback.message.answer(
        f"✅ <b>{title}</b>\nНастройка сохранена. Я буду учитывать её в следующих запросах.",
        parse_mode="HTML"
    )
    await callback.answer()

@dp.callback_query(F.data.in_({"btn_simplify", "btn_save_fav", "btn_word", "btn_continue"}))
async def cb_answer_actions(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    msg_text = callback.message.text or callback.message.caption or ""

    clean_msg = msg_text.split("—\n⚡")[0].strip() if "—\n⚡" in msg_text else msg_text.strip()
    if not clean_msg:
        await callback.answer("⚠️ Нечего обрабатывать!", show_alert=True)
        return

    if callback.data == "btn_word":
        user_data = get_user_data(user_id)
        user_data["last_output"] = clean_msg
        user_data_stats = get_user_stats(user_id)
        user_data_stats["exports"] += 1

        doc = build_professional_word(clean_msg, "MecauAI — ответ")
        bio = io.BytesIO(); doc.save(bio); bio.seek(0)
        path = _artifact_path(user_id, "docx", "Answer")
        with open(path, "wb") as f: f.write(bio.getvalue())
        register_artifact(user_id, "docx", path, {"title": "Ответ"})
        await callback.message.answer_document(
            BufferedInputFile(bio.getvalue(), filename=os.path.basename(path)),
            caption="📄 Готово — оформил ответ как полноценный Word-документ. Можно попросить изменить его."
        )
        await callback.answer()
        return

    if callback.data == "btn_continue":
        await callback.message.answer(
            "🔁 Продолжаем. Напиши следующим сообщением, что изменить, добавить или уточнить — контекст сохранён."
        )
        await callback.answer()
        return

    if callback.data == "btn_save_fav":
        add_to_favorites(user_id, clean_msg)
        await callback.answer("⭐ Успешно сохранено в избранное!", show_alert=True)
    elif callback.data == "btn_simplify":
        await callback.answer("💡 Сжимаю до сути...")
        try:
            response = await call_groq_with_retry(messages=[
                {"role": "system", "content": "Объясни предельно коротко и просто в 3-5 предложениях."},
                {"role": "user", "content": clean_msg}
            ])
            simplified_reply = clean_text_for_html(response.choices[0].message.content)
            full_reply = f"💡 <b>Коротко на пальцах:</b>\n\n{simplified_reply}{AD_FOOTER}"
            await safe_answer(callback.message, full_reply, parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
        except Exception:
            await callback.message.answer("⚠️ Ошибка при обработке. Попробуй ещё раз. Если повторится, пиши - @mecau")

@dp.message(F.text == "📄 Ответ в Word")

def build_professional_word(text, title="MecauAI — документ"):
    """Word pipeline: Heading 1/2/3, lists, TOC field, footer/page numbers, margins."""
    doc=Document()
    sec=doc.sections[0]
    sec.top_margin=Mm(20); sec.bottom_margin=Mm(20); sec.left_margin=Mm(25); sec.right_margin=Mm(20)
    styles=doc.styles
    for name,size,bold in (("Normal",11,False),("Heading 1",16,True),("Heading 2",14,True),("Heading 3",12,True)):
        try:
            st=styles[name]; st.font.name='Aptos'; st.font.size=Pt(size); st.font.bold=bold
        except Exception: pass
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(title[:150]); r.bold=True; r.font.name='Aptos Display'; r.font.size=Pt(20)
    doc.add_paragraph()
    toc=doc.add_paragraph(); toc.add_run('Содержание').bold=True
    fld=toc.add_run(); fld._r.append(__import__('docx').oxml.OxmlElement('w:fldChar')); fld._r[-1].set(__import__('docx').oxml.ns.qn('w:fldCharType'),'begin')
    instr=__import__('docx').oxml.OxmlElement('w:instrText'); instr.set(__import__('docx').oxml.ns.qn('xml:space'),'preserve'); instr.text='TOC \\o "1-3" \\h \\z \\u'; fld._r.append(instr)
    end=__import__('docx').oxml.OxmlElement('w:fldChar'); end.set(__import__('docx').oxml.ns.qn('w:fldCharType'),'end'); fld._r.append(end)
    doc.add_page_break()
    lines=str(text or '').replace('\r','').split('\n')
    for raw in lines:
        line=raw.strip()
        if not line: continue
        if line.startswith('### '): doc.add_paragraph(line[4:].strip(),style='Heading 3'); continue
        if line.startswith('## '): doc.add_paragraph(line[3:].strip(),style='Heading 2'); continue
        if line.startswith('# '): doc.add_paragraph(line[2:].strip(),style='Heading 1'); continue
        if re.match(r'^[-•*]\s+',line):
            p=doc.add_paragraph(re.sub(r'^[-•*]\s+','',line),style='List Bullet'); continue
        if re.match(r'^\d+[.)]\s+',line):
            doc.add_paragraph(re.sub(r'^\d+[.)]\s+','',line),style='List Number'); continue
        p=doc.add_paragraph(line,style='Normal'); p.paragraph_format.space_after=Pt(6)
    # Footer with page number field.
    footer=sec.footer.paragraphs[0]; footer.alignment=WD_ALIGN_PARAGRAPH.CENTER; footer.add_run('MecauAI • ')
    run=footer.add_run(); fld=__import__('docx').oxml.OxmlElement('w:fldSimple'); fld.set(__import__('docx').oxml.ns.qn('w:instr'),'PAGE'); run._r.addnext(fld)
    return doc

async def cmd_download_word(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(user_id)

    user_data = get_user_data(user_id)
    text_to_save = user_data.get("last_output", "").strip()
    if not text_to_save or text_to_save == "Здесь пока нет ответов.":
        await message.answer("⚠️ Нет данных для сохранения. Попробуй ещё раз. Если повторится, пиши - @mecau")
        return

    doc = build_professional_word(text_to_save, "MecauAI — ответ")
    bio = io.BytesIO(); doc.save(bio); bio.seek(0)
    path = _artifact_path(user_id, "docx", "Answer")
    with open(path, "wb") as f: f.write(bio.getvalue())
    register_artifact(user_id, "docx", path, {"title": "Ответ"})
    await message.answer_document(BufferedInputFile(bio.getvalue(), filename=os.path.basename(path)), caption="📄 Вот твой оформленный Word! ✏️ Можно попросить изменить его.")

@dp.message(F.text == "📑 Титульник ГОСТ")
async def cmd_gost_title(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(message.from_user.id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📘 Индивидуальный проект", callback_data="gost_project")],
        [InlineKeyboardButton(text="📗 Курсовая работа",        callback_data="gost_coursework")],
        [InlineKeyboardButton(text="📕 Дипломная работа (ВКР)", callback_data="gost_diploma")],
        [InlineKeyboardButton(text="📙 Отчёт по практике",      callback_data="gost_practice")],
    ])
    await message.answer("📑 Выбери тип работы для титульника:", reply_markup=kb)

@dp.callback_query(F.data.startswith("gost_"))
async def cb_gost_generate(callback: types.CallbackQuery):
    work_type_map = {
        "gost_project":    ("ИНДИВИДУАЛЬНЫЙ ПРОЕКТ",                   "Индивидуальный_проект"),
        "gost_coursework": ("КУРСОВАЯ РАБОТА",                         "Курсовая_работа"),
        "gost_diploma":    ("ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА (ВКР)", "ВКР"),
        "gost_practice":   ("ОТЧЁТ ПО ПРАКТИКЕ",                      "Отчет_по_практике"),
    }
    work_label, filename_base = work_type_map[callback.data]

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Mm(20)
        section.bottom_margin = Mm(20)
        section.left_margin   = Mm(30)
        section.right_margin  = Mm(15)

    def add_centered(text, size=14, bold=False):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.name = 'Times New Roman'
        run.font.size = Pt(size)
        run.bold = bold
        return p

    def add_right(text, size=14):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = p.add_run(text)
        run.font.name = 'Times New Roman'
        run.font.size = Pt(size)
        return p

    def add_empty(count=1):
        for _ in range(count):
            p = doc.add_paragraph()
            p.add_run("").font.name = 'Times New Roman'

    add_centered("ФЕДЕРАЛЬНОЕ ГОСУДАРСТВЕННОЕ БЮДЖЕТНОЕ ПРОФЕССИОНАЛЬНОЕ\nОБРАЗОВАТЕЛЬНОЕ УЧРЕЖДЕНИЕ\n«НАЗВАНИЕ КОЛЛЕДЖА»", size=14)
    add_empty(4)
    add_centered(work_label, size=14, bold=True)
    add_empty(1)
    add_centered("на тему:\n«Введи тему работы здесь»", size=14, bold=True)
    add_empty(6)
    add_right("Выполнил(а): студент(ка) группы ГРУППА\nФамилия Имя Отчество\n\nРуководитель:\nДолжность, Фамилия И.О.")
    add_empty(5)
    add_centered("Город — 2026", size=14)

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)

    file_doc = BufferedInputFile(bio.read(), filename=f"Titulnik_{filename_base}.docx")
    await callback.message.answer_document(file_doc, caption=f"📑 Титульник «{work_label}» готов!")
    await callback.answer()

# Голосовой режим намеренно не добавлен в меню: достаточно отправить voice-сообщение.
# Это сохраняет интерфейс компактным и не перегружает пользователя кнопками.
MENU_BUTTONS = {
    "✨ Возможности", "📂 Файлы", "⭐ Избранное", "⚙️ Настройки",
    "📢 Рассылка", "📊 Статистика",
    # Старые кнопки тоже игнорируем, если они остались у уже открытого клиента.
    "📄 Ответ в Word", "📑 Титульник ГОСТ", "📈 Презентация",
    "📊 Excel-таблица", "📚 Анализ документа", "🔄 Режим ИИ", "ℹ️ О MecauAI"
}


@dp.message(F.voice)
async def handle_voice(message: Message):
    """
    Голосовой режим без отдельной кнопки:
    пользователь отправляет voice -> бот распознаёт речь -> использует
    полученный текст как обычный запрос ИИ.
    """
    user_id = message.from_user.id

    if not await check_subscription(user_id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return

    if not allow_request(user_id, limit=8, window=30):
        await message.answer("🛡 Слишком много запросов подряд. Подожди немного — я защищаю сервер от перегрузки.")
        return

    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    save_user_id(user_id)
    get_user_stats(user_id)["voice"] += 1
    status_msg = await message.answer("🎙 Распознаю голосовое сообщение...")

    busy_users.add(user_id)
    try:
        file_info = await bot.get_file(message.voice.file_id)
        downloaded = await bot.download_file(file_info.file_path)
        audio_bytes = downloaded.read()

        if len(audio_bytes) > MAX_DOC_SIZE_BYTES:
            await status_msg.edit_text("⚠️ Голосовое сообщение слишком большое. Пришли более короткую запись.")
            return

        transcript = await transcribe_voice(audio_bytes)
        if not transcript:
            await status_msg.edit_text("⚠️ Не удалось разобрать речь. Попробуй записать сообщение ещё раз.")
            return

        # Сначала учитываем уже выбранный workflow. Раньше voice обходил
        # ppt_states/excel_states и поэтому тема уходила в обычный AI-чат.
        voice_norm = transcript.lower().strip()

        if user_id in ppt_states:
            topic = transcript.strip()
            if topic:
                ppt_states.discard(user_id)
                user_ppt_topic[user_id] = topic
                await status_msg.edit_text(
                    f"🎙 <b>Распознал тему презентации:</b>\n«{topic[:500]}»\n\n⏳ Начинаю создание…",
                    parse_mode="HTML"
                )
                await generate_presentation_file(message, topic, user_id=user_id)
            else:
                await status_msg.edit_text("⚠️ Не удалось определить тему презентации.")
            return

        if user_id in excel_states:
            topic = transcript.strip()
            if topic:
                excel_states.discard(user_id)
                user_excel_topic[user_id] = topic
                await status_msg.edit_text(
                    f"🎙 <b>Распознал задачу для Excel:</b>\n«{topic[:500]}»\n\n⏳ Начинаю создание…",
                    parse_mode="HTML"
                )
                await generate_excel_file(message, topic, user_id=user_id)
            else:
                await status_msg.edit_text("⚠️ Не удалось определить задачу для Excel.")
            return

        # Естественные голосовые команды: не требуем строго фразу «сделай ...».
        if re.search(r"\b(презентаци[яиюеи]|слайд(?:ы|ов)?)\b", voice_norm) and re.search(r"\b(сделай|создай|подготовь|сформируй|сделать|нужна|нужен|нужно)\b", voice_norm):
            topic = re.sub(r"^.*?\bпрезентаци(?:ю|я|и|й)?\b\s*(?:на тему|по теме|про|о)?\s*", "", transcript, flags=re.I).strip(" .:-") or "Тема не указана"
            user_ppt_topic[user_id] = topic
            ppt_states.discard(user_id)
            user_ppt_design.pop(user_id, None)
            kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text=v["label"], callback_data=k)] for k,v in PPT_THEMES.items()] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
            await status_msg.edit_text(f"🎙 Я понял: сделать презентацию на тему «{topic[:180]}».\nВыбери дизайн:", reply_markup=kb)
            return
        if re.search(r"\b(excel|таблиц(?:а|у|ы)?|xlsx)\b", voice_norm) and re.search(r"\b(сделай|создай|подготовь|сформируй|сделать|нужна|нужен|нужно)\b", voice_norm):
            topic = re.sub(r"^.*?\b(?:excel|таблиц(?:у|а|ы)?|xlsx)\b\s*(?:на тему|по теме|для|с)?\s*", "", transcript, flags=re.I).strip(" .:-") or transcript
            user_excel_topic[user_id] = topic
            excel_states.discard(user_id)
            kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text=v, callback_data=k)] for k,v in EXCEL_MODES.items()] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
            await status_msg.edit_text(f"🎙 Я понял: создать Excel по запросу «{topic[:180]}».\nВыбери формат:", reply_markup=kb)
            return

        user_data = get_user_data(user_id)
        current_mode = user_data["mode"]
        system_prompt = PROMPTS.get(current_mode, PROMPTS["ai"])

        doc_context = user_doc_context.get(user_id)
        if doc_context:
            system_prompt = (
                f"{system_prompt}\n\n"
                f"Пользователь ранее загрузил документ. Используй его содержимое, "
                f"если вопрос с ним связан:\n---\n{doc_context}\n---"
            )

        # Голос обрабатывается как обычный текстовый запрос, поэтому история,
        # режим ИИ и экспорт последнего ответа продолжают работать одинаково.
        user_data["history"].append({"role": "user", "content": transcript})
        if len(user_data["history"]) > 6:
            user_data["history"] = user_data["history"][-6:]

        messages_payload = [{"role": "system", "content": system_prompt}] + user_data["history"]

        await status_msg.edit_text(f"🎙 Распознал:\n\n«{transcript[:3500]}»\n\n🤖 Формирую ответ...")
        await bot.send_chat_action(chat_id=message.chat.id, action="typing")

        response = await call_groq_with_retry(messages=messages_payload)
        ai_reply = clean_text_for_html(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply

        full_message = (
            f"🎙 <b>Твой запрос:</b> {transcript}\n\n"
            f"🤖 <b>Ответ:</b>\n\n{ai_reply}{AD_FOOTER}"
        )

        try:
            await status_msg.delete()
        except Exception:
            pass

        parse_mode = "Markdown" if current_mode == "coder" else "HTML"
        await safe_answer(
            message,
            full_message,
            parse_mode=parse_mode,
            reply_markup=get_answer_inline_keyboard()
        )

    except Exception as e:
        logging.error(f"Ошибка обработки голосового сообщения: {e}", exc_info=True)
        try:
            await status_msg.edit_text(
                "⚠️ Не удалось обработать голосовое сообщение. "
                "Попробуй записать его ещё раз чуть короче и без сильного фонового шума."
            )
        except Exception:
            pass
    finally:
        busy_users.discard(user_id)

@dp.message(F.photo)
async def handle_photo(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    if not allow_request(user_id, limit=6, window=30):
        await message.answer("🛡 Слишком много изображений подряд. Подожди немного, чтобы сервер не перегрузился.")
        return

    save_user_id(user_id)
    get_user_stats(user_id)["images"] += 1
    photo = message.photo[-1]

    # Для презентации храним только Telegram file_id, а не байты картинки.
    # Это резко снижает RAM на слабом VPS.
    if user_id in ppt_states:
        images = user_ppt_images.setdefault(user_id, [])
        if len(images) >= MAX_PPT_IMAGES:
            await message.answer(f"⚠️ Для одной презентации можно добавить максимум {MAX_PPT_IMAGES} картинок.")
            return
        if photo.file_id not in images:
            images.append(photo.file_id)
        await message.answer(f"🖼 Картинка добавлена в презентацию ({len(images)}/{MAX_PPT_IMAGES}). Можешь отправить ещё или продолжить с темой.")
        return

    # Скачивание фото выполняем только когда оно действительно нужно Vision.
    try:
        if photo.file_size and photo.file_size > MAX_DOC_SIZE_BYTES:
            await message.answer("⚠️ Фото слишком большое (лимит Telegram на скачивание — 20 МБ). Пришли файл меньшего размера.")
            return
        file_info = await bot.get_file(photo.file_id)
        downloaded_file = await bot.download_file(file_info.file_path)
        img_bytes = downloaded_file.read()
    except Exception as e:
        logging.error(f"Не удалось скачать фото от пользователя {user_id}: {e}")
        await message.answer("⚠️ Не удалось загрузить фото. Проверь соединение и попробуй отправить ещё раз.")
        return

    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    # Если пользователь в режиме анализа документа — считаем фото сканом страницы
    # и распознаём текст через vision-модель вместо обычного анализа изображения.
    if user_id in doc_analysis_states:
        busy_users.add(user_id)
        status_msg = await message.answer("📚 Распознаю текст на фото...")
        try:
            base64_image = base64.b64encode(img_bytes).decode('utf-8')
            response = await call_groq_with_retry(
                model=VISION_MODEL,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Распознай и выведи весь текст, присутствующий на этом изображении, без комментариев."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                    ]
                }]
            )
            recognized_text = clean_text_for_html(response.choices[0].message.content)
            await process_document_context(message, "фото документа", recognized_text, status_msg)
        except Exception as e:
            logging.error(f"Ошибка распознавания текста с фото: {e}", exc_info=True)
            try:
                await status_msg.edit_text("⚠️ Не удалось распознать текст на фото. Попробуй ещё раз.")
            except Exception:
                pass
            doc_analysis_states.discard(user_id)
        finally:
            busy_users.discard(user_id)
        return

    # Обычный режим анализа картинки через Vision ИИ
    busy_users.add(user_id)
    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        base64_image = base64.b64encode(img_bytes).decode('utf-8')
        response = await call_groq_with_retry(
            model=VISION_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": message.caption or "Подробно проанализируй это изображение, распознай текст если он есть и ответь по сути."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }]
        )
        reply = clean_text_for_html(response.choices[0].message.content)
        user_data = get_user_data(user_id)
        user_data["last_output"] = reply
        await safe_answer(message, f"{reply}{AD_FOOTER}", parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
    except Exception as e:
        logging.error(f"Ошибка обработки фото: {e}")
        await message.answer("⚠️ Ошибка при обработке изображения. Попробуй ещё раз. Если повторится, пиши - @mecau")
    finally:
        busy_users.discard(user_id)
        gc.collect()

@dp.message(F.text)
async def handle_text(message: Message):
    user_id = message.from_user.id

    # Естественные команды: пользователь может запускать функции словами,
    # не изучая меню.
    normalized = (message.text or "").lower().strip()

    # MecauAI 2.0: команды редактирования уже созданного файла и перенос данных между форматами.
    if await handle_artifact_command(message, normalized):
        return

    if normalized in {"/clear", "очисти контекст", "забудь предыдущий диалог", "начать заново"}:
        user_data = get_user_data(user_id)
        user_data["history"] = []
        user_data["last_output"] = "Здесь пока нет ответов."
        user_doc_context.pop(user_id, None)
        await message.answer("🧹 Готово. Контекст очищен — можем начать заново.")
        return

    if normalized in {"/help", "помощь", "что ты умеешь", "что умеешь"}:
        await cmd_capabilities(message)
        return

    if re.search(r"\b(сделай|создай|подготовь|сформируй)\b.*\bпрезентаци", normalized):
        topic = re.sub(r"^.*?\bпрезентаци(?:ю|я|и|й)?\b\s*(?:на тему|по теме|про|о)?\s*", "", message.text, flags=re.I).strip(" .:-")
        user_ppt_topic[user_id] = topic or "Тема не указана"
        await cmd_ppt_prompt(message)
        return

    if re.search(r"\b(сделай|создай|подготовь|сформируй)\b.*\b(excel|таблиц|xlsx)", normalized):
        topic = re.sub(r"^.*?\b(?:excel|таблиц(?:у|а|ы)?|xlsx)\b\s*(?:на тему|по теме|для|с)?\s*", "", message.text, flags=re.I).strip(" .:-")
        user_excel_topic[user_id] = topic or message.text
        await cmd_excel_prompt(message)
        return

    if user_id in excel_states:
        topic = (message.text or "").strip()
        excel_states.discard(user_id)
        if topic:
            await generate_excel_file(message, topic, user_id=user_id)
        else:
            await message.answer("📊 Опиши, какую таблицу создать.")
        return

    if user_id == MY_ADMIN_ID and user_id in broadcast_states:
        broadcast_states.remove(user_id)
        users = load_user_ids()
        success, failed = 0, 0
        status_msg = await message.answer("📢 Начинаю рассылку...")
        for uid in users:
            try:
                await bot.send_message(uid, message.text, disable_web_page_preview=True)
                success += 1
            except TelegramRetryAfter as e:
                # Telegram просит подождать из-за лимита скорости — ждём и пробуем ещё раз
                await asyncio.sleep(e.retry_after)
                try:
                    await bot.send_message(uid, message.text, disable_web_page_preview=True)
                    success += 1
                except Exception:
                    failed += 1
            except TelegramForbiddenError:
                # Пользователь заблокировал бота — это не сбой, просто пропускаем
                failed += 1
            except Exception:
                failed += 1
            await asyncio.sleep(0.05)
        await status_msg.edit_text(f"✅ Рассылка завершена!\n\n👥 Доставлено: {success}\n❌ Ошибок: {failed}")
        return

    # ----------------- ГЕНЕРАЦИЯ ПРЕЗЕНТАЦИИ -----------------
    if user_id in ppt_states:
        ppt_states.discard(user_id)
        topic = (user_ppt_topic.get(user_id) or message.text or "").strip()
        if not topic:
            await message.answer("📈 Напиши тему презентации одним сообщением.")
            return
        await generate_presentation_file(message, topic)
        return

    # ----------------- ГЕНЕРАЦИЯ EXCEL -----------------
    if user_id in excel_states:
        excel_states.discard(user_id)
        topic = (user_excel_topic.get(user_id) or message.text or "").strip()
        if not topic:
            await message.answer("📊 Опиши, какую таблицу создать: столбцы, период или пример данных.")
            return
        await generate_excel_file(message, topic)
        return

    if message.text in MENU_BUTTONS:
        return

    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return

    save_user_id(user_id)
    if user_id in busy_users:
        await message.answer("⏳ Я ещё обрабатываю предыдущий запрос. Подожди немного, чтобы задачи не смешались.")
        return
    busy_users.add(user_id)
    user_data = get_user_data(user_id)
    current_mode = user_data["mode"]

    system_prompt = PROMPTS.get(current_mode, PROMPTS["ai"]) + (
        "\n\nПРАВИЛА ОТВЕТА: "
        "отвечай естественно и полезно. Сначала дай прямой ответ, затем при необходимости "
        "короткое объяснение, пример или конкретные следующие шаги. "
        "Не повторяй вопрос пользователя. Не добавляй пустые вступления. "
        "Если запрос можно решить практическим способом — предложи готовый результат, а не только совет. "
        "Если информации недостаточно — задай один точный уточняющий вопрос."
    ) + "\n\nОтвечай содержательно: если вопрос требует объяснения, дай краткий контекст, основной ответ и 1-2 практических примера или следующих шага. Не растягивай ответ без необходимости."

    # Если ранее был загружен документ через «📚 Анализ документа» — добавляем
    # его содержимое в системный промпт, чтобы отвечать с учётом контекста.
    doc_context = user_doc_context.get(user_id)
    if doc_context:
        system_prompt = (
            f"{system_prompt}\n\n"
            f"Пользователь ранее загрузил документ. Используй его содержимое, если вопрос с ним связан:\n"
            f"---\n{doc_context}\n---"
        )

    user_data["history"].append({"role": "user", "content": message.text})
    if len(user_data["history"]) > 6:
        user_data["history"] = user_data["history"][-6:]

    messages_payload = [{"role": "system", "content": system_prompt}] + user_data["history"]

    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        response = await call_groq_with_retry(messages=messages_payload)
        ai_reply = clean_text_for_html(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply

        if current_mode == "coder":
            full_message = f"{ai_reply}\n\n{AD_FOOTER}"
            await safe_answer(message, full_message, parse_mode="Markdown", reply_markup=get_answer_inline_keyboard())
        else:
            full_message = f"{ai_reply}{AD_FOOTER}"
            await safe_answer(message, full_message, parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
    except Exception as e:
        logging.error(f"Ошибка обычного запроса пользователя {user_id}: {e}", exc_info=True)
        await message.answer("⚠️ Не удалось обработать запрос. Попробуй ещё раз — предыдущий запрос уже освобождён.")
    finally:
        busy_users.discard(user_id)
        gc.collect()

@dp.errors()
async def global_error_handler(event: types.ErrorEvent):
    """
    Глобальный перехватчик непредвиденных ошибок. Раньше необработанное
    исключение в любом хэндлере могло привести к тому, что пользователь
    просто не получал ответа и не понимал, что произошло. Теперь бот
    логирует ошибку и сообщает пользователю, что что-то пошло не так.
    """
    logging.error(f"Необработанная ошибка: {event.exception}", exc_info=True)
    try:
        update = event.update
        if update.message:
            await update.message.answer("⚠️ Произошла непредвиденная ошибка. Попробуй ещё раз позже. Если повторится — пиши @mecau")
        elif update.callback_query:
            await update.callback_query.answer("⚠️ Произошла ошибка, попробуй ещё раз.", show_alert=True)
    except Exception:
        pass
    return True


# ======================= MECAUAI 2.0 QUALITY PASS =======================

def _safe_remove_shape(shape):
    try:
        sp = shape._element
        sp.getparent().remove(sp)
    except Exception:
        pass

async def generate_ai_image(session: aiohttp.ClientSession, prompt: str, width: int = 900, height: int = 650):
    """Надёжный image pipeline: timeout -> retry -> валидация -> уменьшение RAM."""
    prompt = re.sub(r"\s+", " ", str(prompt or "")).strip()[:900]
    if not prompt:
        return None
    for attempt in range(3):
        try:
            encoded = urllib.parse.quote(prompt)
            url = f"https://image.pollinations.ai/prompt/{encoded}?width={width}&height={height}&nologo=true&model=flux"
            timeout = aiohttp.ClientTimeout(total=22 + attempt * 8, connect=8)
            async with session.get(url, timeout=timeout) as resp:
                if resp.status != 200:
                    raise RuntimeError(f"image HTTP {resp.status}")
                data = await resp.read()
                if len(data) < 10_000:
                    raise ValueError("изображение подозрительно маленькое")
                try:
                    from PIL import Image
                    im = Image.open(io.BytesIO(data))
                    im.verify()
                    im = Image.open(io.BytesIO(data))
                    if im.width < 400 or im.height < 300:
                        raise ValueError("низкое разрешение")
                except ImportError:
                    pass
                except Exception as e:
                    raise ValueError(f"invalid image: {e}")
                return data
        except Exception as e:
            logging.warning("Image pipeline attempt %s/3: %s", attempt + 1, e)
            if attempt < 2:
                await asyncio.sleep(1.2 * (2 ** attempt))
    return None


def _ppt_overflow_check(slide):
    # Heuristic overflow guard: shrink dense text instead of letting it spill outside cards.
    for shape in slide.shapes:
        if not hasattr(shape,'text_frame') or not shape.text_frame: continue
        text=shape.text_frame.text or ''
        if len(text)>260:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        cur=r.font.size.pt if r.font.size else 14
                        r.font.size=PptxPt(max(9,cur-2))
                    except Exception: pass
        shape.text_frame.word_wrap=True

def _ppt_add_text(slide, text, x, y, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.clear()
    p = tf.paragraphs[0]; p.text = str(text)[:900]; p.alignment = align
    _set_text_style(p, size, color or RGBColor(0x21,0x21,0x21), bold)
    return box

def _ppt_add_card(slide, x, y, w, h, title, value, theme):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5,0xF7,0xFA); card.line.color.rgb = RGBColor(0xE0,0xE5,0xEB)
    _ppt_add_text(slide, title, x+0.18, y+0.18, w-0.36, 0.35, 11, False, theme["body_text"])
    _ppt_add_text(slide, value, x+0.18, y+0.58, w-0.36, h-0.7, 23, True, theme["primary"])

def build_advanced_ppt_slide(prs, theme, idx, item, img_stream=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = theme["bg"]
    layout = str(item.get("layout") or "bullets").lower()
    title = str(item.get("title") or f"Слайд {idx+1}").strip()[:120]
    _ppt_add_text(s, title, 0.72, 0.52, 11.7, 0.75, 27 if len(title)<65 else 22, True, theme["body_text"])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(13.333), PptxInches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = theme["primary"]; bar.line.fill.background()
    points = [str(x).strip() for x in (item.get("points") or []) if str(x).strip()][:6]
    if layout in ("stats", "statistics", "kpi"):
        stats = item.get("stats") or []
        if not stats:
            stats = [{"label": p.split(":",1)[0][:30], "value": p.split(":",1)[1][:30] if ":" in p else p[:30]} for p in points]
        for i, st in enumerate(stats[:4]):
            _ppt_add_card(s, 0.75 + i*3.05, 1.75, 2.7, 2.1, st.get("label","Показатель"), st.get("value","—"), theme)
        if item.get("note"):
            _ppt_add_text(s, item["note"], 0.85, 4.35, 11.5, 0.9, 16, False, theme["body_text"])
    elif layout in ("comparison", "compare"):
        cols = item.get("columns") or []
        if len(cols) < 2:
            cols = [{"title":"Вариант A","points":points[:3]}, {"title":"Вариант B","points":points[3:6]}]
        for ci, col in enumerate(cols[:3]):
            x = 0.7 + ci*4.15
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PptxInches(x),PptxInches(1.65),PptxInches(3.75),PptxInches(4.7))
            card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF7,0xF8,0xFA); card.line.color.rgb = theme["primary"]
            _ppt_add_text(s,col.get("title","Вариант"),x+0.25,1.95,3.25,0.55,18,True,theme["primary"])
            yy=2.7
            for p in (col.get("points") or [])[:5]:
                _ppt_add_text(s,"• "+str(p),x+0.25,yy,3.2,0.62,13,False,theme["body_text"]); yy+=0.78
    elif layout in ("timeline", "process"):
        steps = item.get("steps") or points
        n=max(1,min(len(steps),6)); y=3.05
        line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,PptxInches(1.0),PptxInches(y+0.18),PptxInches(11.2),PptxInches(0.08))
        line.fill.solid(); line.fill.fore_color.rgb=theme["primary"]; line.line.fill.background()
        for i,step in enumerate(steps[:n]):
            x=1.0 + (10.8/(max(n-1,1)))*i
            circ=s.shapes.add_shape(MSO_SHAPE.OVAL,PptxInches(x),PptxInches(y),PptxInches(0.48),PptxInches(0.48)); circ.fill.solid(); circ.fill.fore_color.rgb=theme["primary"]; circ.line.fill.background()
            p=circ.text_frame.paragraphs[0]; p.text=str(i+1); p.alignment=PP_ALIGN.CENTER; _set_text_style(p,10,theme["title_text"],True)
            _ppt_add_text(s,str(step),x-0.35,y+0.75,1.25,1.3,11,False,theme["body_text"],PP_ALIGN.CENTER)
    elif layout in ("quote", "citation"):
        quote = str(item.get("quote") or (points[0] if points else "Главный тезис темы."))[:500]
        _ppt_add_text(s,"“",0.95,1.65,1.0,1.2,60,True,theme["primary"])
        _ppt_add_text(s,quote,1.75,2.05,9.9,2.2,25,True,theme["body_text"])
        if item.get("author"):
            _ppt_add_text(s,"— "+str(item["author"]),1.8,4.65,8.5,0.5,14,False,theme["primary"])
    elif layout in ("conclusion", "summary"):
        for i,p in enumerate(points[:4]):
            _ppt_add_card(s,0.8+(i%2)*6.1,1.65+(i//2)*2.25,5.55,1.85,f"Вывод {i+1}",p,theme)
    elif layout in ("chart", "graph") and (item.get("labels") and item.get("values")):
        data=ChartData(); data.categories=[str(x)[:30] for x in item["labels"][:8]]; data.add_series(str(item.get("series") or "Значение"), [float(x) for x in item["values"][:8]])
        chart=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,PptxInches(0.75),PptxInches(1.55),PptxInches(11.8),PptxInches(4.95),data).chart
        chart.has_legend=False; chart.has_title=True; chart.chart_title.text_frame.text=str(item.get("chart_title") or title)
    else:
        has_img = img_stream is not None
        left_w=7.15 if has_img else 11.7; yy=1.55
        for i,p in enumerate(points or ["Ключевой тезис будет сформирован автоматически."]):
            h=0.7 if len(p)<110 else 0.9
            card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PptxInches(0.68),PptxInches(yy),PptxInches(left_w),PptxInches(h))
            card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0xF5,0xF7,0xFA); card.line.color.rgb=RGBColor(0xE2,0xE6,0xEC)
            _ppt_add_text(s,f"{i+1}",0.9,yy+0.14,0.35,0.3,9,True,theme["title_text"],PP_ALIGN.CENTER)
            badge=s.shapes[-1]; badge.fill.solid(); badge.fill.fore_color.rgb=theme["primary"]; badge.line.fill.background()
            _ppt_add_text(s,p,1.4,yy+0.1,left_w-0.95,h-0.1,14 if len(p)<120 else 12,False,theme["body_text"])
            yy+=h+0.12
        if has_img:
            frame=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PptxInches(8.15),PptxInches(1.55),PptxInches(4.5),PptxInches(4.95)); frame.fill.solid(); frame.fill.fore_color.rgb=RGBColor(0xEE,0xF1,0xF5); frame.line.color.rgb=RGBColor(0xD9,0xDE,0xE5)
            try: s.shapes.add_picture(img_stream,left=PptxInches(8.28),top=PptxInches(1.68),width=PptxInches(4.24),height=PptxInches(4.69))
            except Exception as e: logging.warning("PPT image insert failed: %s",e)
    _ppt_overflow_check(s)
    _add_footer(s,theme,idx+2)
    return s

async def _ppt_structure(topic, num_slides):
    response=await call_groq_with_retry(messages=[
        {"role":"system","content":f"Создай профессиональную презентацию на русском языке ровно из {num_slides} содержательных слайдов после титульного. Верни только JSON-массив. Каждый объект: title, layout, points, image_prompt, stats, columns, steps, quote, author, labels, values, series, chart_title. layout выбирай из bullets, stats, comparison, timeline, process, quote, chart, conclusion. Не используй chart без labels/values. Слайды должны покрывать введение, ключевые идеи, данные/пример, сравнение/процесс, риски и итог. Тезисы короткие. image_prompt пустой для stats/comparison/timeline/chart/conclusion."},
        {"role":"user","content":f"Тема: {topic[:2500]}"}
    ],temperature=0.35,max_retries=2,timeout=50)
    data=extract_json(clean_text_for_html(response.choices[0].message.content))
    if not isinstance(data,list): raise ValueError("Некорректная структура PPT")
    return data[:num_slides]

async def generate_presentation_file(message: Message, topic: str, user_id: int = None):
    user_id=user_id or message.from_user.id
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется..."); return
    busy_users.add(user_id); heavy_lock=None; status_msg=await message.answer("📈 Анализирую содержание → подбираю макеты → проверяю визуалы…")
    try:
        heavy_lock=await acquire_heavy_job(status_msg)
        theme_key=user_ppt_design.pop(user_id,"ppt_blue"); theme=PPT_THEMES.get(theme_key,PPT_THEMES["ppt_blue"])
        custom=user_ppt_images.pop(user_id,[])[:MAX_PPT_IMAGES]
        n=7 if not custom else max(7,min(10,len(custom)+5))
        slides=await _ppt_structure(topic,n)
        while len(slides)<n:
            slides.append({"title":"Выводы","layout":"conclusion","points":["Главный результат","Практическое значение","Следующий шаг"]})
        prs=Presentation(); prs.slide_width=PptxInches(13.333); prs.slide_height=PptxInches(7.5); build_title_slide(prs,theme,topic)
        async with aiohttp.ClientSession() as session:
            for idx,item in enumerate(slides):
                stream=None
                if idx<len(custom):
                    try:
                        fi=await bot.get_file(custom[idx]); d=await bot.download_file(fi.file_path); b=d.read()
                        if len(b)<=MAX_DOC_SIZE_BYTES: stream=io.BytesIO(b)
                    except Exception as e: logging.warning("custom PPT image %s: %s",idx,e)
                elif item.get("image_prompt") and idx<MAX_PPT_AI_IMAGES:
                    b=await generate_ai_image(session,item["image_prompt"],900,650)
                    if b: stream=io.BytesIO(b)
                build_advanced_ppt_slide(prs,theme,idx,item,stream)
        bio=io.BytesIO(); prs.save(bio); data=bio.getvalue()
        if len(data)>MAX_ARTIFACT_BYTES: raise ValueError("PPT слишком большой")
        path=_artifact_path(user_id,"pptx",topic); open(path,"wb").write(data); register_artifact(user_id,"pptx",path,{"topic":topic,"theme":theme_key,"slides":len(prs.slides)})
        await status_msg.delete(); await message.answer_document(BufferedInputFile(data,filename=os.path.basename(path)),caption=f"📈 Готово: {len(prs.slides)} слайдов • {theme['label']}\n✏️ Можно сразу написать, что изменить.")
        get_user_stats(user_id)["exports"]+=1
    except Exception as e:
        logging.error("PPT 2.0 error: %s",e,exc_info=True)
        try: await status_msg.edit_text("⚠️ Не удалось собрать презентацию. Визуальный блок автоматически пропущен; повтори запрос ещё раз.")
        except Exception: pass
    finally:
        busy_users.discard(user_id); ppt_states.discard(user_id); user_ppt_topic.pop(user_id,None); user_ppt_images.pop(user_id,None); user_ppt_design.pop(user_id,None)
        if heavy_lock: heavy_lock.release()
        gc.collect()

def _excel_detect_format(header, values):
    h=str(header).lower()
    vals=[v for v in values if v not in (None,"")]
    if any(x in h for x in ("%","процент","доля")): return '0.0%'
    if any(x in h for x in ("€","eur","евро","руб","₽","цена","стоим","расход","доход","выруч")): return '#,##0.00'
    if vals and all(isinstance(v,(int,float)) for v in vals): return '#,##0.00'
    return 'General'

async def generate_excel_file(message: Message, topic: str, user_id: int = None):
    user_id=user_id or message.from_user.id
    if user_id in busy_users: await message.answer("⏳ Подожди, предыдущая задача ещё выполняется..."); return
    busy_users.add(user_id); heavy_lock=None; status=await message.answer("📊 Строю данные → формулы → KPI → Dashboard → диаграммы…")
    try:
        heavy_lock=await acquire_heavy_job(status)
        response=await call_groq_with_retry(messages=[
            {"role":"system","content":"Верни только JSON: {title,sheet_name,headers,rows,formulas}. headers 2-12, rows 5-40. Числа числами, даты YYYY-MM-DD. formulas — необязательный список объектов {cell,formula}. Если уместно, данные должны поддерживать KPI и график."},
            {"role":"user","content":topic[:4000]}],temperature=0.2)
        data=extract_json(clean_text_for_html(response.choices[0].message.content)); headers=[str(x).strip() for x in data.get('headers',[]) if str(x).strip()][:12]
        if not headers: raise ValueError('Нет заголовков')
        rows=[]
        for row in (data.get('rows') or [])[:40]:
            r=list(row) if isinstance(row,(list,tuple)) else [row]; rows.append((r+[""]*len(headers))[:len(headers)])
        wb=Workbook(); ws=wb.active; ws.title=re.sub(r'[\\/*?:\[\]]','_',str(data.get('sheet_name') or 'Данные'))[:31] or 'Данные'
        ws.merge_cells(start_row=1,start_column=1,end_row=1,end_column=len(headers)); c=ws.cell(1,1,str(data.get('title') or topic)[:120]); c.font=Font(name='Aptos Display',size=18,bold=True,color='FFFFFF'); c.fill=PatternFill('solid',fgColor='1F4E78'); ws.row_dimensions[1].height=30
        ws.append(headers)
        for cell in ws[2]: cell.font=Font(bold=True,color='FFFFFF'); cell.fill=PatternFill('solid',fgColor='2F75B5'); cell.alignment=Alignment(horizontal='center',wrap_text=True)
        for r in rows: ws.append(r)
        for ci,h in enumerate(headers,1):
            vals=[ws.cell(r,ci).value for r in range(3,ws.max_row+1)]
            fmt=_excel_detect_format(h,vals)
            for r in range(3,ws.max_row+1):
                cell=ws.cell(r,ci); cell.alignment=Alignment(vertical='center',wrap_text=True); cell.number_format=fmt
            maxlen=max([len(str(h))]+[len(str(v or '')) for v in vals]); ws.column_dimensions[get_column_letter(ci)].width=min(max(maxlen+2,12),34)
        ws.freeze_panes='A3'; ws.auto_filter.ref=f'A2:{get_column_letter(len(headers))}{ws.max_row}'
        if ws.max_row>=3:
            tab=Table(displayName='DataTable',ref=f'A2:{get_column_letter(len(headers))}{ws.max_row}'); tab.tableStyleInfo=TableStyleInfo(name='TableStyleMedium2',showRowStripes=True); ws.add_table(tab)
        # Conditional formatting for numeric columns.
        for ci,h in enumerate(headers,1):
            vals=[ws.cell(r,ci).value for r in range(3,ws.max_row+1)]
            if vals and sum(isinstance(v,(int,float)) for v in vals)/len(vals)>=0.7:
                rng=f'{get_column_letter(ci)}3:{get_column_letter(ci)}{ws.max_row}'
                ws.conditional_formatting.add(rng,ColorScaleRule(start_type='min',start_color='F2F2F2',mid_type='percentile',mid_value=50,mid_color='D9EAF7',end_type='max',end_color='5B9BD5'))
        for f in data.get('formulas') or []:
            try: ws[str(f['cell'])]=str(f['formula'])
            except Exception: pass
        # Dashboard with multiple chart types and KPI formulas.
        numeric=[]
        for ci in range(2,len(headers)+1):
            vals=[ws.cell(r,ci).value for r in range(3,ws.max_row+1)]
            if vals and sum(isinstance(v,(int,float)) for v in vals)/len(vals)>=0.7: numeric.append(ci)
        dash=wb.create_sheet('Dashboard'); dash['A1']=str(data.get('title') or topic)[:100]; dash['A1'].font=Font(size=20,bold=True)
        for i,ci in enumerate(numeric[:4]):
            col=get_column_letter(ci); x=1+(i%2)*3; y=3+(i//2)*4
            dash.cell(y,x,'Итого').font=Font(bold=True); dash.cell(y,x+1,f'=SUM(\'{ws.title}\'!{col}3:{col}{ws.max_row})')
            dash.cell(y+1,x,'Среднее').font=Font(bold=True); dash.cell(y+1,x+1,f'=AVERAGE(\'{ws.title}\'!{col}3:{col}{ws.max_row})')
        if numeric:
            chart=BarChart(); chart.type='col'; chart.style=10; chart.title=f'{headers[numeric[0]-1]} по {headers[0]}'; chart.add_data(Reference(ws,min_col=numeric[0],max_col=numeric[0],min_row=2,max_row=ws.max_row),titles_from_data=True); chart.set_categories(Reference(ws,min_col=1,min_row=3,max_row=ws.max_row)); chart.width=18; chart.height=9; dash.add_chart(chart,'A12')
        if len(numeric)>1:
            from openpyxl.chart import LineChart
            line=LineChart(); line.title=f'Динамика: {headers[numeric[1]-1]}'; line.add_data(Reference(ws,min_col=numeric[1],max_col=numeric[1],min_row=2,max_row=ws.max_row),titles_from_data=True); line.set_categories(Reference(ws,min_col=1,min_row=3,max_row=ws.max_row)); line.width=18; line.height=9; dash.add_chart(line,'J12')
        bio=io.BytesIO(); wb.save(bio); raw=bio.getvalue(); path=_artifact_path(user_id,'xlsx',topic); open(path,'wb').write(raw); register_artifact(user_id,'xlsx',path,{'topic':topic,'sheet':ws.title})
        await status.delete(); await message.answer_document(BufferedInputFile(raw,filename=os.path.basename(path)),caption=f'📊 Excel готов: {len(rows)} строк, {len(headers)} столбцов.\n✏️ Можно написать, что изменить.')
        get_user_stats(user_id)['exports']+=1
    except Exception as e:
        logging.error('Excel 2.0 error: %s',e,exc_info=True)
        try: await status.edit_text('⚠️ Excel не удалось собрать. Попробуй указать столбцы, период и единицы измерения.')
        except Exception: pass
    finally:
        busy_users.discard(user_id); excel_states.discard(user_id); user_excel_topic.pop(user_id,None); user_excel_mode.pop(user_id,None)
        if heavy_lock: heavy_lock.release()
        gc.collect()

def _artifact_text(path, limit=14000):
    ext=os.path.splitext(path)[1].lower()
    if ext=='.xlsx':
        wb=__import__('openpyxl').load_workbook(path,read_only=True,data_only=False)
        out=[]
        for ws in wb.worksheets[:3]:
            out.append(f'Лист: {ws.title}')
            for row in ws.iter_rows(min_row=1,max_row=25,values_only=True): out.append(' | '.join(str(v or '') for v in row))
        return '\n'.join(out)[:limit]
    if ext=='.docx':
        d=Document(path); return '\n'.join(p.text for p in d.paragraphs)[:limit]
    if ext=='.pptx':
        prs=Presentation(path); out=[]
        for i,s in enumerate(prs.slides):
            texts=[sh.text for sh in s.shapes if hasattr(sh,'text') and sh.text.strip()]
            out.append(f'Slide {i+1}: '+' | '.join(texts))
        return '\n'.join(out)[:limit]
    return ''

async def _send_artifact(message,user_id,path,caption='Обновлённый файл'):
    with open(path,'rb') as f: raw=f.read()
    await message.answer_document(BufferedInputFile(raw,filename=os.path.basename(path)),caption=caption)

def _edit_ppt_local(path,instruction):
    prs=Presentation(path); low=instruction.lower()
    if re.search(r'\b(убери|удали).*(график|диаграмм)',low):
        for s in prs.slides:
            for sh in list(s.shapes):
                if getattr(sh,'has_chart',False): _safe_remove_shape(sh)
    m=re.search(r'(?:слайд|slide)\s*(\d+)',low); idx=int(m.group(1))-1 if m else None
    if 'добавь слайд' in low or 'добавить слайд' in low:
        s=prs.slides.add_slide(prs.slide_layouts[6]); _ppt_add_text(s,'Выводы',0.8,0.8,11.5,0.7,30,True); _ppt_add_text(s,'• Ключевой результат\n• Практическое значение\n• Следующий шаг',0.9,2.0,10.8,2.5,20); _add_footer(s,PPT_THEMES['ppt_blue'],len(prs.slides));
    if 'подробнее' in low and idx is not None and 0<=idx<len(prs.slides):
        s=prs.slides[idx]; _ppt_add_text(s,'Дополнение',0.8,5.95,2.0,0.35,12,True,PPT_THEMES['ppt_blue']['primary']); _ppt_add_text(s,'Дополнительный контекст, пример и практическое значение по запросу пользователя.',2.3,5.82,9.7,0.7,12)
    if re.search(r'строг(ий|им)|официальн',low):
        theme=PPT_THEMES['ppt_dark']
        for s in prs.slides:
            s.background.fill.solid(); s.background.fill.fore_color.rgb=theme['bg']
            for sh in s.shapes:
                if hasattr(sh,'text_frame'):
                    for p in sh.text_frame.paragraphs:
                        for r in p.runs: r.font.name='Aptos'
    prs.save(path)

async def _edit_excel_local(path,instruction):
    from openpyxl import load_workbook
    wb=load_workbook(path); ws=wb[wb.sheetnames[0]]; low=instruction.lower()
    if 'отсорт' in low:
        # Sort by a named column when present; otherwise by first numeric column.
        col=None
        for i,c in enumerate(ws[2],1):
            if c.value and str(c.value).lower() in low: col=i; break
        if col is None:
            for i in range(2,ws.max_column+1):
                vals=[ws.cell(r,i).value for r in range(3,ws.max_row+1)]
                if vals and any(isinstance(v,(int,float)) for v in vals): col=i; break
        if col:
            rows=list(ws.iter_rows(min_row=3,values_only=True)); rows.sort(key=lambda r:(r[col-1] is None,r[col-1]),reverse='по убыванию' in low or 'убыв' in low)
            for rr,row in enumerate(rows,3):
                for cc,val in enumerate(row,1): ws.cell(rr,cc).value=val
    if 'добавь' in low and ('март' in low or 'апрел' in low or 'май' in low):
        period=next((x for x in ('март','апрель','май','июнь','июль','август','сентябрь','октябрь','ноябрь','декабрь','январь','февраль') if x in low),'новая строка')
        ws.append([period]+['']*(ws.max_column-1))
    if 'диаграмм' in low or 'график' in low:
        from openpyxl.chart import BarChart, Reference
        numeric=2 if ws.max_column>=2 else 1; ch=BarChart(); ch.type='col'; ch.title='Обновлённая диаграмма'; ch.add_data(Reference(ws,min_col=numeric,max_col=numeric,min_row=2,max_row=ws.max_row),titles_from_data=True); ch.set_categories(Reference(ws,min_col=1,min_row=3,max_row=ws.max_row));
        target=wb['Dashboard'] if 'Dashboard' in wb.sheetnames else wb.create_sheet('Dashboard'); target.add_chart(ch,'A30')
    wb.save(path)

async def _edit_word_local(path,instruction):
    d=Document(path); low=instruction.lower()
    if 'заголов' in low or 'структур' in low:
        for p in d.paragraphs:
            if p.text.strip() and len(p.text)<90 and not p.style.name.startswith('Heading'):
                p.style='Heading 2'
    if 'подробнее' in low:
        p=d.add_paragraph('Дополнение: дополнительный контекст, пример и практическое применение по запросу пользователя.')
        p.style='Normal'
    d.save(path)

async def handle_artifact_command(message, normalized):
    user_id=message.from_user.id
    art=last_artifact(user_id)
    if not art: return False
    # Only intercept explicit edit/convert requests, not ordinary conversation.
    edit_words=('измени','изменить','добавь','добавить','убери','удали','удалить','отсортируй','сделай подробнее','сделай более','построй','поставь')
    convert_words=('сделай из него презентацию','сделай из неё презентацию','сделай из него excel','сделай из него word','сделай краткий word','сделай презентацию из excel')
    if not any(w in normalized for w in edit_words+convert_words): return False
    if any(w in normalized for w in convert_words):
        text=_artifact_text(art['path'])
        if 'презента' in normalized:
            await generate_presentation_file(message, 'Создай презентацию по данным последнего файла:\n'+text[:9000], user_id); return True
        if 'excel' in normalized:
            await generate_excel_file(message, 'Создай Excel по данным последнего файла:\n'+text[:9000], user_id); return True
        if 'word' in normalized:
            path=_artifact_path(user_id,'docx','Report'); d=Document();
            for line in text.splitlines()[:80]:
                p=d.add_paragraph(line); p.style='Normal'
            d.save(path); register_artifact(user_id,'docx',path,{'source':art['path']}); await _send_artifact(message,user_id,path,'📄 Word создан из последнего файла.'); return True
    if not art['path'] or not os.path.exists(art['path']): return False
    try:
        if art['kind']=='pptx': _edit_ppt_local(art['path'],normalized)
        elif art['kind']=='xlsx': await _edit_excel_local(art['path'],normalized)
        elif art['kind']=='docx': await _edit_word_local(art['path'],normalized)
        else: return False
        art['created']=time.time(); _save_artifacts(); await _send_artifact(message,user_id,art['path'],'✏️ Готово — изменил уже созданный файл, не начиная всё с нуля.')
        return True
    except Exception as e:
        logging.error('Artifact edit failed: %s',e,exc_info=True); await message.answer('⚠️ Не удалось безопасно изменить файл. Исходный файл сохранён, попробуй описать изменение конкретнее.'); return True

async def main():
    await bot.set_my_commands([
        BotCommand(command="start", description="Перезапустить бота"),
        BotCommand(command="mode",  description="Сменить режим ассистента"),
        BotCommand(command="about", description="О возможностях"),
    ])
    await bot.delete_webhook()
    maintenance = asyncio.create_task(maintenance_loop())
    print("🚀 Бот MecauAI 2.0 запущен: очередь, watchdog, артефакты и редактор активны...")
    try:
        await dp.start_polling(bot)
    finally:
        maintenance.cancel()
        try:
            await maintenance
        except asyncio.CancelledError:
            pass

if __name__ == "__main__":
    asyncio.run(main())
