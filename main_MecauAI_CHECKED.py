import asyncio
import logging
import base64
import html
import io
import json
import os
import re
import time
import gc
import threading
from datetime import datetime, timezone
import urllib.parse
import aiohttp
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import CommandStart, Command
from aiogram.exceptions import TelegramRetryAfter, TelegramForbiddenError, TelegramBadRequest
from aiogram.types import (
    InlineKeyboardMarkup, InlineKeyboardButton, CallbackQuery,
    ReplyKeyboardMarkup, KeyboardButton, Message, BotCommand, BufferedInputFile, FSInputFile
)
from openai import AsyncOpenAI
from docx import Document
from docx.shared import Pt, Inches, Mm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from openpyxl.worksheet.table import Table, TableStyleInfo
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.utils import get_column_letter

# Для вставки картинок в Excel дополнительно нужен пакет Pillow:
#   pip install pillow
# Для чтения PDF в разделе «Анализ документа» нужен пакет pypdf:
#   pip install pypdf
# Оба импортируются лениво (внутри try/except), чтобы их отсутствие
# не ломало весь бот, а просто отключало соответствующую функцию.

from config import (
    BOT_TOKEN, GROQ_API_KEY,
    CHANNEL_1_USERNAME, CHANNEL_1_URL,
    CHANNEL_2_USERNAME, CHANNEL_2_URL,
    TEXT_MODEL, VISION_MODEL, PROMPTS, AD_FOOTER
)

# Groq model resilience:
# the configured model can be unavailable for a particular project/key even when
# the model still exists in Groq's catalog. We therefore fall back automatically
# instead of showing a generic AI error to the user.
TEXT_MODEL_FALLBACKS = [
    os.getenv("MECAUAI_TEXT_MODEL", "").strip() or TEXT_MODEL,
    "openai/gpt-oss-120b",
    "qwen/qwen3.6-27b",
    "openai/gpt-oss-20b",
]
# Llama 4 Scout was deprecated by Groq before this build; Qwen 3.6 27B is the
# current multimodal production-capable fallback documented by Groq.
VISION_MODEL_FALLBACKS = [
    os.getenv("MECAUAI_VISION_MODEL", "").strip() or VISION_MODEL,
    "qwen/qwen3.6-27b",
]
TEXT_MODEL_FALLBACKS = list(dict.fromkeys(x for x in TEXT_MODEL_FALLBACKS if x))
VISION_MODEL_FALLBACKS = list(dict.fromkeys(x for x in VISION_MODEL_FALLBACKS if x))

MY_ADMIN_ID = int(os.getenv("MECAUAI_ADMIN_ID", "1184589026"))
# Единый набор администраторов для всех новых admin-функций.
# Дополнительные ID можно передать через MECAUAI_ADMIN_IDS=123,456,789.
try:
    ADMIN_IDS = {int(x.strip()) for x in os.getenv("MECAUAI_ADMIN_IDS", "").split(",") if x.strip()}
except Exception:
    ADMIN_IDS = set()
ADMIN_IDS.add(MY_ADMIN_ID)

logging.basicConfig(level=logging.INFO)
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher()

groq_client = AsyncOpenAI(
    api_key=GROQ_API_KEY,
    base_url="https://api.groq.com/openai/v1"
)

# Persistent data directory. Set MECAUAI_DATA_DIR to a mounted/persistent volume
# in the deployment platform. If omitted, use /data when available, otherwise
# keep backward compatibility with files next to main.py.
def _choose_data_dir():
    configured = os.getenv("MECAUAI_DATA_DIR", "").strip()
    if configured:
        return os.path.abspath(configured)
    if os.path.isdir("/data") and os.access("/data", os.W_OK):
        return "/data"
    return os.path.abspath(os.path.dirname(__file__) or ".")

DATA_DIR = _choose_data_dir()
os.makedirs(DATA_DIR, exist_ok=True)

# Юридические документы лежат в корне проекта рядом с main.py.
PROJECT_ROOT = os.path.abspath(os.path.dirname(__file__) or ".")
AGREEMENT_FILE = os.path.join(PROJECT_ROOT, "agreement.pdf")
PRIVACY_FILE = os.path.join(PROJECT_ROOT, "privacy.pdf")

# Admin state lives on the persistent volume, not in the ephemeral app folder.
ADMIN_IDS_FILE = os.path.join(DATA_DIR, "admin_ids.json")
try:
    # Migrate an old admin file once, if it exists next to main.py.
    _legacy_admin_file = os.path.join(os.path.abspath(os.path.dirname(__file__) or "."), "admin_ids.json")
    if not os.path.exists(ADMIN_IDS_FILE) and os.path.exists(_legacy_admin_file) and _legacy_admin_file != ADMIN_IDS_FILE:
        import shutil
        shutil.copy2(_legacy_admin_file, ADMIN_IDS_FILE)
    if os.path.exists(ADMIN_IDS_FILE):
        with open(ADMIN_IDS_FILE, "r", encoding="utf-8") as _af:
            _saved_admins = json.load(_af)
        if isinstance(_saved_admins, list):
            ADMIN_IDS.update(int(x) for x in _saved_admins if str(x).lstrip("-").isdigit())
except Exception as exc:
    logging.warning("Не удалось загрузить admin_ids.json: %s", exc)

def _save_admin_ids():
    try:
        tmp = ADMIN_IDS_FILE + ".tmp"
        with open(tmp, "w", encoding="utf-8") as _af:
            json.dump(sorted(ADMIN_IDS), _af, ensure_ascii=False, indent=2)
        os.replace(tmp, ADMIN_IDS_FILE)
    except Exception as exc:
        logging.error("Не удалось сохранить список администраторов: %s", exc)

def _data_file(name):
    return os.path.join(DATA_DIR, name)

USERS_FILE = _data_file("users.json")
FAV_FILE = _data_file("favorites.json")
MODES_FILE = _data_file("user_modes.json")
USER_DATA_FILE = _data_file("users_data.json")
STATS_FILE = _data_file("user_stats.json")

# One-time migration from the old working directory if the persistent file
# does not exist yet. Never overwrites an existing persistent file.
def _migrate_legacy_file(name):
    target = _data_file(name)
    legacy = os.path.abspath(name)
    if target != legacy and not os.path.exists(target) and os.path.exists(legacy):
        try:
            import shutil
            shutil.copy2(legacy, target)
            logging.info("Migrated legacy data file %s -> %s", legacy, target)
        except Exception as exc:
            logging.warning("Could not migrate %s: %s", name, exc)

for _name in ("users.json", "favorites.json", "user_modes.json", "users_data.json", "user_stats.json", "artifacts.json"):
    _migrate_legacy_file(_name)

all_users_cache = set()
favorites_cache = None
favorites_dirty = False
users_db_dirty = False
stats_db_dirty = False
last_gc_at = 0.0
GC_INTERVAL_SECONDS = float(os.getenv("MECAUAI_GC_INTERVAL_SECONDS", "45"))
PERSIST_FLUSH_INTERVAL = float(os.getenv("MECAUAI_PERSIST_FLUSH_INTERVAL", "15"))
users_ids_dirty = False
modes_db_dirty = False
artifacts_dirty = False
last_persist_flush_at = 0.0

def load_user_ids() -> set:
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list):
                    return set(data)
        except Exception:
            pass
    return set()

def save_user_id(user_id: int):
    global all_users_cache, users_ids_dirty
    if user_id in all_users_cache:
        return
    all_users_cache.add(user_id)
    # Первый вход сохраняем сразу: это защищает список пользователей от потери
    # при внезапном restart/rebuild. Дальше новые записи всё равно идут через
    # обычный flush.
    try:
        _atomic_json_write(USERS_FILE, sorted(all_users_cache))
        users_ids_dirty = False
    except Exception as exc:
        users_ids_dirty = True
        logging.error("Не удалось сразу сохранить нового пользователя %s: %s", user_id, exc)

all_users_cache = load_user_ids()

def load_favorites() -> dict:
    global favorites_cache
    if favorites_cache is not None:
        return favorites_cache
    if os.path.exists(FAV_FILE):
        try:
            with open(FAV_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                favorites_cache = data if isinstance(data, dict) else {}
                return favorites_cache
        except Exception:
            pass
    favorites_cache = {}
    return favorites_cache

def save_favorites(favs: dict, immediate: bool = False):
    global favorites_cache, favorites_dirty
    favorites_cache = favs
    if not immediate:
        favorites_dirty = True
        return
    try:
        _atomic_json_write(FAV_FILE, favs)
        favorites_dirty = False
    except Exception as e:
        logging.error("Не удалось сохранить избранное: %s", e, exc_info=True)

def add_to_favorites(user_id: int, text: str):
    favs = load_favorites()
    uid_str = str(user_id)
    if uid_str not in favs:
        favs[uid_str] = []
    if text not in favs[uid_str]:
        favs[uid_str].append(text)
        save_favorites(favs)

def load_user_modes() -> dict:
    if os.path.exists(MODES_FILE):
        try:
            with open(MODES_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, dict):
                    return data
        except Exception:
            pass
    return {}

def save_user_mode(user_id: int, mode: str):
    """Обновляет RAM-кэш; запись на диск выполняется пакетно maintenance_loop."""
    global saved_modes_cache, modes_db_dirty
    saved_modes_cache[str(user_id)] = mode
    modes_db_dirty = True

JSON_WRITE_LOCK = threading.Lock()

def _atomic_json_write(path, data):
    # One-process lock prevents concurrent handlers from racing on the same .tmp file.
    with JSON_WRITE_LOCK:
        tmp = path + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
            f.flush()
            os.fsync(f.fileno())
        os.replace(tmp, path)

def _load_json_dict(path):
    if os.path.exists(path):
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
            return data if isinstance(data, dict) else {}
        except Exception as exc:
            logging.warning("Не удалось загрузить %s: %s", path, exc)
    return {}

users_db = _load_json_dict(USER_DATA_FILE)
# JSON keys are strings after loading; normalize user IDs back to int.
users_db = {int(k): v for k, v in users_db.items()
            if str(k).lstrip("-").isdigit() and isinstance(v, dict)}
saved_modes_cache = load_user_modes()
user_stats = {int(k): v for k, v in _load_json_dict(STATS_FILE).items() if str(k).lstrip("-").isdigit() and isinstance(v, dict)}
ppt_states = set()
excel_states = set()
doc_analysis_states = set()
busy_users = set()
user_ppt_images = {}
user_ppt_design = {}
user_excel_mode = {}
user_ppt_topic = {}
user_ppt_slide_count = {}
user_excel_topic = {}
# Lightweight per-user workflow state. Cleared after a completed/cancelled flow.
user_task_state = {}
# Данные для титульника: город и год вводятся пользователем перед генерацией.
title_page_states = {}

user_last_request = {}
user_custom_images = {}  # temporary presentation assets: user_id -> list Telegram file_id
ppt_creation_users = set()

# Ограничители нагрузки: не даём слабому VPS одновременно запускать несколько
# тяжёлых генераций и десятки обращений к AI.
AI_SEMAPHORE = None
HEAVY_JOB_SEMAPHORE = None
MAX_AI_CONCURRENCY = 2
MAX_HEAVY_CONCURRENCY = 1
MAX_HEAVY_QUEUE = int(os.getenv("MECAUAI_MAX_HEAVY_QUEUE", "3"))
HEAVY_WAITERS = 0
MAX_PPT_IMAGES = 8
MAX_PPT_AI_IMAGES = 4

# ===== MecauAI 2.0: единый registry артефактов + лёгкий persistence =====
ARTIFACT_DIR = os.path.abspath(os.getenv("MECAUAI_ARTIFACT_DIR", os.path.join(DATA_DIR, "generated_files")))
ARTIFACTS_FILE = os.path.abspath(os.getenv("MECAUAI_ARTIFACTS_FILE", os.path.join(DATA_DIR, "artifacts.json")))
ARTIFACT_TTL_SECONDS = int(os.getenv("MECAUAI_ARTIFACT_TTL_SECONDS", str(24 * 3600)))
MAX_ARTIFACTS_PER_USER = 6
MAX_ARTIFACT_BYTES = 25 * 1024 * 1024
RESOURCE_CHECK_INTERVAL = 60

def _detect_memory_limit_mb():
    """Определяет лимит памяти контейнера/VPS, если хост его предоставляет."""
    candidates = ["/sys/fs/cgroup/memory.max", "/sys/fs/cgroup/memory/memory.limit_in_bytes"]
    for path in candidates:
        try:
            if not os.path.exists(path):
                continue
            raw = open(path, "r", encoding="utf-8").read().strip()
            if not raw or raw == "max":
                continue
            value = int(raw)
            if value <= 0 or value >= 2**60:
                continue
            return value / (1024 * 1024)
        except Exception:
            pass
    return None

_CONTAINER_RAM_MB = _detect_memory_limit_mb()
if os.getenv("MECAUAI_RAM_SOFT_LIMIT_MB"):
    RAM_SOFT_LIMIT_MB = int(os.getenv("MECAUAI_RAM_SOFT_LIMIT_MB"))
else:
    RAM_SOFT_LIMIT_MB = int(_CONTAINER_RAM_MB * 0.70) if _CONTAINER_RAM_MB else 700
if os.getenv("MECAUAI_RAM_HARD_LIMIT_MB"):
    RAM_HARD_LIMIT_MB = int(os.getenv("MECAUAI_RAM_HARD_LIMIT_MB"))
else:
    RAM_HARD_LIMIT_MB = int(_CONTAINER_RAM_MB * 0.82) if _CONTAINER_RAM_MB else 900
# Не даём watchdog сработать на совсем маленьком baseline, но и не позволяем
# дефолтному 900 MB превысить реальный лимит контейнера.
if _CONTAINER_RAM_MB:
    RAM_HARD_LIMIT_MB = max(256, min(RAM_HARD_LIMIT_MB, int(_CONTAINER_RAM_MB * 0.88)))
    RAM_SOFT_LIMIT_MB = max(192, min(RAM_SOFT_LIMIT_MB, int(_CONTAINER_RAM_MB * 0.75)))

artifact_registry = {}
artifact_lock = None

def _artifact_lock():
    global artifact_lock
    if artifact_lock is None:
        artifact_lock = asyncio.Lock()
    return artifact_lock

def _load_artifacts():
    global artifact_registry
    try:
        if os.path.exists(ARTIFACTS_FILE):
            with open(ARTIFACTS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            artifact_registry = data if isinstance(data, dict) else {}
    except Exception as e:
        logging.warning("Не удалось загрузить registry артефактов: %s", e)
        artifact_registry = {}

def _save_artifacts(immediate: bool = False):
    global artifacts_dirty
    if not immediate:
        artifacts_dirty = True
        return
    try:
        _atomic_json_write(ARTIFACTS_FILE, artifact_registry)
        artifacts_dirty = False
    except Exception as e:
        logging.warning("Не удалось сохранить registry артефактов: %s", e, exc_info=True)

_load_artifacts()
os.makedirs(ARTIFACT_DIR, exist_ok=True)

def _artifact_path(user_id, ext, stem):
    safe = re.sub(r"[^A-Za-zА-Яа-я0-9_-]+", "_", str(stem))[:45].strip("_") or "file"
    uid_dir = os.path.join(ARTIFACT_DIR, str(user_id))
    os.makedirs(uid_dir, exist_ok=True)
    return os.path.join(uid_dir, f"{int(time.time())}_{safe}.{ext}")

def register_artifact(user_id, kind, path, meta=None):
    if not path or not os.path.exists(path):
        return None
    entry = {"kind": kind, "path": os.path.abspath(path), "created": time.time(), "meta": meta or {}}
    items = artifact_registry.setdefault(str(user_id), [])
    items.append(entry)
    artifact_registry[str(user_id)] = items[-MAX_ARTIFACTS_PER_USER:]
    _save_artifacts()
    return entry

def last_artifact(user_id, kind=None):
    items = artifact_registry.get(str(user_id), [])
    for item in reversed(items):
        if kind is None or item.get("kind") == kind:
            if os.path.exists(item.get("path", "")):
                return item
    return None

def _maybe_gc(force: bool = False):
    global last_gc_at
    now = time.monotonic()
    if not force and now - last_gc_at < GC_INTERVAL_SECONDS:
        return
    ram = _memory_mb()
    if force or ram >= RAM_SOFT_LIMIT_MB:
        gc.collect()
        last_gc_at = now

def _memory_mb():
    """Текущий RSS процесса, а не исторический peak. Linux-first с безопасным fallback."""
    try:
        with open("/proc/self/status", "r", encoding="utf-8") as f:
            for line in f:
                if line.startswith("VmRSS:"):
                    return int(line.split()[1]) / 1024.0
    except Exception:
        pass
    try:
        import resource
        # ru_maxrss — только fallback; это peak, поэтому не используем его как основной watchdog.
        value = resource.getrusage(resource.RUSAGE_SELF).ru_maxrss
        return value / 1024.0 if os.name != "nt" else value / (1024.0 * 1024.0)
    except Exception:
        return 0

def cleanup_artifacts():
    now = time.time()
    changed = False
    for uid, items in list(artifact_registry.items()):
        kept = []
        for item in items:
            path = item.get("path", "")
            age = now - float(item.get("created", now))
            try:
                size = os.path.getsize(path) if os.path.exists(path) else 0
                if path and os.path.exists(path) and (age > ARTIFACT_TTL_SECONDS or size > MAX_ARTIFACT_BYTES):
                    os.remove(path); changed = True
                    continue
            except Exception:
                pass
            if os.path.exists(path):
                kept.append(item)
        if kept:
            artifact_registry[uid] = kept[-MAX_ARTIFACTS_PER_USER:]
        else:
            artifact_registry.pop(uid, None)
    # Empty user directories are cheap to remove.
    try:
        for name in os.listdir(ARTIFACT_DIR):
            full = os.path.join(ARTIFACT_DIR, name)
            if os.path.isdir(full) and not os.listdir(full):
                os.rmdir(full)
    except Exception:
        pass
    if changed:
        _save_artifacts()

def cleanup_stale_tmp_files():
    """Удаляет старые .tmp после аварийного рестарта."""
    now = time.time()
    for root in (DATA_DIR, ARTIFACT_DIR):
        if not os.path.isdir(root):
            continue
        for base, _, files in os.walk(root):
            for name in files:
                if not name.endswith(".tmp"):
                    continue
                path = os.path.join(base, name)
                try:
                    if now - os.path.getmtime(path) > 3600:
                        os.remove(path)
                except Exception:
                    pass

def _disk_free_mb(path=None):
    try:
        import shutil
        target = path or DATA_DIR
        total, used, free = shutil.disk_usage(target)
        return free / (1024 * 1024)
    except Exception:
        return None


def _cleanup_orphan_artifact_files():
    """Удаляет старые файлы генерации, даже если их запись в registry потерялась."""
    now = time.time()
    root = ARTIFACT_DIR
    if not os.path.isdir(root):
        return
    for base, _, files in os.walk(root):
        for name in files:
            if name.endswith('.tmp'):
                continue
            path = os.path.join(base, name)
            try:
                if now - os.path.getmtime(path) > ARTIFACT_TTL_SECONDS:
                    os.remove(path)
            except Exception:
                pass


async def maintenance_loop():
    while True:
        try:
            cleanup_artifacts()
            _cleanup_orphan_artifact_files()
            free_mb = _disk_free_mb()
            if free_mb is not None and free_mb < 250:
                logging.warning("Disk watchdog: only %.0f MB free — cleaning generated files", free_mb)
                cleanup_artifacts()
                _cleanup_orphan_artifact_files()
            # Bounded TTL cache cleanup. Проверку подписки не меняем — чистим только истёкшие записи.
            if len(subscription_cache) > 2000:
                cutoff_sub = time.monotonic() - SUBSCRIPTION_CACHE_TTL
                for uid, ts in list(subscription_cache.items()):
                    if ts < cutoff_sub:
                        subscription_cache.pop(uid, None)
            ram = _memory_mb()
            if ram >= RAM_HARD_LIMIT_MB:
                logging.error("RAM watchdog: %s MB — принудительная очистка временных данных", round(ram))
                _maybe_gc(force=True)
                # Не принимаем новые тяжёлые задачи, пока память не стабилизируется.
                await asyncio.sleep(2)
            elif ram >= RAM_SOFT_LIMIT_MB:
                logging.warning("RAM watchdog: %s MB — мягкая очистка", round(ram))
                _maybe_gc(force=True)
            else:
                _maybe_gc()
            flush_persistent_state(force=False)
        except asyncio.CancelledError:
            raise
        except Exception as e:
            logging.warning("Maintenance loop: %s", e)
        await asyncio.sleep(RESOURCE_CHECK_INTERVAL)

def _release_large_object(*objects):
    """Drop references to large transient objects; GC is intentionally deferred."""
    for obj in objects:
        try:
            del obj
        except Exception:
            pass

def _file_safe_size(path):
    try:
        return os.path.getsize(path) <= MAX_ARTIFACT_BYTES
    except Exception:
        return False


def _atomic_write_bytes(path, data: bytes):
    """Записывает артефакт транзакционно, чтобы сбой не оставил битый файл."""
    if not data or len(data) > MAX_ARTIFACT_BYTES:
        raise ValueError("Артефакт пустой или превышает допустимый размер")
    tmp = path + ".tmp"
    with open(tmp, "wb") as f:
        f.write(data)
        f.flush()
        os.fsync(f.fileno())
    os.replace(tmp, path)
    return path


def set_task_state(user_id: int, **kwargs):
    state = user_task_state.setdefault(user_id, {})
    state.update(kwargs)
    return state

def get_task_state(user_id: int):
    return user_task_state.get(user_id, {})

def clear_task_state(user_id: int):
    user_task_state.pop(user_id, None)
    user_custom_images.pop(user_id, None)
    user_ppt_images.pop(user_id, None)
    ppt_creation_users.discard(user_id)

def ppt_assets(user_id: int):
    return user_custom_images.get(user_id, [])

def save_presentation_image(user_id: int, file_id: str):
    if user_id not in ppt_creation_users:
        return False
    images = user_custom_images.setdefault(user_id, [])
    if file_id not in images:
        images.append(file_id)
    user_custom_images[user_id] = images[-20:]
    return True

def add_custom_image(user_id: int, file_id: str):
    images = user_custom_images.setdefault(user_id, [])
    if file_id not in images:
        images.append(file_id)
    user_custom_images[user_id] = images[-20:]

user_doc_context = {}
user_doc_chunks = {}  # user_id -> {name, chunks:[{page,text}]}

# Лёгкая аналитика поведения — без хранения содержимого сообщений.
user_rate_buckets = {}
SUBSCRIPTION_CACHE_TTL = int(os.getenv("MECAUAI_SUBSCRIPTION_CACHE_TTL", "30"))
subscription_cache = {}

def allow_request(user_id: int, limit: int = 15, window: int = 30) -> bool:
    """Мягкий антиспам: защищает VPS и API, не мешая обычному использованию."""
    now = time.monotonic()
    bucket = [t for t in user_rate_buckets.get(user_id, []) if now - t < window]
    if len(bucket) >= limit:
        user_rate_buckets[user_id] = bucket
        return False
    bucket.append(now)
    user_rate_buckets[user_id] = bucket
    if len(user_rate_buckets) > 5000:
        # Бюджет памяти: сначала удаляем давно неактивные buckets.
        cutoff = now - max(window * 2, 60)
        stale = [uid for uid, ts in user_rate_buckets.items() if not ts or ts[-1] < cutoff]
        for uid in stale[:1200]:
            user_rate_buckets.pop(uid, None)
        if len(user_rate_buckets) > 5200:
            for uid in list(user_rate_buckets)[:500]:
                user_rate_buckets.pop(uid, None)
    return True

def _runtime_health_snapshot():
    """Лёгкая диагностика для защиты VPS: не падает, если psutil недоступен."""
    try:
        import psutil
        vm=psutil.virtual_memory()
        du=psutil.disk_usage(str(DATA_DIR))
        return {'ram_percent':round(vm.percent,1),'ram_available_mb':round(vm.available/1024/1024), 'disk_percent':round(du.percent,1)}
    except Exception:
        return {'ram_percent':None,'ram_available_mb':None,'disk_percent':None}

async def _ingress_allowed(message: Message, limit: int, window: int, label: str) -> bool:
    """
    Асинхронная обёртка над allow_request для входящих хендлеров (документы,
    голосовые, фото, текст). При превышении лимита уведомляет пользователя и
    возвращает False, чтобы хендлер сразу прекратил обработку.
    БАГ: эта функция вызывалась в handle_document/handle_voice/handle_photo/
    handle_text, но нигде не была определена — из-за этого любое входящее
    сообщение падало с NameError ещё до какой-либо обработки.
    """
    user_id = message.from_user.id
    if not allow_request(user_id, limit=limit, window=window):
        try:
            await message.answer(f"⏳ Слишком много запросов ({label}) за короткое время. Подожди немного и попробуй снова.")
        except Exception:
            pass
        return False
    return True

def get_user_stats(user_id: int):
    if user_id not in user_stats:
        user_stats[user_id] = {
            "messages": 0,
            "voice": 0,
            "files": 0,
            "images": 0,
            "exports": 0,
            "last_action": None,
        }
    return user_stats[user_id]

def remember_action(user_id: int, action: str):
    stats = get_user_stats(user_id)
    stats["last_action"] = action
    stats["messages"] += 1

def clean_text_for_html(text: str) -> str:
    """Подготавливает ответ ИИ для Telegram HTML: убирает Markdown-артефакты и экранирует опасные HTML-символы."""
    text = str(text or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    # Telegram HTML допускает только ограниченный набор тегов. Сохраняем простые
    # заголовки/жирный/курсив, а остальные '<' и '>' экранируем.
    placeholders = {}
    allowed = re.compile(r"</?(?:b|strong|i|em|u|s|code|pre|blockquote)(?:\s[^>]*)?>", re.I)
    def protect(m):
        key = f"__MECAU_TAG_{len(placeholders)}__"
        placeholders[key] = m.group(0)
        return key
    text = allowed.sub(protect, text)
    text = html.escape(text, quote=False)
    for key, tag in placeholders.items():
        text = text.replace(key, tag)
    # Базовое Markdown-оформление переводим в Telegram HTML.
    text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text, flags=re.S)
    text = re.sub(r"(?<!\*)\*([^*\n]+)\*(?!\*)", r"<i>\1</i>", text)
    text = re.sub(r"`([^`\n]+)`", r"<code>\1</code>", text)
    return text.strip()


def extract_json(text: str):
    """Извлекает JSON из ответа модели, даже если модель добавила ```json ...``` или пояснение."""
    raw = str(text or "").strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.I)
    raw = re.sub(r"\s*```$", "", raw).strip()
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        # Ищем первый полноценный JSON-массив/объект в тексте.
        decoder = json.JSONDecoder()
        for i, ch in enumerate(raw):
            if ch not in "[{":
                continue
            try:
                value, _ = decoder.raw_decode(raw[i:])
                return value
            except json.JSONDecodeError:
                continue
        raise


def get_answer_inline_keyboard():
    """Кнопки действий под обычным ответом ИИ."""
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="💡 Проще", callback_data="btn_simplify"),
            InlineKeyboardButton(text="🧩 Структура", callback_data="btn_structure"),
            InlineKeyboardButton(text="🔍 Проверить", callback_data="btn_check"),
        ],
        [
            InlineKeyboardButton(text="🧠 Глубже", callback_data="btn_deepen"),
            InlineKeyboardButton(text="💡 Примеры", callback_data="btn_examples"),
            InlineKeyboardButton(text="🎓 Академично", callback_data="btn_academic"),
        ],
        [
            InlineKeyboardButton(text="📄 Word", callback_data="btn_word"),
            InlineKeyboardButton(text="📈 Презентация", callback_data="btn_make_ppt"),
        ],
        [
            InlineKeyboardButton(text="⭐ Сохранить", callback_data="btn_save_fav"),
            InlineKeyboardButton(text="🔁 Продолжить", callback_data="btn_continue"),
        ],
    ])


def get_sub_keyboard():
    """Клавиатура проверки обязательных подписок."""
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📢 MecauAI", url=CHANNEL_1_URL)],
        [InlineKeyboardButton(text="📢 Mecau Info", url=CHANNEL_2_URL)],
        [InlineKeyboardButton(text="✅ Проверить подписку", callback_data="check_sub")],
    ])


def get_quick_actions_keyboard():
    """Основная компактная клавиатура обычного пользователя."""
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="🛠 Создать"), KeyboardButton(text="⭐ Избранное")],
            [KeyboardButton(text="⚙️ Настройки ИИ")],
            [KeyboardButton(text="ℹ️ О MecauAI"), KeyboardButton(text="🛠 Техподдержка")],
        ],
        resize_keyboard=True,
        is_persistent=True,
    )


def get_admin_keyboard():
    """Отдельная клавиатура администратора: пользовательские функции + управление ботом."""
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="🛠 Создать"), KeyboardButton(text="⭐ Избранное")],
            [KeyboardButton(text="⚙️ Настройки ИИ")],
            [KeyboardButton(text="ℹ️ О MecauAI"), KeyboardButton(text="🛠 Техподдержка")],
            [KeyboardButton(text="🛡 Админ-панель")],
        ],
        resize_keyboard=True,
        is_persistent=True,
    )


def keyboard_for(user_id: int):
    """Возвращает пользовательскую или административную клавиатуру."""
    if user_id in ADMIN_IDS:
        return get_admin_keyboard()
    return get_quick_actions_keyboard()


def get_cancel_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]
    ])


def get_create_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📈 Презентация", callback_data="create_ppt")],
        [InlineKeyboardButton(text="📊 Excel", callback_data="create_excel")],
        [InlineKeyboardButton(text="📄 Word", callback_data="create_word")],
        [InlineKeyboardButton(text="📑 Титульник ГОСТ", callback_data="create_title")],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")],
    ])


async def safe_answer(message: Message, text: str, parse_mode: str = None, reply_markup=None, disable_web_page_preview: bool = True):
    """
    Безопасная отправка сообщения: если Telegram не может распарсить
    HTML/Markdown в ответе ИИ (например, из-за случайных символов < > _ *),
    бот отправит тот же текст обычным сообщением вместо падения с ошибкой.
    """
    try:
        return await message.answer(
            text, parse_mode=parse_mode, reply_markup=reply_markup,
            disable_web_page_preview=disable_web_page_preview
        )
    except Exception as e:
        logging.warning(f"Не удалось отправить с parse_mode={parse_mode}, отправляю без форматирования: {e}")
        plain_text = re.sub(r'<[^>]+>', '', text) if parse_mode == "HTML" else text
        try:
            return await message.answer(
                plain_text, reply_markup=reply_markup,
                disable_web_page_preview=disable_web_page_preview
            )
        except Exception as e2:
            logging.error(f"Повторная ошибка при отправке сообщения: {e2}")
            return None


def get_help_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="🎙 Голос", callback_data="help_voice"),
            InlineKeyboardButton(text="📎 Файл", callback_data="help_file"),
        ],
        [
            InlineKeyboardButton(text="📈 Презентация", callback_data="help_ppt"),
            InlineKeyboardButton(text="📊 Excel", callback_data="help_excel"),
        ],
        [
            InlineKeyboardButton(text="⚙️ Настройки ИИ", callback_data="help_mode"),
        ]
    ])

def get_settings_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🧠 Изменить режим ИИ", callback_data="settings_mode")],
        [InlineKeyboardButton(text="🧹 Очистить контекст", callback_data="settings_clear")],
        [InlineKeyboardButton(text="❌ Закрыть", callback_data="settings_close")],
    ])

# БАГ: использовалась в _trim_user_record, но нигде не была определена —
# NameError при первом же сохранении/подрезке пользовательской записи.
MAX_USER_LAST_OUTPUT_CHARS = 4000

def _trim_user_record(record):
    if not isinstance(record, dict):
        return {"mode": "ai", "history": [], "last_output": ""}
    history = record.get("history") or []
    if isinstance(history, list):
        record["history"] = history[-6:]
    else:
        record["history"] = []
    record["last_output"] = str(record.get("last_output") or "")[-MAX_USER_LAST_OUTPUT_CHARS:]
    return record

# Bound resident user history/output after startup so old oversized records
# cannot keep unnecessary RAM alive for the entire process lifetime.
for _uid, _record in list(users_db.items()):
    users_db[_uid] = _trim_user_record(_record)

def save_users_db(immediate: bool = False):
    global users_db_dirty
    for uid, record in list(users_db.items()):
        users_db[uid] = _trim_user_record(record)
    if not immediate:
        users_db_dirty = True
        return
    try:
        _atomic_json_write(USER_DATA_FILE, {str(k): v for k, v in users_db.items()})
        users_db_dirty = False
    except Exception as exc:
        logging.error("Не удалось сохранить пользовательские данные: %s", exc)

def save_stats_db(immediate: bool = False):
    global stats_db_dirty
    if not immediate:
        stats_db_dirty = True
        return
    try:
        _atomic_json_write(STATS_FILE, {str(k): v for k, v in user_stats.items()})
        stats_db_dirty = False
    except Exception as exc:
        logging.error("Не удалось сохранить статистику: %s", exc)

def flush_persistent_state(force: bool = False):
    global users_db_dirty, stats_db_dirty, favorites_dirty, users_ids_dirty, modes_db_dirty, artifacts_dirty, last_persist_flush_at, activity_dirty
    now = time.monotonic()
    if not force and now - last_persist_flush_at < PERSIST_FLUSH_INTERVAL:
        return
    if users_ids_dirty:
        try:
            _atomic_json_write(USERS_FILE, sorted(all_users_cache))
            users_ids_dirty = False
        except Exception as exc:
            logging.error("Не удалось сохранить список пользователей: %s", exc)
    if modes_db_dirty:
        try:
            _atomic_json_write(MODES_FILE, saved_modes_cache)
            modes_db_dirty = False
        except Exception as exc:
            logging.error("Не удалось сохранить режимы: %s", exc)
    if users_db_dirty:
        save_users_db(immediate=True)
    if stats_db_dirty:
        save_stats_db(immediate=True)
    if favorites_dirty:
        save_favorites(load_favorites(), immediate=True)
    if artifacts_dirty:
        _save_artifacts(immediate=True)
    if activity_dirty:
        _save_user_activity(immediate=True)
    last_persist_flush_at = now

def get_user_data(user_id: int):
    if user_id not in users_db:
        users_db[user_id] = {
            "mode": saved_modes_cache.get(str(user_id), "ai"),
            "history": [],
            "last_output": "Здесь пока нет ответов."
        }
        save_users_db()
    return users_db[user_id]

def clear_pending_states(user_id: int):
    """
    Сбрасывает все 'ожидающие ввода' состояния пользователя (создание презентации,
    таблицы, анализ документа). Нужно, чтобы если пользователь передумал и нажал
    другую кнопку меню, бот не принял её текст за тему презентации/таблицы —
    раньше это приводило к путанице и «зависшим» состояниям.
    """
    ppt_states.discard(user_id)
    excel_states.discard(user_id)
    doc_analysis_states.discard(user_id)
    user_ppt_images.pop(user_id, None)
    user_ppt_design.pop(user_id, None)
    user_ppt_slide_count.pop(user_id, None)
    user_excel_mode.pop(user_id, None)
    user_ppt_topic.pop(user_id, None)
    user_excel_topic.pop(user_id, None)
    user_doc_chunks.pop(user_id, None)

def _get_semaphores():
    global AI_SEMAPHORE, HEAVY_JOB_SEMAPHORE
    if AI_SEMAPHORE is None:
        AI_SEMAPHORE = asyncio.Semaphore(MAX_AI_CONCURRENCY)
    if HEAVY_JOB_SEMAPHORE is None:
        HEAVY_JOB_SEMAPHORE = asyncio.Semaphore(MAX_HEAVY_CONCURRENCY)
    return AI_SEMAPHORE, HEAVY_JOB_SEMAPHORE


async def acquire_heavy_job(status_msg=None):
    global HEAVY_WAITERS
    _, heavy = _get_semaphores()
    current_ram = _memory_mb()
    if current_ram >= RAM_HARD_LIMIT_MB:
        _maybe_gc()
        current_ram = _memory_mb()
        if current_ram >= RAM_HARD_LIMIT_MB:
            raise RuntimeError(f"Сервер временно перегружен по памяти ({current_ram:.0f} MB). Повтори задачу через минуту.")
    queued=False
    if heavy.locked():
        if HEAVY_WAITERS >= MAX_HEAVY_QUEUE:
            raise RuntimeError("Сервер временно занят: очередь тяжёлых задач заполнена")
        HEAVY_WAITERS += 1; queued=True
        if status_msg is not None:
            try: await status_msg.edit_text("⏳ Сервер занят. Задача поставлена в короткую очередь — я не перегружаю VPS.")
            except Exception: pass
    try:
        await heavy.acquire()
        return heavy
    finally:
        if queued: HEAVY_WAITERS=max(0,HEAVY_WAITERS-1)


def _looks_like_model_access_error(exc: Exception) -> bool:
    text = str(exc).lower()
    return (
        "model_not_found" in text
        or "does not exist" in text
        or "do not have access" in text
        or "model is not available" in text
        or "invalid model" in text
        or "error code: 404" in text
        or "error code: 403" in text
    )


async def _probe_groq_model(model: str) -> bool:
    """Проверяет, доступна ли модель именно этому API-ключу/проекту."""
    if not model:
        return False
    try:
        await asyncio.wait_for(groq_client.models.retrieve(model), timeout=10)
        return True
    except Exception as exc:
        logging.warning("Groq model probe failed for %s: %s", model, exc)
        return False


async def _refresh_groq_candidates():
    """Убирает из очереди очевидно недоступные/устаревшие модели до первого запроса."""
    global TEXT_MODEL_FALLBACKS, VISION_MODEL_FALLBACKS
    for attr in ("TEXT_MODEL_FALLBACKS", "VISION_MODEL_FALLBACKS"):
        candidates = list(dict.fromkeys(globals().get(attr, [])))
        if not candidates:
            continue
        checked = await asyncio.gather(*[_probe_groq_model(x) for x in candidates], return_exceptions=True)
        available = [m for m, ok in zip(candidates, checked) if ok is True]
        # Если проверка сети временно не удалась для всех моделей, не обнуляем
        # список: runtime fallback всё равно сможет повторить запрос.
        if available:
            globals()[attr] = available
            logging.info("Groq available %s models: %s", "text" if attr.startswith("TEXT") else "vision", available)
        else:
            logging.warning("Groq model probe returned no confirmed models for %s; keeping fallback list", attr)


async def call_groq_with_retry(messages, model: str = None, temperature: float = 0.7, max_retries: int = 2, timeout: int = 45):
    """
    Надёжный вызов Groq.
    Если конкретная модель недоступна для ключа/проекта, автоматически
    переключается на доступную резервную модель. Сетевые/временные ошибки
    по-прежнему получают повторные попытки.
    """
    requested = model or TEXT_MODEL
    is_vision = model == VISION_MODEL
    candidates = VISION_MODEL_FALLBACKS if is_vision else TEXT_MODEL_FALLBACKS
    if requested not in candidates:
        candidates = [requested] + candidates

    last_err = None
    for candidate_index, used_model in enumerate(dict.fromkeys(candidates)):
        for attempt in range(max_retries + 1):
            try:
                ai_sem, _ = _get_semaphores()
                async with ai_sem:
                    result = await asyncio.wait_for(
                        groq_client.chat.completions.create(
                            model=used_model,
                            messages=messages,
                            temperature=temperature
                        ),
                        timeout=timeout
                    )
                if used_model != requested:
                    logging.info("Groq: модель %s недоступна, успешно использована резервная %s", requested, used_model)
                return result
            except Exception as e:
                last_err = e
                if _looks_like_model_access_error(e):
                    logging.warning(
                        "Groq: модель %s недоступна (попытка %s/%s): %s",
                        used_model, attempt + 1, max_retries + 1, e
                    )
                    # Нет смысла трижды повторять гарантированно недоступную модель.
                    break
                logging.warning(
                    "Groq API — модель %s, попытка %s/%s не удалась: %s",
                    used_model, attempt + 1, max_retries + 1, e
                )
                if attempt < max_retries:
                    await asyncio.sleep(1.5 * (attempt + 1))
        if candidate_index < len(candidates) - 1 and last_err is not None and _looks_like_model_access_error(last_err):
            logging.warning("Groq: переключаюсь с %s на следующую резервную модель", used_model)
            continue
        if last_err is not None and not _looks_like_model_access_error(last_err):
            break

    raise last_err


async def transcribe_voice(file_bytes: bytes, filename: str = "voice.ogg") -> str:
    """Надёжное распознавание Telegram voice через Groq Whisper."""
    last_err = None
    for attempt in range(3):
        try:
            ai_sem, _ = _get_semaphores()
            async with ai_sem:
                result = await asyncio.wait_for(
                    groq_client.audio.transcriptions.create(
                        model="whisper-large-v3-turbo",
                        file=(filename, file_bytes),
                        response_format="text",
                        language="ru",
                        temperature=0
                    ),
                    timeout=60
                )
            transcript = result if isinstance(result, str) else getattr(result, "text", "")
            transcript = re.sub(r"\s+", " ", (transcript or "")).strip()
            if transcript:
                return transcript
            raise ValueError("Пустая расшифровка")
        except Exception as e:
            last_err = e
            logging.warning(f"Whisper attempt {attempt + 1}/3 failed: {e}")
            if attempt < 2:
                await asyncio.sleep(1.5 * (attempt + 1))
    raise last_err

def _normalize_channel_ref(value):
    """Нормализует публичный username, t.me URL или числовой chat_id."""
    if value is None:
        return None
    raw = str(value).strip()
    if not raw:
        return None
    if raw.lstrip("-").isdigit():
        return int(raw)
    raw = re.sub(r"^https?://(?:www\.)?(?:t\.me|telegram\.me)/", "", raw, flags=re.I)
    raw = raw.split("?", 1)[0].split("/", 1)[0].strip()
    if raw.startswith("+") or raw.startswith("joinchat/"):
        raise ValueError("Для проверки подписки нужен публичный @username или числовой chat_id, а не invite-ссылка")
    return raw.lstrip("@").strip()


async def _resolve_channel(value):
    """Получает реальный объект канала. Публичные username всегда передаются с @."""
    ref = _normalize_channel_ref(value)
    if ref is None:
        raise ValueError(f"Пустой идентификатор канала: {value!r}")
    lookup = ref if isinstance(ref, int) else f"@{ref}"
    chat = await bot.get_chat(chat_id=lookup)
    if getattr(chat, "type", None) != "channel":
        raise ValueError(f"{value!r} разрешился не в канал, а в chat type={getattr(chat, 'type', None)!r}")
    return chat


def _member_is_subscribed(member) -> bool:
    status = getattr(member, "status", None)
    if status in {"creator", "administrator", "member"}:
        return True
    return status == "restricted" and bool(getattr(member, "is_member", False))


async def _read_member(chat_id, user_id: int, attempts: int = 4):
    """Читает статус несколько раз: Telegram иногда обновляет membership не мгновенно."""
    last_member = None
    last_error = None
    for attempt in range(attempts):
        try:
            member = await bot.get_chat_member(chat_id=chat_id, user_id=user_id)
            last_member = member
            if _member_is_subscribed(member):
                return member
            if attempt + 1 < attempts:
                await asyncio.sleep(0.7)
        except Exception as exc:
            last_error = exc
            if attempt + 1 < attempts:
                await asyncio.sleep(0.7)
    if last_member is not None:
        return last_member
    raise last_error


async def check_subscription_details(user_id: int, force: bool = False):
    """Проверка подписки с коротким TTL успешного результата."""
    if user_id in ADMIN_IDS:
        return True, None, "admin_bypass"
    now = time.monotonic()
    cached_at = subscription_cache.get(user_id)
    if not force and cached_at and now - cached_at < SUBSCRIPTION_CACHE_TTL:
        return True, None, "cached_ok"
    if force:
        subscription_cache.pop(user_id, None)

    channels = [
        ("CHANNEL_1", CHANNEL_1_USERNAME, CHANNEL_1_URL),
        ("CHANNEL_2", CHANNEL_2_USERNAME, CHANNEL_2_URL),
    ]

    for label, configured_ref, public_url in channels:
        try:
            chat = await _resolve_channel(configured_ref)
            # get_chat_member is the authoritative check. Retry transient Telegram
            # failures and membership propagation, but never cache a positive result.
            member = None
            last_exc = None
            for attempt in range(5):
                try:
                    member = await bot.get_chat_member(chat_id=chat.id, user_id=user_id)
                    if _member_is_subscribed(member):
                        break
                    if attempt < 4:
                        await asyncio.sleep(0.6)
                except Exception as exc:
                    last_exc = exc
                    if attempt < 4:
                        await asyncio.sleep(0.8)
                    else:
                        raise
            if member is None:
                raise last_exc or RuntimeError("Telegram did not return a member object")

            status = getattr(member, "status", None)
            is_member = getattr(member, "is_member", None)
            ok = _member_is_subscribed(member)
            logging.info(
                "SUB_CHECK user=%s channel=%s configured=%r lookup=%r chat_id=%s "
                "chat_username=%r chat_type=%r status=%r is_member=%r result=%s",
                user_id, label, configured_ref,
                f"@{_normalize_channel_ref(configured_ref)}" if not str(configured_ref).lstrip("-").isdigit() else configured_ref,
                chat.id, getattr(chat, "username", None), getattr(chat, "type", None),
                status, is_member, ok
            )
            if not ok:
                return False, (label, configured_ref, public_url, chat.id, status, is_member), "not_member"
        except Exception as exc:
            logging.exception(
                "SUB_CHECK_ERROR user=%s channel=%s configured=%r url=%r error=%s",
                user_id, label, configured_ref, public_url, exc
            )
            return False, (label, configured_ref, public_url, None, None, None), f"error:{type(exc).__name__}:{exc}"

    subscription_cache[user_id] = time.monotonic()
    if len(subscription_cache) > 10000:
        cutoff = time.monotonic() - SUBSCRIPTION_CACHE_TTL
        for uid, ts in list(subscription_cache.items()):
            if ts < cutoff:
                subscription_cache.pop(uid, None)
    return True, None, "ok"

async def check_subscription(user_id: int, force: bool = False) -> bool:
    ok, _, _ = await check_subscription_details(user_id, force=force)
    return ok

async def require_subscription_callback(callback: types.CallbackQuery) -> bool:
    ok, failed, reason = await check_subscription_details(callback.from_user.id, force=True)
    if ok:
        return True
    if reason == "not_member" and failed:
        label, configured_ref, public_url, chat_id, status, is_member = failed
        await callback.answer(f"❌ Нет подписки на {label}: {public_url}", show_alert=True)
    else:
        await callback.answer("⚠️ Telegram временно не дал проверить подписку. Попробуй ещё раз.", show_alert=True)
    return False

@dp.message(CommandStart())
async def cmd_start(message: Message):
    user_id = message.from_user.id
    save_user_id(user_id)
    _track_user_event(user_id, message=message, kind="start", status="received", prompt="/start")

    if not await check_subscription(user_id):
        await message.answer(
            f"🔒 Доступ заблокирован!\n\n"
            f"Чтобы пользоваться MecauAI, подпишись на оба наших канала:\n"
            f" 👉 {CHANNEL_1_URL}\n"
            f" 👉 {CHANNEL_2_URL}\n\n"
            f"После подписки нажми кнопку ниже 👇",
            reply_markup=get_sub_keyboard()
        )
        return

    start_text = (
        f"Привет, {message.from_user.first_name}! Ты активировал MecauAI 🚀\n\n"
        "Я твой карманный помощник для учебы и разработки. Можешь отправлять вопросы, задачи, код, изображения, документы и голосовые сообщения — я распознаю запрос, разберусь в нём и дам готовый ответ. Также могу подготовить Word, презентацию, Excel или помочь разобрать документ. Если задача большая, просто опиши её своими словами — я сам выберу подходящий способ обработки."
    )
    await message.answer(start_text, reply_markup=keyboard_for(user_id))

@dp.callback_query(F.data == "check_sub")
async def cb_check_sub(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    ok, failed, reason = await check_subscription_details(user_id, force=True)
    if ok:
        save_user_id(user_id)
        await callback.answer("✅ Подписка подтверждена!")
        try:
            await callback.message.delete()
        except Exception:
            pass
        await callback.message.answer(
            "🎉 Подписка на оба канала подтверждена! Добро пожаловать в MecauAI 🚀",
            reply_markup=keyboard_for(user_id)
        )
        return
    if failed and reason == "not_member":
        await callback.answer(f"❌ Нет подписки на канал: {failed[2]}", show_alert=True)
    else:
        await callback.answer("⚠️ Telegram временно не дал проверить подписку. Попробуй ещё раз.", show_alert=True)


@dp.message(F.text == "✨ Возможности")
async def cmd_capabilities(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    clear_pending_states(message.from_user.id)
    await message.answer(
        "✨ <b>MecauAI умеет больше, чем кажется</b>\n\n"
        "💬 <b>Общение</b> — вопросы, объяснения, тексты, код, идеи.\n"
        "🎙 <b>Голос</b> — отправь голосовое, я сам превращу его в задачу.\n"
        "🖼 <b>Изображения</b> — распознаю текст, задачи, схемы и содержимое фото.\n"
        "📚 <b>Документы</b> — PDF, DOCX, TXT и фото страниц с вопросами по содержимому.\n"
        "📄 <b>Экспорт</b> — ответ можно сохранить в Word.\n"
        "📈 <b>Презентации</b> — попроси сделать презентацию на нужную тему.\n"
        "📊 <b>Excel</b> — попроси создать таблицу с данными и диаграммой.\n\n"
        "💡 Главное: <b>не нужно искать правильную кнопку</b>. Напиши или скажи задачу своими словами.",
        parse_mode="HTML",
        reply_markup=get_help_keyboard()
    )

@dp.message(Command("help"))
async def cmd_help(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🤖 <b>MecauAI умеет работать не только с текстом.</b>\n\n"
        "🎙 Голос — говори задачу обычными словами.\n"
        "🖼 Фото — отправляй задачи, скриншоты и изображения.\n"
        "📁 Документы — загружай PDF/DOCX/TXT и задавай вопросы.\n"
        "📈 Презентации — создание с выбором содержания и оформления.\n"
        "📊 Excel — таблицы, формулы, расчёты и визуализация.\n"
        "📄 Word — оформление готового материала в документ.\n\n"
        "💡 Самый простой способ: просто напиши, что тебе нужно.",
        parse_mode="HTML",
        reply_markup=get_quick_actions_keyboard()
    )

@dp.message(F.text == "🛠 Техподдержка")
async def cmd_support(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🛠 <b>Техподдержка MecauAI</b>\n\n"
        "Если бот работает неправильно, нашёл ошибку или есть вопрос — напиши в поддержку.\n\n"
        "👤 Поддержка: <b>@mecau</b>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="💬 Написать в поддержку", url="https://t.me/mecau")]
        ])
    )

@dp.message(F.text == "📁 Документы")
async def cmd_documents_menu(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "📁 <b>Работа с документами</b>\n\n"
        "Просто отправь PDF, DOCX или TXT — дальше можно:\n"
        "• сделать конспект;\n"
        "• найти нужную информацию;\n"
        "• задать вопросы по документу;\n"
        "• подготовить материал к экзамену;\n"
        "• сравнить документы;\n"
        "• сделать основу для презентации.\n\n"
        "💡 Можно сразу написать запрос вместе с файлом.",
        parse_mode="HTML",
        reply_markup=get_cancel_keyboard()
    )

@dp.message(F.text == "🛠 Создать")
async def cmd_create_menu(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    await message.answer(
        "🛠 <b>Что создать?</b>\n\n"
        "Выбери готовый вариант. Если нужного нет — просто опиши задачу обычными словами.\n\n"
        "📈 В презентации я отдельно помогу с темой, структурой, стилем, титульником и твоими картинками.",
        parse_mode="HTML",
        reply_markup=get_create_keyboard()
    )

@dp.message(F.text == "⚙️ Настройки ИИ")
async def cmd_settings(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    clear_pending_states(message.from_user.id)
    user_data = get_user_data(message.from_user.id)
    names = {
        "coder": "💻 ИИ-Программист",
        "ai": "🧠 Академический ассистент",
        "friend": "🫂 Лучший друг"
    }
    await message.answer(
        f"⚙️ <b>Настройки ИИ</b>\n\n"
        f"Текущий режим: <b>{names.get(user_data['mode'], 'Стандартный')}</b>\n\n"
        "Здесь собраны все настройки ИИ в одном месте.\n\nМожно сменить режим общения или очистить контекст текущего диалога.",
        parse_mode="HTML",
        reply_markup=get_settings_keyboard()
    )

@dp.callback_query(F.data.in_({"help_voice", "help_file", "help_ppt", "help_excel", "help_mode",
                               "settings_mode", "settings_clear", "settings_close"}))
async def cb_help_and_settings(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    user_id = callback.from_user.id

    if callback.data == "help_voice":
        await callback.message.answer("🎙 Просто отправь мне голосовое сообщение. Ничего включать не нужно.")
    elif callback.data == "help_file":
        await callback.message.answer("📎 Просто отправь PDF, DOCX, TXT или фото страницы — я предложу, что можно сделать.")
    elif callback.data == "help_ppt":
        await callback.message.answer("📈 Напиши: «Сделай презентацию на тему ...» — бот сам запустит создание презентации.")
    elif callback.data == "help_excel":
        await callback.message.answer("📊 Напиши: «Сделай Excel-таблицу ...» — бот предложит нужный формат и создаст файл.")
    elif callback.data == "help_mode":
        await cmd_settings(callback.message)
    elif callback.data == "settings_mode":
        await cmd_mode(callback.message)
    elif callback.data == "settings_close":
        await callback.answer()
        try:
            await callback.message.delete()
        except Exception:
            pass
        await callback.message.answer("🏠 <b>Главное меню</b>", parse_mode="HTML", reply_markup=keyboard_for(user_id))
        return
    elif callback.data == "settings_clear":
        user_data = get_user_data(user_id)
        user_data["history"] = []
        user_data["last_output"] = "Здесь пока нет ответов."
        save_users_db()
        user_doc_context.pop(user_id, None)
        await callback.message.answer("🧹 Контекст очищен. Начинаем с чистого листа.")

    await callback.answer()

@dp.message(F.text == "ℹ️ О MecauAI")
@dp.message(Command("about"))
async def cmd_about(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return

    about_text = (
        "ℹ️ <b>О MecauAI</b>\n\n"
        "MecauAI — помощник для учёбы и разработки. "
        "Просто отправляй задачу текстом, голосом, фото или файлом — "
        "бот сам выбирает подходящий способ обработки.\n\n"
        "Возможности и юридические документы собраны здесь, чтобы главное меню оставалось компактным."
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✨ Возможности", callback_data="about_features")],
        [InlineKeyboardButton(text="📄 Документы", callback_data="about_docs")],
        [InlineKeyboardButton(text="⬅️ Назад", callback_data="nav_main")],
    ])
    await message.answer(about_text, parse_mode="HTML", reply_markup=kb)

@dp.callback_query(F.data == "about_features")
async def cb_about_features(callback: CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    text = (
        "✨ <b>Возможности MecauAI</b>\n\n"
        "💬 Вопросы, объяснения, тексты, код и идеи\n"
        "🎙 Голос → расшифровка → обработка\n"
        "🖼 Фото → распознавание и анализ\n"
        "📚 PDF/DOCX/TXT → анализ и вопросы по содержимому\n"
        "📄 Word → экспорт ответа\n"
        "📈 Презентации → структура, дизайн и изображения\n"
        "📊 Excel → таблицы, формулы, KPI и диаграммы\n"
        "⭐ Избранное и 🧠 режимы ИИ\n\n"
        "💡 Не обязательно искать функцию в меню — просто опиши задачу своими словами."
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="⬅️ Назад", callback_data="about_back")]
    ])
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)

@dp.callback_query(F.data == "about_back")
async def cb_about_back(callback: CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    uid = callback.from_user.id
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    await callback.message.answer(
        "ℹ️ <b>О MecauAI</b>\n\n"
        "MecauAI — помощник для учёбы и разработки. "
        "Просто отправляй задачу текстом, голосом, фото или файлом — "
        "бот сам выбирает подходящий способ обработки.",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="✨ Возможности", callback_data="about_features")],
            [InlineKeyboardButton(text="📄 Документы", callback_data="about_docs")],
            [InlineKeyboardButton(text="⬅️ Назад", callback_data="nav_main")],
        ])
    )

@dp.callback_query(F.data == "about_docs")
async def cb_about_docs(callback: CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    rows = []
    if os.path.isfile(AGREEMENT_FILE):
        rows.append([InlineKeyboardButton(text="📄 Пользовательское соглашение", callback_data="doc_agreement")])
    if os.path.isfile(PRIVACY_FILE):
        rows.append([InlineKeyboardButton(text="🔐 Политика конфиденциальности", callback_data="doc_privacy")])
    if not rows:
        text = "📄 <b>Документы</b>\n\nДокументы пока не загружены в корень проекта."
    else:
        text = "📄 <b>Документы MecauAI</b>\n\nВыбери нужный документ:"
    rows.append([InlineKeyboardButton(text="⬅️ Назад", callback_data="about_back")])
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    await callback.message.answer(text, parse_mode="HTML", reply_markup=InlineKeyboardMarkup(inline_keyboard=rows))

async def _send_legal_pdf(callback: CallbackQuery, path: str, caption: str):
    if not os.path.isfile(path):
        await callback.answer("⚠️ Файл документа не найден на сервере.", show_alert=True)
        logging.error("Legal document missing: %s", path)
        return
    try:
        await callback.message.delete()
    except Exception:
        pass
    await callback.message.answer_document(FSInputFile(path), caption=caption)
    await callback.message.answer(
        "📄 Документы",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="⬅️ Назад", callback_data="about_back")]
        ])
    )
    await callback.answer()

@dp.callback_query(F.data == "doc_agreement")
async def cb_doc_agreement(callback: CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    await _send_legal_pdf(callback, AGREEMENT_FILE, "📄 Пользовательское соглашение MecauAI")

@dp.callback_query(F.data == "doc_privacy")
async def cb_doc_privacy(callback: CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    await _send_legal_pdf(callback, PRIVACY_FILE, "🔐 Политика конфиденциальности MecauAI")

@dp.callback_query(F.data == "nav_main")
async def cb_nav_main(callback: CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    uid = callback.from_user.id
    clear_pending_states(uid)
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    await callback.message.answer("🏠 <b>Главное меню</b>", parse_mode="HTML", reply_markup=keyboard_for(uid))

@dp.message(F.text == "🔄 Режим ИИ")
@dp.message(Command("mode"))
async def cmd_mode(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(message.from_user.id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💻 ИИ-Программист",       callback_data="set_mode_coder")],
        [InlineKeyboardButton(text="🧠 Академический ассистент", callback_data="set_mode_ai")],
        [InlineKeyboardButton(text="🫂 Лучший друг",             callback_data="set_mode_friend")]
    ])
    await message.answer("Выбери режим работы бота:", reply_markup=kb)

@dp.callback_query(F.data.startswith("set_mode_"))
async def cb_set_mode(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    mode = callback.data.split("_")[-1]
    user_id = callback.from_user.id
    user_data = get_user_data(user_id)
    user_data["mode"] = mode
    save_user_mode(user_id, mode)
    saved_modes_cache[str(user_id)] = mode
    save_users_db()
    clear_pending_states(user_id)

    names = {
        "coder": "ИИ-Программист 💻",
        "ai": "Академический ассистент 🧠",
        "friend": "Лучший друг 🫂"
    }
    await callback.message.edit_text(f"Режим переключен на: {names.get(mode, 'Стандартный')}")
    await callback.answer()

@dp.message(Command("checksub"))
async def cmd_checksub(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    parts = (message.text or "").split(maxsplit=1)
    if len(parts) != 2 or not parts[1].lstrip("-").isdigit():
        await message.answer("Использование: /checksub USER_ID")
        return
    uid = int(parts[1])
    ok, failed, reason = await check_subscription_details(uid)
    if ok:
        await message.answer(f"✅ USER_ID {uid}: подписан на оба канала.")
    else:
        await message.answer(f"❌ USER_ID {uid}: проблема с {failed[0] if failed else 'каналом'}.\nПричина: {reason}")

# ======================= ПРЕЗЕНТАЦИИ =======================

PPT_THEMES = {
    "ppt_blue": {
        "label": "🔵 Классический синий",
        "primary": RGBColor(0x1F, 0x4E, 0x79),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xD6, 0xE4, 0xF0),
        "body_text": RGBColor(0x21, 0x21, 0x21),
    },
    "ppt_dark": {
        "label": "⚫ Тёмный модерн",
        "primary": RGBColor(0x18, 0x18, 0x18),
        "bg": RGBColor(0x2B, 0x2B, 0x2B),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xCC, 0xCC, 0xCC),
        "body_text": RGBColor(0xF2, 0xF2, 0xF2),
    },
    "ppt_green": {
        "label": "🟢 Изумрудный",
        "primary": RGBColor(0x0E, 0x6B, 0x4F),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xD2, 0xF0, 0xE4),
        "body_text": RGBColor(0x1A, 0x1A, 0x1A),
    },
    "ppt_orange": {
        "label": "🟠 Тёплый оранжевый",
        "primary": RGBColor(0xC1, 0x5A, 0x11),
        "bg": RGBColor(0xFF, 0xF8, 0xF2),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle_text": RGBColor(0xFC, 0xE2, 0xC8),
        "body_text": RGBColor(0x2B, 0x1A, 0x0D),
    },
}

def _set_text_style(paragraph, size, color, bold=False, font="Aptos"):
    paragraph.font.name = font
    paragraph.font.size = PptxPt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold


def _add_footer(slide, theme, slide_no=None):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0.55), PptxInches(7.08), PptxInches(12.2), PptxInches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = theme["primary"]; line.line.fill.background()
    if slide_no is not None:
        tb = slide.shapes.add_textbox(PptxInches(11.9), PptxInches(7.1), PptxInches(0.7), PptxInches(0.25))
        p = tb.text_frame.paragraphs[0]; p.text = str(slide_no); p.alignment = PP_ALIGN.RIGHT
        _set_text_style(p, 9, theme["body_text"])


def build_title_slide(prs, theme, topic_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = theme["primary"]

    # Декоративные блоки делают титульник похожим на современный шаблон, а не на обычный текстовый лист.
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(0.22), PptxInches(7.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = theme["subtitle_text"]; accent.line.fill.background()

    label = slide.shapes.add_textbox(PptxInches(0.9), PptxInches(1.0), PptxInches(3.0), PptxInches(0.4))
    p = label.text_frame.paragraphs[0]; p.text = "MECAUAI • ПРЕЗЕНТАЦИЯ"
    _set_text_style(p, 11, theme["subtitle_text"], True)

    title = slide.shapes.add_textbox(PptxInches(0.9), PptxInches(2.0), PptxInches(11.2), PptxInches(2.2))
    tf = title.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(topic_text).strip()[:180]; p.alignment = PP_ALIGN.LEFT
    _set_text_style(p, 34 if len(str(topic_text)) < 80 else 28, theme["title_text"], True)

    sub = slide.shapes.add_textbox(PptxInches(0.92), PptxInches(4.55), PptxInches(8.8), PptxInches(0.8))
    p = sub.text_frame.paragraphs[0]; p.text = "Структурировано • Визуально • Готово к защите"
    _set_text_style(p, 16, theme["subtitle_text"])

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(9.9), PptxInches(5.55), PptxInches(2.25), PptxInches(0.7))
    badge.fill.solid(); badge.fill.fore_color.rgb = theme["bg"]; badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]; p.text = "AI • 2026"; p.alignment = PP_ALIGN.CENTER
    _set_text_style(p, 13, theme["body_text"], True)
    return slide


def build_content_slide(prs, theme, idx, item, img_stream):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = theme["bg"]

    # Небольшая верхняя полоса вместо огромного синего заголовка.
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), prs.slide_width, PptxInches(0.18))
    top.fill.solid(); top.fill.fore_color.rgb = theme["primary"]; top.line.fill.background()

    num = s.shapes.add_shape(MSO_SHAPE.OVAL, PptxInches(0.62), PptxInches(0.55), PptxInches(0.55), PptxInches(0.55))
    num.fill.solid(); num.fill.fore_color.rgb = theme["primary"]; num.line.fill.background()
    p = num.text_frame.paragraphs[0]; p.text = str(idx + 1); p.alignment = PP_ALIGN.CENTER
    _set_text_style(p, 13, theme["title_text"], True)

    tb_title = s.shapes.add_textbox(PptxInches(1.35), PptxInches(0.47), PptxInches(10.9), PptxInches(0.75))
    p_title = tb_title.text_frame.paragraphs[0]; p_title.text = str(item.get("title", f"Слайд {idx + 1}"))[:120]
    _set_text_style(p_title, 27 if len(p_title.text) < 65 else 22, theme["body_text"], True)

    points = [str(x).strip() for x in (item.get("points") or []) if str(x).strip()]
    points = points[:6] or ["Ключевая информация по теме будет добавлена автоматически."]
    has_img = img_stream is not None
    left_w = 7.05 if has_img else 11.7

    # Карточки с тезисами вместо длинного списка маркеров.
    card_h = 0.72 if len(points) <= 4 else 0.62
    y = 1.55
    for i, pt in enumerate(points):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.62), PptxInches(y), PptxInches(left_w), PptxInches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA) if theme["bg"] == RGBColor(0xFF,0xFF,0xFF) else theme["bg"]
        card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xEC)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, PptxInches(0.82), PptxInches(y + 0.16), PptxInches(0.32), PptxInches(0.32))
        badge.fill.solid(); badge.fill.fore_color.rgb = theme["primary"]; badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]; bp.text = str(i + 1); bp.alignment = PP_ALIGN.CENTER
        _set_text_style(bp, 8, theme["title_text"], True)
        tb = s.shapes.add_textbox(PptxInches(1.3), PptxInches(y + 0.10), PptxInches(left_w - 0.85), PptxInches(card_h - 0.15))
        tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
        pp = tf.paragraphs[0]; pp.text = pt[:220]
        _set_text_style(pp, 14 if len(pt) < 130 else 12, theme["body_text"])
        y += card_h + 0.13

    if has_img:
        frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(8.05), PptxInches(1.55), PptxInches(4.65), PptxInches(4.95))
        frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(0xEE,0xF1,0xF5); frame.line.color.rgb = RGBColor(0xD9,0xDE,0xE5)
        try:
            s.shapes.add_picture(img_stream, left=PptxInches(8.18), top=PptxInches(1.68), width=PptxInches(4.39), height=PptxInches(4.69))
        except Exception as img_err:
            logging.error(f"Не удалось вставить картинку на слайд {idx}: {img_err}")

    _add_footer(s, theme, idx + 2)
    return s


PPT_SLIDE_COUNT_LABELS = {6:"📑 6 слайдов",8:"📑 8 слайдов",10:"📑 10 слайдов",12:"📑 12 слайдов"}

def _ppt_design_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=v["label"],callback_data=k)] for k,v in PPT_THEMES.items()
    ]+[[InlineKeyboardButton(text="❌ Отмена",callback_data="cancel_flow")]])

def _ppt_count_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=PPT_SLIDE_COUNT_LABELS[6],callback_data="ppt_count_6"),InlineKeyboardButton(text=PPT_SLIDE_COUNT_LABELS[8],callback_data="ppt_count_8")],
        [InlineKeyboardButton(text=PPT_SLIDE_COUNT_LABELS[10],callback_data="ppt_count_10"),InlineKeyboardButton(text=PPT_SLIDE_COUNT_LABELS[12],callback_data="ppt_count_12")],
        [InlineKeyboardButton(text="❌ Отмена",callback_data="cancel_flow")]
    ])

async def _start_ppt_flow(message,user_id,topic=None):
    clear_pending_states(user_id)
    if topic: user_ppt_topic[user_id]=topic.strip()[:2500]
    await message.answer("🎨 <b>Шаг 1/3 — выбери дизайн презентации</b>\n\nПосле дизайна я предложу количество слайдов, затем запрошу тему.",parse_mode="HTML",reply_markup=_ppt_design_keyboard())

@dp.message(F.text == "📈 Презентация")
async def cmd_ppt_prompt(message: Message):
    user_id=message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",reply_markup=get_sub_keyboard()); return
    if user_id in busy_users: await message.answer("⏳ Подожди, предыдущая задача ещё выполняется..."); return
    topic=user_ppt_topic.get(user_id)
    await _start_ppt_flow(message,user_id,topic)

@dp.callback_query(F.data.in_(set(PPT_THEMES.keys())))
async def cb_ppt_design(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    uid=callback.from_user.id; user_ppt_design[uid]=callback.data; ppt_states.add(uid)
    try: await callback.message.edit_text(f"🎨 <b>Шаг 1/3 — дизайн:</b> {PPT_THEMES[callback.data]['label']}",parse_mode="HTML")
    except Exception: pass
    await callback.message.answer("📑 <b>Шаг 2/3 — сколько содержательных слайдов?</b>\n\nТитульный слайд добавится автоматически.",parse_mode="HTML",reply_markup=_ppt_count_keyboard())
    await callback.answer()

@dp.callback_query(F.data.startswith("ppt_count_"))
async def cb_ppt_count(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    uid=callback.from_user.id
    try: count=int(callback.data.rsplit("_",1)[1])
    except Exception: await callback.answer("Некорректное количество",show_alert=True); return
    if count not in PPT_SLIDE_COUNTS: await callback.answer("Недоступное количество",show_alert=True); return
    user_ppt_slide_count[uid]=count; ppt_states.add(uid); topic=user_ppt_topic.get(uid)
    try: await callback.message.edit_text(f"📑 <b>Шаг 2/3 — {count} содержательных слайдов</b>",parse_mode="HTML")
    except Exception: pass
    if topic:
        await callback.message.answer(f"📌 <b>Шаг 3/3 — тема уже получена:</b> «{topic[:300]}»\n\nОтправь свои картинки сейчас или нажми «Готово», чтобы начать.",parse_mode="HTML",reply_markup=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="🚀 Начать генерацию",callback_data="ppt_start")],[InlineKeyboardButton(text="❌ Отмена",callback_data="cancel_flow")]]))
    else:
        await callback.message.answer("📌 <b>Шаг 3/3 — отправь тему презентации.</b>\n\nПосле темы можно отправить свои изображения. Когда закончишь — напиши <b>«готово»</b>.",parse_mode="HTML")
    await callback.answer()

@dp.callback_query(F.data == "ppt_wait_images")
async def cb_ppt_wait_images(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    uid = callback.from_user.id
    if not user_ppt_topic.get(uid):
        await callback.answer("Сначала отправь тему", show_alert=True)
        return
    await callback.answer("Жду картинки")
    await callback.message.answer(
        "🖼 <b>Режим добавления картинок включён.</b>\n\n"
        "Отправь одну или несколько фотографий. После этого нажми «🚀 Создать» или напиши «готово».\n"
        "Если картинки не нужны — просто нажми «Создать без моих картинок»."
        , parse_mode="HTML", reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🚀 Создать", callback_data="ppt_start")],
            [InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")],
        ])
    )

@dp.callback_query(F.data == "ppt_start")
async def cb_ppt_start(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    uid=callback.from_user.id; topic=user_ppt_topic.get(uid)
    if not topic: await callback.answer("Сначала отправь тему",show_alert=True); return
    await callback.answer("Запускаю генерацию…"); await generate_presentation_file(callback.message,topic,uid)

@dp.callback_query(F.data == "cancel_flow")
async def cb_cancel_flow(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    uid = callback.from_user.id
    clear_pending_states(uid); user_ppt_slide_count.pop(uid,None)
    await callback.answer("Отменено")
    try:
        await callback.message.delete()
    except Exception:
        pass
    await callback.message.answer("🏠 <b>Главное меню</b>", parse_mode="HTML", reply_markup=keyboard_for(uid))


# ======================= EXCEL-ТАБЛИЦЫ =======================

EXCEL_MODES = {
    "excel_plain": "📋 Только таблица",
    "excel_chart": "📊 Таблица + диаграмма",
    "excel_full": "🖼 Таблица + диаграмма + ИИ-картинка",
}

@dp.message(F.text == "📊 Excel-таблица")
async def cmd_excel_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    topic = user_excel_topic.get(user_id)
    clear_pending_states(user_id)
    if topic:
        user_excel_topic[user_id] = topic
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=v, callback_data=k)] for k, v in EXCEL_MODES.items()
    ] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
    prompt = "📊 Выбери формат Excel:" if not topic else f"📊 Запрос: «{topic[:180]}»\nВыбери формат:"
    await message.answer(prompt, reply_markup=kb)

@dp.callback_query(F.data.in_(set(EXCEL_MODES.keys())))
async def cb_excel_mode(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    user_id = callback.from_user.id
    user_excel_mode[user_id] = callback.data
    topic = user_excel_topic.get(user_id)
    excel_states.add(user_id)
    mode_label = EXCEL_MODES[callback.data]
    try:
        await callback.message.edit_text(f"📊 Выбрано: {mode_label}")
    except Exception:
        pass
    if topic:
        await callback.message.answer(f"✅ Формат выбран: {mode_label}. Начинаю генерацию Excel…")
        await generate_excel_file(callback.message, topic, user_id=user_id)
    else:
        await callback.message.answer(
            "Опиши задачу или тему для таблицы следующим сообщением "
            "(например: «Таблица успеваемости студентов с оценками»), "
            "и я сформирую готовый Excel-файл (.xlsx)!"
        )
    await callback.answer()

# ======================= АНАЛИЗ ДОКУМЕНТОВ =======================

MAX_DOC_CONTEXT_CHARS = 12000
MAX_DOC_SIZE_BYTES = 20 * 1024 * 1024  # ограничение Telegram Bot API на скачивание файла

@dp.message(F.text == "📚 Анализ документа")
async def cmd_doc_analysis_prompt(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return
    clear_pending_states(user_id)
    doc_analysis_states.add(user_id)
    await message.answer(
        "📚 Пришли файл в формате <b>.pdf</b>, <b>.docx</b> или <b>.txt</b> (можно также прислать фото страницы) — "
        "я прочитаю его, сделаю краткое содержание, и дальше сможешь задавать вопросы по документу прямо в чате.",
        parse_mode="HTML"
    )

async def extract_text_from_document(ext: str, file_bytes: bytes):
    """Извлекает текст и структурированные фрагменты документа."""
    if ext == "txt":
        raw=file_bytes.decode("utf-8",errors="ignore")
        return raw,[{"page":None,"text":raw}],None
    if ext == "docx":
        try:
            d=Document(io.BytesIO(file_bytes)); blocks=[]
            for p in d.paragraphs:
                if p.text.strip(): blocks.append(p.text.strip())
            for table in d.tables:
                rows=[]
                for row in table.rows:
                    vals=[c.text.replace("\n"," ").strip() for c in row.cells]
                    if any(vals): rows.append(" | ".join(vals))
                if rows: blocks.append("[ТАБЛИЦА]\n"+"\n".join(rows))
            raw="\n\n".join(blocks)
            return raw,[{"page":None,"text":raw}],None
        except Exception as e:
            return None,None,f"⚠️ Не удалось прочитать DOCX: {e}"
    if ext == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return None,None,"⚠️ Для чтения PDF нужна библиотека pypdf. Установи: pip install pypdf"
        try:
            reader=PdfReader(io.BytesIO(file_bytes)); pages=[]
            for n,page in enumerate(reader.pages,1):
                t=(page.extract_text() or "").strip()
                if t: pages.append({"page":n,"text":t})
            raw="\n\n".join(f"[Страница {x['page']}]\n{x['text']}" for x in pages)
            return raw,pages,None
        except Exception as e:
            return None,None,f"⚠️ Не удалось прочитать PDF: {e}"
    return None,None,"⚠️ Поддерживаются .pdf, .docx и .txt."

def _chunk_document(pages, size=5000, overlap=600):
    out=[]
    for page in pages or []:
        raw=str(page.get("text") or "").strip()
        pos=0
        while raw and pos < len(raw):
            part=raw[pos:pos+size].strip()
            if part: out.append({"page":page.get("page"),"text":part})
            if pos+size>=len(raw): break
            pos += size-overlap
    return out

def _select_document_context(user_id, query="", limit=14000):
    data=user_doc_chunks.get(user_id) or {}; chunks=data.get("chunks") or []
    if not chunks: return user_doc_context.get(user_id,"")[:limit]
    words=re.findall(r"[a-zа-яё0-9]{3,}",(query or "").lower())[:50]
    scored=[]
    for i,ch in enumerate(chunks):
        low=ch["text"].lower(); score=sum(low.count(w) for w in words)
        score += max(0,2-i*0.02)
        scored.append((score,i,ch))
    scored.sort(key=lambda x:x[0],reverse=True)
    result=[]; total=0
    for _,_,ch in scored:
        block=(f"[Страница {ch['page']}]\n" if ch.get("page") else "")+ch["text"]
        if total+len(block)>limit: continue
        result.append(block); total+=len(block)
        if total>=limit*0.9: break
    return "\n\n---\n\n".join(result)

async def process_document_context(message: Message, filename: str, text: str, pages, status_msg: Message):
    user_id=message.from_user.id; text=(text or "").strip()
    if not text:
        await status_msg.edit_text("⚠️ Текстовый слой не найден. Если это скан, пришли страницы как фото — я разберу их через Vision.")
        doc_analysis_states.discard(user_id); return
    chunks=_chunk_document(pages or [{"page":None,"text":text}])
    user_doc_chunks[user_id]={"name":filename,"chunks":chunks}
    user_doc_context[user_id]=_select_document_context(user_id,"",14000)
    doc_analysis_states.discard(user_id)
    try:
        source=_select_document_context(user_id,"",18000)
        r=await call_groq_with_retry(messages=[
            {"role":"system","content":"Сделай структурированное резюме документа: тема, 5-8 главных тезисов, важные даты/числа, вывод и практические действия. Не выдумывай."},
            {"role":"user","content":source}
        ])
        summary=clean_text_for_html(r.choices[0].message.content)
    except Exception as e:
        logging.error("Document summary failed: %s",e,exc_info=True); summary="Резюме не сформировано, но документ доступен для вопросов."
    page_count=len({x.get("page") for x in pages or [] if x.get("page")})
    await status_msg.edit_text(f"✅ <b>{html.escape(filename)}</b> прочитан.\n📄 Страниц: {page_count or '—'}\n🧩 Фрагментов для поиска: {len(chunks)}",parse_mode="HTML")
    await safe_answer(message,f"📋 <b>Содержание:</b>\n\n{summary}\n\n💬 Теперь спрашивай конкретно — я ищу релевантные фрагменты по всему документу, а не только в его начале.",parse_mode="HTML")

@dp.message(F.document)
async def handle_document(message: Message):
    user_id = message.from_user.id
    _track_user_event(user_id, message=message, kind="document", status="received", prompt=message.document.file_name if message.document else "[document]")
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if not await _ingress_allowed(message, 3, 60, "документы"): return
    # Отдельный режим больше не обязателен: отправка поддержанного файла сама запускает анализ.
    doc_analysis_states.add(user_id)
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    get_user_stats(user_id)["files"] += 1
    save_stats_db()
    doc = message.document
    filename = doc.file_name or "документ"
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in ("pdf", "docx", "txt"):
        await message.answer("⚠️ Поддерживаются только форматы .pdf, .docx и .txt")
        return
    if doc.file_size and doc.file_size > MAX_DOC_SIZE_BYTES:
        await message.answer("⚠️ Файл слишком большой (лимит Telegram на скачивание — 20 МБ).")
        return

    busy_users.add(user_id)
    status_msg = await message.answer("📚 Читаю документ...")
    try:
        file_info = await bot.get_file(doc.file_id)
        downloaded = await bot.download_file(file_info.file_path)
        file_bytes = downloaded.read()

        text, pages, err = await extract_text_from_document(ext, file_bytes)
        file_bytes = None
        downloaded = None
        if err:
            await status_msg.edit_text(err)
            doc_analysis_states.discard(user_id)
            return

        await process_document_context(message, filename, text, pages, status_msg)
    except Exception as e:
        logging.error(f"Ошибка при обработке документа: {e}", exc_info=True)
        try:
            await status_msg.edit_text("⚠️ Ошибка при чтении документа. Попробуй ещё раз.")
        except Exception:
            pass
        doc_analysis_states.discard(user_id)
    finally:
        busy_users.discard(user_id)

@dp.message(F.text == "⭐ Избранное")
async def cmd_favorites(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(user_id)

    favs = load_favorites()
    uid_str = str(user_id)
    user_favs = favs.get(uid_str, [])

    if not user_favs:
        await message.answer("⭐ У тебя пока нет сохраненных ответов в избранном.")
        return

    await message.answer(f"⭐ Твои сохраненные ответы ({len(user_favs)}):")
    for idx, fav_text in enumerate(user_favs, 1):
        display_text = f"Сохранение #{idx}:\n\n{fav_text}"
        if len(display_text) > 4000:
            display_text = display_text[:3997] + "..."
        await message.answer(display_text, disable_web_page_preview=True)

@dp.callback_query(F.data.in_({
    "guide_voice", "guide_image", "guide_doc", "guide_ppt", "guide_excel", "guide_word", "create_ppt", "create_excel", "create_word", "create_title"
}))
async def cb_product_navigation(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    data = callback.data
    messages = {
        "guide_voice": "🎙 Отправь голосовое — ничего включать не нужно.",
        "guide_image": "🖼 Отправь фото — я могу разобрать задачу, текст или изображение. Для презентации свои картинки можно добавить прямо в процессе её создания.",
        "guide_doc": "📁 Отправь PDF, DOCX или TXT — после загрузки задавай вопросы.",
        "guide_ppt": "📈 Нажми «Создать → Презентация». Я буду задавать вопросы по одному.",
        "guide_excel": "📊 Нажми «Создать → Excel» и опиши, какую таблицу хочешь.",
        "guide_word": "📄 Под ответом нажми «В Word», чтобы получить файл.",
        "create_ppt": "📈 Хорошо. Я помогу выбрать назначение, стиль и титульный лист — без сложных терминов.",
        "create_excel": "📊 Опиши нужную таблицу или пришли пример.",
        "create_word": "📄 Получи ответ и нажми «В Word».",
        "create_title": "📑 Я сначала уточню, для чего титульник нужен, и предложу понятные варианты."
    }
    if data == "create_ppt":
        user_id = callback.from_user.id
        await _start_ppt_flow(callback.message,user_id)
    elif data == "create_excel":
        user_id = callback.from_user.id
        clear_pending_states(user_id)
        kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text=v, callback_data=k)] for k, v in EXCEL_MODES.items()] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
        await callback.message.answer(
            "📊 <b>Создание Excel</b>\n\n"
            "Выбери формат. После выбора просто напиши, какую таблицу нужно сделать.",
            parse_mode="HTML", reply_markup=kb
        )
    elif data == "create_title":
        await callback.message.answer("📑 Выбери тип титульного листа:", reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📘 Индивидуальный проект", callback_data="gost_project")],
            [InlineKeyboardButton(text="📗 Курсовая работа", callback_data="gost_coursework")],
            [InlineKeyboardButton(text="📕 Дипломная работа (ВКР)", callback_data="gost_diploma")],
            [InlineKeyboardButton(text="📙 Отчёт по практике", callback_data="gost_practice")]
        ]))
    else:
        await callback.message.answer(messages.get(data, "Готово."))
    await callback.answer()

@dp.callback_query(F.data == "flow_restart")
async def cb_flow_restart(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    clear_task_state(callback.from_user.id)
    await callback.message.answer(
        "🔄 Начинаем заново.\n\n"
        "Просто напиши, что хочешь сделать, или выбери действие ниже.",
        reply_markup=get_quick_actions_keyboard()
    )
    await callback.answer()

@dp.callback_query(F.data.in_({"btn_simplify", "btn_structure", "btn_check", "btn_deepen", "btn_examples", "btn_academic", "btn_save_fav", "btn_word", "btn_make_ppt", "btn_continue"}))
async def cb_answer_actions(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback):
        return
    user_id = callback.from_user.id
    msg_text = callback.message.text or callback.message.caption or ""

    clean_msg = msg_text.split("—\n⚡")[0].strip() if "—\n⚡" in msg_text else msg_text.strip()
    if not clean_msg:
        await callback.answer("⚠️ Нечего обрабатывать!", show_alert=True)
        return

    if callback.data == "btn_word":
        user_data = get_user_data(user_id)
        user_data["last_output"] = clean_msg
        user_data_stats = get_user_stats(user_id)
        user_data_stats["exports"] += 1

        doc = build_answer_only_word(clean_msg)
        bio = io.BytesIO(); doc.save(bio); raw = bio.getvalue(); bio.close()
        path = _artifact_path(user_id, "docx", "Answer")
        _atomic_write_bytes(path, raw)
        register_artifact(user_id, "docx", path, {"title": "Ответ"})
        await callback.message.answer_document(
            BufferedInputFile(raw, filename=os.path.basename(path)),
            caption="📄 Готово — это только ответ ИИ, без титульника и содержания."
        )
        await callback.answer()
        return

    if callback.data == "btn_continue":
        await callback.message.answer(
            "🔁 Продолжаем. Напиши следующим сообщением, что изменить, добавить или уточнить — контекст сохранён."
        )
        await callback.answer()
        return

    if callback.data == "btn_make_ppt":
        user_ppt_topic[user_id]=clean_msg[:2500]
        await _start_ppt_flow(callback.message,user_id,clean_msg[:2500])
        await callback.answer("📈 Запускаю создание презентации…")
        return

    if callback.data in {"btn_structure","btn_check","btn_deepen","btn_examples","btn_academic"}:
        await callback.answer("🧠 Обрабатываю…")
        instruction={"btn_check":"Проведи строгую проверку текста: фактические ошибки, логические пробелы, противоречия, сомнительные утверждения и что стоит уточнить. Не выдумывай ошибки.","btn_structure":"Перестрой текст в максимально понятную структуру: краткий вывод, основные тезисы, пошаговое решение/алгоритм, примеры и итог. Сохрани исходный смысл.","btn_deepen":"Существенно углуби ответ: добавь детали, причины, последствия, ограничения, нюансы и полезные выводы. Не выдумывай факты.","btn_examples":"Добавь 3-5 конкретных примеров, мини-кейсов или аналогий, которые помогают понять материал. Не меняй исходные факты.","btn_academic":"Перепиши ответ в строгом академическом стиле: точные формулировки, логичная аргументация, определения, нейтральный тон и вывод. Не добавляй неподтверждённых фактов."}[callback.data]
        try:
            r=await call_groq_with_retry(messages=[{"role":"system","content":instruction},{"role":"user","content":clean_msg}])
            result=clean_text_for_html(r.choices[0].message.content)
            title={"btn_check":"🔍 <b>Проверка ответа</b>","btn_structure":"🧩 <b>Структурированная версия</b>","btn_deepen":"🧠 <b>Углублённая версия</b>","btn_examples":"💡 <b>Версия с примерами</b>","btn_academic":"🎓 <b>Академическая версия</b>"}[callback.data]
            await safe_answer(callback.message,f"{title}\n\n{result}{AD_FOOTER}",parse_mode="HTML",reply_markup=get_answer_inline_keyboard())
        except Exception as e:
            logging.error("Answer transform failed: %s",e,exc_info=True)
            await callback.message.answer("⚠️ Не удалось выполнить преобразование. Попробуй ещё раз.")
        return

    if callback.data == "btn_save_fav":
        add_to_favorites(user_id, clean_msg)
        await callback.answer("⭐ Успешно сохранено в избранное!", show_alert=True)
    elif callback.data == "btn_simplify":
        await callback.answer("💡 Сжимаю до сути...")
        try:
            response = await call_groq_with_retry(messages=[
                {"role": "system", "content": "Объясни предельно коротко и просто в 3-5 предложениях."},
                {"role": "user", "content": clean_msg}
            ])
            simplified_reply = clean_text_for_html(response.choices[0].message.content)
            full_reply = f"💡 <b>Коротко на пальцах:</b>\n\n{simplified_reply}{AD_FOOTER}"
            await safe_answer(callback.message, full_reply, parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
        except Exception:
            await callback.message.answer("⚠️ Ошибка при обработке. Попробуй ещё раз. Если повторится, пиши - @mecau")

def build_answer_only_word(text):
    """Экспорт ответа ИИ: только содержание ответа, без титульника/содержания/брендинга."""
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Mm(20); sec.bottom_margin = Mm(20)
    sec.left_margin = Mm(25); sec.right_margin = Mm(20)
    styles = doc.styles
    try:
        normal = styles['Normal']
        normal.font.name = 'Times New Roman'
        normal.font.size = Pt(14)
    except Exception:
        pass
    clean = re.sub(r'<br\s*/?>', '\\n', text or '', flags=re.I)
    clean = re.sub(r'</?(?:b|strong|i|em|u|code|pre)>', '', clean, flags=re.I)
    clean = re.sub(r'<[^>]+>', '', clean)
    clean = html.unescape(clean)
    clean = clean.replace(AD_FOOTER, '').strip()
    for block in re.split(r'\\n\\s*\\n', clean):
        block=block.strip()
        if not block: continue
        lines=block.splitlines()
        for line in lines:
            line=line.strip()
            if not line: continue
            if re.match(r'^#{1,3}\\s+', line):
                level=len(line)-len(line.lstrip('#'))
                doc.add_paragraph(line[level:].strip(), style=f'Heading {min(level,3)}')
            elif re.match(r'^[-•*]\\s+', line):
                doc.add_paragraph(re.sub(r'^[-•*]\\s+', '', line), style='List Bullet')
            elif re.match(r'^\d+[.)]\\s+', line):
                doc.add_paragraph(re.sub(r'^\d+[.)]\\s+', '', line), style='List Number')
            else:
                p=doc.add_paragraph(line)
                p.paragraph_format.space_after=Pt(6)
    return doc

def build_professional_word(text, title="MecauAI — документ"):
    """Word pipeline: Heading 1/2/3, lists, TOC field, footer/page numbers, margins."""
    doc=Document()
    sec=doc.sections[0]
    sec.top_margin=Mm(20); sec.bottom_margin=Mm(20); sec.left_margin=Mm(25); sec.right_margin=Mm(20)
    styles=doc.styles
    for name,size,bold in (("Normal",11,False),("Heading 1",16,True),("Heading 2",14,True),("Heading 3",12,True)):
        try:
            st=styles[name]; st.font.name='Aptos'; st.font.size=Pt(size); st.font.bold=bold
        except Exception: pass
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=p.add_run(title[:150]); r.bold=True; r.font.name='Aptos Display'; r.font.size=Pt(20)
    doc.add_paragraph()
    toc=doc.add_paragraph(); toc.add_run('Содержание').bold=True
    fld=toc.add_run(); fld._r.append(OxmlElement('w:fldChar')); fld._r[-1].set(qn('w:fldCharType'),'begin')
    instr=OxmlElement('w:instrText'); instr.set(qn('xml:space'),'preserve'); instr.text='TOC \\o "1-3" \\h \\z \\u'; fld._r.append(instr)
    end=OxmlElement('w:fldChar'); end.set(qn('w:fldCharType'),'end'); fld._r.append(end)
    doc.add_page_break()
    lines=str(text or '').replace('\r','').split('\n')
    for raw in lines:
        line=raw.strip()
        if not line: continue
        if line.startswith('### '): doc.add_paragraph(line[4:].strip(),style='Heading 3'); continue
        if line.startswith('## '): doc.add_paragraph(line[3:].strip(),style='Heading 2'); continue
        if line.startswith('# '): doc.add_paragraph(line[2:].strip(),style='Heading 1'); continue
        if re.match(r'^[-•*]\s+',line):
            p=doc.add_paragraph(re.sub(r'^[-•*]\s+','',line),style='List Bullet'); continue
        if re.match(r'^\d+[.)]\s+',line):
            doc.add_paragraph(re.sub(r'^\d+[.)]\s+','',line),style='List Number'); continue
        p=doc.add_paragraph(line,style='Normal'); p.paragraph_format.space_after=Pt(6)
    # Footer with page number field.
    footer=sec.footer.paragraphs[0]; footer.alignment=WD_ALIGN_PARAGRAPH.CENTER; footer.add_run('MecauAI • ')
    run=footer.add_run(); fld=OxmlElement('w:fldSimple'); fld.set(qn('w:instr'),'PAGE'); run._r.addnext(fld)
    return doc

@dp.message(F.text == "📄 Ответ в Word")
async def cmd_download_word(message: Message):
    user_id = message.from_user.id
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(user_id)

    user_data = get_user_data(user_id)
    text_to_save = user_data.get("last_output", "").strip()
    if not text_to_save or text_to_save == "Здесь пока нет ответов.":
        await message.answer("⚠️ Нет данных для сохранения. Попробуй ещё раз. Если повторится, пиши - @mecau")
        return

    doc = build_answer_only_word(text_to_save)
    bio = io.BytesIO(); doc.save(bio); raw = bio.getvalue(); bio.close()
    path = _artifact_path(user_id, "docx", "Answer")
    _atomic_write_bytes(path, raw)
    register_artifact(user_id, "docx", path, {"title": "Ответ"})
    await message.answer_document(BufferedInputFile(raw, filename=os.path.basename(path)), caption="📄 Готово — это только ответ ИИ, без титульника и содержания.")

@dp.message(F.text == "📑 Титульник ГОСТ")
async def cmd_gost_title(message: Message):
    if not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    clear_pending_states(message.from_user.id)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📘 Индивидуальный проект", callback_data="gost_project")],
        [InlineKeyboardButton(text="📗 Курсовая работа",        callback_data="gost_coursework")],
        [InlineKeyboardButton(text="📕 Дипломная работа (ВКР)", callback_data="gost_diploma")],
        [InlineKeyboardButton(text="📙 Отчёт по практике",      callback_data="gost_practice")],
    ])
    await message.answer("📑 Выбери тип работы для титульника:", reply_markup=kb)

async def _generate_gost_title_for_user(message: Message, state: dict):
    user_id = message.from_user.id
    work_label = state["work_label"]
    filename_base = state["filename_base"]
    institution = state["institution"].strip()
    student = state["student"].strip()
    teacher = state["teacher"].strip()
    topic = state["topic"].strip()
    city = state["city"].strip()
    year = state["year"]

    doc = Document()
    sec = doc.sections[0]
    sec.top_margin=Mm(20); sec.bottom_margin=Mm(20); sec.left_margin=Mm(30); sec.right_margin=Mm(15)
    try:
        normal=doc.styles['Normal']; normal.font.name='Times New Roman'; normal.font.size=Pt(14)
    except Exception: pass

    def para(text='', align=WD_ALIGN_PARAGRAPH.CENTER, bold=False, size=14, space_after=0):
        p=doc.add_paragraph(); p.alignment=align; p.paragraph_format.space_after=Pt(space_after)
        r=p.add_run(text); r.font.name='Times New Roman'; r.font.size=Pt(size); r.bold=bold
        return p
    def empty(n=1):
        for _ in range(n): doc.add_paragraph()

    # Учебное заведение
    para(institution, bold=True)
    empty(4)
    para(work_label, bold=True)
    empty(1)
    para(f'на тему:\n«{topic}»', bold=True)
    empty(6)

    # Исполнитель/руководитель блок — без выдумывания группы или должности.
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.RIGHT
    for label, value in (("Выполнил(а):", student),("Руководитель:", teacher)):
        r=p.add_run(label+'\n'); r.font.name='Times New Roman'; r.font.size=Pt(14)
        r2=p.add_run(value+'\n\n'); r2.font.name='Times New Roman'; r2.font.size=Pt(14)
    empty(4)
    para(f'{city}, {year}', size=14)

    bio=io.BytesIO(); doc.save(bio); raw=bio.getvalue(); bio.close()
    filename=f'Titulnik_{filename_base}.docx'
    await message.answer_document(BufferedInputFile(raw, filename=filename), caption=f'📑 Титульник «{work_label}» готов!')
    title_page_states.pop(user_id, None)


def _title_preview_text(state: dict) -> str:
    return (
        '📑 <b>Предпросмотр титульника</b>\n\n'
        f'🏫 <b>Учебное заведение:</b> {html.escape(state["institution"])}\n'
        f'👨‍🎓 <b>Студент:</b> {html.escape(state["student"])}\n'
        f'👨‍🏫 <b>Преподаватель:</b> {html.escape(state["teacher"])}\n'
        f'📚 <b>Тема:</b> {html.escape(state["topic"])}\n'
        f'📍 <b>Город, год:</b> {html.escape(state["city"])}, {state["year"]}\n\n'
        f'📘 <b>Тип:</b> {html.escape(state["work_label"])}\n\n'
        'Проверь данные перед созданием файла.'
    )


def _title_preview_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text='✅ Всё верно — создать', callback_data='title_confirm')],
        [InlineKeyboardButton(text='✏️ Заполнить заново', callback_data='title_edit')],
        [InlineKeyboardButton(text='❌ Отмена', callback_data='cancel_flow')],
    ])


# Голосовой режим намеренно не добавлен в меню: достаточно отправить voice-сообщение.
# Это сохраняет интерфейс компактным и не перегружает пользователя кнопками.
MENU_BUTTONS = {
    "✨ Возможности", "📁 Документы", "🛠 Техподдержка", "⭐ Избранное", "⚙️ Настройки ИИ",
    "📢 Рассылка", "📊 Статистика", "🖥 Сервер", "👥 Пользователи", "📜 История",
    # Старые кнопки тоже игнорируем, если они остались у уже открытого клиента.
    "📄 Ответ в Word", "📑 Титульник ГОСТ", "📈 Презентация",
    "📊 Excel-таблица", "📚 Анализ документа", "ℹ️ О MecauAI"
}


async def _handle_title_page_text(message: Message):
    user_id=message.from_user.id
    state=title_page_states.get(user_id)
    if not state: return False
    step=state.get('step')
    raw=(message.text or '').strip()

    if step=='institution':
        if len(raw)<4:
            await message.answer('⚠️ Укажи полное название колледжа/вуза.'); return True
        state['institution']=raw[:500]; state['step']='student'
        await message.answer('2️⃣ <b>ФИО студента</b>\n\nНапример: <code>Иванов Иван Иванович</code>',parse_mode='HTML'); return True
    if step=='student':
        if len(raw)<5:
            await message.answer('⚠️ Укажи полное ФИО студента.'); return True
        state['student']=raw[:250]; state['step']='teacher'
        await message.answer('3️⃣ <b>ФИО преподавателя</b>\n\nНапример: <code>Петров Пётр Петрович</code>',parse_mode='HTML'); return True
    if step=='teacher':
        if len(raw)<5:
            await message.answer('⚠️ Укажи полное ФИО преподавателя.'); return True
        state['teacher']=raw[:250]; state['step']='topic'
        await message.answer('4️⃣ <b>Тема работы</b>\n\nНапиши точное название темы.',parse_mode='HTML'); return True
    if step=='topic':
        if len(raw)<3:
            await message.answer('⚠️ Тема слишком короткая. Укажи полное название.'); return True
        state['topic']=raw[:1000]; state['step']='location_year'
        await message.answer('5️⃣ <b>Город и год</b>\n\nНапиши в формате:\n<code>Москва, 2026</code>',parse_mode='HTML'); return True
    if step=='location_year':
        m=re.match(r'^\s*(.+?)\s*[,;]\s*(\d{4})\s*$',raw) or re.match(r'^\s*(.+?)\s+(\d{4})\s*$',raw)
        if not m:
            await message.answer('⚠️ Не понял формат.\nПример: <code>Москва, 2026</code>',parse_mode='HTML'); return True
        city=m.group(1).strip(' ,;.-'); year=int(m.group(2))
        if not city or not (1900<=year<=2200):
            await message.answer('⚠️ Проверь город и год.\nПример: <code>Москва, 2026</code>',parse_mode='HTML'); return True
        state['city']=city; state['year']=year; state['step']='confirm'
        await message.answer(_title_preview_text(state),parse_mode='HTML',reply_markup=_title_preview_keyboard()); return True
    if step=='confirm':
        await message.answer('ℹ️ Проверь предпросмотр выше и нажми кнопку подтверждения.'); return True
    return False

@dp.callback_query(F.data.startswith('gost_'))
async def cb_gost_generate(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback): return
    work_type_map={
        'gost_project':('ИНДИВИДУАЛЬНЫЙ ПРОЕКТ','Индивидуальный_проект'),
        'gost_coursework':('КУРСОВАЯ РАБОТА','Курсовая_работа'),
        'gost_diploma':('ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА (ВКР)','ВКР'),
        'gost_practice':('ОТЧЁТ ПО ПРАКТИКЕ','Отчет_по_практике')}
    if callback.data not in work_type_map:
        await callback.answer('Неизвестный тип работы.',show_alert=True); return
    work_label,filename_base=work_type_map[callback.data]
    title_page_states[callback.from_user.id]={'work_label':work_label,'filename_base':filename_base,'step':'institution'}
    await callback.message.answer(
        f'📑 <b>{html.escape(work_label)}</b>\n\n'
        '1️⃣ <b>Полное название колледжа / вуза</b>\n\n'
        'Например: <code>Московский государственный технический университет</code>',parse_mode='HTML')
    await callback.answer('🏫 Начинаем')

@dp.callback_query(F.data=='title_confirm')
async def cb_title_confirm(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback): return
    state=title_page_states.get(callback.from_user.id)
    if not state or state.get('step')!='confirm':
        await callback.answer('Сначала заполни данные титульника.',show_alert=True); return
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception: pass
    await callback.answer('Создаю титульник…')
    try:
        await _generate_gost_title_for_user(callback.message,state)
    except Exception as exc:
        logging.error('Title generation failed: %s',exc,exc_info=True)
        title_page_states.pop(callback.from_user.id,None)
        await callback.message.answer('⚠️ Не удалось создать титульник. Попробуй ещё раз.')

@dp.callback_query(F.data=='title_edit')
async def cb_title_edit(callback: types.CallbackQuery):
    if not await require_subscription_callback(callback): return
    state=title_page_states.get(callback.from_user.id)
    if not state:
        await callback.answer('Сценарий уже завершён.',show_alert=True); return
    work_label=state.get('work_label'); filename_base=state.get('filename_base')
    state.clear(); state.update({'work_label':work_label,'filename_base':filename_base,'step':'institution'})
    await callback.answer('Заполняем заново')
    await callback.message.answer('🏫 <b>Полное название колледжа / вуза</b>\n\nНапример: <code>Московский государственный технический университет</code>',parse_mode='HTML')
    # Restore type from callback message is not available; preserve a local backup before clear is safer.

@dp.message(F.voice)
async def handle_voice(message: Message):
    """
    Голосовой режим без отдельной кнопки:
    пользователь отправляет voice -> бот распознаёт речь -> использует
    полученный текст как обычный запрос ИИ.
    """
    user_id = message.from_user.id
    _track_user_event(user_id, message=message, kind="voice", status="received", prompt="[голосовое сообщение]")

    if not await check_subscription(user_id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    if not await _ingress_allowed(message, 6, 60, "голосовые"): return

    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    save_user_id(user_id)
    get_user_stats(user_id)["voice"] += 1
    save_stats_db()
    status_msg = await message.answer("🎙 Распознаю голосовое сообщение...")

    busy_users.add(user_id)
    try:
        file_info = await bot.get_file(message.voice.file_id)
        downloaded = await bot.download_file(file_info.file_path)
        audio_bytes = downloaded.read()

        if len(audio_bytes) > MAX_DOC_SIZE_BYTES:
            await status_msg.edit_text("⚠️ Голосовое сообщение слишком большое. Пришли более короткую запись.")
            return

        transcript = await transcribe_voice(audio_bytes)
        if not transcript:
            await status_msg.edit_text("⚠️ Не удалось разобрать речь. Попробуй записать сообщение ещё раз.")
            return

        # Сначала учитываем уже выбранный workflow. Раньше voice обходил
        # ppt_states/excel_states и поэтому тема уходила в обычный AI-чат.
        voice_norm = transcript.lower().strip()

        if user_id in ppt_states:
            topic = transcript.strip()
            if not topic:
                await status_msg.edit_text("⚠️ Не удалось определить тему презентации.")
                return
            user_ppt_topic[user_id] = topic[:2500]
            if user_id not in user_ppt_slide_count:
                await status_msg.delete()
                await _start_ppt_flow(message,user_id,topic)
                return
            ppt_states.discard(user_id)
            await status_msg.edit_text(f"🎙 <b>Распознал тему:</b> «{topic[:500]}»\n\n⏳ Начинаю создание…",parse_mode="HTML")
            await generate_presentation_file(message, topic, user_id=user_id)
            return

        if user_id in excel_states:
            topic = transcript.strip()
            if topic:
                excel_states.discard(user_id)
                user_excel_topic[user_id] = topic
                await status_msg.edit_text(
                    f"🎙 <b>Распознал задачу для Excel:</b>\n«{topic[:500]}»\n\n⏳ Начинаю создание…",
                    parse_mode="HTML"
                )
                await generate_excel_file(message, topic, user_id=user_id)
            else:
                await status_msg.edit_text("⚠️ Не удалось определить задачу для Excel.")
            return

        # Естественные голосовые команды: не требуем строго фразу «сделай ...».
        if re.search(r"\b(презентаци[яиюеи]|слайд(?:ы|ов)?)\b", voice_norm) and re.search(r"\b(сделай|создай|подготовь|сформируй|сделать|нужна|нужен|нужно)\b", voice_norm):
            topic = re.sub(r"^.*?\bпрезентаци(?:ю|я|и|й)?\b\s*(?:на тему|по теме|про|о)?\s*", "", transcript, flags=re.I).strip(" .:-") or "Тема не указана"
            user_ppt_topic[user_id] = topic
            ppt_states.clear()
            user_ppt_design.pop(user_id, None); user_ppt_slide_count.pop(user_id,None)
            await status_msg.delete()
            await _start_ppt_flow(message,user_id,topic)
            return
        if re.search(r"\b(excel|таблиц(?:а|у|ы)?|xlsx)\b", voice_norm) and re.search(r"\b(сделай|создай|подготовь|сформируй|сделать|нужна|нужен|нужно)\b", voice_norm):
            topic = re.sub(r"^.*?\b(?:excel|таблиц(?:у|а|ы)?|xlsx)\b\s*(?:на тему|по теме|для|с)?\s*", "", transcript, flags=re.I).strip(" .:-") or transcript
            user_excel_topic[user_id] = topic
            excel_states.discard(user_id)
            kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text=v, callback_data=k)] for k,v in EXCEL_MODES.items()] + [[InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")]])
            await status_msg.edit_text(f"🎙 Я понял: создать Excel по запросу «{topic[:180]}».\nВыбери формат:", reply_markup=kb)
            return

        user_data = get_user_data(user_id)
        current_mode = user_data["mode"]
        system_prompt = PROMPTS.get(current_mode, PROMPTS["ai"])

        doc_context = user_doc_context.get(user_id)
        if doc_context:
            system_prompt = (
                f"{system_prompt}\n\n"
                f"Пользователь ранее загрузил документ. Используй его содержимое, "
                f"если вопрос с ним связан:\n---\n{doc_context}\n---"
            )

        # Голос обрабатывается как обычный текстовый запрос, поэтому история,
        # режим ИИ и экспорт последнего ответа продолжают работать одинаково.
        user_data["history"].append({"role": "user", "content": transcript})
        if len(user_data["history"]) > 6:
            user_data["history"] = user_data["history"][-6:]

        messages_payload = [{"role": "system", "content": system_prompt}] + user_data["history"]

        await status_msg.edit_text(f"🎙 Распознал:\n\n«{transcript[:3500]}»\n\n🤖 Формирую ответ...")
        await bot.send_chat_action(chat_id=message.chat.id, action="typing")

        response = await call_groq_with_retry(messages=messages_payload)
        ai_reply = clean_text_for_html(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply
        save_users_db()

        full_message = (
            f"🎙 <b>Твой запрос:</b> {transcript}\n\n"
            f"🤖 <b>Ответ:</b>\n\n{ai_reply}{AD_FOOTER}"
        )

        try:
            await status_msg.delete()
        except Exception:
            pass

        parse_mode = "Markdown" if current_mode == "coder" else "HTML"
        await safe_answer(
            message,
            full_message,
            parse_mode=parse_mode,
            reply_markup=get_answer_inline_keyboard()
        )

    except Exception as e:
        logging.error(f"Ошибка обработки голосового сообщения: {e}", exc_info=True)
        try:
            await status_msg.edit_text(
                "⚠️ Не удалось обработать голосовое сообщение. "
                "Попробуй записать его ещё раз чуть короче и без сильного фонового шума."
            )
        except Exception:
            pass
    finally:
        busy_users.discard(user_id)

@dp.message(F.photo)
async def handle_photo(message: Message):
    user_id = message.from_user.id
    _track_user_event(user_id, message=message, kind="vision", status="received", prompt=message.caption or "[изображение]")
    if not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}", reply_markup=get_sub_keyboard())
        return
    if not await _ingress_allowed(message, 8, 60, "изображения"): return

    save_user_id(user_id)
    get_user_stats(user_id)["images"] += 1
    save_stats_db()
    photo = message.photo[-1]

    # Для презентации храним только Telegram file_id, а не байты картинки.
    # Это резко снижает RAM на слабом VPS.
    if user_id in ppt_states:
        images = user_ppt_images.setdefault(user_id, [])
        if len(images) >= MAX_PPT_IMAGES:
            await message.answer(f"⚠️ Для одной презентации можно добавить максимум {MAX_PPT_IMAGES} картинок.")
            return
        if photo.file_id not in images:
            images.append(photo.file_id)
        await message.answer(f"🖼 Картинка добавлена в презентацию ({len(images)}/{MAX_PPT_IMAGES}). Можешь отправить ещё или продолжить с темой.")
        return

    # Скачивание фото выполняем только когда оно действительно нужно Vision.
    try:
        if photo.file_size and photo.file_size > MAX_DOC_SIZE_BYTES:
            await message.answer("⚠️ Фото слишком большое (лимит Telegram на скачивание — 20 МБ). Пришли файл меньшего размера.")
            return
        file_info = await bot.get_file(photo.file_id)
        downloaded_file = await bot.download_file(file_info.file_path)
        img_bytes = downloaded_file.read()
    except Exception as e:
        logging.error(f"Не удалось скачать фото от пользователя {user_id}: {e}")
        await message.answer("⚠️ Не удалось загрузить фото. Проверь соединение и попробуй отправить ещё раз.")
        return

    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется...")
        return

    # Если пользователь в режиме анализа документа — считаем фото сканом страницы
    # и распознаём текст через vision-модель вместо обычного анализа изображения.
    if user_id in doc_analysis_states:
        busy_users.add(user_id)
        status_msg = await message.answer("📚 Распознаю текст на фото...")
        try:
            base64_image = base64.b64encode(img_bytes).decode('utf-8')
            response = await call_groq_with_retry(
                model=VISION_MODEL,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Распознай и выведи весь текст, присутствующий на этом изображении, без комментариев."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                    ]
                }]
            )
            recognized_text = clean_text_for_html(response.choices[0].message.content)
            base64_image = None
            img_bytes = None
            await process_document_context(message, "фото документа", recognized_text, status_msg)
        except Exception as e:
            logging.error(f"Ошибка распознавания текста с фото: {e}", exc_info=True)
            try:
                await status_msg.edit_text("⚠️ Не удалось распознать текст на фото. Попробуй ещё раз.")
            except Exception:
                pass
            doc_analysis_states.discard(user_id)
        finally:
            busy_users.discard(user_id)
        return

    # Обычный режим анализа картинки через Vision ИИ
    busy_users.add(user_id)
    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        # Base64 temporarily doubles memory; reject unexpectedly huge images early.
        if len(img_bytes) > 8 * 1024 * 1024:
            raise ValueError("image too large")
        base64_image = base64.b64encode(img_bytes).decode('ascii')
        response = await call_groq_with_retry(
            model=VISION_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": message.caption or "Подробно проанализируй это изображение, распознай текст если он есть и ответь по сути."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }]
        )
        reply = clean_text_for_html(response.choices[0].message.content)
        base64_image = None
        img_bytes = None
        user_data = get_user_data(user_id)
        user_data["last_output"] = reply
        save_users_db()
        await safe_answer(message, f"{reply}{AD_FOOTER}", parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
    except Exception as e:
        logging.error(f"Ошибка обработки фото: {e}")
        await message.answer("⚠️ Ошибка при обработке изображения. Попробуй ещё раз. Если повторится, пиши - @mecau")
    finally:
        busy_users.discard(user_id)

# ВАЖНО: эти обработчики зарегистрированы ДО общего @dp.message(F.text).
# Иначе /admin и кнопка "🛡 Админ-панель" попадали в ИИ как обычный текст.
@dp.message(Command("admin"))
async def admin_panel(message: Message):
    await _admin_panel_message(message)

@dp.message(F.text == "🛡 Админ-панель")
async def admin_panel_button(message: Message):
    await _admin_panel_message(message)

@dp.message(F.text)
async def handle_text(message: Message):
    user_id = message.from_user.id
    _track_user_event(user_id, message=message, kind="text", status="received", prompt=message.text or "")

    # Естественные команды: пользователь может запускать функции словами,
    # не изучая меню.
    normalized = (message.text or "").lower().strip()

    # Админские вводы обрабатываем ПЕРЕД проверкой подписки и любым ИИ.
    # Иначе ввод Telegram ID/текста рассылки мог уйти в обычный AI-flow.
    if user_id in admin_grant_states or -user_id in admin_grant_states:
        if not _admin_is_allowed(user_id) or message.chat.type != "private":
            admin_grant_states.discard(user_id)
            admin_grant_states.discard(-user_id)
            await message.answer("⛔ Доступ к этому действию запрещён.")
            return
        is_revoke = -user_id in admin_grant_states
        admin_grant_states.discard(user_id)
        admin_grant_states.discard(-user_id)
        raw = (message.text or "").strip()
        if not raw.lstrip("-").isdigit():
            await message.answer("⚠️ Нужен числовой Telegram ID.")
            return
        target = int(raw)
        if target == MY_ADMIN_ID:
            await message.answer("ℹ️ Главного владельца изменить нельзя.", reply_markup=_admin_back_keyboard())
            return
        if is_revoke:
            ADMIN_IDS.discard(target)
            _save_admin_ids()
            _admin_log(user_id, "revoke_admin", str(target))
            await message.answer(
                f"✅ Доступ администратора забран у <code>{target}</code>.",
                parse_mode="HTML", reply_markup=_admin_back_keyboard()
            )
        else:
            ADMIN_IDS.add(target)
            _save_admin_ids()
            _admin_log(user_id, "grant_admin", str(target))
            await message.answer(
                f"✅ Пользователь <code>{target}</code> теперь администратор.",
                parse_mode="HTML", reply_markup=_admin_back_keyboard()
            )
        return

    if user_id in admin_broadcast_states:
        if not _admin_is_allowed(user_id) or message.chat.type != "private":
            admin_broadcast_states.discard(user_id)
            await message.answer("⛔ Доступ к рассылке запрещён.")
            return
        admin_broadcast_states.discard(user_id)
        status_msg = await message.answer("📢 Начинаю рассылку...")
        success = failed = 0
        for uid in sorted(_all_known_user_ids()):
            try:
                await bot.send_message(uid, message.text, disable_web_page_preview=True)
                success += 1
            except TelegramRetryAfter as e:
                await asyncio.sleep(e.retry_after)
                try:
                    await bot.send_message(uid, message.text, disable_web_page_preview=True)
                    success += 1
                except Exception:
                    failed += 1
            except Exception:
                failed += 1
            await asyncio.sleep(0.05)
        await status_msg.edit_text(
            f"✅ Рассылка завершена!\n\n👥 Доставлено: {success}\n❌ Ошибок: {failed}",
            reply_markup=_admin_back_keyboard()
        )
        return

    # Доступ к редактированию/конвертации артефактов тоже защищён подпиской.
    if not await check_subscription(user_id):
        await message.answer(
            f"🔒 Сначала подпишись на каналы:\n1️⃣ {CHANNEL_1_URL}\n2️⃣ {CHANNEL_2_URL}",
            reply_markup=get_sub_keyboard()
        )
        return
    if not await _ingress_allowed(message, 15, 30, "сообщения"): return

    # Сначала обслуживаем многошаговый титульник — отдельного catch-all handler больше нет.
    if await _handle_title_page_text(message):
        return

    # Если вопрос относится к загруженному документу, подмешиваем релевантные фрагменты,
    # а не фиксированные первые 12k символов.
    doc_context=_select_document_context(user_id,message.text,14000)
    if doc_context:
        set_task_state(user_id, document_context=True)

    # Умный анализ последнего созданного файла.
    if await _analyze_last_artifact(message, normalized): return

    # MecauAI 2.0: команды редактирования уже созданного файла и перенос данных между форматами.
    if await handle_artifact_command(message, normalized):
        return

    if normalized in {"/clear", "очисти контекст", "забудь предыдущий диалог", "начать заново"}:
        user_data = get_user_data(user_id)
        user_data["history"] = []
        user_data["last_output"] = "Здесь пока нет ответов."
        save_users_db()
        user_doc_context.pop(user_id, None)
        user_doc_chunks.pop(user_id, None)
        title_page_states.pop(user_id, None)
        await message.answer("🧹 Готово. Контекст очищен — можем начать заново.")
        return

    if normalized in {"/help", "помощь", "что ты умеешь", "что умеешь"}:
        await cmd_capabilities(message)
        return

    if re.search(r"\b(проверь|проверить|найди ошибки|есть ли ошибки|проведи проверку)\b", normalized) and get_user_data(user_id).get("last_output"):
        source=get_user_data(user_id).get("last_output","")
        try:
            r=await call_groq_with_retry(messages=[{"role":"system","content":"Проверь текст на фактические ошибки, логические пробелы, противоречия и сомнительные утверждения. Не выдумывай."},{"role":"user","content":source}])
            reply=clean_text_for_html(r.choices[0].message.content); await safe_answer(message,f"🔍 <b>Проверка:</b>\n\n{reply}{AD_FOOTER}",parse_mode="HTML",reply_markup=get_answer_inline_keyboard())
        except Exception as e:
            logging.error("smart check failed: %s",e,exc_info=True); await message.answer("⚠️ Не удалось выполнить проверку.")
        return

    if re.search(r"\b(структурируй|структурировать|сделай структуру|разбей по пунктам)\b", normalized) and get_user_data(user_id).get("last_output"):
        source=get_user_data(user_id).get("last_output","")
        try:
            r=await call_groq_with_retry(messages=[{"role":"system","content":"Структурируй текст: вывод, тезисы, шаги, примеры, итог. Не меняй факты."},{"role":"user","content":source}])
            reply=clean_text_for_html(r.choices[0].message.content); await safe_answer(message,f"🧩 <b>Структура:</b>\n\n{reply}{AD_FOOTER}",parse_mode="HTML",reply_markup=get_answer_inline_keyboard())
        except Exception as e:
            logging.error("smart structure failed: %s",e,exc_info=True); await message.answer("⚠️ Не удалось структурировать ответ.")
        return

    if re.search(r"\b(сделай|создай|подготовь|сформируй)\b.*\bпрезентаци", normalized):
        topic = re.sub(r"^.*?\bпрезентаци(?:ю|я|и|й)?\b\s*(?:на тему|по теме|про|о)?\s*", "", message.text, flags=re.I).strip(" .:-")
        user_ppt_topic[user_id] = topic or "Тема не указана"
        await cmd_ppt_prompt(message)
        return

    if re.search(r"\b(сделай|создай|подготовь|сформируй)\b.*\b(excel|таблиц|xlsx)", normalized):
        topic = re.sub(r"^.*?\b(?:excel|таблиц(?:у|а|ы)?|xlsx)\b\s*(?:на тему|по теме|для|с)?\s*", "", message.text, flags=re.I).strip(" .:-")
        user_excel_topic[user_id] = topic or message.text
        await cmd_excel_prompt(message)
        return

    if user_id in excel_states:
        topic = (message.text or "").strip()
        excel_states.discard(user_id)
        if topic:
            await generate_excel_file(message, topic, user_id=user_id)
        else:
            await message.answer("📊 Опиши, какую таблицу создать.")
        return


    # ----------------- ГЕНЕРАЦИЯ ПРЕЗЕНТАЦИИ -----------------
    if user_id in ppt_states and normalized in {"готово", "готово!", "начать", "создавай", "генерируй"}:
        topic=(user_ppt_topic.get(user_id) or "").strip()
        if not topic:
            await message.answer("📌 Сначала отправь тему презентации.")
            return
        if user_id not in user_ppt_slide_count:
            await message.answer("📑 Сначала выбери количество слайдов через кнопки выше.")
            return
        await generate_presentation_file(message,topic,user_id=user_id); return

    if user_id in ppt_states:
        # ШАГ 3 презентации: первое текстовое сообщение после выбора
        # количества слайдов — это ТЕМА. Состояние сохраняем, чтобы после
        # темы пользователь мог отправить собственные изображения.
        current_topic = (user_ppt_topic.get(user_id) or "").strip()
        incoming = (message.text or "").strip()

        if not current_topic:
            if not incoming:
                await message.answer("📈 Напиши тему презентации одним сообщением.")
                return
            user_ppt_topic[user_id] = incoming[:2500]
            await message.answer(
                "📌 <b>Тема принята:</b> «" + html.escape(incoming[:500]) + "»\n\n"
                "🖼 Теперь можешь отправить свои изображения для презентации — "
                "я использую их в слайдах.\n\n"
                "Когда картинки будут добавлены, нажми «🚀 Создать». Если картинки не нужны — "
                "сразу запускай создание.",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="🚀 Создать без моих картинок", callback_data="ppt_start")],
                    [InlineKeyboardButton(text="🖼 Добавить свои картинки", callback_data="ppt_wait_images")],
                    [InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")],
                ])
            )
            return

        await message.answer(
            "🖼 Тема уже принята. Отправь свои изображения или нажми "
            "<b>«🚀 Начать генерацию»</b>.",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="🚀 Начать генерацию", callback_data="ppt_start")],
                [InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_flow")],
            ])
        )
        return

    # ----------------- ГЕНЕРАЦИЯ EXCEL -----------------
    if user_id in excel_states:
        excel_states.discard(user_id)
        topic = (user_excel_topic.get(user_id) or message.text or "").strip()
        if not topic:
            await message.answer("📊 Опиши, какую таблицу создать: столбцы, период или пример данных.")
            return
        await generate_excel_file(message, topic)
        return

    if message.text in MENU_BUTTONS:
        return

    save_user_id(user_id)
    if user_id in busy_users:
        await message.answer("⏳ Я ещё обрабатываю предыдущий запрос. Подожди немного, чтобы задачи не смешались.")
        return
    busy_users.add(user_id)
    user_data = get_user_data(user_id)
    current_mode = user_data["mode"]

    system_prompt = PROMPTS.get(current_mode, PROMPTS["ai"]) + (
        "\n\nПРАВИЛА ОТВЕТА: "
        "отвечай естественно и полезно. Сначала дай прямой ответ, затем при необходимости "
        "короткое объяснение, пример или конкретные следующие шаги. "
        "Не повторяй вопрос пользователя. Не добавляй пустые вступления. "
        "Если запрос можно решить практическим способом — предложи готовый результат, а не только совет. "
        "Если информации недостаточно — задай один точный уточняющий вопрос."
    ) + "\n\nОтвечай содержательно: если вопрос требует объяснения, дай краткий контекст, основной ответ и 1-2 практических примера или следующих шага. Не растягивай ответ без необходимости."

    # Если ранее был загружен документ через «📚 Анализ документа» — добавляем
    # его содержимое в системный промпт, чтобы отвечать с учётом контекста.
    doc_context = _select_document_context(user_id,message.text,14000)
    if doc_context:
        system_prompt = (
            f"{system_prompt}\n\n"
            f"У пользователя есть загруженный документ. Если вопрос относится к нему, отвечай по релевантным фрагментам ниже. Если данных нет — скажи это прямо, не выдумывай.\n---\n{doc_context}\n---"
        )

    user_data["history"].append({"role": "user", "content": message.text})
    if len(user_data["history"]) > 6:
        user_data["history"] = user_data["history"][-6:]

    messages_payload = [{"role": "system", "content": system_prompt}] + user_data["history"]

    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        response = await call_groq_with_retry(messages=messages_payload)
        ai_reply = clean_text_for_html(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply
        save_users_db()

        if current_mode == "coder":
            full_message = f"{ai_reply}\n\n{AD_FOOTER}"
            await safe_answer(message, full_message, parse_mode="Markdown", reply_markup=get_answer_inline_keyboard())
        else:
            full_message = f"{ai_reply}{AD_FOOTER}"
            await safe_answer(message, full_message, parse_mode="HTML", reply_markup=get_answer_inline_keyboard())
    except Exception as e:
        logging.error(f"Ошибка обычного запроса пользователя {user_id}: {e}", exc_info=True)
        _track_user_event(user_id, message=message, kind="text", status="error", prompt=message.text or "", error=str(e))
        await message.answer("⚠️ Сейчас ИИ временно недоступен. Я сохранил ошибку в журнале и освободил задачу. Попробуй ещё раз через несколько секунд.")
    finally:
        busy_users.discard(user_id)
        _maybe_gc()

@dp.errors()
async def global_error_handler(event: types.ErrorEvent):
    """
    Глобальный перехватчик непредвиденных ошибок. Раньше необработанное
    исключение в любом хэндлере могло привести к тому, что пользователь
    просто не получал ответа и не понимал, что произошло. Теперь бот
    логирует ошибку и сообщает пользователю, что что-то пошло не так.
    """
    logging.error(f"Необработанная ошибка: {event.exception}", exc_info=True)
    try:
        update = event.update
        if update.message:
            await update.message.answer("⚠️ Произошла непредвиденная ошибка. Попробуй ещё раз позже. Если повторится — пиши @mecau")
        elif update.callback_query:
            await update.callback_query.answer("⚠️ Произошла ошибка, попробуй ещё раз.", show_alert=True)
    except Exception:
        pass
    return True


# ======================= MECAUAI 2.0 QUALITY PASS =======================

PPT_SLIDE_COUNTS = (6, 8, 10, 12)
PPT_LAYOUTS = ("bullets", "stats", "comparison", "timeline", "process", "quote", "chart", "conclusion")


def _safe_remove_shape(shape):
    try:
        sp = shape._element
        sp.getparent().remove(sp)
    except Exception:
        pass


def _ppt_normalize_item(raw, idx=0):
    """Presentation-grade normalization: dense enough to be useful, safe enough to render."""
    if not isinstance(raw, dict):
        raw = {"title": str(raw or f"Слайд {idx+1}"), "points": [str(raw or "Ключевой тезис.")]}
    def as_list(v):
        if v is None: return []
        if isinstance(v, list): return v
        return [v]
    def clean(v, limit=500):
        return re.sub(r'\s+', ' ', str(v or '')).strip()[:limit]
    points=[clean(x,360) for x in as_list(raw.get("points")) if clean(x,360)][:7]
    stats=[]
    for st in as_list(raw.get("stats"))[:6]:
        if isinstance(st, dict):
            stats.append({"label": clean(st.get("label") or st.get("title") or "Показатель",55), "value": clean(st.get("value") or st.get("number") or "—",80)})
        else:
            text=clean(st,120)
            if text:
                parts=text.split(":",1); stats.append({"label":clean(parts[0],55),"value":clean(parts[1] if len(parts)>1 else text,80)})
    cols=[]
    for col in as_list(raw.get("columns"))[:3]:
        if isinstance(col, dict):
            cols.append({"title":clean(col.get("title") or "Вариант",70),"points":[clean(x,260) for x in as_list(col.get("points")) if clean(x,260)][:6]})
        else: cols.append({"title":clean(col,70),"points":[]})
    steps=[clean(x,240) for x in as_list(raw.get("steps")) if clean(x,240)][:7]
    labels=[clean(x,45) for x in as_list(raw.get("labels")) if clean(x,45)][:10]
    values=[]
    for x in as_list(raw.get("values"))[:10]:
        try: values.append(float(str(x).replace(",",".").replace("%","").strip()))
        except Exception: pass
    layout=str(raw.get("layout") or "bullets").lower().strip()
    if layout not in PPT_LAYOUTS: layout="bullets"
    # Never allow a chart to be empty. The renderer will fall back to a data slide.
    if layout in ("chart","graph") and (len(labels)<2 or len(values)<2 or len(labels)!=len(values)):
        layout="bullets"
        if not points:
            points=["Для этой темы не удалось получить надёжный числовой ряд.","Лучше показать проверяемые факты и качественные выводы, чем пустой график."]
    return {
        "title":clean(raw.get("title") or f"Слайд {idx+1}",120),
        "layout":layout,"points":points,"stats":stats,"columns":cols,"steps":steps,
        "quote":clean(raw.get("quote"),700),"author":clean(raw.get("author"),120),
        "labels":labels,"values":values,"series":clean(raw.get("series") or "Значение",70),
        "chart_title":clean(raw.get("chart_title") or raw.get("title") or "Данные",110),
        "chart_kind":clean(raw.get("chart_kind") or "column",30),
        "subtitle":clean(raw.get("subtitle"),180),
        "takeaway":clean(raw.get("takeaway"),260),
        "source":clean(raw.get("source"),240),
        "image_prompt":clean(raw.get("image_prompt"),1000),
        "note":clean(raw.get("note"),700),
    }

async def generate_ai_image(session: aiohttp.ClientSession, prompt: str, width: int = 900, height: int = 650):
    """Надёжный image pipeline: timeout -> retry -> валидация -> уменьшение RAM."""
    prompt = re.sub(r"\s+", " ", str(prompt or "")).strip()[:900]
    if not prompt:
        return None
    for attempt in range(3):
        try:
            encoded = urllib.parse.quote(prompt)
            url = f"https://image.pollinations.ai/prompt/{encoded}?width={width}&height={height}&nologo=true&model=flux"
            timeout = aiohttp.ClientTimeout(total=22 + attempt * 8, connect=8)
            async with session.get(url, timeout=timeout) as resp:
                if resp.status != 200:
                    raise RuntimeError(f"image HTTP {resp.status}")
                data = await resp.read()
                if len(data) < 10_000:
                    raise ValueError("изображение подозрительно маленькое")
                try:
                    from PIL import Image
                    im = Image.open(io.BytesIO(data))
                    im.verify()
                    im = Image.open(io.BytesIO(data))
                    if im.width < 400 or im.height < 300:
                        raise ValueError("низкое разрешение")
                except ImportError:
                    pass
                except Exception as e:
                    raise ValueError(f"invalid image: {e}")
                return data
        except Exception as e:
            logging.warning("Image pipeline attempt %s/3: %s", attempt + 1, e)
            if attempt < 2:
                await asyncio.sleep(1.2 * (2 ** attempt))
    return None


def _ppt_overflow_check(slide):
    # Heuristic overflow guard: shrink dense text instead of letting it spill outside cards.
    for shape in slide.shapes:
        if not hasattr(shape,'text_frame') or not shape.text_frame: continue
        text=shape.text_frame.text or ''
        if len(text)>260:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        cur=r.font.size.pt if r.font.size else 14
                        r.font.size=PptxPt(max(9,cur-2))
                    except Exception: pass
        shape.text_frame.word_wrap=True

def _ppt_add_text(slide, text, x, y, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.clear()
    p = tf.paragraphs[0]; p.text = str(text)[:900]; p.alignment = align
    _set_text_style(p, size, color or RGBColor(0x21,0x21,0x21), bold)
    return box


def _ppt_add_card(slide, x, y, w, h, title, value, theme):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5,0xF7,0xFA); card.line.color.rgb = RGBColor(0xE0,0xE5,0xEB)
    _ppt_add_text(slide, title, x+0.18, y+0.18, w-0.36, 0.35, 11, False, theme["body_text"])
    _ppt_add_text(slide, value, x+0.18, y+0.58, w-0.36, h-0.7, 23, True, theme["primary"])


def _ppt_add_header(slide,theme,idx,title):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,PptxInches(0),PptxInches(0),PptxInches(13.333),PptxInches(.12)); bar.fill.solid(); bar.fill.fore_color.rgb=theme["primary"]; bar.line.fill.background()
    _ppt_add_text(slide,f"{idx:02d}",.68,.48,.5,.35,10,True,theme["primary"])
    _ppt_add_text(slide,title,1.28,.42,11.3,.72,27 if len(title)<60 else 21,True,theme["body_text"])


def _ppt_add_image(slide,img_stream,x=8.15,y=1.55,w=4.45,h=4.95):
    if img_stream is None: return
    try: slide.shapes.add_picture(img_stream,left=PptxInches(x),top=PptxInches(y),width=PptxInches(w),height=PptxInches(h))
    except Exception as e: logging.warning("PPT image insert failed: %s",e)


def _ppt_add_takeaway(slide, text, theme, y=6.18):
    if not text:
        return
    _ppt_add_text(slide, "КЛЮЧЕВАЯ МЫСЛЬ", .72, y, 1.45, .28, 8, True, theme["primary"])
    _ppt_add_text(slide, text, 2.08, y-.04, 10.45, .42, 11, True, theme["body_text"])

def _ppt_add_source(slide, source, theme):
    if source:
        _ppt_add_text(slide, "Источник: "+source, .72, 6.92, 11.8, .25, 8, False, theme.get("muted_text", theme.get("subtitle_text", theme.get("body_text", "666666"))))

def _ppt_add_subtitle(slide, subtitle, theme):
    if subtitle:
        _ppt_add_text(slide, subtitle, .72, 1.16, 11.7, .4, 11, False, theme.get("muted_text", theme.get("subtitle_text", theme.get("body_text", "666666"))))

def build_advanced_ppt_slide(prs, theme, idx, item, img_stream=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = theme["bg"]
    layout = str(item.get("layout") or "bullets").lower()
    title = str(item.get("title") or f"Слайд {idx+1}").strip()[:120]
    _ppt_add_text(s, title, 0.72, 0.52, 11.7, 0.75, 27 if len(title)<65 else 22, True, theme["body_text"])
    _ppt_add_subtitle(s, item.get("subtitle"), theme)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(13.333), PptxInches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = theme["primary"]; bar.line.fill.background()
    points = [str(x).strip() for x in (item.get("points") or []) if str(x).strip()][:6]
    if layout in ("stats", "statistics", "kpi"):
        stats = item.get("stats") or []
        if not stats:
            stats = [{"label": p.split(":",1)[0][:30], "value": p.split(":",1)[1][:30] if ":" in p else p[:30]} for p in points]
        for i, st in enumerate(stats[:4]):
            _ppt_add_card(s, 0.75 + i*3.05, 1.75, 2.7, 2.1, st.get("label","Показатель"), st.get("value","—"), theme)
        if item.get("note"):
            _ppt_add_text(s, item["note"], 0.85, 4.35, 11.5, 0.9, 16, False, theme["body_text"])
    elif layout in ("comparison", "compare"):
        cols = item.get("columns") or []
        if len(cols) < 2:
            cols = [{"title":"Вариант A","points":points[:3]}, {"title":"Вариант B","points":points[3:6]}]
        for ci, col in enumerate(cols[:3]):
            x = 0.7 + ci*4.15
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PptxInches(x),PptxInches(1.65),PptxInches(3.75),PptxInches(4.7))
            card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF7,0xF8,0xFA); card.line.color.rgb = theme["primary"]
            _ppt_add_text(s,col.get("title","Вариант"),x+0.25,1.95,3.25,0.55,18,True,theme["primary"])
            yy=2.7
            for p in (col.get("points") or [])[:5]:
                _ppt_add_text(s,"• "+str(p),x+0.25,yy,3.2,0.62,13,False,theme["body_text"]); yy+=0.78
    elif layout in ("timeline", "process"):
        steps = item.get("steps") or points
        n=max(1,min(len(steps),6)); y=3.05
        line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,PptxInches(1.0),PptxInches(y+0.18),PptxInches(11.2),PptxInches(0.08))
        line.fill.solid(); line.fill.fore_color.rgb=theme["primary"]; line.line.fill.background()
        for i,step in enumerate(steps[:n]):
            x=1.0 + (10.8/(max(n-1,1)))*i
            circ=s.shapes.add_shape(MSO_SHAPE.OVAL,PptxInches(x),PptxInches(y),PptxInches(0.48),PptxInches(0.48)); circ.fill.solid(); circ.fill.fore_color.rgb=theme["primary"]; circ.line.fill.background()
            p=circ.text_frame.paragraphs[0]; p.text=str(i+1); p.alignment=PP_ALIGN.CENTER; _set_text_style(p,10,theme["title_text"],True)
            _ppt_add_text(s,str(step),x-0.35,y+0.75,1.25,1.3,11,False,theme["body_text"],PP_ALIGN.CENTER)
    elif layout in ("quote", "citation"):
        quote = str(item.get("quote") or (points[0] if points else "Главный тезис темы."))[:500]
        _ppt_add_text(s,"“",0.95,1.65,1.0,1.2,60,True,theme["primary"])
        _ppt_add_text(s,quote,1.75,2.05,9.9,2.2,25,True,theme["body_text"])
        if item.get("author"):
            _ppt_add_text(s,"— "+str(item["author"]),1.8,4.65,8.5,0.5,14,False,theme["primary"])
    elif layout in ("conclusion", "summary"):
        for i,p in enumerate(points[:4]):
            _ppt_add_card(s,0.8+(i%2)*6.1,1.65+(i//2)*2.25,5.55,1.85,f"Вывод {i+1}",p,theme)
    elif layout in ("chart", "graph") and (item.get("labels") and item.get("values")):
        data=ChartData(); data.categories=[str(x)[:30] for x in item["labels"][:8]]; data.add_series(str(item.get("series") or "Значение"), [float(x) for x in item["values"][:8]])
        chart=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,PptxInches(0.75),PptxInches(1.55),PptxInches(11.8),PptxInches(4.95),data).chart
        chart.has_legend=False; chart.has_title=True; chart.chart_title.text_frame.text=str(item.get("chart_title") or title)
    else:
        has_img = img_stream is not None
        left_w=7.15 if has_img else 11.7; yy=1.55
        for i,p in enumerate(points or ["Ключевой тезис будет сформирован автоматически."]):
            h=0.7 if len(p)<110 else 0.9
            card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PptxInches(0.68),PptxInches(yy),PptxInches(left_w),PptxInches(h))
            card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0xF5,0xF7,0xFA); card.line.color.rgb=RGBColor(0xE2,0xE6,0xEC)
            _ppt_add_text(s,f"{i+1}",0.9,yy+0.14,0.35,0.3,9,True,theme["title_text"],PP_ALIGN.CENTER)
            badge=s.shapes[-1]; badge.fill.solid(); badge.fill.fore_color.rgb=theme["primary"]; badge.line.fill.background()
            _ppt_add_text(s,p,1.4,yy+0.1,left_w-0.95,h-0.1,14 if len(p)<120 else 12,False,theme["body_text"])
            yy+=h+0.12
        if has_img:
            frame=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PptxInches(8.15),PptxInches(1.55),PptxInches(4.5),PptxInches(4.95)); frame.fill.solid(); frame.fill.fore_color.rgb=RGBColor(0xEE,0xF1,0xF5); frame.line.color.rgb=RGBColor(0xD9,0xDE,0xE5)
            try: s.shapes.add_picture(img_stream,left=PptxInches(8.28),top=PptxInches(1.68),width=PptxInches(4.24),height=PptxInches(4.69))
            except Exception as e: logging.warning("PPT image insert failed: %s",e)
    _ppt_add_takeaway(s, item.get("takeaway"), theme)
    _ppt_add_source(s, item.get("source"), theme)
    _ppt_overflow_check(s)
    _add_footer(s,theme,idx+2)
    return s


async def _ppt_structure(topic, num_slides):
    response=await call_groq_with_retry(messages=[
        {"role":"system","content":f"Создай визуально сильную, содержательную презентацию на русском языке ровно из {num_slides} содержательных слайдов после титульного. Это должна быть презентация уровня современного pitch deck / Canva: один главный тезис на слайд, сильная визуальная иерархия, мало текста, конкретные факты, примеры и выводы. Верни ТОЛЬКО JSON-массив. Каждый объект: title, subtitle, layout, points, image_prompt, stats, columns, steps, quote, author, labels, values, series, chart_title, takeaway, source, note. layout выбирай из bullets, stats, comparison, timeline, process, quote, chart, conclusion. Не используй chart без labels/values и не выдумывай числовые данные: если надёжных чисел нет, используй comparison/process/bullets. Структура должна логично рассказывать историю: контекст → проблема/вопрос → ключевые идеи → данные или пример → сравнение/процесс → риски/ограничения → практический смысл → итог. Не повторяй одни и те же тезисы. Для большинства смысловых слайдов дай короткий image_prompt для тематического визуала; для chart и чистых comparison/stats можно оставить пустым. points — 2–5 коротких, конкретных тезисов. takeaway — одна сильная фраза. source — только если источник реально известен из темы/контекста."},
        {"role":"user","content":f"Тема: {topic[:2500]}"}
    ],temperature=0.35,max_retries=2,timeout=50)
    data=extract_json(clean_text_for_html(response.choices[0].message.content))
    if not isinstance(data,list): raise ValueError("Некорректная структура PPT")
    return data[:num_slides]


def _ppt_quality_audit(prs):
    """Deterministic QA. No rasterization or duplicate in-memory presentations."""
    issues=[]
    for i,slide in enumerate(prs.slides,1):
        texts=[sh.text.strip() for sh in slide.shapes if hasattr(sh,'text') and sh.text.strip()]
        if len(texts)<3:
            issues.append(f"slide {i}: too little content")
        for sh in slide.shapes:
            if getattr(sh,'has_chart',False):
                try:
                    chart=sh.chart
                    if not chart.series or len(chart.series)==0:
                        issues.append(f"slide {i}: empty chart")
                except Exception:
                    issues.append(f"slide {i}: invalid chart")
    return issues


async def generate_presentation_file(message: Message, topic: str, user_id: int = None):
    user_id=user_id or message.from_user.id
    if user_id in busy_users:
        await message.answer("⏳ Подожди, предыдущая задача ещё выполняется..."); return
    busy_users.add(user_id); heavy_lock=None; status_msg=await message.answer("📈 Анализирую содержание → подбираю макеты → проверяю визуалы…")
    try:
        heavy_lock=await acquire_heavy_job(status_msg)
        theme_key=user_ppt_design.pop(user_id,"ppt_blue"); theme=PPT_THEMES.get(theme_key,PPT_THEMES["ppt_blue"])
        custom=user_ppt_images.pop(user_id,[])[:MAX_PPT_IMAGES]
        # Сохраняем выбранное пользователем количество содержательных слайдов.
        n=int(user_ppt_slide_count.pop(user_id, 8) or 8)
        if n not in PPT_SLIDE_COUNTS:
            n=8
        try:
            slides=await _ppt_structure(topic,n)
        except Exception as structure_error:
            logging.error("PPT structure generation failed, using fallback: %s", structure_error, exc_info=True)
            fallback_titles = [
                "Введение и контекст", "Ключевые факты", "Основные тенденции",
                "Примеры и данные", "Сравнение подходов", "Проблемы и риски",
                "Практическое значение", "Итоги", "Выводы", "Рекомендации",
                "Перспективы", "Заключение"
            ]
            slides=[]
            for i in range(n):
                title=fallback_titles[i % len(fallback_titles)]
                slides.append({
                    "title": title,
                    "layout": "bullets" if i < n-1 else "conclusion",
                    "points": [
                        f"Основной аспект: {title.lower()}",
                        "Факты и практический контекст по теме",
                        "Что это означает на практике",
                    ],
                    "image_prompt": ""
                })
        while len(slides)<n:
            slides.append({"title":"Выводы","layout":"conclusion","points":["Главный результат","Практическое значение","Следующий шаг"]})
        prs=Presentation(); prs.slide_width=PptxInches(13.333); prs.slide_height=PptxInches(7.5); build_title_slide(prs,theme,topic)
        async with aiohttp.ClientSession() as session:
            for idx,item in enumerate(slides):
                stream=None
                if idx<len(custom):
                    try:
                        fi=await bot.get_file(custom[idx]); d=await bot.download_file(fi.file_path); b=d.read()
                        if len(b)<=MAX_DOC_SIZE_BYTES: stream=io.BytesIO(b)
                    except Exception as e: logging.warning("custom PPT image %s: %s",idx,e)
                elif item.get("image_prompt") and idx<MAX_PPT_AI_IMAGES:
                    try:
                        b=await generate_ai_image(session,item["image_prompt"],900,650)
                        if b: stream=io.BytesIO(b)
                    except Exception as e:
                        logging.warning("AI visual skipped on PPT slide %s: %s", idx + 1, e)
                        stream=None
                try:
                    build_advanced_ppt_slide(prs,theme,idx,item,stream)
                except Exception as e:
                    logging.error("PPT slide %s build failed: %s", idx + 1, e, exc_info=True)
                    fallback = {
                        "title": str(item.get("title") or f"Слайд {idx + 1}")[:120],
                        "layout": "bullets",
                        "points": [str(x) for x in (item.get("points") or []) if str(x).strip()][:5] or [
                            "Ключевой тезис по теме.", "Практическое значение.", "Основной вывод."
                        ],
                    }
                    build_advanced_ppt_slide(prs,theme,idx,fallback,None)
        bio=io.BytesIO(); prs.save(bio); data=bio.getvalue()
        if len(data)>MAX_ARTIFACT_BYTES: raise ValueError("PPT слишком большой")
        path=_artifact_path(user_id,"pptx",topic); _atomic_write_bytes(path, data); register_artifact(user_id,"pptx",path,{"topic":topic,"theme":theme_key,"slides":len(prs.slides)})
        await status_msg.delete(); await message.answer_document(BufferedInputFile(data,filename=os.path.basename(path)),caption=f"📈 Готово: {len(prs.slides)} слайдов • {theme['label']}\n✏️ Можно сразу написать, что изменить.")
        get_user_stats(user_id)["exports"]+=1
    except Exception as e:
        logging.error("PPT 2.0 error: %s", e, exc_info=True)
        # Последний резерв: даже если AI/визуалы/сложный макет сломались,
        # пользователь должен получить настоящий .pptx, а не только ошибку.
        try:
            fallback_prs = Presentation()
            fallback_prs.slide_width = PptxInches(13.333)
            fallback_prs.slide_height = PptxInches(7.5)
            build_title_slide(fallback_prs, PPT_THEMES.get("ppt_blue", theme), topic)
            fallback_items = [
                {"title": "Главное о теме", "points": [f"Тема: {topic[:180]}", "Ключевые понятия и контекст", "Основной смысл"], "layout": "bullets"},
                {"title": "Ключевые аспекты", "points": ["Основные факты", "Важные тенденции", "Практическое значение"], "layout": "bullets"},
                {"title": "Вывод", "points": ["Главный результат", "Что важно запомнить", "Практический следующий шаг"], "layout": "conclusion"},
            ]
            for idx, item in enumerate(fallback_items):
                build_advanced_ppt_slide(fallback_prs, PPT_THEMES.get("ppt_blue", theme), idx, item, None)
            fb = io.BytesIO(); fallback_prs.save(fb); fb_data = fb.getvalue()
            fb_path = _artifact_path(user_id, "pptx", topic); _atomic_write_bytes(fb_path, fb_data)
            register_artifact(user_id, "pptx", fb_path, {"topic": topic, "fallback": True, "slides": len(fallback_prs.slides)})
            await status_msg.edit_text("⚠️ Визуалы/сложный макет не загрузились, поэтому я собрал упрощённую версию презентации.")
            await message.answer_document(BufferedInputFile(fb_data, filename=os.path.basename(fb_path)), caption="📈 Презентация готова в упрощённом режиме. Можно написать, что изменить.")
        except Exception as fallback_error:
            logging.error("PPT fallback also failed: %s", fallback_error, exc_info=True)
            try:
                await status_msg.edit_text("⚠️ Не удалось создать даже резервную презентацию. Подробность сохранена в журнале ошибок.")
            except Exception:
                pass
    finally:
        busy_users.discard(user_id); ppt_states.discard(user_id); user_ppt_topic.pop(user_id,None); user_ppt_images.pop(user_id,None); user_ppt_design.pop(user_id,None)
        if heavy_lock: heavy_lock.release()
        gc.collect()



def _excel_detect_format(header, values):
    h=str(header).lower()
    vals=[v for v in values if v not in (None,"")]
    if any(x in h for x in ("%","процент","доля")): return '0.0%'
    if any(x in h for x in ("€","eur","евро","руб","₽","цена","стоим","расход","доход","выруч")): return '#,##0.00'
    if vals and all(isinstance(v,(int,float)) for v in vals): return '#,##0.00'
    return 'General'

async def generate_excel_file(message: Message, topic: str, user_id: int = None):
    user_id=user_id or message.from_user.id
    if user_id in busy_users: await message.answer("⏳ Подожди, предыдущая задача ещё выполняется..."); return
    busy_users.add(user_id); heavy_lock=None; status=await message.answer("📊 Строю данные → формулы → KPI → Dashboard → диаграммы…")
    try:
        heavy_lock=await acquire_heavy_job(status)
        response=await call_groq_with_retry(messages=[
            {"role":"system","content":"Ты генератор Excel-данных. Верни ТОЛЬКО валидный JSON без Markdown и пояснений. Строго объект с ключами title, sheet_name, headers, rows, formulas. headers — массив из 2-12 строк. rows — массив массивов, 5-40 строк. formulas — массив объектов {cell,formula} или []. Не используй комментарии вне JSON."},
            {"role":"user","content":topic[:4000]}],temperature=0.1)
        raw_ai = response.choices[0].message.content
        try:
            data=extract_json(clean_text_for_html(raw_ai))
        except Exception:
            # Второй, более короткий запрос часто спасает модели, которые добавили пояснение
            # к JSON в первом ответе.
            retry=await call_groq_with_retry(messages=[
                {"role":"system","content":'Верни только JSON. Формат: {"title":"...","sheet_name":"Данные","headers":["Столбец 1","Столбец 2"],"rows":[["значение 1","значение 2"]],"formulas":[]}. Никаких ``` и текста вокруг JSON.'},
                {"role":"user","content":topic[:3000]}],temperature=0)
            data=extract_json(retry.choices[0].message.content)
        if not isinstance(data, dict):
            raise ValueError("AI вернул не объект JSON")
        headers=[str(x).strip() for x in data.get('headers',[]) if str(x).strip()][:12]
        if not headers: raise ValueError('Нет заголовков')
        rows=[]
        for row in (data.get('rows') or [])[:40]:
            r=list(row) if isinstance(row,(list,tuple)) else [row]; rows.append((r+[""]*len(headers))[:len(headers)])
        wb=Workbook(); ws=wb.active; ws.title=re.sub(r'[\\/*?:\[\]]','_',str(data.get('sheet_name') or 'Данные'))[:31] or 'Данные'
        ws.merge_cells(start_row=1,start_column=1,end_row=1,end_column=len(headers)); c=ws.cell(1,1,str(data.get('title') or topic)[:120]); c.font=Font(name='Aptos Display',size=18,bold=True,color='FFFFFF'); c.fill=PatternFill('solid',fgColor='1F4E78'); ws.row_dimensions[1].height=30
        ws.append(headers)
        for cell in ws[2]: cell.font=Font(bold=True,color='FFFFFF'); cell.fill=PatternFill('solid',fgColor='2F75B5'); cell.alignment=Alignment(horizontal='center',wrap_text=True)
        for r in rows: ws.append(r)
        for ci,h in enumerate(headers,1):
            vals=[ws.cell(r,ci).value for r in range(3,ws.max_row+1)]
            fmt=_excel_detect_format(h,vals)
            for r in range(3,ws.max_row+1):
                cell=ws.cell(r,ci); cell.alignment=Alignment(vertical='center',wrap_text=True); cell.number_format=fmt
            maxlen=max([len(str(h))]+[len(str(v or '')) for v in vals]); ws.column_dimensions[get_column_letter(ci)].width=min(max(maxlen+2,12),34)
        ws.freeze_panes='A3'; ws.auto_filter.ref=f'A2:{get_column_letter(len(headers))}{ws.max_row}'
        if ws.max_row>=3:
            tab=Table(displayName='DataTable',ref=f'A2:{get_column_letter(len(headers))}{ws.max_row}'); tab.tableStyleInfo=TableStyleInfo(name='TableStyleMedium2',showRowStripes=True); ws.add_table(tab)
        # Conditional formatting for numeric columns.
        for ci,h in enumerate(headers,1):
            vals=[ws.cell(r,ci).value for r in range(3,ws.max_row+1)]
            if vals and sum(isinstance(v,(int,float)) for v in vals)/len(vals)>=0.7:
                rng=f'{get_column_letter(ci)}3:{get_column_letter(ci)}{ws.max_row}'
                ws.conditional_formatting.add(rng,ColorScaleRule(start_type='min',start_color='F2F2F2',mid_type='percentile',mid_value=50,mid_color='D9EAF7',end_type='max',end_color='5B9BD5'))
        for f in data.get('formulas') or []:
            try: ws[str(f['cell'])]=str(f['formula'])
            except Exception: pass
        # Dashboard with multiple chart types and KPI formulas.
        numeric=[]
        for ci in range(2,len(headers)+1):
            vals=[ws.cell(r,ci).value for r in range(3,ws.max_row+1)]
            if vals and sum(isinstance(v,(int,float)) for v in vals)/len(vals)>=0.7: numeric.append(ci)
        dash=wb.create_sheet('Dashboard'); dash['A1']=str(data.get('title') or topic)[:100]; dash['A1'].font=Font(size=20,bold=True)
        for i,ci in enumerate(numeric[:4]):
            col=get_column_letter(ci); x=1+(i%2)*3; y=3+(i//2)*4
            dash.cell(y,x,'Итого').font=Font(bold=True); dash.cell(y,x+1,f'=SUM(\'{ws.title}\'!{col}3:{col}{ws.max_row})')
            dash.cell(y+1,x,'Среднее').font=Font(bold=True); dash.cell(y+1,x+1,f'=AVERAGE(\'{ws.title}\'!{col}3:{col}{ws.max_row})')
        if numeric:
            chart=BarChart(); chart.type='col'; chart.style=10; chart.title=f'{headers[numeric[0]-1]} по {headers[0]}'; chart.add_data(Reference(ws,min_col=numeric[0],max_col=numeric[0],min_row=2,max_row=ws.max_row),titles_from_data=True); chart.set_categories(Reference(ws,min_col=1,min_row=3,max_row=ws.max_row)); chart.width=18; chart.height=9; dash.add_chart(chart,'A12')
        if len(numeric)>1:
            from openpyxl.chart import LineChart
            line=LineChart(); line.title=f'Динамика: {headers[numeric[1]-1]}'; line.add_data(Reference(ws,min_col=numeric[1],max_col=numeric[1],min_row=2,max_row=ws.max_row),titles_from_data=True); line.set_categories(Reference(ws,min_col=1,min_row=3,max_row=ws.max_row)); line.width=18; line.height=9; dash.add_chart(line,'J12')
        bio=io.BytesIO(); wb.save(bio); raw=bio.getvalue();
        if not raw or len(raw) > MAX_ARTIFACT_BYTES: raise ValueError('Excel слишком большой или пустой')
        path=_artifact_path(user_id,'xlsx',topic); _atomic_write_bytes(path, raw); register_artifact(user_id,'xlsx',path,{'topic':topic,'sheet':ws.title})
        await status.delete(); await message.answer_document(BufferedInputFile(raw,filename=os.path.basename(path)),caption=f'📊 Excel готов: {len(rows)} строк, {len(headers)} столбцов.\n✏️ Можно написать, что изменить.')
        get_user_stats(user_id)['exports']+=1
    except Exception as e:
        logging.error('Excel 2.0 error: %s', e, exc_info=True)
        # Резервный режим: создаём валидный XLSX даже при сбое AI/JSON.
        try:
            wb = Workbook(); ws = wb.active; ws.title = 'Данные'
            ws.append(['Параметр', 'Значение'])
            ws.append(['Запрос', topic[:300]])
            ws.append(['Статус', 'Создано в резервном режиме'])
            ws.append(['Дата', datetime.now().strftime('%Y-%m-%d')])
            ws.append(['Комментарий', 'AI не вернул корректную структуру; файл всё равно создан.'])
            for cell in ws[1]:
                cell.font = Font(bold=True, color='FFFFFF'); cell.fill = PatternFill('solid', fgColor='2F75B5')
            ws.column_dimensions['A'].width = 28; ws.column_dimensions['B'].width = 70
            ws.freeze_panes = 'A2'
            bio = io.BytesIO(); wb.save(bio); raw = bio.getvalue()
            fb_path = _artifact_path(user_id, 'xlsx', topic); _atomic_write_bytes(fb_path, raw)
            register_artifact(user_id, 'xlsx', fb_path, {'topic': topic, 'fallback': True})
            await status.edit_text('⚠️ AI не смог разложить запрос по столбцам, поэтому я создал резервный Excel. Его можно попросить дополнить.')
            await message.answer_document(BufferedInputFile(raw, filename=os.path.basename(fb_path)), caption='📊 Excel готов в резервном режиме. Напиши, какие столбцы/данные добавить.')
        except Exception as fallback_error:
            logging.error('Excel fallback also failed: %s', fallback_error, exc_info=True)
            try: await status.edit_text('⚠️ Не удалось создать Excel-файл. Подробность сохранена в журнале ошибок.')
            except Exception: pass
    finally:
        busy_users.discard(user_id); excel_states.discard(user_id); user_excel_topic.pop(user_id,None); user_excel_mode.pop(user_id,None)
        if heavy_lock: heavy_lock.release()
        gc.collect()

def _artifact_text(path, limit=14000):
    ext=os.path.splitext(path)[1].lower()
    if ext=='.xlsx':
        wb=__import__('openpyxl').load_workbook(path,read_only=True,data_only=False)
        out=[]
        for ws in wb.worksheets[:3]:
            out.append(f'Лист: {ws.title}')
            for row in ws.iter_rows(min_row=1,max_row=25,values_only=True): out.append(' | '.join(str(v or '') for v in row))
        return '\n'.join(out)[:limit]
    if ext=='.docx':
        d=Document(path); return '\n'.join(p.text for p in d.paragraphs)[:limit]
    if ext=='.pptx':
        prs=Presentation(path); out=[]
        for i,s in enumerate(prs.slides):
            texts=[sh.text for sh in s.shapes if hasattr(sh,'text') and sh.text.strip()]
            out.append(f'Slide {i+1}: '+' | '.join(texts))
        return '\n'.join(out)[:limit]
    return ''

async def _send_artifact(message,user_id,path,caption='Обновлённый файл'):
    with open(path,'rb') as f: raw=f.read()
    await message.answer_document(BufferedInputFile(raw,filename=os.path.basename(path)),caption=caption)

def _edit_ppt_local(path,instruction):
    prs=Presentation(path); low=instruction.lower()
    if re.search(r'\b(убери|удали).*(график|диаграмм)',low):
        for s in prs.slides:
            for sh in list(s.shapes):
                if getattr(sh,'has_chart',False): _safe_remove_shape(sh)
    m=re.search(r'(?:слайд|slide)\s*(\d+)',low); idx=int(m.group(1))-1 if m else None
    if 'добавь слайд' in low or 'добавить слайд' in low:
        s=prs.slides.add_slide(prs.slide_layouts[6]); _ppt_add_text(s,'Выводы',0.8,0.8,11.5,0.7,30,True); _ppt_add_text(s,'• Ключевой результат\n• Практическое значение\n• Следующий шаг',0.9,2.0,10.8,2.5,20); _add_footer(s,PPT_THEMES['ppt_blue'],len(prs.slides));
    if 'подробнее' in low and idx is not None and 0<=idx<len(prs.slides):
        s=prs.slides[idx]; _ppt_add_text(s,'Дополнение',0.8,5.95,2.0,0.35,12,True,PPT_THEMES['ppt_blue']['primary']); _ppt_add_text(s,'Дополнительный контекст, пример и практическое значение по запросу пользователя.',2.3,5.82,9.7,0.7,12)
    if re.search(r'строг(ий|им)|официальн',low):
        theme=PPT_THEMES['ppt_dark']
        for s in prs.slides:
            s.background.fill.solid(); s.background.fill.fore_color.rgb=theme['bg']
            for sh in s.shapes:
                if hasattr(sh,'text_frame'):
                    for p in sh.text_frame.paragraphs:
                        for r in p.runs: r.font.name='Aptos'
    prs.save(path)

async def _edit_excel_local(path,instruction):
    from openpyxl import load_workbook
    wb=load_workbook(path); ws=wb[wb.sheetnames[0]]; low=instruction.lower()
    if 'отсорт' in low:
        # Sort by a named column when present; otherwise by first numeric column.
        col=None
        for i,c in enumerate(ws[2],1):
            if c.value and str(c.value).lower() in low: col=i; break
        if col is None:
            for i in range(2,ws.max_column+1):
                vals=[ws.cell(r,i).value for r in range(3,ws.max_row+1)]
                if vals and any(isinstance(v,(int,float)) for v in vals): col=i; break
        if col:
            rows=list(ws.iter_rows(min_row=3,values_only=True)); rows.sort(key=lambda r:(r[col-1] is None,r[col-1]),reverse='по убыванию' in low or 'убыв' in low)
            for rr,row in enumerate(rows,3):
                for cc,val in enumerate(row,1): ws.cell(rr,cc).value=val
    if 'добавь' in low and ('март' in low or 'апрел' in low or 'май' in low):
        period=next((x for x in ('март','апрель','май','июнь','июль','август','сентябрь','октябрь','ноябрь','декабрь','январь','февраль') if x in low),'новая строка')
        ws.append([period]+['']*(ws.max_column-1))
    if 'диаграмм' in low or 'график' in low:
        from openpyxl.chart import BarChart, Reference
        numeric=2 if ws.max_column>=2 else 1; ch=BarChart(); ch.type='col'; ch.title='Обновлённая диаграмма'; ch.add_data(Reference(ws,min_col=numeric,max_col=numeric,min_row=2,max_row=ws.max_row),titles_from_data=True); ch.set_categories(Reference(ws,min_col=1,min_row=3,max_row=ws.max_row));
        target=wb['Dashboard'] if 'Dashboard' in wb.sheetnames else wb.create_sheet('Dashboard'); target.add_chart(ch,'A30')
    wb.save(path)

async def _edit_word_local(path,instruction):
    d=Document(path); low=instruction.lower()
    if 'заголов' in low or 'структур' in low:
        for p in d.paragraphs:
            if p.text.strip() and len(p.text)<90 and not p.style.name.startswith('Heading'):
                p.style='Heading 2'
    if 'подробнее' in low:
        p=d.add_paragraph('Дополнение: дополнительный контекст, пример и практическое применение по запросу пользователя.')
        p.style='Normal'
    d.save(path)

async def _analyze_last_artifact(message, normalized):
    user_id=message.from_user.id; art=last_artifact(user_id)
    triggers=("проанализируй последний","проанализируй файл","разбери последний файл","сделай выводы из файла","найди ошибки в файле","проверь файл","сделай конспект файла","сделай кратко из файла")
    if not art or not art.get("path") or not os.path.exists(art["path"]) or not any(x in normalized for x in triggers): return False
    text=_artifact_text(art["path"],limit=18000)
    if not text.strip(): await message.answer("⚠️ Не удалось извлечь текст из последнего файла."); return True
    if "ошиб" in normalized or "проверь" in normalized: task="Проверь материал на логические противоречия, очевидные ошибки, пропуски и слабые места. Отдельно укажи, что требует проверки человеком."; title="🔍 Проверка файла"
    elif "конспект" in normalized or "кратко" in normalized: task="Сделай структурированный конспект: тема, ключевые идеи, важные цифры/даты, выводы и что важно запомнить."; title="📝 Конспект файла"
    else: task="Проанализируй материал: краткий вывод, ключевые факты, проблемы, сильные стороны, риски и практические рекомендации."; title="🧠 Анализ файла"
    try:
        r=await call_groq_with_retry(messages=[{"role":"system","content":task+" Не выдумывай сведения, которых нет в материале."},{"role":"user","content":text}])
        result=clean_text_for_html(r.choices[0].message.content)
        await safe_answer(message,f"{title}:\n\n{result}{AD_FOOTER}",parse_mode="HTML",reply_markup=get_answer_inline_keyboard())
    except Exception as e: logging.error("Artifact analysis failed: %s",e,exc_info=True); await message.answer("⚠️ Не удалось проанализировать файл.")
    return True

async def handle_artifact_command(message, normalized):
    user_id=message.from_user.id
    art=last_artifact(user_id)
    if not art: return False
    # Only intercept explicit edit/convert requests, not ordinary conversation.
    edit_words=('измени презентац','измени файл','изменить файл','добавь слайд','убери слайд','удали слайд','убери график','удали график','добавь график','диаграмм','отсортируй','добавь столбец','добавь строк','удали столбец','удали строк','переименуй столбец','расходы за','март','апрель','май','сделай подробнее','сделай более строг','официальн','добавь заголов')
    convert_words=('сделай из него презентацию','сделай из неё презентацию','сделай из него excel','сделай из него word','сделай краткий word','сделай презентацию из excel')
    if not any(w in normalized for w in edit_words+convert_words): return False
    if any(w in normalized for w in convert_words):
        text=_artifact_text(art['path'])
        if 'презента' in normalized:
            await generate_presentation_file(message, 'Создай презентацию по данным последнего файла:\n'+text[:9000], user_id); return True
        if 'excel' in normalized:
            await generate_excel_file(message, 'Создай Excel по данным последнего файла:\n'+text[:9000], user_id); return True
        if 'word' in normalized:
            path=_artifact_path(user_id,'docx','Report'); d=Document();
            for line in text.splitlines()[:80]:
                p=d.add_paragraph(line); p.style='Normal'
            d.save(path); register_artifact(user_id,'docx',path,{'source':art['path']}); await _send_artifact(message,user_id,path,'📄 Word создан из последнего файла.'); return True
    if not art['path'] or not os.path.exists(art['path']): return False
    try:
        if art['kind']=='pptx': _edit_ppt_local(art['path'],normalized)
        elif art['kind']=='xlsx': await _edit_excel_local(art['path'],normalized)
        elif art['kind']=='docx': await _edit_word_local(art['path'],normalized)
        else: return False
        art['created']=time.time(); _save_artifacts(); await _send_artifact(message,user_id,art['path'],'✏️ Готово — изменил уже созданный файл, не начиная всё с нуля.')
        return True
    except Exception as e:
        logging.error('Artifact edit failed: %s',e,exc_info=True); await message.answer('⚠️ Не удалось безопасно изменить файл. Исходный файл сохранён, попробуй описать изменение конкретнее.'); return True

def _install_asyncio_exception_handler():
    """Логирует исключения фоновых asyncio-задач, которые иначе легко потерять."""
    loop = asyncio.get_running_loop()

    def _handler(loop, context):
        exc = context.get("exception")
        message = context.get("message", "Unhandled asyncio exception")
        if exc is not None:
            logging.error("ASYNCIO_BACKGROUND_ERROR: %s", message, exc_info=(type(exc), exc, exc.__traceback__))
        else:
            logging.error("ASYNCIO_BACKGROUND_ERROR: %s | context=%r", message, context)

    loop.set_exception_handler(_handler)


async def _startup_telegram_call(name, func, attempts=4):
    """Retries non-critical Telegram startup calls so a transient API/network error does not keep the bot offline."""
    last_err = None
    for attempt in range(attempts):
        try:
            return await func()
        except Exception as exc:
            last_err = exc
            logging.warning("Startup Telegram call %s failed (%s/%s): %s", name, attempt + 1, attempts, exc)
            if attempt + 1 < attempts:
                await asyncio.sleep(1.5 * (attempt + 1))
    raise last_err



def _maintenance_notice() -> str:
    return "🛠 MecauAI временно находится на техническом обслуживании. Попробуйте немного позже."



# ========================= ЕДИНАЯ АДМИН-ПАНЕЛЬ =========================
ADMIN_AUDIT_LOG = {}
ADMIN_MAINTENANCE = False
admin_grant_states = set()
admin_broadcast_states = set()

def _admin_find_user(value):
    """Resolve a Telegram ID safely for admin history tools."""
    try:
        uid = int(str(value).strip())
    except (TypeError, ValueError):
        return None
    if uid <= 0:
        return None
    return uid

def _admin_is_allowed(user_id: int) -> bool:
    return int(user_id) in ADMIN_IDS

def _admin_user_allowed(user_id: int) -> bool:
    return _admin_is_allowed(user_id)

def _admin_log(user_id: int, action: str, details: str = ""):
    ADMIN_AUDIT_LOG.setdefault(int(user_id), []).append({
        "action": action, "details": str(details)[:1000],
        "ts": datetime.now().isoformat(timespec="seconds")
    })
    ADMIN_AUDIT_LOG[int(user_id)] = ADMIN_AUDIT_LOG[int(user_id)][-100:]

def _admin_runtime_report() -> str:
    ram = _memory_mb()
    disk = _disk_free_mb()
    disk_text = f"{disk:.0f} MB" if disk is not None else "н/д"
    return (
        "🛡 <b>ЕДИНАЯ АДМИН-ПАНЕЛЬ</b>\n\n"
        f"🟢 Бот: <b>{'ТЕХРАБОТЫ' if ADMIN_MAINTENANCE else 'ONLINE'}</b>\n"
        f"👥 Пользователей: <b>{len(_all_known_user_ids())}</b>\n"
        f"👑 Администраторов: <b>{len(ADMIN_IDS)}</b>\n"
        f"⚙️ Активных задач: <b>{len(busy_users)}</b>\n"
        f"💾 RAM процесса: <b>{ram:.0f} MB</b>\n"
        f"💿 Свободно на диске: <b>{disk_text}</b>\n"
        f"⏳ Очередь тяжёлых задач: <b>{HEAVY_WAITERS}</b>"
    )

def _admin_stats_report() -> str:
    total_messages = sum(v.get("messages",0) for v in user_stats.values())
    total_voice = sum(v.get("voice",0) for v in user_stats.values())
    total_files = sum(v.get("files",0) for v in user_stats.values())
    total_images = sum(v.get("images",0) for v in user_stats.values())
    total_exports = sum(v.get("exports",0) for v in user_stats.values())
    activity_events = sum(len(v) for v in USER_REQUEST_HISTORY.values())
    # Если старый user_stats ещё пуст, новая история всё равно даёт реальные числа.
    if not total_messages and activity_events:
        total_messages = activity_events
    return (
        "📊 <b>СТАТИСТИКА</b>\n\n"
        f"👥 Пользователей: <b>{len(_all_known_user_ids())}</b>\n"
        f"💬 Сообщений: <b>{total_messages}</b>\n"
        f"🎙 Голосовых: <b>{total_voice}</b>\n"
        f"📎 Файлов: <b>{total_files}</b>\n"
        f"🖼 Изображений: <b>{total_images}</b>\n"
        f"📤 Экспортов: <b>{total_exports}</b>\n"
        f"⏳ Сейчас занято: <b>{len(busy_users)}</b>"
    )

def _admin_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📊 Статистика", callback_data="adm:stats"), InlineKeyboardButton(text="🖥 Сервер", callback_data="adm:server")],
        [InlineKeyboardButton(text="👥 Пользователи", callback_data="adm:users"), InlineKeyboardButton(text="📋 История", callback_data="adm:audit")],
        [InlineKeyboardButton(text="👑 Администраторы", callback_data="adm:admins"), InlineKeyboardButton(text="📢 Рассылка", callback_data="adm:broadcast")],
        [InlineKeyboardButton(text="🚨 Ошибки/состояние", callback_data="adm:health"), InlineKeyboardButton(text="🛠 Техобслуживание", callback_data="adm:maintenance")],
        [InlineKeyboardButton(text="🔄 Обновить", callback_data="adm:dashboard")],
    ])

def _admin_back_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="⬅️ Админ-панель", callback_data="adm:back")]])

async def _admin_replace_with_dashboard(call: CallbackQuery):
    try:
        await call.message.delete()
    except Exception:
        pass
    await call.message.answer(_admin_runtime_report(), parse_mode="HTML", reply_markup=_admin_keyboard())

def _admin_users_text(page=0, per_page=8):
    ids = sorted(_all_known_user_ids(), reverse=True)
    total_pages = max(1, (len(ids) + per_page - 1) // per_page)
    page = max(0, min(int(page), total_pages - 1))
    chunk = ids[page * per_page:(page + 1) * per_page]
    if not chunk:
        body = "Пользователей пока нет."
    else:
        lines = []
        for uid in chunk:
            p = USER_PROFILE_CACHE.get(int(uid), {})
            name = " ".join(x for x in [p.get("first_name"), p.get("last_name")] if x).strip()
            username = f"@{p['username']}" if p.get("username") else "без username"
            last = USER_ACTIVITY_INDEX.get(int(uid), "нет истории")
            lines.append(
                f"👤 <b>{html.escape(name or username)}</b>\n"
                f"ID: <code>{uid}</code> · {html.escape(username)}\n"
                f"Последняя активность: <b>{html.escape(str(last))}</b>"
            )
    return f"👥 <b>ПОЛЬЗОВАТЕЛИ</b>\n\n{body}\n\nСтраница {page+1}/{total_pages}"

def _admin_users_keyboard(page=0):
    ids = sorted(_all_known_user_ids(), reverse=True)
    pages = max(1, (len(ids) + 7) // 8)
    buttons = []
    nav = []
    if page > 0:
        nav.append(InlineKeyboardButton(text="⬅️", callback_data=f"au:{page-1}"))
    if page < pages - 1:
        nav.append(InlineKeyboardButton(text="➡️", callback_data=f"au:{page+1}"))
    if nav:
        buttons.append(nav)
    chunk = ids[page*8:(page+1)*8]
    for uid in chunk:
        p = USER_PROFILE_CACHE.get(int(uid), {})
        label = " ".join(x for x in [p.get("first_name"), p.get("last_name")] if x).strip()
        if not label:
            label = f"ID {uid}"
        buttons.append([InlineKeyboardButton(text=f"👤 {label[:40]}", callback_data=f"up:{uid}")])
    buttons.append([InlineKeyboardButton(text="⬅️ Админ-панель", callback_data="adm:back")])
    return InlineKeyboardMarkup(inline_keyboard=buttons)

def _admin_admins_text():
    ids = sorted(ADMIN_IDS)
    lines = [f"• <code>{uid}</code>{' — владелец' if uid == MY_ADMIN_ID else ''}" for uid in ids]
    return "👑 <b>АДМИНИСТРАТОРЫ</b>\n\n" + "\n".join(lines)

def _admin_admins_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="➕ Выдать доступ", callback_data="adm:grant")],
        [InlineKeyboardButton(text="➖ Забрать доступ", callback_data="adm:revoke")],
        [InlineKeyboardButton(text="⬅️ Админ-панель", callback_data="adm:back")],
    ])

async def _admin_panel_message(message: Message):
    user_id = message.from_user.id
    # Админ-панель никогда не открывается в группах/каналах:
    # так её содержимое и управляющие кнопки не могут попасть другим людям.
    if message.chat.type != "private":
        await message.answer("⛔ Админ-панель доступна только в личном чате с ботом.")
        return
    admin_grant_states.discard(user_id)
    admin_grant_states.discard(-user_id)
    admin_broadcast_states.discard(user_id)
    if not _admin_is_allowed(user_id):
        # Если у пользователя осталась старая клавиатура, сразу скрываем
        # административную кнопку и возвращаем обычное меню.
        await message.answer("⛔ Доступ запрещён.", reply_markup=get_quick_actions_keyboard())
        return
    _admin_log(user_id, "open_admin")
    await message.answer(_admin_runtime_report(), parse_mode="HTML", reply_markup=_admin_keyboard())

@dp.callback_query(F.data.startswith("au:"))
async def admin_users_page_callback(call: CallbackQuery):
    if not _admin_is_allowed(call.from_user.id):
        await call.answer("Доступ запрещён", show_alert=True)
        return
    try:
        page = int(call.data.split(":", 1)[1])
    except Exception:
        page = 0
    await call.answer()
    await call.message.edit_text(_admin_users_text(page), parse_mode="HTML", reply_markup=_admin_users_keyboard(page))

@dp.callback_query(F.data.startswith("adm:"))
async def admin_callback(call: CallbackQuery):
    if not _admin_is_allowed(call.from_user.id):
        await call.answer("Доступ запрещён", show_alert=True)
        return
    if not call.message or call.message.chat.type != "private":
        await call.answer("Админ-панель доступна только в личном чате.", show_alert=True)
        return
    action = call.data.split(":",1)[1]
    _admin_log(call.from_user.id, "callback", action)
    if action in {"dashboard", "back"}:
        await _admin_replace_with_dashboard(call)
    elif action == "stats":
        await call.message.edit_text(_admin_stats_report(), parse_mode="HTML", reply_markup=_admin_back_keyboard())
    elif action == "server":
        ram = _memory_mb(); disk = _disk_free_mb()
        disk_text = f"{disk:.0f} MB" if disk is not None else "н/д"
        server_text = (
            "🖥 <b>СЕРВЕР</b>\n\n"
            f"RAM процесса: <b>{ram:.0f} MB</b>\n"
            f"Свободно на диске: <b>{disk_text}</b>\n"
            f"Активные задачи: <b>{len(busy_users)}</b>\n"
            f"Очередь тяжёлых задач: <b>{HEAVY_WAITERS}</b>\n"
            f"AI concurrency: <b>{MAX_AI_CONCURRENCY}</b>\n"
            f"Heavy concurrency: <b>{MAX_HEAVY_CONCURRENCY}</b>\n"
            f"DATA_DIR: <code>{html.escape(DATA_DIR)}</code>"
        )
        await call.message.edit_text(server_text, parse_mode="HTML", reply_markup=_admin_back_keyboard())
    elif action == "users":
        await call.message.edit_text(_admin_users_text(0), parse_mode="HTML", reply_markup=_admin_users_keyboard(0))
    elif action == "audit":
        rows = []
        for uid, events in USER_REQUEST_HISTORY.items():
            for e in events[-5:]:
                rows.append((e.get("ts",""), uid, e))
        rows.sort(key=lambda x: x[0], reverse=True)
        lines = []
        for ts, uid, e in rows[:25]:
            prompt = html.escape((e.get("prompt") or "—")[:160])
            lines.append(f"{'✅' if e.get('status') == 'ok' else '⚠️'} <b>{html.escape(ts)}</b> · <code>{uid}</code> · {html.escape(e.get('kind','request'))}\n💬 {prompt}")
        body = "\n\n".join(lines) or "История действий пока пуста."
        await call.message.edit_text("📋 <b>ИСТОРИЯ ПОЛЬЗОВАТЕЛЕЙ</b>\n\n" + body, parse_mode="HTML", reply_markup=_admin_back_keyboard())
    elif action == "admins":
        await call.message.edit_text(_admin_admins_text(), parse_mode="HTML", reply_markup=_admin_admins_keyboard())
    elif action == "grant":
        admin_grant_states.add(call.from_user.id)
        await call.message.edit_text("➕ <b>Выдать доступ администратора</b>\n\nОтправь следующим сообщением Telegram ID пользователя.\n\nНапример: <code>123456789</code>", parse_mode="HTML", reply_markup=_admin_back_keyboard())
    elif action == "revoke":
        await call.message.edit_text(_admin_admins_text() + "\n\nЧтобы забрать доступ, нажми кнопку ниже и введи ID.", parse_mode="HTML", reply_markup=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="➖ Ввести ID", callback_data="adm:revoke_input")],[InlineKeyboardButton(text="⬅️ Назад", callback_data="adm:admins")]]))
    elif action == "revoke_input":
        admin_grant_states.discard(call.from_user.id)
        admin_grant_states.add(-call.from_user.id)
        await call.message.edit_text("➖ <b>Забрать доступ</b>\n\nОтправь Telegram ID пользователя.", parse_mode="HTML", reply_markup=_admin_back_keyboard())
    elif action == "broadcast":
        admin_broadcast_states.add(call.from_user.id)
        await call.message.edit_text("📢 <b>Рассылка</b>\n\nОтправь следующим сообщением текст рассылки.\n\n❌ Отмена — нажми «Админ-панель».", parse_mode="HTML", reply_markup=_admin_back_keyboard())
    elif action == "health":
        await call.message.edit_text(_admin_runtime_report(), parse_mode="HTML", reply_markup=_admin_back_keyboard())
    elif action == "maintenance":
        global ADMIN_MAINTENANCE
        ADMIN_MAINTENANCE = not ADMIN_MAINTENANCE
        _admin_log(call.from_user.id, "maintenance", str(ADMIN_MAINTENANCE))
        await call.message.edit_text(_admin_runtime_report(), parse_mode="HTML", reply_markup=_admin_keyboard())
    await call.answer()

# ================= USER ANALYTICS & REQUEST HISTORY =================
USER_ANALYTICS_HISTORY = True
USER_REQUEST_HISTORY = {}
USER_PROFILE_CACHE = {}
USER_ACTIVITY_INDEX = {}
USER_ACTIVITY_FILE = _data_file("user_activity.json")
activity_dirty = False

def _load_user_activity():
    global USER_REQUEST_HISTORY, USER_PROFILE_CACHE, USER_ACTIVITY_INDEX
    # One-time migration from old app root if present.
    legacy = os.path.join(PROJECT_ROOT, "user_activity.json")
    if not os.path.exists(USER_ACTIVITY_FILE) and os.path.exists(legacy) and legacy != USER_ACTIVITY_FILE:
        try:
            import shutil
            shutil.copy2(legacy, USER_ACTIVITY_FILE)
        except Exception as exc:
            logging.warning("Не удалось мигрировать user_activity.json: %s", exc)
    data = _load_json_dict(USER_ACTIVITY_FILE)
    for uid_s, item in data.items():
        try:
            uid = int(uid_s)
        except Exception:
            continue
        if not isinstance(item, dict):
            continue
        USER_PROFILE_CACHE[uid] = item.get("profile") or {}
        rows = item.get("history") or []
        USER_REQUEST_HISTORY[uid] = rows[-500:] if isinstance(rows, list) else []
        if USER_REQUEST_HISTORY[uid]:
            USER_ACTIVITY_INDEX[uid] = USER_REQUEST_HISTORY[uid][-1].get("ts", "—")

def _save_user_activity(immediate=False):
    global activity_dirty
    if not immediate:
        activity_dirty = True
        return
    data = {}
    all_ids = set(USER_PROFILE_CACHE) | set(USER_REQUEST_HISTORY)
    for uid in all_ids:
        data[str(uid)] = {
            "profile": USER_PROFILE_CACHE.get(uid, {}),
            "history": USER_REQUEST_HISTORY.get(uid, [])[-500:]
        }
    try:
        _atomic_json_write(USER_ACTIVITY_FILE, data)
        activity_dirty = False
    except Exception as exc:
        logging.error("Не удалось сохранить историю пользователей: %s", exc)

_load_user_activity()


def _all_known_user_ids() -> set:
    """Единый источник для админки: объединяет все локальные индексы пользователей."""
    ids = set(all_users_cache)
    ids.update(int(uid) for uid in users_db.keys() if str(uid).lstrip("-").isdigit())
    ids.update(int(uid) for uid in user_stats.keys() if str(uid).lstrip("-").isdigit())
    ids.update(int(uid) for uid in USER_PROFILE_CACHE.keys() if str(uid).lstrip("-").isdigit())
    ids.update(int(uid) for uid in USER_REQUEST_HISTORY.keys() if str(uid).lstrip("-").isdigit())
    return {uid for uid in ids if uid > 0}


def _reconcile_user_storage():
    """Восстанавливает индекс пользователей из имеющихся persistent JSON-файлов."""
    global all_users_cache, users_ids_dirty
    known = _all_known_user_ids()
    if known - all_users_cache:
        all_users_cache.update(known)
        try:
            _atomic_json_write(USERS_FILE, sorted(all_users_cache))
            users_ids_dirty = False
            logging.info("User index reconciled: %d users", len(all_users_cache))
        except Exception as exc:
            users_ids_dirty = True
            logging.warning("Не удалось восстановить users.json: %s", exc)


_reconcile_user_storage()


def _track_user_event(user_id: int, message=None, kind="request", status="ok",
                      prompt="", artifact=None, error=None):
    """Lightweight per-user audit/history. Stores metadata, not full generated files."""
    try:
        if message is not None:
            u = getattr(message, "from_user", None)
            profile = USER_PROFILE_CACHE.setdefault(user_id, {})
            if u:
                profile.update({
                    "id": user_id,
                    "username": getattr(u, "username", None),
                    "first_name": getattr(u, "first_name", None),
                    "last_name": getattr(u, "last_name", None),
                    "language_code": getattr(u, "language_code", None),
                })
        history = USER_REQUEST_HISTORY.setdefault(user_id, [])
        event = {
            "ts": datetime.now().isoformat(timespec="seconds"),
            "kind": kind,
            "status": status,
            "prompt": (prompt or "")[:1500],
            "artifact": (artifact or "")[:300],
            "error": (error or "")[:500],
        }
        history.append(event)
        # Hard cap per user to protect disk/RAM.
        if len(history) > 500:
            del history[:-500]
        USER_ACTIVITY_INDEX[user_id] = event["ts"]
        global activity_dirty
        activity_dirty = True
    except Exception:
        pass

def _user_profile_text(uid: int) -> str:
    p = USER_PROFILE_CACHE.get(uid, {})
    history = USER_REQUEST_HISTORY.get(uid, [])
    total = len(history)
    ok = sum(1 for x in history if x.get("status") == "ok")
    errors = total - ok
    by_kind = {}
    for x in history:
        by_kind[x.get("kind", "request")] = by_kind.get(x.get("kind", "request"), 0) + 1
    return (
        f"👤 <b>{html.escape(str(p.get('first_name') or 'Без имени'))}</b>\n"
        f"Username: <b>{('@' + p['username']) if p.get('username') else '—'}</b>\n"
        f"Telegram ID: <code>{uid}</code>\n"
        f"Язык: <b>{html.escape(str(p.get('language_code') or '—'))}</b>\n"
        f"Первая запись: <b>{history[0]['ts'] if history else '—'}</b>\n"
        f"Последняя активность: <b>{USER_ACTIVITY_INDEX.get(uid, '—')}</b>\n\n"
        f"📊 Запросов: <b>{total}</b>\n"
        f"✅ Успешных: <b>{ok}</b>\n"
        f"❌ Ошибок: <b>{errors}</b>\n"
        f"🎙 Voice: <b>{by_kind.get('voice', 0)}</b>\n"
        f"🖼 Vision: <b>{by_kind.get('vision', 0)}</b>\n"
        f"📄 Word: <b>{by_kind.get('word', 0)}</b>\n"
        f"📊 Excel: <b>{by_kind.get('excel', 0)}</b>\n"
        f"📈 PPTX: <b>{by_kind.get('pptx', 0)}</b>\n"
        f"📚 Documents: <b>{by_kind.get('document', 0)}</b>"
    )

def _user_history_text(uid: int, page=0, per_page=8) -> str:
    rows = USER_REQUEST_HISTORY.get(uid, [])
    total_pages = max(1, (len(rows) + per_page - 1) // per_page)
    page = max(0, min(page, total_pages - 1))
    chunk = list(reversed(rows))[page * per_page:(page + 1) * per_page]
    if not chunk:
        body = "История пока пуста."
    else:
        lines = []
        for e in chunk:
            prompt = html.escape(e.get("prompt") or "—")
            if len(prompt) > 180:
                prompt = prompt[:177] + "..."
            icon = "✅" if e.get("status") == "ok" else "❌"
            lines.append(
                f"{icon} <b>{e.get('ts')}</b> · {html.escape(e.get('kind','request'))}\n"
                f"💬 {prompt}"
            )
        body = "\n\n".join(lines)
    return f"📜 <b>История пользователя {uid}</b>\n\n{body}\n\nСтраница {page+1}/{total_pages}"

def _user_history_keyboard(uid: int, page=0):
    rows = USER_REQUEST_HISTORY.get(uid, [])
    pages = max(1, (len(rows) + 7) // 8)
    buttons = []
    nav = []
    if page > 0:
        nav.append(InlineKeyboardButton(text="⬅️", callback_data=f"uh:{uid}:{page-1}"))
    if page < pages - 1:
        nav.append(InlineKeyboardButton(text="➡️", callback_data=f"uh:{uid}:{page+1}"))
    if nav:
        buttons.append(nav)
    buttons.append([InlineKeyboardButton(text="👤 Профиль", callback_data=f"up:{uid}")])
    return InlineKeyboardMarkup(inline_keyboard=buttons)

@dp.message(Command("history"))
async def admin_user_history(message: Message):
    if not _admin_user_allowed(message.from_user.id):
        await message.answer("⛔ Доступ запрещён.")
        return
    parts = (message.text or "").split(maxsplit=1)
    if len(parts) != 2:
        await message.answer("Использование: <code>/history 123456789</code>", parse_mode="HTML")
        return
    uid = _admin_find_user(parts[1])
    if uid is None:
        await message.answer("⚠️ Некорректный ID пользователя.")
        return
    _admin_log(message.from_user.id, "history", str(uid))
    await message.answer(_user_history_text(uid), parse_mode="HTML",
                         reply_markup=_user_history_keyboard(uid))

@dp.callback_query(F.data.startswith("uh:"))
async def admin_user_history_callback(call: CallbackQuery):
    if not _admin_user_allowed(call.from_user.id):
        await call.answer("Доступ запрещён", show_alert=True)
        return
    _, uid_s, page_s = call.data.split(":", 2)
    uid, page = int(uid_s), int(page_s)
    await call.answer()
    await call.message.edit_text(_user_history_text(uid, page), parse_mode="HTML",
                                 reply_markup=_user_history_keyboard(uid, page))

@dp.callback_query(F.data.startswith("up:"))
async def admin_user_profile_callback(call: CallbackQuery):
    if not _admin_user_allowed(call.from_user.id):
        await call.answer("Доступ запрещён", show_alert=True)
        return
    uid = int(call.data.split(":", 1)[1])
    await call.answer()
    await call.message.edit_text(
        _user_profile_text(uid), parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📜 История", callback_data=f"uh:{uid}:0")],
            [InlineKeyboardButton(text="⬅️ Админ-панель", callback_data="adm:back")]
        ])
    )


async def main():
    _install_asyncio_exception_handler()
    cleanup_stale_tmp_files()
    cleanup_artifacts()
    _cleanup_orphan_artifact_files()
    logging.info("Legal PDFs: agreement=%s (%s), privacy=%s (%s)", os.path.isfile(AGREEMENT_FILE), AGREEMENT_FILE, os.path.isfile(PRIVACY_FILE), PRIVACY_FILE)
    await _refresh_groq_candidates()
    logging.info("Groq models: configured_text=%s candidates=%s configured_vision=%s candidates=%s", TEXT_MODEL, TEXT_MODEL_FALLBACKS, VISION_MODEL, VISION_MODEL_FALLBACKS)
    logging.info("Persistent users: %s (%d users)", USERS_FILE, len(_all_known_user_ids()))
    logging.info("Persistent user history: %s (%d users)", USER_ACTIVITY_FILE, len(USER_REQUEST_HISTORY))
    logging.info(
        "Startup: DATA_DIR=%s ARTIFACT_DIR=%s RAM=%.1fMB LIMIT=%sMB SOFT=%sMB HARD=%sMB "
        "DISK_FREE=%sMB AI_CONCURRENCY=%s HEAVY_CONCURRENCY=%s",
        DATA_DIR, ARTIFACT_DIR, _memory_mb(),
        round(_CONTAINER_RAM_MB) if _CONTAINER_RAM_MB else "unknown",
        RAM_SOFT_LIMIT_MB, RAM_HARD_LIMIT_MB,
        round(_disk_free_mb()) if _disk_free_mb() is not None else "unknown",
        MAX_AI_CONCURRENCY, MAX_HEAVY_CONCURRENCY
    )
    await _startup_telegram_call("set_my_commands", lambda: bot.set_my_commands([
        BotCommand(command="start", description="Перезапустить бота"),
        BotCommand(command="mode",  description="Сменить режим ассистента"),
        BotCommand(command="about", description="О возможностях"),
        BotCommand(command="checksub", description="Проверить подписку пользователя (админ)"),
    ]))
    await _startup_telegram_call("delete_webhook", bot.delete_webhook)
    maintenance = asyncio.create_task(maintenance_loop())
    print("🚀 Бот MecauAI 2.0 запущен: оптимизированные ресурсы, очередь, watchdog, артефакты и редактор активны...")
    try:
        await dp.start_polling(bot)
    finally:
        flush_persistent_state(force=True)
        maintenance.cancel()
        try:
            await maintenance
        except asyncio.CancelledError:
            pass
        except Exception:
            logging.exception("Maintenance shutdown error")
        try:
            await bot.session.close()
        except Exception:
            logging.exception("Bot session shutdown error")

if __name__ == "__main__":
    asyncio.run(main())
